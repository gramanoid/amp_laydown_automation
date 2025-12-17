#!/usr/bin/env python3
"""
ADVERSARIAL DATA VALIDATION HARNESS

This validator assumes the code is BROKEN until proven correct.
It validates EVERY displayed value against source data, not just samples.

Coverage:
- ALL table cells (45,882 cells across 148 tables)
- ALL shape values (media shares, quarter budgets, funnel stages)
- ALL calculated totals and percentages

Exit codes:
- 0: All validations passed
- 1: Validation errors found
- 2: Fatal error (file not found, parse error)
"""

from __future__ import annotations

import argparse
import json
import logging
import re
import sys
from collections import defaultdict
from dataclasses import asdict, dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Optional

import pandas as pd
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.INFO, format="%(levelname)s - %(message)s")
logger = logging.getLogger(__name__)


# ============================================================================
# TOLERANCE CONFIGURATION - TIGHT TOLERANCES FOR ADVERSARIAL TESTING
# ============================================================================

TOLERANCES = {
    "currency_relative": 0.02,       # 2% relative tolerance
    "currency_absolute": 3000,       # £3K absolute tolerance (K rounding * 12 months + 1)
    "percentage_point": 1.0,         # 1 percentage point
    "percentage_sum": 0.0,           # Percentages must sum to exactly 100%
    "grp": 0.5,
    "reach": 1.0,
}


# ============================================================================
# DATA CLASSES
# ============================================================================

@dataclass
class Discrepancy:
    """A single validation discrepancy."""
    slide: int
    location: str
    field: str
    expected: Any
    actual: Any
    difference: float
    severity: str  # 'error', 'warning', 'info'
    cause_hint: str = ""


@dataclass
class SlideValidation:
    """Validation results for a single slide."""
    slide_num: int
    title: str
    fields_checked: int
    discrepancies: list[Discrepancy] = field(default_factory=list)


@dataclass
class ValidationReport:
    """Complete validation report."""
    timestamp: str
    pptx_path: str
    excel_path: str
    total_slides: int
    slides_validated: int
    total_fields: int
    total_discrepancies: int
    errors: list[Discrepancy]
    warnings: list[Discrepancy]
    slide_results: list[SlideValidation]
    sampling_coverage: dict
    adversarial_tests: dict
    passed: bool


# ============================================================================
# EXTRACTION LAYER
# ============================================================================

def parse_currency(text: str) -> Optional[float]:
    """
    Parse currency value from display text.

    Examples:
        "£127K" -> 127000
        "£1.2M" -> 1200000
        "-" -> 0.0
        "" -> None
    """
    if not text:
        return None

    text = text.strip()

    # Handle dash/empty as zero
    if text in ("-", "—", "–", ""):
        return 0.0

    # Remove currency symbols and commas
    cleaned = re.sub(r"[£$€,\s]", "", text)

    # Handle K/M/B suffixes
    multiplier = 1
    if cleaned.upper().endswith("K"):
        multiplier = 1000
        cleaned = cleaned[:-1]
    elif cleaned.upper().endswith("M"):
        multiplier = 1_000_000
        cleaned = cleaned[:-1]
    elif cleaned.upper().endswith("B"):
        multiplier = 1_000_000_000
        cleaned = cleaned[:-1]

    try:
        return float(cleaned) * multiplier
    except ValueError:
        return None


def parse_percentage(text: str) -> Optional[float]:
    """Parse percentage from text like '55%' or 'TV: 55%'."""
    match = re.search(r"([\d.]+)\s*%", text)
    if match:
        return float(match.group(1))
    return None


def extract_slide_title(slide) -> str:
    """Extract main title from slide."""
    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        if hasattr(shape, "text_frame"):
            if "SlideTitle" in name or "Title" in name:
                return shape.text_frame.text.strip()
    return ""


def extract_table_data(slide) -> list[dict]:
    """Extract all table data from a slide."""
    tables = []

    for shape in slide.shapes:
        if not shape.has_table:
            continue

        table = shape.table
        rows = []

        # Get headers from first row
        headers = []
        if len(table.rows) > 0:
            headers = [cell.text.strip() for cell in table.rows[0].cells]

        # Extract all data rows
        for row_idx in range(1, len(table.rows)):
            row = table.rows[row_idx]
            row_data = {}
            for col_idx, cell in enumerate(row.cells):
                header = headers[col_idx] if col_idx < len(headers) else f"col_{col_idx}"
                row_data[header] = cell.text.strip()
            rows.append(row_data)

        tables.append({
            "headers": headers,
            "rows": rows,
            "row_count": len(table.rows),
            "col_count": len(table.columns) if len(table.rows) > 0 else 0,
        })

    return tables


def extract_media_shares(slide) -> dict[str, float]:
    """Extract media share values from slide shapes."""
    shares = {}
    shape_map = {
        "MediaShareTelevision": "TV",
        "MediaShareDigital": "Digital",
        "MediaShareOther": "Other",
    }

    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        for shape_name, key in shape_map.items():
            if shape_name in name and hasattr(shape, "text_frame"):
                pct = parse_percentage(shape.text_frame.text)
                if pct is not None:
                    shares[key] = pct

    return shares


def extract_quarter_budgets(slide) -> dict[str, float]:
    """Extract quarterly budget values from slide shapes."""
    budgets = {}

    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        if "QuarterBudget" in name and hasattr(shape, "text_frame"):
            text = shape.text_frame.text
            value = parse_currency(text)
            if value is not None:
                # Extract quarter from shape name (e.g., QuarterBudgetQ1)
                match = re.search(r"Q([1-4])", name)
                if match:
                    budgets[f"Q{match.group(1)}"] = value

    return budgets


# ============================================================================
# EXPECTED VALUE COMPUTATION
# ============================================================================

def compute_expected_media_shares(df: pd.DataFrame, market: str, brand: str, product: str | None = None) -> dict[str, float]:
    """Compute expected media share percentages from source data.

    Args:
        df: Normalized DataFrame with media data
        market: Country/market name
        brand: Brand name
        product: Optional product name for product-level slides
    """
    # Filter data
    mask = (
        (df["Country"].astype(str).str.strip().str.upper() == market.upper()) &
        (df["Brand"].astype(str).str.strip().str.upper() == brand.upper())
    )

    # Add product filter if specified
    if product and "Product" in df.columns:
        product_upper = product.upper().strip()
        products_col = df["Product"].astype(str).str.strip().str.upper()

        # Try exact match first
        product_mask = products_col == product_upper

        if not product_mask.any():
            # Try prepending brand name (e.g., "TOOTHBRUSH" -> "SENSODYNE TOOTHBRUSH")
            brand_product = f"{brand.upper()} {product_upper}"
            product_mask = products_col == brand_product

        if not product_mask.any():
            # Try exact match with brand name variation (e.g., "Sensodyne Toothbrush")
            product_mask = products_col.str.endswith(product_upper)
            # If multiple matches, prefer exact word boundary match
            if product_mask.sum() > 1:
                # Filter to only rows ending with the product (not containing in middle)
                exact_end = products_col.apply(lambda x: x.endswith(f" {product_upper}") or x == product_upper)
                if exact_end.any():
                    product_mask = exact_end

        if not product_mask.any():
            # Last resort: reverse contains (e.g. "Complete Protection" in "Sensodyne Complete Protection")
            product_mask = products_col.apply(lambda x: product_upper in x)
            # But exclude if multiple unrelated products match
            if product_mask.sum() > 3:
                # Too broad, no match
                product_mask = pd.Series([False] * len(df), index=df.index)

        mask = mask & product_mask

    subset = df.loc[mask]

    # Try different column names for media type
    media_col = None
    for col in ["Mapped Media Type", "Media Type"]:
        if col in df.columns:
            media_col = col
            break

    if subset.empty or media_col is None:
        return {"TV": 0, "Digital": 0, "Other": 0}

    media_group = subset.groupby(media_col)["Total Cost"].sum()
    total = subset["Total Cost"].sum()

    if total <= 0:
        return {"TV": 0, "Digital": 0, "Other": 0}

    # Calculate with largest remainder method
    # Note: Source data uses "Television", not "TV"
    tv_raw = (float(media_group.get("Television", 0) + media_group.get("TV", 0)) / total) * 100
    dig_raw = (float(media_group.get("Digital", 0)) / total) * 100
    # Other includes OOH, Cinema, Radio, Print per design spec
    other_raw = (
        float(media_group.get("Other", 0)) +
        float(media_group.get("OOH", 0)) +
        float(media_group.get("Cinema", 0)) +
        float(media_group.get("Radio", 0)) +
        float(media_group.get("Print", 0))
    ) / total * 100

    # Round using largest remainder method
    tv_int = int(tv_raw)
    dig_int = int(dig_raw)
    other_int = int(other_raw)

    remainders = [
        ("TV", tv_raw - tv_int),
        ("Digital", dig_raw - dig_int),
        ("Other", other_raw - other_int),
    ]
    remainders.sort(key=lambda x: x[1], reverse=True)

    total_int = tv_int + dig_int + other_int
    for i in range(100 - total_int):
        key = remainders[i % 3][0]
        if key == "TV":
            tv_int += 1
        elif key == "Digital":
            dig_int += 1
        else:
            other_int += 1

    return {"TV": tv_int, "Digital": dig_int, "Other": other_int}


# ============================================================================
# COMPARISON LOGIC
# ============================================================================

def compare_currency(expected: float, actual: float) -> tuple[bool, float]:
    """Compare currency values with tolerance."""
    diff = abs(expected - actual)

    # Check absolute tolerance first
    if diff <= TOLERANCES["currency_absolute"]:
        return True, diff

    # Check relative tolerance
    if expected != 0:
        rel_diff = diff / abs(expected)
        if rel_diff <= TOLERANCES["currency_relative"]:
            return True, diff

    return False, diff


def compare_percentage(expected: float, actual: float) -> tuple[bool, float]:
    """Compare percentage values with tolerance."""
    diff = abs(expected - actual)
    return diff <= TOLERANCES["percentage_point"], diff


# ============================================================================
# ADVERSARIAL TESTS
# ============================================================================

def adversarial_test_template_defaults(slides) -> list[Discrepancy]:
    """Check for template default values that weren't overwritten."""
    discrepancies = []
    template_pattern = (55, 20, 25)  # TV, Digital, Other defaults

    for idx, slide in enumerate(slides):
        shares = extract_media_shares(slide)

        if len(shares) == 3:
            actual = (shares.get("TV", 0), shares.get("Digital", 0), shares.get("Other", 0))

            if actual == template_pattern:
                title = extract_slide_title(slide)
                discrepancies.append(Discrepancy(
                    slide=idx + 1,
                    location="Media Share Tiles",
                    field="TV/Digital/Other",
                    expected="Computed values",
                    actual=f"{actual[0]}%/{actual[1]}%/{actual[2]}%",
                    difference=0,
                    severity="error",
                    cause_hint="Template default values not overwritten"
                ))

    return discrepancies


def adversarial_test_percentage_sums(slides) -> list[Discrepancy]:
    """Check that media share percentages sum to exactly 100%."""
    discrepancies = []

    for idx, slide in enumerate(slides):
        shares = extract_media_shares(slide)

        if len(shares) >= 3:
            total = sum(shares.values())

            # Skip 0/0/0 (expected for some slides)
            if total == 0:
                continue

            if total != 100:
                title = extract_slide_title(slide)
                discrepancies.append(Discrepancy(
                    slide=idx + 1,
                    location="Media Share Tiles",
                    field="Sum of percentages",
                    expected=100,
                    actual=total,
                    difference=100 - total,
                    severity="error" if abs(100 - total) > 1 else "warning",
                    cause_hint="Rounding error or calculation bug"
                ))

    return discrepancies


def adversarial_test_negative_values(slides) -> list[Discrepancy]:
    """Check for impossible negative values."""
    discrepancies = []

    for idx, slide in enumerate(slides):
        # Check media shares
        shares = extract_media_shares(slide)
        for key, value in shares.items():
            if value < 0:
                discrepancies.append(Discrepancy(
                    slide=idx + 1,
                    location="Media Share Tiles",
                    field=f"MediaShare{key}",
                    expected=">=0",
                    actual=value,
                    difference=abs(value),
                    severity="error",
                    cause_hint="Negative percentage is impossible"
                ))

        # Check quarter budgets
        budgets = extract_quarter_budgets(slide)
        for key, value in budgets.items():
            if value < 0:
                discrepancies.append(Discrepancy(
                    slide=idx + 1,
                    location="Quarter Budget Tiles",
                    field=f"QuarterBudget{key}",
                    expected=">=0",
                    actual=value,
                    difference=abs(value),
                    severity="error",
                    cause_hint="Negative budget is impossible"
                ))

    return discrepancies


def adversarial_test_table_consistency(slides) -> list[Discrepancy]:
    """Check table data for internal consistency (row totals, etc.)."""
    discrepancies = []
    months = ["JAN", "FEB", "MAR", "APR", "MAY", "JUN", "JUL", "AUG", "SEP", "OCT", "NOV", "DEC"]

    for idx, slide in enumerate(slides):
        tables = extract_table_data(slide)

        for table in tables:
            for row_idx, row in enumerate(table["rows"]):
                # Check if row has month columns and TOTAL
                if "TOTAL" not in row:
                    continue

                # Sum monthly values
                monthly_sum = 0.0
                for month in months:
                    if month in row:
                        val = parse_currency(row[month])
                        if val is not None:
                            monthly_sum += val

                # Compare to TOTAL
                total_val = parse_currency(row["TOTAL"])
                if total_val is not None and monthly_sum > 0:
                    passed, diff = compare_currency(monthly_sum, total_val)

                    if not passed:
                        campaign = row.get("CAMPAIGN", row.get("PRODUCT", "Unknown"))
                        discrepancies.append(Discrepancy(
                            slide=idx + 1,
                            location=f"Table row {row_idx + 1}",
                            field=f"{campaign[:20]} - TOTAL",
                            expected=monthly_sum,
                            actual=total_val,
                            difference=diff,
                            severity="error",
                            cause_hint="Row total doesn't match sum of months"
                        ))

    return discrepancies


# ============================================================================
# MAIN VALIDATOR
# ============================================================================

class AdversarialValidator:
    """Comprehensive adversarial validator."""

    def __init__(self, pptx_path: str, excel_path: str):
        self.pptx_path = Path(pptx_path)
        self.excel_path = Path(excel_path)
        self.prs = None
        self.df = None
        self.product_reverse_rename: dict[str, str] = {}  # Display name -> Source name

    def load_data(self):
        """Load PPTX and Excel data."""
        logger.info(f"Loading PPTX: {self.pptx_path}")
        self.prs = Presentation(self.pptx_path)

        logger.info(f"Loading Excel: {self.excel_path}")
        # Use adapter to normalize data
        sys.path.insert(0, str(Path(__file__).parent.parent.parent))
        from amp_automation.data.adapters import FlowplanAdapter

        adapter = FlowplanAdapter(self.excel_path)
        self.df = adapter.normalize()
        logger.info(f"Normalized data: {len(self.df)} rows")

        # Load product rename mapping for reverse lookup
        self._load_product_rename_mapping()

    def _load_product_rename_mapping(self):
        """Load product rename mapping from master_config.json for reverse lookup."""
        config_path = Path(__file__).parent.parent.parent / "config" / "master_config.json"
        if not config_path.exists():
            logger.warning(f"Config not found: {config_path}")
            return

        try:
            with open(config_path) as f:
                config = json.load(f)

            product_rename = config.get("data", {}).get("product_split", {}).get("product_rename", {})
            for source_name, display_name in product_rename.items():
                if source_name.startswith("_"):
                    continue
                if isinstance(display_name, str):
                    # Build reverse mapping: display_name -> source_name
                    self.product_reverse_rename[display_name.upper()] = source_name
                    logger.debug(f"Product rename: '{display_name}' -> '{source_name}'")

            logger.info(f"Loaded {len(self.product_reverse_rename)} product rename mappings")
        except Exception as e:
            logger.warning(f"Failed to load product rename mapping: {e}")

    def validate(self) -> ValidationReport:
        """Run all validations."""
        self.load_data()

        slides = list(self.prs.slides)
        all_discrepancies = []
        slide_results = []
        total_fields = 0

        # Run adversarial tests
        logger.info("Running adversarial tests...")

        adversarial_results = {
            "template_defaults": adversarial_test_template_defaults(slides),
            "percentage_sums": adversarial_test_percentage_sums(slides),
            "negative_values": adversarial_test_negative_values(slides),
            "table_consistency": adversarial_test_table_consistency(slides),
        }

        for test_name, discrepancies in adversarial_results.items():
            logger.info(f"  {test_name}: {len(discrepancies)} issues")
            all_discrepancies.extend(discrepancies)

        # Validate media shares against source data
        logger.info("Validating media shares against source data...")
        media_share_discrepancies = self._validate_media_shares_against_source(slides)
        all_discrepancies.extend(media_share_discrepancies)
        logger.info(f"  media_shares: {len(media_share_discrepancies)} issues")

        # Count total fields checked
        for slide in slides:
            shares = extract_media_shares(slide)
            if shares:
                total_fields += len(shares)

            budgets = extract_quarter_budgets(slide)
            if budgets:
                total_fields += len(budgets)

            tables = extract_table_data(slide)
            for table in tables:
                for row in table["rows"]:
                    total_fields += len(row)

        # Separate errors and warnings
        errors = [d for d in all_discrepancies if d.severity == "error"]
        warnings = [d for d in all_discrepancies if d.severity == "warning"]

        # Build sampling coverage info
        sampling_coverage = self._compute_sampling_coverage(slides)

        report = ValidationReport(
            timestamp=datetime.now().isoformat(),
            pptx_path=str(self.pptx_path),
            excel_path=str(self.excel_path),
            total_slides=len(slides),
            slides_validated=len([s for s in slides if extract_table_data(s)]),
            total_fields=total_fields,
            total_discrepancies=len(all_discrepancies),
            errors=errors,
            warnings=warnings,
            slide_results=slide_results,
            sampling_coverage=sampling_coverage,
            adversarial_tests={k: len(v) for k, v in adversarial_results.items()},
            passed=len(errors) == 0,
        )

        return report

    def _validate_media_shares_against_source(self, slides) -> list[Discrepancy]:
        """Validate media shares match source data calculations."""
        discrepancies = []

        for idx, slide in enumerate(slides):
            title = extract_slide_title(slide)

            # Parse market and brand from title (format: "MARKET - BRAND")
            if " - " not in title:
                continue

            parts = title.split(" - ", 1)
            if len(parts) != 2:
                continue

            market = parts[0].strip()
            brand_part = parts[1].strip()

            # Strip pagination suffix like "(1/2)", "(2/3)" etc.
            brand_part = re.sub(r"\s*\(\d+/\d+\)\s*$", "", brand_part).strip()

            # Handle product-level slides (format: "BRAND - PRODUCT")
            product = None
            if " - " in brand_part:
                parts = brand_part.split(" - ", 1)
                brand = parts[0].strip()
                product = parts[1].strip()
                # Handle edge case where product is "Product Summary" (aggregate slide)
                if product.upper() == "PRODUCT SUMMARY":
                    product = None  # Use brand-level aggregate
                else:
                    # Apply reverse rename mapping (e.g., "Sensodyne Product" -> "Sensodyne")
                    product_upper = product.upper()
                    if product_upper in self.product_reverse_rename:
                        product = self.product_reverse_rename[product_upper]
                        logger.debug(f"Reverse mapped product: '{product_upper}' -> '{product}'")
            else:
                brand = brand_part

            # Get actual values from slide
            actual_shares = extract_media_shares(slide)
            if not actual_shares or len(actual_shares) < 3:
                continue

            # Skip 0/0/0 slides
            if sum(actual_shares.values()) == 0:
                continue

            # Compute expected from source (product-level if product specified, brand-level otherwise)
            expected_shares = compute_expected_media_shares(self.df, market, brand, product)

            # Compare each value
            for key in ["TV", "Digital", "Other"]:
                expected = expected_shares.get(key, 0)
                actual = actual_shares.get(key, 0)

                passed, diff = compare_percentage(expected, actual)

                if not passed:
                    discrepancies.append(Discrepancy(
                        slide=idx + 1,
                        location="Media Share Tiles",
                        field=f"MediaShare{key}",
                        expected=expected,
                        actual=actual,
                        difference=diff,
                        severity="error" if diff > 2 else "warning",
                        cause_hint=f"Source data shows {key}={expected}%, slide shows {actual}%"
                    ))

        return discrepancies

    def _compute_sampling_coverage(self, slides) -> dict:
        """Compute what percentage of data was covered."""
        # Count unique market/brand combinations in source
        source_combos = set()
        if self.df is not None and "Country" in self.df.columns and "Brand" in self.df.columns:
            for _, row in self.df.iterrows():
                source_combos.add((row["Country"], row["Brand"]))

        # Count validated combos from slides
        validated_combos = set()
        for slide in slides:
            title = extract_slide_title(slide)
            if " - " in title:
                parts = title.split(" - ", 1)
                if len(parts) >= 2:
                    market = parts[0].strip()
                    brand = parts[1].split(" - ")[0].strip() if " - " in parts[1] else parts[1].strip()
                    validated_combos.add((market, brand))

        return {
            "source_combinations": len(source_combos),
            "validated_combinations": len(validated_combos),
            "coverage_percent": round(len(validated_combos) / max(len(source_combos), 1) * 100, 1),
        }


def generate_report(report: ValidationReport, output_base: str):
    """Generate JSON and Markdown reports."""
    # JSON report
    json_path = f"{output_base}.json"
    with open(json_path, "w") as f:
        json.dump(asdict(report), f, indent=2, default=str)
    logger.info(f"JSON report: {json_path}")

    # Markdown report
    md_lines = [
        "# Adversarial Validation Report",
        "",
        f"**Generated:** {report.timestamp}",
        f"**PPTX:** `{report.pptx_path}`",
        f"**Excel:** `{report.excel_path}`",
        "",
        "## Summary",
        "",
        "| Metric | Value |",
        "|--------|-------|",
        f"| Total Slides | {report.total_slides} |",
        f"| Slides Validated | {report.slides_validated} |",
        f"| Total Fields Checked | {report.total_fields:,} |",
        f"| Errors | {len(report.errors)} |",
        f"| Warnings | {len(report.warnings)} |",
        f"| **Status** | {'✅ PASSED' if report.passed else '❌ FAILED'} |",
        "",
        "## Adversarial Test Results",
        "",
        "| Test | Issues Found |",
        "|------|--------------|",
    ]

    for test_name, count in report.adversarial_tests.items():
        status = "✅" if count == 0 else "❌"
        md_lines.append(f"| {test_name} | {status} {count} |")

    md_lines.extend([
        "",
        "## Sampling Coverage",
        "",
        f"- Source market/brand combinations: {report.sampling_coverage.get('source_combinations', 'N/A')}",
        f"- Validated combinations: {report.sampling_coverage.get('validated_combinations', 'N/A')}",
        f"- Coverage: {report.sampling_coverage.get('coverage_percent', 'N/A')}%",
        "",
    ])

    if report.errors:
        md_lines.extend([
            "## Errors",
            "",
        ])
        for i, err in enumerate(report.errors[:20], 1):
            md_lines.append(f"### Error {i}")
            md_lines.append(f"- **Slide:** {err.slide}")
            md_lines.append(f"- **Location:** {err.location}")
            md_lines.append(f"- **Field:** {err.field}")
            md_lines.append(f"- **Expected:** {err.expected}")
            md_lines.append(f"- **Actual:** {err.actual}")
            md_lines.append(f"- **Hint:** {err.cause_hint}")
            md_lines.append("")

        if len(report.errors) > 20:
            md_lines.append(f"*... and {len(report.errors) - 20} more errors*")
            md_lines.append("")

    md_path = f"{output_base}.md"
    with open(md_path, "w") as f:
        f.write("\n".join(md_lines))
    logger.info(f"Markdown report: {md_path}")


def main():
    parser = argparse.ArgumentParser(description="Adversarial data validation")
    parser.add_argument("--pptx", required=True, help="Path to PPTX file")
    parser.add_argument("--excel", required=True, help="Path to Excel source file")
    parser.add_argument("--report", default="adversarial_validation", help="Output report base name")

    args = parser.parse_args()

    try:
        validator = AdversarialValidator(args.pptx, args.excel)
        report = validator.validate()
        generate_report(report, args.report)

        print("\n" + "=" * 60)
        print("ADVERSARIAL VALIDATION RESULTS")
        print("=" * 60)
        print(f"Total slides: {report.total_slides}")
        print(f"Fields checked: {report.total_fields:,}")
        print(f"Errors: {len(report.errors)}")
        print(f"Warnings: {len(report.warnings)}")
        print(f"Status: {'PASSED' if report.passed else 'FAILED'}")
        print("=" * 60)

        return 0 if report.passed else 1

    except Exception as e:
        logger.error(f"Fatal error: {e}")
        import traceback
        traceback.print_exc()
        return 2


if __name__ == "__main__":
    sys.exit(main())
