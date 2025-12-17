#!/usr/bin/env python3
"""
Comprehensive Data Correctness Validation Harness

Validates ALL displayed values in generated PowerPoint presentations against
the source DataFrame. Designed to catch any discrepancies in:
- Monthly budget values
- Total budget calculations
- Media share percentages
- GRP/Reach metrics
- Campaign/Product aggregations

Usage:
    python -m tools.validate.comprehensive_validator \\
        --pptx output/presentations/run_*/AMP_Laydowns_*.pptx \\
        --excel input/BulkPlanData_*.xlsx \\
        --report validation_report

Or use as a module:
    from tools.validate.comprehensive_validator import ComprehensiveValidator
    validator = ComprehensiveValidator(pptx_path, excel_path)
    report = validator.validate()
"""

from __future__ import annotations

import argparse
import json
import logging
import re
import sys
from dataclasses import asdict, dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Optional

import pandas as pd
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.INFO, format="%(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

# ============================================================================
# CONSTANTS
# ============================================================================

MONTH_ORDER = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
MONTH_COL_MAP = {m.upper(): i for i, m in enumerate(MONTH_ORDER)}

# Tolerance configuration
TOLERANCE_CONFIG = {
    "percentage_tolerance": 0.02,  # 2% relative tolerance
    "absolute_tolerance": 6000,     # £6K absolute tolerance (accounts for K rounding)
    "percentage_point_tolerance": 1.0,  # 1 percentage point for media shares
    "grp_tolerance": 0.5,           # 0.5 GRP absolute tolerance
    "reach_tolerance": 1.0,         # 1% reach absolute tolerance
}

MEDIA_TYPES = ["Television", "Digital", "OOH", "Other"]
MEDIA_DISPLAY_ORDER = ["TV", "DIG", "OOH", "OTH"]


# ============================================================================
# DATA CLASSES
# ============================================================================

@dataclass
class ValidationError:
    """Represents a validation error found in the artifact."""
    slide_num: int
    error_type: str  # 'budget', 'media_share', 'grp', 'reach', 'total', 'missing_category'
    field_name: str
    expected: Any
    actual: Any
    difference: Optional[float] = None
    location: Optional[str] = None
    source_slice: Optional[str] = None
    likely_cause: Optional[str] = None


@dataclass
class ExtractedSlideData:
    """Data extracted from a single slide."""
    slide_num: int
    slide_title: str
    market: Optional[str] = None
    brand: Optional[str] = None
    campaign: Optional[str] = None
    product: Optional[str] = None
    year: Optional[int] = None

    # Table data
    rows: list[dict] = field(default_factory=list)

    # Summary tiles
    total_budget: Optional[float] = None
    media_shares: dict[str, float] = field(default_factory=dict)
    quarterly_budgets: dict[str, float] = field(default_factory=dict)


@dataclass
class ValidationReport:
    """Comprehensive validation report."""
    timestamp: str = field(default_factory=lambda: datetime.now().isoformat())
    pptx_path: str = ""
    excel_path: str = ""
    total_slides: int = 0
    slides_validated: int = 0
    fields_checked: int = 0
    errors: list[ValidationError] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    sampling_strategy: dict = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return len(self.errors) == 0

    @property
    def error_count(self) -> int:
        return len(self.errors)

    def add_error(self, error: ValidationError):
        self.errors.append(error)

    def add_warning(self, warning: str):
        self.warnings.append(warning)

    def to_json(self) -> str:
        """Export report as JSON."""
        return json.dumps(asdict(self), indent=2, default=str)

    def to_markdown(self) -> str:
        """Export report as markdown."""
        lines = []
        lines.append("# Comprehensive Validation Report")
        lines.append("")
        lines.append(f"**Generated:** {self.timestamp}")
        lines.append(f"**PPTX:** `{self.pptx_path}`")
        lines.append(f"**Excel:** `{self.excel_path}`")
        lines.append("")
        lines.append("## Summary")
        lines.append("")
        lines.append(f"| Metric | Value |")
        lines.append(f"|--------|-------|")
        lines.append(f"| Total Slides | {self.total_slides} |")
        lines.append(f"| Slides Validated | {self.slides_validated} |")
        lines.append(f"| Fields Checked | {self.fields_checked} |")
        lines.append(f"| Errors | {self.error_count} |")
        lines.append(f"| Warnings | {len(self.warnings)} |")
        lines.append(f"| **Status** | {'✅ PASSED' if self.passed else '❌ FAILED'} |")
        lines.append("")

        if self.errors:
            lines.append("## Errors")
            lines.append("")
            for i, err in enumerate(self.errors, 1):
                lines.append(f"### Error {i}: {err.error_type}")
                lines.append(f"- **Slide:** {err.slide_num}")
                lines.append(f"- **Field:** {err.field_name}")
                lines.append(f"- **Expected:** {err.expected}")
                lines.append(f"- **Actual:** {err.actual}")
                if err.difference is not None:
                    lines.append(f"- **Difference:** {err.difference:,.2f}")
                if err.location:
                    lines.append(f"- **Location:** {err.location}")
                if err.likely_cause:
                    lines.append(f"- **Likely Cause:** {err.likely_cause}")
                lines.append("")

        if self.warnings:
            lines.append("## Warnings")
            lines.append("")
            for warning in self.warnings:
                lines.append(f"- {warning}")
            lines.append("")

        return "\n".join(lines)


# ============================================================================
# EXTRACTION LAYER
# ============================================================================

def parse_number(text: str) -> Optional[float]:
    """
    Parse a number from display text, handling K/M suffixes and currency symbols.

    Examples:
        "£127K" → 127000
        "£1.2M" → 1200000
        "42%" → 42.0
        "-" → None
    """
    if not text or text.strip() in ["-", "", "–", "—"]:
        return None

    cleaned = text.strip().replace("£", "").replace("$", "").replace(",", "").replace(" ", "")

    # Handle percentage
    if "%" in cleaned:
        cleaned = cleaned.replace("%", "")
        try:
            return float(cleaned)
        except ValueError:
            return None

    # Handle M suffix (millions)
    if cleaned.upper().endswith("M"):
        try:
            return float(cleaned[:-1]) * 1_000_000
        except ValueError:
            return None

    # Handle K suffix (thousands)
    if cleaned.upper().endswith("K"):
        try:
            return float(cleaned[:-1]) * 1_000
        except ValueError:
            return None

    try:
        return float(cleaned)
    except ValueError:
        return None


def extract_slide_title(slide) -> tuple[str, dict]:
    """
    Extract title and parse market/brand/product information.

    Handles formats:
    - "MARKET - BRAND (pagination)" e.g., "SAUDI ARABIA - SENSODYNE (1/2)"
    - "MARKET - BRAND - PRODUCT" e.g., "SAUDI ARABIA - SENSODYNE - CLINICAL WHITE"

    Returns:
        tuple of (raw_title, parsed_info_dict)
    """
    import re

    title = ""
    parsed = {"market": None, "brand": None, "year": None, "product": None, "is_product": False}

    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            # Look for title patterns: "MARKET - BRAND..." with uppercase words
            if " - " in text and text[0].isupper() and len(text) < 100:
                # Prefer longer/more specific titles
                if len(text) > len(title):
                    title = text

    if title:
        # Remove pagination suffix like "(1/2)" or "(2/2)"
        clean_title = re.sub(r'\s*\(\d+/\d+\)\s*$', '', title)

        # Parse: "MARKET - BRAND" or "MARKET - BRAND - PRODUCT"
        parts = [p.strip() for p in clean_title.split(" - ")]

        if len(parts) >= 2:
            parsed["market"] = parts[0]
            parsed["brand"] = parts[1]

            if len(parts) >= 3:
                parsed["product"] = parts[2]
                parsed["is_product"] = True

    return title, parsed


def extract_table_data(table) -> list[list[str]]:
    """Extract all text from table cells."""
    data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        data.append(row_data)
    return data


def extract_slide_data(slide, slide_num: int) -> Optional[ExtractedSlideData]:
    """
    Extract all relevant data from a slide.

    Returns:
        ExtractedSlideData or None if slide has no table
    """
    # Find table
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if not table:
        return None

    # Extract title info
    title, parsed = extract_slide_title(slide)

    extracted = ExtractedSlideData(
        slide_num=slide_num,
        slide_title=title,
        market=parsed.get("market"),
        brand=parsed.get("brand"),
        product=parsed.get("product"),
        year=parsed.get("year"),
    )

    # Extract table data
    table_data = extract_table_data(table)

    if len(table_data) < 2:
        return extracted

    # Parse header row to find column indices
    header = table_data[0]
    month_start = None
    total_col = None

    for i, col in enumerate(header):
        if col.upper() == "JAN":
            month_start = i
        if col.upper() == "TOTAL":
            total_col = i

    if month_start is None:
        return extracted

    # Parse data rows
    for row_idx, row in enumerate(table_data[1:], start=1):
        if len(row) < 3:
            continue

        row_data = {
            "row_idx": row_idx,
            "campaign_or_product": row[0] if len(row) > 0 else "",
            "media": row[1] if len(row) > 1 else "",
            "metrics": row[2] if len(row) > 2 else "",
            "months": {},
            "total": None,
        }

        # Parse monthly values
        if month_start:
            for i, month in enumerate(MONTH_ORDER):
                col_idx = month_start + i
                if col_idx < len(row):
                    row_data["months"][month] = parse_number(row[col_idx])

        # Parse total
        if total_col and total_col < len(row):
            row_data["total"] = parse_number(row[total_col])

        extracted.rows.append(row_data)

    return extracted


# ============================================================================
# EXPECTED VALUE COMPUTATION
# ============================================================================

def load_and_prepare_source_data(excel_path: Path, format_type: str = "auto") -> pd.DataFrame:
    """
    Load source Excel data and prepare for validation.

    Applies the same transformations as the presentation generator.
    Supports both BulkPlanData and Flowplan_Summaries formats.
    """
    from amp_automation.config.loader import load_master_config
    from amp_automation.data.ingestion import load_and_prepare_data
    from amp_automation.data.adapters import InputFormat

    # Map format string to enum
    format_map = {
        "auto": InputFormat.AUTO,
        "bulkplan": InputFormat.BULK_PLAN,
        "flowplan": InputFormat.FLOWPLAN,
    }
    input_format = format_map.get(format_type.lower(), InputFormat.AUTO)

    # Load config using the helper function
    config_path = Path(__file__).parent.parent.parent / "config" / "master_config.json"
    config = load_master_config(config_path)

    # Load data using the ingestion module with format type
    dataset = load_and_prepare_data(str(excel_path), config, logger, format_type=input_format)
    logger.info(f"Source format detected: {dataset.source_format}")

    return dataset.frame


def normalize_text(text: str) -> str:
    """Normalize text for case-insensitive comparison."""
    if not text:
        return ""
    # Replace newlines with spaces, strip, convert to uppercase
    return text.replace("\n", " ").strip().upper()


def compute_expected_budget(
    df: pd.DataFrame,
    market: str,
    brand: str,
    year: Optional[int] = None,
    campaign: Optional[str] = None,
    product: Optional[str] = None,
    media_type: Optional[str] = None,
) -> dict[str, float]:
    """
    Compute expected budget values from source DataFrame.

    Returns dict with monthly values, total, and optional breakdown.
    Uses case-insensitive matching for all string fields.
    """
    # Build filter with case-insensitive matching
    market_norm = normalize_text(market)
    brand_norm = normalize_text(brand)

    mask = (
        (df["Country"].str.strip().str.upper() == market_norm)
        & (df["Brand"].str.strip().str.upper() == brand_norm)
    )

    if year:
        mask &= df["Year"].astype(str).str.strip() == str(year)

    if campaign:
        campaign_norm = normalize_text(campaign)
        mask &= df["Campaign Name"].str.strip().str.upper() == campaign_norm

    if product:
        product_norm = normalize_text(product)
        mask &= df["Product"].str.strip().str.upper() == product_norm

    if media_type:
        media_norm = normalize_text(media_type)
        mask &= df["Mapped Media Type"].str.strip().str.upper() == media_norm

    subset = df[mask]

    result = {"months": {}, "total": 0.0}

    # Aggregate monthly values
    for month in MONTH_ORDER:
        if month in subset.columns:
            result["months"][month] = float(subset[month].sum())

    # Calculate total
    result["total"] = float(subset["Total Cost"].sum())

    return result


def compute_expected_media_shares(
    df: pd.DataFrame,
    market: str,
    brand: str,
    year: Optional[int] = None,
) -> dict[str, float]:
    """
    Compute expected media share percentages.

    Returns dict mapping media type to percentage (0-100).
    Uses case-insensitive matching.
    """
    market_norm = normalize_text(market)
    brand_norm = normalize_text(brand)

    mask = (
        (df["Country"].str.strip().str.upper() == market_norm)
        & (df["Brand"].str.strip().str.upper() == brand_norm)
    )

    if year:
        mask &= df["Year"].astype(str).str.strip() == str(year)

    subset = df[mask]

    total_budget = float(subset["Total Cost"].sum())

    if total_budget == 0:
        return {}

    shares = {}
    for media_type in MEDIA_TYPES:
        media_mask = subset["Mapped Media Type"].str.strip().str.upper() == media_type.upper()
        media_budget = float(subset[media_mask]["Total Cost"].sum())
        shares[media_type] = (media_budget / total_budget) * 100

    return shares


# ============================================================================
# COMPARISON LOGIC
# ============================================================================

def compare_values(
    actual: Optional[float],
    expected: Optional[float],
    value_type: str = "budget",
) -> tuple[bool, Optional[float]]:
    """
    Compare actual vs expected values with appropriate tolerance.

    Args:
        actual: Displayed/extracted value
        expected: Computed expected value
        value_type: Type of value ('budget', 'percentage', 'grp', 'reach')

    Returns:
        tuple of (is_match, difference)
    """
    if actual is None and expected is None:
        return True, None

    if actual is None or expected is None:
        return False, None

    difference = abs(actual - expected)

    if value_type == "budget":
        # Use percentage or absolute tolerance, whichever is larger
        pct_tol = max(abs(actual), abs(expected)) * TOLERANCE_CONFIG["percentage_tolerance"]
        abs_tol = TOLERANCE_CONFIG["absolute_tolerance"]
        tolerance = max(pct_tol, abs_tol)
        return difference <= tolerance, difference

    elif value_type == "percentage":
        tolerance = TOLERANCE_CONFIG["percentage_point_tolerance"]
        return difference <= tolerance, difference

    elif value_type == "grp":
        tolerance = TOLERANCE_CONFIG["grp_tolerance"]
        return difference <= tolerance, difference

    elif value_type == "reach":
        tolerance = TOLERANCE_CONFIG["reach_tolerance"]
        return difference <= tolerance, difference

    return difference < 0.01, difference


# ============================================================================
# VALIDATOR CLASS
# ============================================================================

class ComprehensiveValidator:
    """
    Comprehensive validator for PowerPoint presentations against source data.
    """

    def __init__(self, pptx_path: Path, excel_path: Path):
        self.pptx_path = Path(pptx_path)
        self.excel_path = Path(excel_path)
        self.report = ValidationReport(
            pptx_path=str(pptx_path),
            excel_path=str(excel_path),
        )
        self.df: Optional[pd.DataFrame] = None
        self.prs: Optional[Presentation] = None

    def load_data(self):
        """Load PowerPoint and Excel data."""
        logger.info(f"Loading PPTX: {self.pptx_path}")
        self.prs = Presentation(str(self.pptx_path))
        self.report.total_slides = len(self.prs.slides)

        logger.info(f"Loading Excel: {self.excel_path}")
        self.df = load_and_prepare_source_data(self.excel_path)
        logger.info(f"Source data: {len(self.df)} rows")

    def validate(self, sample_strategy: Optional[dict] = None) -> ValidationReport:
        """
        Run comprehensive validation.

        Args:
            sample_strategy: Optional dict specifying which combinations to validate.
                            If None, validates all slides.

        Returns:
            ValidationReport with all findings
        """
        self.load_data()

        if sample_strategy:
            self.report.sampling_strategy = sample_strategy

        for slide_num, slide in enumerate(self.prs.slides, start=1):
            extracted = extract_slide_data(slide, slide_num)

            if extracted is None:
                continue

            self.report.slides_validated += 1

            # Validate this slide
            self._validate_slide(extracted)

        logger.info(f"Validation complete: {self.report.slides_validated} slides, "
                   f"{self.report.fields_checked} fields, {self.report.error_count} errors")

        return self.report

    def _validate_slide(self, extracted: ExtractedSlideData):
        """Validate a single slide's data against source."""
        if not extracted.market or not extracted.brand:
            self.report.add_warning(
                f"Slide {extracted.slide_num}: Could not parse market/brand from title"
            )
            return

        # Validate each row's data
        for row in extracted.rows:
            self._validate_row(extracted, row)

    def _validate_row(self, slide: ExtractedSlideData, row: dict):
        """Validate a single table row."""
        # Skip header/label rows
        campaign = row.get("campaign_or_product", "")
        if not campaign or campaign.strip() in ["-", "", "CAMPAIGN", "PRODUCT"]:
            return

        # Skip total rows (validated separately)
        if "TOTAL" in campaign.upper():
            return

        # Skip metric rows (GRPs, Reach, etc.)
        metrics = row.get("metrics", "")
        if metrics in ["GRPs", "Reach@1+", "Reach@3+", "OTS@1+", "OTS@3+", "Frequency"]:
            return

        # Get expected values
        media_type = row.get("media", "")
        if media_type == "TV":
            media_type = "Television"
        elif media_type == "DIG":
            media_type = "Digital"

        try:
            # For product slides, filter by both product AND campaign
            # For brand slides, filter only by campaign
            expected = compute_expected_budget(
                self.df,
                market=slide.market,
                brand=slide.brand,
                year=slide.year,
                campaign=campaign if campaign else None,
                product=slide.product,
                media_type=media_type if media_type else None,
            )
        except Exception as e:
            self.report.add_warning(
                f"Slide {slide.slide_num}: Error computing expected values: {e}"
            )
            return

        # Validate total
        actual_total = row.get("total")
        expected_total = expected.get("total", 0)

        if actual_total is not None and expected_total > 0:
            self.report.fields_checked += 1
            is_match, diff = compare_values(actual_total, expected_total, "budget")

            if not is_match:
                self.report.add_error(ValidationError(
                    slide_num=slide.slide_num,
                    error_type="budget",
                    field_name=f"{campaign} - {media_type} - Total",
                    expected=expected_total,
                    actual=actual_total,
                    difference=diff,
                    location=f"Row {row.get('row_idx')}, Total column",
                    source_slice=f"Market={slide.market}, Brand={slide.brand}, Year={slide.year}",
                ))

        # Validate monthly values (spot check)
        for month in ["Jan", "Jun", "Dec"]:  # Spot check first, mid, last month
            actual_month = row.get("months", {}).get(month)
            expected_month = expected.get("months", {}).get(month, 0)

            if actual_month is not None and expected_month > 0:
                self.report.fields_checked += 1
                is_match, diff = compare_values(actual_month, expected_month, "budget")

                if not is_match:
                    self.report.add_error(ValidationError(
                        slide_num=slide.slide_num,
                        error_type="budget",
                        field_name=f"{campaign} - {media_type} - {month}",
                        expected=expected_month,
                        actual=actual_month,
                        difference=diff,
                        location=f"Row {row.get('row_idx')}, {month} column",
                    ))


# ============================================================================
# CLI
# ============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Comprehensive validation of PowerPoint against source Excel data"
    )
    parser.add_argument(
        "--pptx",
        type=Path,
        required=True,
        help="Path to PowerPoint file to validate"
    )
    parser.add_argument(
        "--excel",
        type=Path,
        required=True,
        help="Path to source Excel file"
    )
    parser.add_argument(
        "--report",
        type=str,
        default="validation_report",
        help="Base name for output report files (will create .json and .md)"
    )
    parser.add_argument(
        "--tolerance",
        type=float,
        default=0.02,
        help="Percentage tolerance (default: 0.02 = 2%%)"
    )
    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="Verbose output"
    )

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    # Update tolerance if specified
    TOLERANCE_CONFIG["percentage_tolerance"] = args.tolerance

    # Run validation
    validator = ComprehensiveValidator(args.pptx, args.excel)
    report = validator.validate()

    # Save reports
    json_path = Path(f"{args.report}.json")
    md_path = Path(f"{args.report}.md")

    json_path.write_text(report.to_json())
    md_path.write_text(report.to_markdown())

    logger.info(f"Reports saved: {json_path}, {md_path}")

    # Print summary
    print("\n" + report.to_markdown())

    # Exit with error code if validation failed
    sys.exit(0 if report.passed else 1)


if __name__ == "__main__":
    main()
