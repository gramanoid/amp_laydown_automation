#!/usr/bin/env python
"""Verify that unmerge operation removed ALL merge attributes."""

import sys
from pptx import Presentation

def verify_no_merges(pptx_path: str, slide_idx: int) -> bool:
    """Verify that a specific slide has NO merge attributes."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]

    # Find table
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if not table:
        print(f"❌ No table found on slide {slide_idx + 1}")
        return False

    # Check every cell
    merge_count = 0
    merge_details = []

    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            tc = cell._tc

            # Check all merge attributes
            row_span = tc.get('rowSpan')
            grid_span = tc.get('gridSpan')
            v_merge = tc.get('vMerge')
            h_merge = tc.get('hMerge')

            if row_span or grid_span or v_merge or h_merge:
                merge_count += 1
                attrs = []
                if row_span:
                    attrs.append(f"rowSpan={row_span}")
                if grid_span:
                    attrs.append(f"gridSpan={grid_span}")
                if v_merge:
                    attrs.append(f"vMerge={v_merge}")
                if h_merge:
                    attrs.append(f"hMerge={h_merge}")

                merge_details.append(f"  Cell ({row_idx}, {col_idx}): {', '.join(attrs)}")

    print(f"\n{'='*60}")
    print(f"Slide {slide_idx + 1} Verification Results")
    print(f"{'='*60}")
    print(f"Table size: {len(table.rows)} rows × {len(table.columns)} cols")
    print(f"Total cells: {len(table.rows) * len(table.columns)}")
    print(f"Cells with merge attributes: {merge_count}")

    if merge_count == 0:
        print(f"✅ SUCCESS: All merge attributes removed!")
        return True
    else:
        print(f"❌ FAILURE: {merge_count} cells still have merge attributes:")
        for detail in merge_details[:20]:  # Show first 20
            print(detail)
        if len(merge_details) > 20:
            print(f"  ... and {len(merge_details) - 20} more")
        return False

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python verify_unmerge.py <pptx_path> <slide_number>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_num = int(sys.argv[2])
    slide_idx = slide_num - 1  # Convert to 0-based index

    success = verify_no_merges(pptx_path, slide_idx)
    sys.exit(0 if success else 1)
