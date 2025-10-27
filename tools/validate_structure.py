"""Structural validator for generated AMP decks.

This script inspects a generated deck to ensure the layout, shapes, and metric
rows adhere to the Template V4 structural contract documented on 20 Oct 2025.
It is intentionally data-agnostic: values may vary, but geometry, labels, and
required elements must match the template.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional

from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_CONTRACT_PATH = PROJECT_ROOT / "config" / "structural_contract.json"


@dataclass(slots=True)
class Issue:
    slide_index: int
    message: str

    def __str__(self) -> str:
        return f"Slide {self.slide_index}: {self.message}"


def _normalize(text: str) -> str:
    """Normalize text for comparisons (uppercase, trim, replace newlines)."""

    if text is None:
        return ""
    cleaned = text.replace("\uFFFD", "£").replace("\xa0", " ").replace("\n", " ")
    return re.sub(r"\s+", " ", cleaned).strip()


def _normalize_metric(metric: str) -> str:
    return _normalize(metric).upper()


def _load_contract(path: Path) -> dict[str, object]:
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def _derive_export_date_from_excel(excel_path: Optional[Path]) -> Optional[str]:
    if not excel_path:
        return None
    match = re.search(r"(\d{4})_(\d{2})_(\d{2})", excel_path.stem)
    if not match:
        return None
    year, month, day = match.groups()
    return f"{day}{month}{year[-2:]}"


def _iter_table_rows(table) -> Iterable:
    return list(table.rows)


def _validate_required_shapes(slide, required_shapes: List[str]) -> List[str]:
    available = {shape.name for shape in slide.shapes}
    missing = [shape for shape in required_shapes if shape not in available]
    return missing


def _table_has_grand_total(table, grand_total_label: str) -> bool:
    """Check if the table contains a grand total row (indicates last slide of brand/market)."""
    grand_total_normalized = _normalize_metric(grand_total_label)
    for row in _iter_table_rows(table)[1:]:  # Skip header
        cells = row.cells
        if cells:
            campaign_cell = _normalize_metric(cells[0].text)
            if campaign_cell == grand_total_normalized:
                return True
    return False


def _validate_last_slide_shapes(slide, last_slide_only_shapes: List[str]) -> List[str]:
    """Validate that all last-slide-only shapes exist on the slide."""
    available = {shape.name for shape in slide.shapes}
    missing = [shape for shape in last_slide_only_shapes if shape not in available]
    return missing


def _validate_table_header(table, expected_header: List[str]) -> Optional[str]:
    header_cells = [cell for cell in table.rows[0].cells]
    observed = [_normalize(cell.text).upper() for cell in header_cells]
    if observed != [item.upper() for item in expected_header]:
        diff = ", ".join(observed)
        return f"Header row mismatch. Observed: [{diff}]"
    return None


def _validate_media_sections(
    table,
    contract: dict[str, object],
    is_final_slide: bool = True,
) -> List[str]:
    issues: List[str] = []
    expected_order = {media: idx for idx, media in enumerate(contract["media_order"])}
    media_metrics = {
        media: {metric.upper() for metric in metrics}
        for media, metrics in contract["media_metrics"].items()
    }
    campaign_total_label = _normalize_metric(contract["campaign_total_label"])
    grand_total_label = _normalize_metric(contract["grand_total_label"])

    rows = _iter_table_rows(table)
    current_media: Optional[str] = None
    last_media_position = -1
    seen_grand_total = False

    for idx, row in enumerate(rows[1:], start=1):
        cells = row.cells
        media_cell = _normalize(cells[1].text)
        metric_cell = _normalize_metric(cells[2].text)
        campaign_cell = _normalize_metric(cells[0].text)

        if not any(_normalize(cell.text) for cell in cells):
            continue

        if campaign_cell == grand_total_label:
            seen_grand_total = True
            continue

        if campaign_cell == campaign_total_label:
            continue

        if media_cell and media_cell != "-":
            normalized_media = media_cell.upper()
            if campaign_cell not in {"-", campaign_total_label, grand_total_label}:
                last_media_position = -1
            if normalized_media not in expected_order:
                issues.append(
                    f"Unexpected media '{media_cell}' in row {idx + 1}. "
                    "Update structural contract if this is intentional."
                )
            else:
                position = expected_order[normalized_media]
                if position < last_media_position:
                    issues.append(
                        f"Media '{media_cell}' appears out of order in row {idx + 1}."
                    )
                last_media_position = position
            current_media = normalized_media

        if metric_cell in {"", "-", campaign_total_label, grand_total_label}:
            continue

        if metric_cell == "MONTHLY TOTAL (£ 000)":
            continue

        if current_media is None:
            issues.append(
                f"Metric '{metric_cell}' found in row {idx + 1} before any media header."
            )
            continue

        allowed_metrics = media_metrics.get(current_media, set())
        if metric_cell not in allowed_metrics:
            issues.append(
                f"Metric '{metric_cell}' in row {idx + 1} is not allowed for media '{current_media}'."
            )

    # Only check for grand total on final slides (indicated by is_final_slide parameter)
    if is_final_slide and not seen_grand_total:
        issues.append("Grand total row not found in table.")

    return issues


def _validate_footer(footer_shape, expected_date: Optional[str]) -> Optional[str]:
    if footer_shape is None:
        return "FooterNotes shape missing."

    text = _normalize(footer_shape.text).upper()
    if "SOURCE:" not in text:
        return "Footer missing 'Source:' prefix."

    if expected_date:
        normalized = text.replace("_", "").replace(" ", "")
        target = expected_date.upper()
        if target not in normalized:
            return (
                f"Source line does not include expected export date '{expected_date}'. "
                "Update generator or footnote."
            )

    return None


def validate_presentation(
    presentation_path: Path,
    contract_path: Path = DEFAULT_CONTRACT_PATH,
    excel_path: Optional[Path] = None,
) -> List[Issue]:
    contract = _load_contract(contract_path)
    presentation = Presentation(str(presentation_path))
    expected_date = _derive_export_date_from_excel(excel_path)

    issues: List[Issue] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        table_shape = next(
            (shape for shape in slide.shapes if shape.name == contract["table_shape_name"]),
            None,
        )
        if table_shape is None:
            continue

        missing_shapes = _validate_required_shapes(slide, contract["required_shapes"])
        if missing_shapes:
            issues.append(
                Issue(
                    slide_index,
                    f"Missing required shapes: {', '.join(missing_shapes)}",
                )
            )

        header_error = _validate_table_header(table_shape.table, contract["table_header"])
        if header_error:
            issues.append(Issue(slide_index, header_error))

        # Determine if this is a final slide (has grand total row)
        is_final_slide = _table_has_grand_total(table_shape.table, contract["grand_total_label"])

        for message in _validate_media_sections(table_shape.table, contract, is_final_slide):
            issues.append(Issue(slide_index, message))

        if is_final_slide:
            # Check for last-slide-only shapes if they're defined in the contract
            if "last_slide_only_shapes" in contract:
                missing_last_slide_shapes = _validate_last_slide_shapes(
                    slide, contract["last_slide_only_shapes"]
                )
                if missing_last_slide_shapes:
                    issues.append(
                        Issue(
                            slide_index,
                            f"Missing last-slide-only shapes: {', '.join(missing_last_slide_shapes)}",
                        )
                    )

            # Validate footer (which appears on final slides)
            footer_shape = next(
                (shape for shape in slide.shapes if shape.name == contract["footer_shape"]),
                None,
            )
            footer_error = _validate_footer(footer_shape, expected_date)
            if footer_error:
                issues.append(Issue(slide_index, footer_error))

    return issues


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Validate structural fidelity of a generated AMP deck.")
    parser.add_argument("presentation", type=Path, help="Path to the generated PPTX deck.")
    parser.add_argument(
        "--contract",
        type=Path,
        default=DEFAULT_CONTRACT_PATH,
        help="Path to structural contract JSON (default: config/structural_contract.json).",
    )
    parser.add_argument(
        "--excel",
        type=Path,
        default=None,
        help="Raw Excel path used for generation (enables Source date validation).",
    )
    return parser.parse_args()


def main() -> None:
    args = _parse_args()
    issues = validate_presentation(args.presentation, args.contract, args.excel)

    if issues:
        print("Structural validation failed:")
        for issue in issues:
            print(f" - {issue}")
        sys.exit(1)

    print(f"Structural validation passed for {args.presentation}.")


if __name__ == "__main__":  # pragma: no cover
    main()
