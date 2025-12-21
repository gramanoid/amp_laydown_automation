#!/usr/bin/env python3
"""
Comprehensive Split Box Validation Tool

Validates CEJ (AWA/CON/PUR), MEDIA (TV/DIG/OTHER), and QUARTER (Q1-Q4) split boxes
displayed at the bottom of slides against source DataFrame calculations.

Usage:
    python tools/validate_split_boxes.py [--pptx PATH] [--excel PATH] [--tolerance 0.01]
"""

import argparse
import json
import re
import sys
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation

# Add project root to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent))
from amp_automation.data.adapters import get_adapter, InputFormat
import logging

# Media type mapping (same as config)
MEDIA_TYPE_MAPPING = {
    "Television": "TV",
    "TV": "TV",
    "Digital": "Digital",
    "OOH": "OOH",
    "Other": "Other",
    "Print": "Other",
    "Radio": "Other",
    "Cinema": "Other",
}

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# =============================================================================
# CONFIGURATION
# =============================================================================

DEFAULT_PPTX = Path("output/presentations/run_20251217_130600/AMP_Laydowns_171225.pptx")
DEFAULT_EXCEL = Path("input/Flowplan_Summaries_MEA_2025_12_17.xlsx")

# Shape name patterns
SHAPE_PATTERNS = {
    "media": {
        "tv": "MediaShareTelevision",
        "digital": "MediaShareDigital",
        "other": "MediaShareOther",
    },
    "funnel": {
        "awareness": "FunnelShareAwareness",
        "consideration": "FunnelShareConsideration",
        "purchase": "FunnelSharePurchase",
    },
    "quarter": {
        "q1": "QuarterBudgetQ1",
        "q2": "QuarterBudgetQ2",
        "q3": "QuarterBudgetQ3",
        "q4": "QuarterBudgetQ4",
    },
}

# Tolerance settings
PCT_TOLERANCE = 1  # 1 percentage point for percentages
CURRENCY_TOLERANCE = 0.05  # 5% relative tolerance for currency values


# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class ExtractedValue:
    """Represents a value extracted from a slide."""
    slide_num: int
    slide_title: str
    category: str  # "media", "funnel", "quarter"
    field_name: str  # "tv", "digital", "awareness", "q1", etc.
    raw_text: str
    parsed_value: float | None
    unit: str  # "%", "£", etc.


@dataclass
class ExpectedValue:
    """Represents expected computed value from source data."""
    category: str
    field_name: str
    value: float
    unit: str
    source_filter: dict  # market, brand, year used to compute


@dataclass
class Discrepancy:
    """Represents a mismatch between actual and expected."""
    slide_num: int
    slide_title: str
    category: str
    field_name: str
    actual_value: float | None
    expected_value: float
    diff: float
    unit: str
    severity: str  # "error", "warning"
    source_filter: dict


@dataclass
class ValidationResult:
    """Complete validation result."""
    timestamp: str
    pptx_path: str
    excel_path: str
    total_slides: int
    slides_with_split_boxes: int
    fields_checked: int
    errors: list[Discrepancy] = field(default_factory=list)
    warnings: list[Discrepancy] = field(default_factory=list)
    passed: list[dict] = field(default_factory=list)
    sampling_coverage: dict = field(default_factory=dict)


# =============================================================================
# EXTRACTION LAYER
# =============================================================================

def extract_slide_title(slide) -> str:
    """Extract slide title text."""
    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        if hasattr(shape, "text_frame"):
            if "SlideTitle" in name or "Title" in name or "TitlePlaceholder" in name:
                return shape.text_frame.text.strip()
    # Fallback: look for any text in title-like position
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame.text:
            text = shape.text_frame.text.strip()
            if " - " in text and len(text) < 100:  # Likely a title
                return text
    return ""


def parse_slide_context(title: str) -> dict:
    """Parse market, brand, product from slide title."""
    context = {"market": "", "brand": "", "product": "", "is_product_slide": False, "is_continuation": False}

    if not title:
        return context

    # Remove continuation indicator like "(2/2)", "(Continued)", etc.
    clean_title = re.sub(r'\s*\(\d+/\d+\)\s*$', '', title)
    clean_title = re.sub(r'\s*\(Continued\)\s*$', '', clean_title, flags=re.IGNORECASE)
    
    if clean_title != title:
        context["is_continuation"] = True

    # Pattern: "MARKET - BRAND" or "MARKET - BRAND - PRODUCT"
    parts = clean_title.split(" - ")
    if len(parts) >= 2:
        context["market"] = parts[0].strip()
        context["brand"] = parts[1].strip()
        if len(parts) >= 3:
            context["product"] = parts[2].strip()
            context["is_product_slide"] = True

    return context


def extract_shape_text(slide, shape_name: str) -> str | None:
    """Extract text from a shape by name."""
    for shape in slide.shapes:
        if getattr(shape, "name", "") == shape_name:
            if hasattr(shape, "text_frame"):
                return shape.text_frame.text.strip()
    return None


def parse_percentage(text: str) -> tuple[float | None, str]:
    """Parse percentage from text like 'TV: 55%' -> (55.0, '%')."""
    if not text:
        return None, ""
    match = re.search(r"(\d+(?:\.\d+)?)\s*%", text)
    if match:
        return float(match.group(1)), "%"
    return None, ""


def parse_currency(text: str) -> tuple[float | None, str]:
    """Parse currency from text like 'Q1: £659K' or 'Q2: £1.2M'."""
    if not text:
        return None, ""

    # Match currency symbol followed by number and optional K/M suffix
    # Pattern: £262K, £1M, £1.2M, £1,234K
    match = re.search(r"[£$€]\s*([\d,]+(?:\.\d+)?)\s*([KkMm])?", text)
    if match:
        value_str = match.group(1).replace(",", "")
        multiplier = match.group(2)

        value = float(value_str)

        if multiplier:
            if multiplier.upper() == "K":
                value *= 1000
            elif multiplier.upper() == "M":
                value *= 1_000_000

        return value, "£"

    return None, ""


def extract_all_split_box_values(slide, slide_num: int) -> list[ExtractedValue]:
    """Extract all split box values from a slide."""
    title = extract_slide_title(slide)
    values = []

    # Extract media shares (TV, Digital, Other)
    for field_key, shape_name in SHAPE_PATTERNS["media"].items():
        raw_text = extract_shape_text(slide, shape_name)
        if raw_text:
            parsed, unit = parse_percentage(raw_text)
            values.append(ExtractedValue(
                slide_num=slide_num,
                slide_title=title,
                category="media",
                field_name=field_key,
                raw_text=raw_text,
                parsed_value=parsed,
                unit=unit
            ))

    # Extract funnel shares (AWA, CON, PUR)
    for field_key, shape_name in SHAPE_PATTERNS["funnel"].items():
        raw_text = extract_shape_text(slide, shape_name)
        if raw_text:
            parsed, unit = parse_percentage(raw_text)
            values.append(ExtractedValue(
                slide_num=slide_num,
                slide_title=title,
                category="funnel",
                field_name=field_key,
                raw_text=raw_text,
                parsed_value=parsed,
                unit=unit
            ))

    # Extract quarter budgets (Q1-Q4)
    for field_key, shape_name in SHAPE_PATTERNS["quarter"].items():
        raw_text = extract_shape_text(slide, shape_name)
        if raw_text:
            parsed, unit = parse_currency(raw_text)
            values.append(ExtractedValue(
                slide_num=slide_num,
                slide_title=title,
                category="quarter",
                field_name=field_key,
                raw_text=raw_text,
                parsed_value=parsed,
                unit=unit
            ))

    return values


# =============================================================================
# EXPECTED VALUE COMPUTATION
# =============================================================================

def load_excel_data(excel_path: Path) -> pd.DataFrame:
    """Load and normalize Excel source data using the project's adapter pattern."""
    # Use the adapter pattern to normalize data
    adapter = get_adapter(excel_path, InputFormat.AUTO, logger)
    df = adapter.normalize()
    
    # Apply media type mapping
    if "Media Type" in df.columns:
        df["Mapped Media Type"] = df["Media Type"].map(lambda m: MEDIA_TYPE_MAPPING.get(m, m))
    
    return df


def filter_data_for_slide(df: pd.DataFrame, context: dict) -> pd.DataFrame:
    """Filter DataFrame based on slide context (market, brand, product)."""
    market = context.get("market", "").strip()
    brand = context.get("brand", "").strip()
    product = context.get("product", "").strip()
    is_product_slide = context.get("is_product_slide", False)

    if not market or not brand:
        return pd.DataFrame()

    # Build base filter
    mask = (
        (df["Country"].astype(str).str.strip().str.upper() == market.upper()) &
        (df["Brand"].astype(str).str.strip().str.upper() == brand.upper())
    )

    # Add year filter if available (default to 2025)
    if "Year" in df.columns:
        mask = mask & (df["Year"].astype(str).str.strip() == "2025")

    # Handle product-level filtering
    if is_product_slide and product:
        # Skip "PRODUCT SUMMARY" - it aggregates all products
        if product.upper() != "PRODUCT SUMMARY":
            # Try direct match first
            product_col = df["Product"].astype(str).str.strip().str.upper()
            product_mask = (product_col == product.upper())

            # If no direct match, try with brand prefix
            if not product_mask.any():
                prefixed = f"{brand} {product}".upper()
                product_mask = (product_col == prefixed)

            # Also try reverse mapping for renamed products
            rename_map = {
                "PARODONTAX PRODUCT": "PARODONTAX",
                "SENSODYNE PRODUCT": "SENSODYNE",
                "CALPOL PRODUCT": "CALPOL",
            }
            if product.upper() in rename_map:
                original = rename_map[product.upper()]
                product_mask = product_mask | (product_col == original)

            mask = mask & product_mask

    return df.loc[mask].copy()


def compute_expected_media_shares(subset: pd.DataFrame) -> dict[str, float]:
    """Compute expected media share percentages using largest remainder rounding."""
    if subset.empty or "Mapped Media Type" not in subset.columns:
        return {"tv": 0, "digital": 0, "other": 0}

    media_group = subset.groupby("Mapped Media Type")["Total Cost"].sum()
    total_cost = subset["Total Cost"].sum()

    if total_cost <= 0:
        return {"tv": 0, "digital": 0, "other": 0}

    tv_value = float(media_group.get("TV", 0.0))
    digital_value = float(media_group.get("Digital", 0.0))
    # Other includes OOH
    other_value = float(media_group.get("Other", 0.0)) + float(media_group.get("OOH", 0.0))

    # Calculate raw percentages
    tv_pct_raw = (tv_value / total_cost) * 100
    digital_pct_raw = (digital_value / total_cost) * 100
    other_pct_raw = (other_value / total_cost) * 100

    # Floor values
    tv_pct = int(tv_pct_raw)
    digital_pct = int(digital_pct_raw)
    other_pct = int(other_pct_raw)

    # Largest remainder method to sum to 100%
    remainders = [
        ("tv", tv_pct_raw - tv_pct, tv_pct),
        ("digital", digital_pct_raw - digital_pct, digital_pct),
        ("other", other_pct_raw - other_pct, other_pct),
    ]
    remainders.sort(key=lambda x: x[1], reverse=True)

    current_sum = tv_pct + digital_pct + other_pct
    result = {"tv": tv_pct, "digital": digital_pct, "other": other_pct}

    for i in range(100 - current_sum):
        key = remainders[i % 3][0]
        result[key] += 1

    return result


def compute_expected_funnel_shares(subset: pd.DataFrame) -> dict[str, float]:
    """Compute expected funnel stage percentages using largest remainder rounding."""
    if subset.empty or "Funnel Stage" not in subset.columns:
        return {"awareness": 0, "consideration": 0, "purchase": 0}

    funnel_group = subset.groupby("Funnel Stage")["Total Cost"].sum()
    total_cost = subset["Total Cost"].sum()

    if total_cost <= 0:
        return {"awareness": 0, "consideration": 0, "purchase": 0}

    awa_value = float(funnel_group.get("Awareness", 0.0))
    con_value = float(funnel_group.get("Consideration", 0.0))
    pur_value = float(funnel_group.get("Purchase", 0.0))

    # Calculate raw percentages
    awa_pct_raw = (awa_value / total_cost) * 100
    con_pct_raw = (con_value / total_cost) * 100
    pur_pct_raw = (pur_value / total_cost) * 100

    # Floor values
    awa_pct = int(awa_pct_raw)
    con_pct = int(con_pct_raw)
    pur_pct = int(pur_pct_raw)

    # Largest remainder method
    remainders = [
        ("awareness", awa_pct_raw - awa_pct, awa_pct),
        ("consideration", con_pct_raw - con_pct, con_pct),
        ("purchase", pur_pct_raw - pur_pct, pur_pct),
    ]
    remainders.sort(key=lambda x: x[1], reverse=True)

    current_sum = awa_pct + con_pct + pur_pct
    result = {"awareness": awa_pct, "consideration": con_pct, "purchase": pur_pct}

    for i in range(100 - current_sum):
        key = remainders[i % 3][0]
        result[key] += 1

    return result


def compute_expected_quarter_budgets(subset: pd.DataFrame) -> dict[str, float]:
    """Compute expected quarter budget totals."""
    if subset.empty:
        return {"q1": 0, "q2": 0, "q3": 0, "q4": 0}

    months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    month_values = {}
    for month in months:
        if month in subset.columns:
            month_values[month] = float(subset[month].sum())
        else:
            month_values[month] = 0.0

    return {
        "q1": month_values["Jan"] + month_values["Feb"] + month_values["Mar"],
        "q2": month_values["Apr"] + month_values["May"] + month_values["Jun"],
        "q3": month_values["Jul"] + month_values["Aug"] + month_values["Sep"],
        "q4": month_values["Oct"] + month_values["Nov"] + month_values["Dec"],
    }


# =============================================================================
# COMPARISON LOGIC
# =============================================================================

def compare_values(
    extracted: ExtractedValue,
    expected: dict[str, float],
    source_filter: dict
) -> Discrepancy | None:
    """Compare extracted value against expected and return discrepancy if any."""

    expected_val = expected.get(extracted.field_name)
    if expected_val is None:
        return None

    actual_val = extracted.parsed_value

    # Handle missing extracted value
    if actual_val is None:
        return Discrepancy(
            slide_num=extracted.slide_num,
            slide_title=extracted.slide_title,
            category=extracted.category,
            field_name=extracted.field_name,
            actual_value=None,
            expected_value=expected_val,
            diff=expected_val,
            unit=extracted.unit,
            severity="error",
            source_filter=source_filter
        )

    # Compute difference based on category
    if extracted.category in ("media", "funnel"):
        # Percentage comparison - absolute difference
        diff = abs(actual_val - expected_val)
        if diff > PCT_TOLERANCE:
            return Discrepancy(
                slide_num=extracted.slide_num,
                slide_title=extracted.slide_title,
                category=extracted.category,
                field_name=extracted.field_name,
                actual_value=actual_val,
                expected_value=expected_val,
                diff=diff,
                unit="%",
                severity="error" if diff > 5 else "warning",
                source_filter=source_filter
            )
    elif extracted.category == "quarter":
        # Currency comparison - relative difference
        if expected_val > 0:
            rel_diff = abs(actual_val - expected_val) / expected_val
            if rel_diff > CURRENCY_TOLERANCE:
                return Discrepancy(
                    slide_num=extracted.slide_num,
                    slide_title=extracted.slide_title,
                    category=extracted.category,
                    field_name=extracted.field_name,
                    actual_value=actual_val,
                    expected_value=expected_val,
                    diff=actual_val - expected_val,
                    unit="£",
                    severity="error" if rel_diff > 0.1 else "warning",
                    source_filter=source_filter
                )
        elif actual_val > 0:
            # Expected 0 but got something
            return Discrepancy(
                slide_num=extracted.slide_num,
                slide_title=extracted.slide_title,
                category=extracted.category,
                field_name=extracted.field_name,
                actual_value=actual_val,
                expected_value=expected_val,
                diff=actual_val,
                unit="£",
                severity="warning",
                source_filter=source_filter
            )

    return None


# =============================================================================
# MAIN VALIDATION
# =============================================================================

def validate_presentation(pptx_path: Path, excel_path: Path) -> ValidationResult:
    """Run comprehensive validation of all split boxes."""

    prs = Presentation(pptx_path)
    df = load_excel_data(excel_path)

    result = ValidationResult(
        timestamp=datetime.now().isoformat(),
        pptx_path=str(pptx_path),
        excel_path=str(excel_path),
        total_slides=len(prs.slides),
        slides_with_split_boxes=0,
        fields_checked=0,
    )

    # Track unique combinations validated
    validated_combos = set()

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1

        # Extract all split box values from this slide
        extracted_values = extract_all_split_box_values(slide, slide_num)

        if not extracted_values:
            continue

        result.slides_with_split_boxes += 1

        # Parse slide context
        title = extract_slide_title(slide)
        context = parse_slide_context(title)

        if not context["market"] or not context["brand"]:
            # Can't validate without market/brand
            for ev in extracted_values:
                result.warnings.append(Discrepancy(
                    slide_num=slide_num,
                    slide_title=title,
                    category=ev.category,
                    field_name=ev.field_name,
                    actual_value=ev.parsed_value,
                    expected_value=0,
                    diff=0,
                    unit=ev.unit,
                    severity="warning",
                    source_filter={"error": "Could not parse market/brand from title"}
                ))
            continue

        # Filter source data
        subset = filter_data_for_slide(df, context)
        source_filter = {
            "market": context["market"],
            "brand": context["brand"],
            "product": context.get("product", ""),
            "rows_matched": len(subset),
        }

        if subset.empty:
            # No data for this combination
            for ev in extracted_values:
                if ev.parsed_value and ev.parsed_value > 0:
                    result.warnings.append(Discrepancy(
                        slide_num=slide_num,
                        slide_title=title,
                        category=ev.category,
                        field_name=ev.field_name,
                        actual_value=ev.parsed_value,
                        expected_value=0,
                        diff=ev.parsed_value,
                        unit=ev.unit,
                        severity="warning",
                        source_filter=source_filter
                    ))
            continue

        # Compute expected values
        expected_media = compute_expected_media_shares(subset)
        expected_funnel = compute_expected_funnel_shares(subset)
        expected_quarters = compute_expected_quarter_budgets(subset)

        # Compare each extracted value
        for ev in extracted_values:
            result.fields_checked += 1

            if ev.category == "media":
                expected = expected_media
            elif ev.category == "funnel":
                expected = expected_funnel
            elif ev.category == "quarter":
                expected = expected_quarters
            else:
                continue

            discrepancy = compare_values(ev, expected, source_filter)

            if discrepancy:
                if discrepancy.severity == "error":
                    result.errors.append(discrepancy)
                else:
                    result.warnings.append(discrepancy)
            else:
                result.passed.append({
                    "slide_num": ev.slide_num,
                    "slide_title": ev.slide_title,
                    "category": ev.category,
                    "field_name": ev.field_name,
                    "actual": ev.parsed_value,
                    "expected": expected.get(ev.field_name),
                    "source_filter": source_filter,
                })

        # Track validated combination
        combo_key = f"{context['market']}|{context['brand']}|{context.get('product', '')}"
        validated_combos.add(combo_key)

    result.sampling_coverage = {
        "unique_combinations": len(validated_combos),
        "combinations": list(validated_combos)[:20],  # First 20 for report
    }

    return result


def generate_report(result: ValidationResult) -> tuple[str, str]:
    """Generate JSON and Markdown reports."""

    # JSON report
    json_report = {
        "timestamp": result.timestamp,
        "pptx_path": result.pptx_path,
        "excel_path": result.excel_path,
        "summary": {
            "total_slides": result.total_slides,
            "slides_with_split_boxes": result.slides_with_split_boxes,
            "fields_checked": result.fields_checked,
            "errors": len(result.errors),
            "warnings": len(result.warnings),
            "passed": len(result.passed),
            "status": "FAILED" if result.errors else "PASSED",
        },
        "errors": [
            {
                "slide": d.slide_num,
                "title": d.slide_title,
                "category": d.category,
                "field": d.field_name,
                "actual": d.actual_value,
                "expected": d.expected_value,
                "diff": d.diff,
                "unit": d.unit,
                "source": d.source_filter,
            }
            for d in result.errors
        ],
        "warnings": [
            {
                "slide": d.slide_num,
                "title": d.slide_title,
                "category": d.category,
                "field": d.field_name,
                "actual": d.actual_value,
                "expected": d.expected_value,
                "diff": d.diff,
                "unit": d.unit,
                "source": d.source_filter,
            }
            for d in result.warnings
        ],
        "sampling_coverage": result.sampling_coverage,
    }

    # Markdown report
    status_emoji = "❌ FAILED" if result.errors else "✅ PASSED"

    md_lines = [
        "# Split Box Validation Report",
        "",
        f"**Generated:** {result.timestamp}",
        f"**PPTX:** `{result.pptx_path}`",
        f"**Excel:** `{result.excel_path}`",
        "",
        "## Summary",
        "",
        "| Metric | Value |",
        "|--------|-------|",
        f"| Total Slides | {result.total_slides} |",
        f"| Slides with Split Boxes | {result.slides_with_split_boxes} |",
        f"| Fields Checked | {result.fields_checked} |",
        f"| Errors | {len(result.errors)} |",
        f"| Warnings | {len(result.warnings)} |",
        f"| Passed | {len(result.passed)} |",
        f"| **Status** | {status_emoji} |",
        "",
    ]

    if result.errors:
        md_lines.extend([
            "## ❌ Errors",
            "",
            "| Slide | Title | Category | Field | Actual | Expected | Diff |",
            "|-------|-------|----------|-------|--------|----------|------|",
        ])
        for e in result.errors:
            actual = f"{e.actual_value:.0f}{e.unit}" if e.actual_value is not None else "N/A"
            expected = f"{e.expected_value:.0f}{e.unit}"
            diff_str = f"{e.diff:+.1f}" if e.category in ("media", "funnel") else f"{e.diff:+,.0f}"
            md_lines.append(
                f"| {e.slide_num} | {e.slide_title[:40]} | {e.category} | {e.field_name} | {actual} | {expected} | {diff_str} |"
            )
        md_lines.append("")

    if result.warnings:
        md_lines.extend([
            "## ⚠️ Warnings",
            "",
            f"Found {len(result.warnings)} warnings (showing first 10):",
            "",
        ])
        for w in result.warnings[:10]:
            actual = f"{w.actual_value:.0f}{w.unit}" if w.actual_value is not None else "N/A"
            expected = f"{w.expected_value:.0f}{w.unit}"
            md_lines.append(f"- **Slide {w.slide_num}** ({w.slide_title[:30]}): {w.category}/{w.field_name} - actual={actual}, expected={expected}")
        md_lines.append("")

    md_lines.extend([
        "## Coverage",
        "",
        f"Validated {result.sampling_coverage.get('unique_combinations', 0)} unique market/brand/product combinations.",
        "",
    ])

    return json.dumps(json_report, indent=2), "\n".join(md_lines)


def main():
    parser = argparse.ArgumentParser(description="Validate split box values in PPTX against source data")
    parser.add_argument("--pptx", type=Path, default=DEFAULT_PPTX, help="Path to PPTX file")
    parser.add_argument("--excel", type=Path, default=DEFAULT_EXCEL, help="Path to Excel source data")
    parser.add_argument("--output-json", type=Path, default=Path("split_box_validation.json"))
    parser.add_argument("--output-md", type=Path, default=Path("split_box_validation.md"))
    args = parser.parse_args()

    print(f"Validating: {args.pptx}")
    print(f"Against source: {args.excel}")

    result = validate_presentation(args.pptx, args.excel)
    json_report, md_report = generate_report(result)

    args.output_json.write_text(json_report)
    args.output_md.write_text(md_report)

    print(f"\n{'='*60}")
    print(f"VALIDATION {'FAILED' if result.errors else 'PASSED'}")
    print(f"{'='*60}")
    print(f"Fields checked: {result.fields_checked}")
    print(f"Errors: {len(result.errors)}")
    print(f"Warnings: {len(result.warnings)}")
    print(f"Passed: {len(result.passed)}")
    print(f"\nReports saved to:")
    print(f"  - {args.output_json}")
    print(f"  - {args.output_md}")

    if result.errors:
        print(f"\n{'='*60}")
        print("ERRORS FOUND:")
        print(f"{'='*60}")
        for e in result.errors[:10]:
            actual = f"{e.actual_value:.0f}{e.unit}" if e.actual_value is not None else "N/A"
            expected = f"{e.expected_value:.0f}{e.unit}"
            print(f"  Slide {e.slide_num}: {e.category}/{e.field_name} - got {actual}, expected {expected}")

    return 1 if result.errors else 0


if __name__ == "__main__":
    exit(main())
