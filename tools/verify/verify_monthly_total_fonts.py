#!/usr/bin/env python
"""Verify MONTHLY TOTAL fonts across multiple slides."""

import sys
from pptx import Presentation

def verify_monthly_total_fonts(pptx_path: str, slide_indices: list):
    """Check MONTHLY TOTAL rows have correct fonts."""
    prs = Presentation(pptx_path)

    print(f"\n{'='*80}")
    print(f"MONTHLY TOTAL Font Verification")
    print(f"{'='*80}\n")

    issues = []
    correct_count = 0

    for slide_idx in slide_indices:
        if slide_idx >= len(prs.slides):
            continue

        slide = prs.slides[slide_idx]

        # Find table
        table = None
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                break

        if not table:
            continue

        # Check MONTHLY TOTAL rows
        for row_idx in range(1, len(table.rows)):
            cell = table.cell(row_idx, 0)
            text = cell.text_frame.text.strip() if cell.text_frame else ""

            if "MONTHLY" in text.upper() and "TOTAL" in text.upper():
                # Check font
                if cell.text_frame.paragraphs:
                    para = cell.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        font_name = run.font.name or "N/A"
                        font_size = f"{run.font.size.pt:.0f}pt" if run.font.size else "N/A"

                        if font_name == "Verdana" and font_size == "6pt":
                            correct_count += 1
                        else:
                            issues.append(f"Slide {slide_idx + 1}, Row {row_idx}: {font_name} {font_size}")

    print(f"Total MONTHLY TOTAL rows checked: {correct_count + len(issues)}")
    print(f"Correct (Verdana 6pt): {correct_count}")
    print(f"Issues found: {len(issues)}")

    if issues:
        print(f"\n{'='*80}")
        print("ISSUES:")
        print(f"{'='*80}")
        for issue in issues[:10]:  # Show first 10
            print(f"  {issue}")
        if len(issues) > 10:
            print(f"  ... and {len(issues) - 10} more")
    else:
        print(f"\n{'='*80}")
        print("ALL MONTHLY TOTAL ROWS CORRECT!")
        print(f"{'='*80}")

if __name__ == "__main__":
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else "output/presentations/run_20251024_172953/deck_test_fonts.pptx"

    # Check slides 2-10
    slide_indices = list(range(1, 11))  # 0-indexed: slides 2-11

    verify_monthly_total_fonts(pptx_path, slide_indices)
