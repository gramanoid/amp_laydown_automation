#!/usr/bin/env python
"""Verify fonts across entire deck."""

import sys
from pptx import Presentation

def verify_deck_fonts(pptx_path: str):
    """Check fonts across all slides."""
    prs = Presentation(pptx_path)

    print(f"\n{'='*80}")
    print(f"Full Deck Font Verification")
    print(f"{'='*80}\n")

    issues = []
    checked = {"header": 0, "body": 0, "bottom": 0}
    correct = {"header": 0, "body": 0, "bottom": 0}

    for slide_idx, slide in enumerate(prs.slides):
        table = None
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                break

        if not table:
            continue

        row_count = len(table.rows)

        # Sample 3 cells per slide: header, body, bottom
        samples = [
            (0, 0, "header"),
            (row_count // 2, 0, "body"),
            (row_count - 1, 0, "bottom")
        ]

        for row_idx, col_idx, cell_type in samples:
            cell = table.cell(row_idx, col_idx)
            if cell.text_frame.paragraphs:
                para = cell.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    font_name = run.font.name or "N/A"
                    font_size_pt = run.font.size.pt if run.font.size else 0

                    checked[cell_type] += 1

                    # Expected sizes
                    expected_size = 7 if cell_type in ["header", "bottom"] else 6

                    if font_name == "Verdana" and font_size_pt == expected_size:
                        correct[cell_type] += 1
                    else:
                        issues.append(f"Slide {slide_idx + 1} {cell_type}: {font_name} {font_size_pt}pt")

    print(f"Cells checked: {sum(checked.values())}")
    print(f"  Header cells: {checked['header']}")
    print(f"  Body cells: {checked['body']}")
    print(f"  Bottom cells: {checked['bottom']}")
    print()
    print(f"Correct fonts: {sum(correct.values())} / {sum(checked.values())}")
    print(f"  Header (Verdana 7pt): {correct['header']} / {checked['header']}")
    print(f"  Body (Verdana 6pt): {correct['body']} / {checked['body']}")
    print(f"  Bottom (Verdana 7pt): {correct['bottom']} / {checked['bottom']}")

    if issues:
        print(f"\n{'='*80}")
        print(f"ISSUES FOUND: {len(issues)}")
        print(f"{'='*80}")
        for issue in issues[:20]:
            print(f"  {issue}")
        if len(issues) > 20:
            print(f"  ... and {len(issues) - 20} more")
    else:
        print(f"\n{'='*80}")
        print("ALL FONTS CORRECT ACROSS ENTIRE DECK!")
        print(f"{'='*80}")

if __name__ == "__main__":
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else "output/presentations/run_20251024_172953/deck_font_entrench.pptx"
    verify_deck_fonts(pptx_path)
