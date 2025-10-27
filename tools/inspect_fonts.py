#!/usr/bin/env python
"""Inspect font properties in a table."""

import sys
from pptx import Presentation
from collections import defaultdict

def inspect_fonts(pptx_path: str, slide_idx: int):
    """Inspect font properties for a specific slide."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]

    # Find table
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if not table:
        print(f"No table found on slide {slide_idx + 1}")
        return

    print(f"\n{'='*80}")
    print(f"Slide {slide_idx + 1} Font Inspection")
    print(f"{'='*80}")
    print(f"Table: {len(table.rows)} rows x {len(table.columns)} cols\n")

    # Sample cells to inspect
    row_count = len(table.rows)
    samples = [
        (0, 0, "Header row, col 1"),
        (0, 5, "Header row, col 6"),
        (1, 0, "Body row 1, col 1"),
        (1, 5, "Body row 1, col 6"),
        (10, 0, "Body row 10, col 1"),
        (10, 10, "Body row 10, col 11"),
        (row_count - 1, 0, "Bottom row, col 1"),
        (row_count - 1, 5, "Bottom row, col 6"),
    ]

    print("Sample Cell Font Inspection:")
    print(f"{'Location':<30s} {'Font':<15s} {'Size':<10s} {'Text Preview':<30s}")
    print("-" * 85)

    for row_idx, col_idx, description in samples:
        if row_idx >= len(table.rows) or col_idx >= len(table.columns):
            continue

        cell = table.cell(row_idx, col_idx)
        text = cell.text_frame.text[:25] if cell.text_frame.text else "(empty)"

        # Get font info from first run
        font_name = "N/A"
        font_size = "N/A"
        if cell.text_frame.paragraphs:
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
                font_name = run.font.name or "N/A"
                font_size = f"{run.font.size.pt:.1f}pt" if run.font.size else "N/A"

        print(f"{description:<30s} {font_name:<15s} {font_size:<10s} {text:<30s}")

    # Aggregate statistics
    print(f"\n{'='*80}")
    print("Font Statistics by Row Type:")
    print(f"{'='*80}\n")

    font_stats = defaultdict(lambda: defaultdict(int))

    for row_idx in range(len(table.rows)):
        row_type = "header" if row_idx == 0 else ("bottom" if row_idx == row_count - 1 else "body")

        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            if cell.text_frame.paragraphs:
                para = cell.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    font_name = run.font.name or "Unknown"
                    font_size = f"{run.font.size.pt:.0f}pt" if run.font.size else "Unknown"
                    key = f"{font_name} {font_size}"
                    font_stats[row_type][key] += 1

    for row_type in ["header", "body", "bottom"]:
        if row_type in font_stats:
            print(f"{row_type.upper()} ROWS:")
            for font_key, count in sorted(font_stats[row_type].items(), key=lambda x: -x[1]):
                print(f"  {font_key}: {count} cells")
            print()

    # Check for issues
    issues = []
    for row_type, expected_size in [("header", 7), ("body", 6), ("bottom", 7)]:
        if row_type in font_stats:
            for font_key, count in font_stats[row_type].items():
                if "Verdana" not in font_key:
                    issues.append(f"{row_type}: {count} cells not using Verdana ({font_key})")
                elif f"{expected_size}pt" not in font_key:
                    issues.append(f"{row_type}: {count} cells not using {expected_size}pt ({font_key})")

    if issues:
        print(f"{'='*80}")
        print("ISSUES FOUND:")
        print(f"{'='*80}")
        for issue in issues:
            print(f"  ❌ {issue}")
    else:
        print(f"{'='*80}")
        print("✅ ALL FONTS CORRECT: Verdana 6pt (body), Verdana 7pt (header/bottom)")
        print(f"{'='*80}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python inspect_fonts.py <pptx_path> <slide_number>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_num = int(sys.argv[2])
    slide_idx = slide_num - 1

    inspect_fonts(pptx_path, slide_idx)
