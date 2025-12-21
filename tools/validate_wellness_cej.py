#!/usr/bin/env python3
"""
Validate CEJ calculations specifically for WELLNESS brands.

WELLNESS = Nutrition and Digestive Health category:
- Centrum
- ENO
- CAC / Cac-1000
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

import pandas as pd
from pptx import Presentation
from amp_automation.data.adapters import FlowplanAdapter
import re

# WELLNESS brands
WELLNESS_BRANDS = ["CENTRUM", "ENO", "CAC", "CAC-1000"]

def parse_slide_context(title: str) -> dict:
    """Parse market, brand, product from slide title."""
    context = {"market": "", "brand": "", "product": "", "is_product_slide": False}
    if not title:
        return context

    clean_title = re.sub(r'\s*\(\d+/\d+\)\s*$', '', title)
    clean_title = re.sub(r'\s*\(Continued\)\s*$', '', clean_title, flags=re.IGNORECASE)

    parts = clean_title.split(" - ")
    if len(parts) >= 2:
        context["market"] = parts[0].strip()
        context["brand"] = parts[1].strip()
        if len(parts) >= 3:
            context["product"] = parts[2].strip()
            context["is_product_slide"] = True

    return context

def extract_cej_values(slide) -> dict:
    """Extract CEJ percentages from slide shapes."""
    values = {}
    shape_map = {
        "FunnelShareAwareness": "awareness",
        "FunnelShareConsideration": "consideration",
        "FunnelSharePurchase": "purchase"
    }

    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        if name in shape_map and hasattr(shape, "text_frame"):
            text = shape.text_frame.text.strip()
            match = re.search(r"(\d+(?:\.\d+)?)\s*%", text)
            if match:
                values[shape_map[name]] = int(float(match.group(1)))

    return values

def compute_expected_cej(df: pd.DataFrame, market: str, brand: str, product: str = None) -> dict:
    """Compute expected CEJ percentages from source data."""
    # Filter data
    mask = (
        (df["Country"].astype(str).str.strip().str.upper() == market.upper()) &
        (df["Brand"].astype(str).str.strip().str.upper() == brand.upper())
    )

    if "Year" in df.columns:
        mask = mask & (df["Year"].astype(str).str.strip() == "2025")

    # Apply product filter if product-level slide
    if product and product.upper() != "PRODUCT SUMMARY":
        product_col = df["Product"].astype(str).str.strip().str.upper()
        product_mask = (product_col == product.upper())

        # Try with brand prefix
        if not product_mask.any():
            prefixed = f"{brand} {product}".upper()
            product_mask = (product_col == prefixed)

        mask = mask & product_mask

    subset = df.loc[mask]

    if subset.empty:
        return {"awareness": 0, "consideration": 0, "purchase": 0}

    funnel_group = subset.groupby("Funnel Stage")["Total Cost"].sum()
    total_cost = subset["Total Cost"].sum()

    if total_cost <= 0:
        return {"awareness": 0, "consideration": 0, "purchase": 0}

    result = {}
    for stage, key in [("Awareness", "awareness"), ("Consideration", "consideration"), ("Purchase", "purchase")]:
        value = float(funnel_group.get(stage, 0.0))
        pct = (value / total_cost) * 100
        result[key] = round(pct)

    return result

def main():
    pptx_path = Path("output/presentations/run_20251217_172239/AMP_Laydowns_171225.pptx")
    excel_path = Path("input/Flowplan_Summaries_MEA_2025_12_17.xlsx")

    # Load data
    adapter = FlowplanAdapter(str(excel_path))
    df = adapter.normalize()

    # Apply media type mapping
    media_type_mapping = {"TV": "Television", "Digital": "Digital", "OOH": "OOH", "Radio": "Other", "Cinema": "Other", "Print": "Other"}
    df["Mapped Media Type"] = df["Media Type"].map(media_type_mapping).fillna("Other")

    # Load presentation
    prs = Presentation(str(pptx_path))

    print("=" * 80)
    print("WELLNESS BRANDS CEJ VALIDATION")
    print("=" * 80)
    print(f"Brands: {', '.join(WELLNESS_BRANDS)}")
    print()

    results = []
    wellness_slides = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Get title
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "name") and "title" in shape.name.lower():
                if hasattr(shape, "text_frame"):
                    title = shape.text_frame.text.strip()
                    break

        if not title:
            continue

        context = parse_slide_context(title)
        brand_upper = context["brand"].upper()

        # Check if this is a WELLNESS brand
        is_wellness = any(wb in brand_upper for wb in WELLNESS_BRANDS)

        if not is_wellness:
            continue

        wellness_slides.append(slide_num)

        # Extract actual CEJ values
        actual = extract_cej_values(slide)

        if not actual:
            print(f"Slide {slide_num}: {title}")
            print(f"  WARNING: No CEJ values found on slide")
            print()
            continue

        # Compute expected values
        expected = compute_expected_cej(
            df,
            context["market"],
            context["brand"],
            context["product"] if context["is_product_slide"] else None
        )

        # Compare
        all_match = True
        mismatches = []

        for key in ["awareness", "consideration", "purchase"]:
            actual_val = actual.get(key, 0)
            expected_val = expected.get(key, 0)

            if abs(actual_val - expected_val) > 1:  # Allow 1% tolerance
                all_match = False
                mismatches.append(f"{key.upper()}: actual={actual_val}%, expected={expected_val}%")

        status = "✅ PASS" if all_match else "❌ FAIL"

        results.append({
            "slide": slide_num,
            "title": title,
            "actual": actual,
            "expected": expected,
            "match": all_match
        })

        print(f"Slide {slide_num}: {title}")
        print(f"  Actual:   AWA={actual.get('awareness', '?')}%, CON={actual.get('consideration', '?')}%, PUR={actual.get('purchase', '?')}%")
        print(f"  Expected: AWA={expected['awareness']}%, CON={expected['consideration']}%, PUR={expected['purchase']}%")
        print(f"  Status: {status}")
        if mismatches:
            for m in mismatches:
                print(f"    ⚠️ {m}")
        print()

    # Summary
    print("=" * 80)
    print("SUMMARY")
    print("=" * 80)

    total = len(results)
    passed = sum(1 for r in results if r["match"])
    failed = total - passed

    print(f"WELLNESS slides found: {len(wellness_slides)}")
    print(f"Slides with CEJ: {total}")
    print(f"Passed: {passed}")
    print(f"Failed: {failed}")
    print()

    if failed == 0:
        print("✅ ALL WELLNESS CEJ CALCULATIONS ARE CORRECT!")
    else:
        print("❌ SOME WELLNESS CEJ CALCULATIONS HAVE MISMATCHES")
        print("\nFailed slides:")
        for r in results:
            if not r["match"]:
                print(f"  - Slide {r['slide']}: {r['title']}")

if __name__ == "__main__":
    main()
