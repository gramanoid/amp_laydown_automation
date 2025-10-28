"""Find which column contains percentages."""
from pptx import Presentation

# Open the latest presentation
prs = Presentation(r"D:\Drive\projects\work\AMP Laydowns Automation\output\presentations\run_20251028_160313\presentations.pptx")

# Get slide 3 (index 2)
slide = prs.slides[2]

# Find the table on the slide
table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break

if table is None:
    print("No table found on slide 3")
    exit(1)

# Check header row (row 0)
print("Header row content:\n")
for col_idx in range(table.columns.__len__()):
    cell = table.cell(0, col_idx)
    cell_text = cell.text.strip()
    print(f"Column {col_idx + 1}: '{cell_text}'")

print("\n\nChecking for percentage columns (looking for '%' in header):")
percent_columns = []
for col_idx in range(table.columns.__len__()):
    cell = table.cell(0, col_idx)
    cell_text = cell.text.strip()
    if '%' in cell_text or 'percent' in cell_text.lower():
        percent_columns.append((col_idx, cell_text))
        print(f"Found percentage column: Column {col_idx + 1} (index {col_idx}): '{cell_text}'")

if percent_columns:
    print(f"\n\nChecking all rows in percentage column (Column {percent_columns[0][0] + 1}):\n")
    col_idx = percent_columns[0][0]

    all_percentage_cells = []

    for row_idx in range(table.rows.__len__()):
        cell = table.cell(row_idx, col_idx)
        row_bold_status = []
        cell_text = ""

        if cell.text_frame:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        cell_text += run.text
                        row_bold_status.append(run.font.bold)
                        bold_status = run.font.bold if run.font.bold is not None else False
                        print(f"Row {row_idx + 1}, Col {col_idx + 1} - Text: '{run.text}', Bold: {bold_status}")

        if row_bold_status:
            # Only consider cells with actual percentage values (not '-')
            if '%' in cell_text:
                row_is_bold = all([b if b is not None else False for b in row_bold_status])
                all_percentage_cells.append(row_is_bold)

                # Highlight specific rows
                if row_idx == 1:  # Row 2
                    print(f"  > Row 2, Col {col_idx + 1}: {'YES' if row_is_bold else 'NO'}\n")
                elif row_idx == 7:  # Row 8
                    print(f"  > Row 8, Col {col_idx + 1}: {'YES' if row_is_bold else 'NO'}\n")

    # Final result
    print("\n" + "="*50)
    if all_percentage_cells and all(all_percentage_cells):
        print("Percentage cells bold: YES")
    else:
        print("Percentage cells bold: NO")
