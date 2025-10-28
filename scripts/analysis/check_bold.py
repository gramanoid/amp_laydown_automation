"""Check if percentage cells in presentation are bold."""
from pptx import Presentation

# Open the latest presentation
prs = Presentation(r"D:\Drive\projects\work\AMP Laydowns Automation\output\presentations\run_20251028_160313\presentations.pptx")

# Get slide 3 (index 2)
slide = prs.slides[2]

# Find the table on the slide
table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break

if table is None:
    print("No table found on slide 3")
    exit(1)

# Check column 17 (index 16)
col_idx = 16

# Check multiple rows to find percentage cells
print("Checking all rows in column 17:\n")
all_cells_bold = []

for row_idx in range(table.rows.__len__()):
    cell = table.cell(row_idx, col_idx)
    row_bold_status = []
    cell_text = ""

    if cell.text_frame:
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():  # Only check non-empty runs
                    cell_text += run.text
                    row_bold_status.append(run.font.bold)
                    print(f"Row {row_idx + 1}, Col 17 - Text: '{run.text}', Bold: {run.font.bold}")

    if row_bold_status:
        row_is_bold = all(row_bold_status)
        all_cells_bold.append(row_is_bold)

        # Highlight specific rows
        if row_idx == 1:  # Row 2
            print(f"  > Row 2, Col 17: {'YES' if row_is_bold else 'NO'}\n")
        elif row_idx == 7:  # Row 8
            print(f"  > Row 8, Col 17: {'YES' if row_is_bold else 'NO'}\n")

# Final result
if all_cells_bold and all(all_cells_bold):
    print("\nPercentage cells bold: YES")
else:
    print("\nPercentage cells bold: NO")
