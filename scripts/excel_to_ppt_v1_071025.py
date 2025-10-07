"""
Excel to PowerPoint Automation Script - GEOGRAPHY VERSION
This version uses the 'Plan - Geography' column instead of '**Country' column
to extract market/country information from the hierarchical geography string.

Changes from FINAL_PRODUCTION.py:
- Uses 'Plan - Geography' column (column K) instead of '**Country'
- Same extraction logic: splits by " | " and takes the last part
- Example: "Global | EMEA | MEA | Pakistan" → "Pakistan"
"""
import os
import logging
import traceback
from datetime import datetime
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL, MSO_FILL_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import argparse
import sys
import ast, pathlib, inspect, textwrap
from pathlib import Path

# Define the constant we need (EXACTLY = 1 is the standard value)
class WD_ROW_HEIGHT_RULE:
    AT_LEAST = 0  # Add this for more flexible row height
    EXACTLY = 1

TABLE_HEIGHT_RULE_AVAILABLE = True

# --- OXML-derived Constants for Cell Styling ---
MARGIN_EMU_LR = 45720  # Approx Pt(3.6), from template analysis for left/right cell margins
ZERO_THRESHOLD = 0.01  # Values below this (absolute) are treated as zero for display/coloring

# --- Constants for Named Shapes (from Template_Refactoring_Guide.md) ---
SHAPE_NAME_TITLE = "TitlePlaceholder"
SHAPE_NAME_TABLE = "MainDataTable"
SHAPE_NAME_COMMENTS_TITLE = "CommentsTitle"
SHAPE_NAME_COMMENTS_BOX = "CommentsBox"
SHAPE_NAME_FUNNEL_CHART = "FunnelChart"
SHAPE_NAME_MEDIA_TYPE_CHART = "MediaTypeChart"
SHAPE_NAME_CAMPAIGN_TYPE_CHART = "CampaignTypeChart"
SHAPE_NAME_TV_LEGEND_COLOR = "TelevisionLegendColor"
SHAPE_NAME_TV_LEGEND_TEXT = "TelevisionLegendText"
SHAPE_NAME_DIGITAL_LEGEND_COLOR = "DigitalLegendColor"
SHAPE_NAME_DIGITAL_LEGEND_TEXT = "DigitalLegendText"
SHAPE_NAME_OOH_LEGEND_COLOR = "OOHLegendColor"
SHAPE_NAME_OOH_LEGEND_TEXT = "OOHLegendText"
SHAPE_NAME_OTHER_LEGEND_COLOR = "OtherLegendColor"
SHAPE_NAME_OTHER_LEGEND_TEXT = "OtherLegendText"

# Must exactly match the placeholder name in Slide Master
TABLE_PLACEHOLDER_NAME = "Table Placeholder 1"

# --- Color Constants (QA checklist verified RGB values) ---
CLR_TELEVISION = RGBColor(113, 212, 141)      # #71D48D - TV swatch/rows
CLR_DIGITAL = RGBColor(253, 242, 183)         # #FDF2B7 - Digital swatch/rows  
CLR_OOH = RGBColor(255, 191, 0)               # #FFBF00 - OOH swatch/rows
CLR_OTHER = RGBColor(176, 211, 255)           # #B0D3FF - Other swatch/rows
CLR_TABLE_GRAY = RGBColor(191, 191, 191)      # Light grey for headers
CLR_HEADER_GREEN = RGBColor(56, 236, 4)       # #38EC04 - Title bar bright green
CLR_COMMENTS_GRAY = RGBColor(242, 242, 242)   # #F2F2F2 - Comments box fill
CLR_SUBTOTAL_GRAY = RGBColor(217, 217, 217)   # #D9D9D9 - Subtotal row
CLR_BLACK = RGBColor(0, 0, 0)
CLR_WHITE = RGBColor(255, 255, 255)
CLR_LIGHT_GRAY_TEXT = RGBColor(191, 191, 191)  # Light grey for empty cell dashes

# --- PIXEL-PERFECT FONT CONSTANTS ---
DEFAULT_FONT_NAME = "Calibri"           # Standard font for all text
FONT_SIZE_HEADER = Pt(7.5)             # Header rows - 7.5pt as requested
FONT_SIZE_BODY = Pt(7)                 # Body rows - 7pt as requested
FONT_SIZE_CHART_TITLE = Pt(8)          # Chart titles - small, non-bold
FONT_SIZE_CHART_LABELS = Pt(6)         # Chart data labels - minimum 6pt

# --- TABLE SPLITTING CONSTANTS ---
# Calculate based on available height: 2.34" total height, with header at 0.139" and body rows at 0.118"
# Available for body rows: 2.34 - 0.139 = 2.201"
# Number of body rows that fit: 2.201 / 0.118 = ~18.6, so 18 body rows + 1 header = 19 total
# Reduced to 17 to ensure proper fit with margins and prevent overflow
MAX_ROWS_PER_SLIDE = 17                # Maximum rows per slide (including header and total)
SPLIT_STRATEGY = "by_campaign"         # Split by complete campaigns
SHOW_CHARTS_ON_SPLITS = "all"          # Show charts on all split slides
SHOW_CARRIED_SUBTOTAL = True           # Show carried forward subtotal on continuation slides
CONTINUATION_INDICATOR = " (Continued)" # Text to append to campaign names on split slides

# --- GEOSPATIAL 2D COORDINATE SYSTEM ---
# Precise positioning for 10" × 5.625" PowerPoint slide (FINAL QA VERIFIED)
# Canvas: 16:9 aspect ratio with origin (0,0) at top-left corner
# All coordinates verified against "Egypt - Centrum" reference slide + Final QA checklist
ELEMENT_COORDINATES = {
    # Title bar - FINAL QA: exactly 2.952 in width (ends at X=3.136)
    'title': {'left': 0.184, 'top': 0.308, 'width': 2.952, 'height': 0.370},
    
    # Main data table - Height limited to stay above COMMENTS/charts (Y=3.300 - 0.15" gap = 3.150 max bottom)
    # With top at 0.812, max height = 3.150 - 0.812 = 2.338"
    'main_table': {'left': 0.184, 'top': 0.812, 'width': 9.299, 'height': 2.338},
    
    # Comments block - coordinates verified
    'comments_title': {'left': 1.097, 'top': 3.697, 'width': 0.640, 'height': 0.151},
    'comments_box': {'left': 0.184, 'top': 3.886, 'width': 2.466, 'height': 1.489},
    
    # Pie charts - FINAL QA: Y=3.300 (moved up 0.20" from 3.500), titles clear table by ≥0.15"
    'chart_1': {'left': 2.650, 'top': 3.300, 'width': 2.466, 'height': 2.000},  # Funnel
    'chart_2': {'left': 4.725, 'top': 3.300, 'width': 2.647, 'height': 2.000},  # Media Type
    'chart_3': {'left': 6.985, 'top': 3.300, 'width': 2.647, 'height': 2.000},  # Campaign Type
    
    # Legend strip - FINAL QA: labels baseline Y=0.416 (nudged up 0.02"), first swatch X=6.645
    'tv_legend_color': {'left': 6.645, 'top': 0.438, 'width': 0.259, 'height': 0.139},
    'tv_legend_text': {'left': 6.841, 'top': 0.416, 'width': 0.612, 'height': 0.219},  # FINAL QA: Y=0.416 exact
    'digital_legend_color': {'left': 7.463, 'top': 0.449, 'width': 0.259, 'height': 0.139},
    'digital_legend_text': {'left': 7.658, 'top': 0.416, 'width': 0.467, 'height': 0.219},  # FINAL QA: Y=0.416 exact
    'ooh_legend_color': {'left': 8.196, 'top': 0.449, 'width': 0.259, 'height': 0.139},
    'ooh_legend_text': {'left': 8.392, 'top': 0.416, 'width': 0.393, 'height': 0.219},  # FINAL QA: Y=0.416 exact
    'other_legend_color': {'left': 8.866, 'top': 0.449, 'width': 0.259, 'height': 0.139},
    'other_legend_text': {'left': 9.061, 'top': 0.416, 'width': 0.439, 'height': 0.219}  # FINAL QA: Y=0.416 exact
}

def get_element_position(element_name):
    """
    Get positioning coordinates for a named element in inches.
    
    Args:
        element_name: Key from ELEMENT_COORDINATES dictionary
        
    Returns:
        dict: Dictionary with left, top, width, height in Inches objects
    """
    if element_name not in ELEMENT_COORDINATES:
        logger.error(f"Element '{element_name}' not found in coordinate system")
        return None
    
    coords = ELEMENT_COORDINATES[element_name]
    return {
        'left': Inches(coords['left']),
        'top': Inches(coords['top']), 
        'width': Inches(coords['width']),
        'height': Inches(coords['height'])
    }

# --- Font Constants (FINAL QA-verified specifications) ---
FONT_SIZE_TITLE = Pt(11)      # Updated: Calibri 11pt bold for title
FONT_SIZE_LEGEND = Pt(6)      # FINAL QA: Calibri 6pt for legend labels (was 7pt)
FONT_SIZE_COMMENTS = Pt(9)    # FINAL QA: Calibri 9pt bold for "COMMENTS" heading

# PRIORITY 3: Enhanced Font System with validation and fallbacks
FALLBACK_FONT_NAME = "Arial"  # Fallback if Calibri unavailable
FONT_VALIDATION_ENABLED = True  # Enable font validation checks

# --- Logging Setup ---
def setup_logging(log_path_base="excel_to_ppt_v2_log"):
    """Sets up logging to file and console."""
    # Create a timestamped log file name
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    log_file = f"{log_path_base}_{timestamp}.log"

    # Get the root logger
    logger = logging.getLogger()
    logger.setLevel(logging.DEBUG)  # Set root logger to DEBUG

    # Clear existing handlers (if any from previous runs in the same session, e.g. in a notebook)
    if logger.hasHandlers():
        logger.handlers.clear()

    # File Handler - logs INFO and up (reduced from DEBUG for smaller log files)
    fh = logging.FileHandler(log_file)
    fh.setLevel(logging.INFO)
    file_formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(filename)s:%(lineno)d - %(message)s')
    fh.setFormatter(file_formatter)
    logger.addHandler(fh)

    # Console Handler - logs INFO and up (or change to DEBUG for more console verbosity)
    ch = logging.StreamHandler()
    ch.setLevel(logging.INFO) # Keep console less verbose, or set to DEBUG if needed
    console_formatter = logging.Formatter('%(levelname)s - %(filename)s:%(lineno)d - %(message)s')
    ch.setFormatter(console_formatter)
    logger.addHandler(ch)

    # Get a specific logger for this module if preferred, to avoid impacting other library loggers
    # For simplicity, we're using the root logger configured above.
    # module_logger = logging.getLogger(__name__)
    # module_logger.setLevel(logging.DEBUG) # Ensure this module's logger is also at DEBUG

    logger.info(f"Logging setup complete. Log file: {log_file}")
    return logger # Though typically we use logging.getLogger(__name__) in modules

# Initialize logger variable but don't set it up yet
logger = None

# --- CORRECTED Excel Column Indices for YOUR ACTUAL FILE ---
# Updated to match YOUR BULK_PLAN_EXPORT_2025_08_25.xlsx structure
COLUMN_GEOGRAPHY = 10              # Plan - Geography (extract country from hierarchy) - NO CHANGE
COLUMN_GLOBAL_MASTERBRAND = 17     # Plan - Brand - NO CHANGE
COLUMN_MEDIA_TYPE = 20             # Media Type - NO CHANGE
COLUMN_CAMPAIGN_NAMES = 83         # **Campaign Name(s) - WAS 85
COLUMN_CAMPAIGN_TYPE = 84          # **Campaign Type - WAS 86
COLUMN_FUNNEL_STAGE = 95           # **Funnel Stage - WAS 97
COLUMN_YEAR = 15                   # Plan - Year - NO CHANGE
COLUMN_FLIGHT_START_DATE = 110     # **Flight Start Date (for month extraction) - NO CHANGE
COLUMN_NET_COST = 71               # *Net Cost (budget data) - WAS 73
COLUMN_NATIONAL_GRP = 55           # GRP - WAS 56

# Additional Lumina columns for enhanced TV data (CORRECTED)
COLUMN_REACH_1PLUS = 104           # Reach 1+ - WAS 106
COLUMN_REACH_3PLUS = 64            # Reach 3+ - NO CHANGE
COLUMN_FREQUENCY = 105             # Frequency - WAS 107

# Legacy column references (kept for backwards compatibility)
COLUMN_COUNTRY = 10                # Now uses COLUMN_GEOGRAPHY
COLUMN_MONTH_START_DATE = 110      # Now uses COLUMN_FLIGHT_START_DATE

# Actual column name for GRP data, used if reading by name
COLUMN_NATIONAL_GRP_NAME = 'National GRP [Current]'

# --- Data Processing Constants (from v1.5) ---
MEDIA_TYPE_MAPPING = {
    'Television': 'TV',
    'TV': 'TV',
    'Digital': 'Digital',
    'OOH': 'OOH',
    'Other': 'Other',
    'Print': 'Other', # Group Print under Other
    'Radio': 'Other', # Group Radio under Other
    'Cinema': 'Other' # Group Cinema under Other
}

MONTH_MAP = {
    1: 'Jan', 2: 'Feb', 3: 'Mar', 4: 'Apr', 5: 'May', 6: 'Jun',
    7: 'Jul', 8: 'Aug', 9: 'Sep', 10: 'Oct', 11: 'Nov', 12: 'Dec'
}

EXPECTED_COLUMNS = [
    COLUMN_GLOBAL_MASTERBRAND,
    COLUMN_MEDIA_TYPE,
    COLUMN_CAMPAIGN_NAMES,
    COLUMN_COUNTRY,
    COLUMN_NET_COST,
    COLUMN_MONTH_START_DATE # Used to derive month
]

# --- Helper Functions ---
def _get_shape_by_name(slide, name):
    """Find a shape on a slide by its name."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    logger.warning(f"Shape with name '{name}' not found on slide.")
    return None

def _copy_text_box(source_shape, target_slide, new_name=None, new_text=None):
    """Copies a source text box shape to a target slide, with optional new name and text."""
    if not source_shape.has_text_frame:
        logger.warning(f"Source shape '{source_shape.name}' does not have a text frame. Cannot copy as text box.")
        return None

    left, top, width, height = source_shape.left, source_shape.top, source_shape.width, source_shape.height
    new_shape = target_slide.shapes.add_textbox(left, top, width, height)
    new_shape.name = new_name if new_name else source_shape.name

    # Copy text frame properties
    source_tf = source_shape.text_frame
    new_tf = new_shape.text_frame

    new_tf.word_wrap = source_tf.word_wrap
    # MSO_AUTO_SIZE enum: NONE = 0, SHAPE_TO_FIT_TEXT = 1, TEXT_TO_FIT_SHAPE = 2
    # Ensure auto_size is handled correctly; direct copy might be problematic if enum values differ or not available.
    # For simplicity, let's try direct copy if it's simple. Otherwise, might need specific handling.
    try:
        new_tf.auto_size = source_tf.auto_size
    except AttributeError:
        logger.debug(f"Could not directly copy auto_size for shape '{new_shape.name}'. Defaulting.")
        # Default or no action if direct copy fails

    if new_text is not None:
        new_tf.text = new_text
    else:
        new_tf.text = source_tf.text

    # Copy paragraph formatting for the first paragraph (if text exists)
    if source_tf.paragraphs and new_tf.paragraphs:
        source_para = source_tf.paragraphs[0]
        new_para = new_tf.paragraphs[0]
        new_para.alignment = source_para.alignment
        
        if source_para.runs and new_para.runs:
            source_run = source_para.runs[0]
            new_run = new_para.runs[0]
            new_run.font.name = source_run.font.name
            new_run.font.size = source_run.font.size
            if source_run.font.color.type == MSO_COLOR_TYPE.RGB:
                 if hasattr(source_run.font.color, 'rgb'):
                    new_run.font.color.rgb = source_run.font.color.rgb
            new_run.font.bold = source_run.font.bold
            new_run.font.italic = source_run.font.italic
        elif not new_para.runs and new_text == "": # Handle empty text box case for font styling
            # If new_text is empty, add a run to apply font style from source
            if source_para.runs:
                source_run = source_para.runs[0]
                new_run = new_para.add_run()
                new_run.font.name = source_run.font.name
                new_run.font.size = source_run.font.size
                if source_run.font.color.type == MSO_COLOR_TYPE.RGB:
                     if hasattr(source_run.font.color, 'rgb'):
                        new_run.font.color.rgb = source_run.font.color.rgb
                new_run.font.bold = source_run.font.bold
                new_run.font.italic = source_run.font.italic
                new_para.text = "" # Clear the run's default text if any

    logger.debug(f"Copied text box '{source_shape.name}' to '{new_shape.name}' on target slide.")
    return new_shape

def _copy_shape(source_shape, target_slide, new_name=None):
    """Copies a basic source auto shape to a target slide, with optional new name."""
    if source_shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        logger.warning(f"Source shape '{source_shape.name}' is not an AUTO_SHAPE (type: {source_shape.shape_type}). Skipping copy.")
        return None

    left, top, width, height = source_shape.left, source_shape.top, source_shape.width, source_shape.height
    # Add the shape with the same auto_shape_type
    new_shape = target_slide.shapes.add_shape(
        source_shape.auto_shape_type, left, top, width, height
    )
    new_shape.name = new_name if new_name else source_shape.name

    # Copy fill properties
    source_fill = source_shape.fill
    new_fill = new_shape.fill
    # Assuming solid fill for legend color boxes
    new_fill.solid()
    if hasattr(source_fill.fore_color, 'rgb'):
        new_fill.fore_color.rgb = source_fill.fore_color.rgb
    else:
        logger.debug(f"Source shape '{source_shape.name}' fill fore_color has no rgb. Defaulting fill.")

    # Copy line properties
    source_line = source_shape.line
    new_line = new_shape.line
    if hasattr(source_line.color, 'rgb') and source_line.color.type != 0: # Check if color is set and not NO_FILL implicitly by type
        new_line.color.rgb = source_line.color.rgb
    elif source_line.fill.type == 0: # NO_FILL (0)
        new_line.fill.background() # Or new_line.fill.solid() with transparency, or just leave as default if that's no line
        logger.debug("Source shape '{source_shape.name}' line has NO_FILL. Applying background fill to line of new shape.")
    else:
        logger.debug(f"Source shape '{source_shape.name}' line color has no rgb or complex fill. Defaulting line color.")
    
    new_line.width = source_line.width # Width is in EMUs, direct copy is fine

    logger.debug(f"Copied shape '{source_shape.name}' to '{new_shape.name}' on target slide.")
    return new_shape

def normalize_media_type(media_type):
    """
    Standardize media type strings for consistent comparison.
    Converts various TV/Television representations to a standard format.
    """
    if not media_type:
        return ""
    
    media_type_str = str(media_type).strip().upper()
    if media_type_str in ['TV', 'TELEVISION']:
        return "Television"
    elif media_type_str == 'DIGITAL':
        return "Digital"
    else:
        return media_type_str

def is_empty_formatted_value(formatted_value):
    """Check if a formatted value represents empty/zero data that should not be displayed"""
    if not formatted_value:
        return True
    
    # Standard empty indicators
    if formatted_value in ["-", "", "0.0%"]:
        return True
    
    # Budget-specific empty indicators
    if formatted_value in ["£0K", "£0", "£0.00K"]:
        return True
        
    return False

def format_number(value, is_budget=False, is_percentage=False, is_grp=False, is_monthly_column=False):
    """Formats numbers for display in tables, with improved accuracy for budgets and enhanced zero handling."""
    if pd.isna(value):
        if is_percentage:
            return "0.0%"
        return ""  # Return empty string for NaN values to trigger "-" display
    
    numeric_value = 0
    try:
        numeric_value = float(value)
    except ValueError:
        if is_percentage:
            return "0.0%"
        return ""  # Return empty string for invalid values to trigger "-" display

    if numeric_value == 0 or abs(numeric_value) < ZERO_THRESHOLD:  # Enhanced zero detection including near-zero values
        if is_percentage:
            return "0.0%"
        # For budget and GRP values, return empty string to trigger "-" display in cell styling
        if is_budget or is_grp:
            return ""  # This will trigger "-" display in _apply_table_cell_styling
        return ""

    # Handle percentages first
    if is_percentage:
        return f"{numeric_value:.1f}%"

    # Handle GRPs with K suffix for thousands
    if is_grp:
        abs_value = abs(numeric_value)
        if abs_value >= 1_000:
            formatted_val = numeric_value / 1_000.0
            # Show one decimal place for accuracy (e.g., 23.5K)
            return f"{formatted_val:.1f}K"
        else:
            # For values less than 1000, show as integer
            return f"{int(numeric_value)}"

    # IMPROVED BUDGET FORMATTING: More accurate representation
    if is_budget:
        abs_value = abs(numeric_value)
        
        # For values >= 1M, use millions
        if abs_value >= 1_000_000:
            formatted_val = numeric_value / 1_000_000.0
            # For monthly columns, always use whole numbers (no decimals)
            if is_monthly_column:
                return f"£{formatted_val:.0f}M"
            else:
                # Use 1 decimal place for millions to preserve accuracy in main Budget column
                if formatted_val == int(formatted_val):
                    return f"£{formatted_val:.0f}M"
                else:
                    return f"£{formatted_val:.1f}M"
        
        # For values >= 1K, use thousands
        elif abs_value >= 1_000:
            formatted_val = numeric_value / 1_000.0
            # For monthly columns, always use whole numbers (no decimals)
            if is_monthly_column:
                return f"£{formatted_val:.0f}K"
            else:
                # Use 1 decimal place for thousands when needed for accuracy in main Budget column
                if formatted_val == int(formatted_val):
                    return f"£{formatted_val:.0f}K"
                else:
                    return f"£{formatted_val:.1f}K"
        
        # For values < 1K, show in pounds with K suffix for consistency
        else:
            # CRITICAL FIX: For monthly columns with values < 1K, show actual amount in pounds
            # This prevents £500 from displaying as "£0K" which appears empty
            if is_monthly_column and abs_value > 0:
                return f"£{int(numeric_value)}"
            
            formatted_val = numeric_value / 1_000.0
            # For monthly columns, always round to whole K units
            if is_monthly_column:
                return f"£{formatted_val:.0f}K"
            else:
                return f"£{formatted_val:.2f}K"
    
    # Non-budget formatting (fallback)
    abs_value = abs(numeric_value)
    if abs_value >= 1_000_000:
        suffix = 'M'
        divisor = 1_000_000.0
    elif abs_value >= 1_000:
        suffix = 'K'
        divisor = 1_000.0
    else:
        suffix = 'K'
        divisor = 1_000.0

    formatted_val = numeric_value / divisor
    formatted_str = f"{formatted_val:.1f}{suffix}"
    return formatted_str

def get_month_specific_tv_metrics(raw_excel_path, country, brand, campaign, year, month):
    """
    Get month-specific TV metrics (GRP, Frequency, Reach) for a specific campaign and month.
    This function aggregates raw data on-the-fly for the requested month.
    """
    
    # Cache the raw data to avoid reloading for each call
    if not hasattr(get_month_specific_tv_metrics, '_cached_data'):
        df = pd.read_excel(raw_excel_path, header=0)
        
        # Check if Month column exists, if not extract from Flight Start Date
        if 'Month' not in df.columns or df['Month'].isna().all():
            if '**Flight Start Date' in df.columns:
                # Convert to datetime and extract month name
                df['**Flight Start Date'] = pd.to_datetime(df['**Flight Start Date'])
                df['Month'] = df['**Flight Start Date'].dt.strftime('%b')  # Short month names (Jan, Feb, etc.)
        
        get_month_specific_tv_metrics._cached_data = df
    else:
        df = get_month_specific_tv_metrics._cached_data
    
    # Extract country from hierarchical format
    def extract_country(country_str):
        if pd.isna(country_str):
            return None
        parts = str(country_str).split(' | ')
        return parts[-1].strip()
    
    # Clean brand name
    def clean_brand(brand_str):
        if pd.isna(brand_str):
            return ""
        brand = str(brand_str)
        if " | " in brand:
            return brand.split(" | ")[-1].strip()
        return brand.strip()
    
    # Month name mapping (handle both upper and proper case)
    month_mapping = {
        'Jan': 'Jan', 'Feb': 'Feb', 'Mar': 'Mar', 'Apr': 'Apr', 
        'May': 'May', 'Jun': 'Jun', 'Jul': 'Jul', 'Aug': 'Aug',
        'Sep': 'Sep', 'Sept': 'Sep', 'Oct': 'Oct', 'Nov': 'Nov', 'Dec': 'Dec',
        'JAN': 'Jan', 'FEB': 'Feb', 'MAR': 'Mar', 'APR': 'Apr',
        'MAY': 'May', 'JUN': 'Jun', 'JUL': 'Jul', 'AUG': 'Aug',
        'SEP': 'Sep', 'SEPT': 'Sep', 'OCT': 'Oct', 'NOV': 'Nov', 'DEC': 'Dec'
    }
    
    # Normalize month name
    month = month_mapping.get(month, month)
    
    # Filter data for this specific country/brand/campaign/year/month combination
    filtered_data = df[
        (df['Plan - Geography'].apply(extract_country) == country) &
        (df['Plan - Brand'].apply(clean_brand) == brand) &
        (df['**Campaign Name(s)'] == campaign) &
        (df['Plan - Year'] == year) &
        (df['Month'] == month) &
        (df['Media Type'] == 'Television')
    ]
    
    # Apply GNE Television filter - exclude "Pan Asian TV" from Flight Comments
    # Check if this is a GNE campaign (GNE in the raw geography string)
    gne_mask = filtered_data['Plan - Geography'].astype(str).str.contains('GNE', na=False)
    pan_asian_mask = filtered_data['Flight Comments'].astype(str).str.contains('Pan Asian TV', na=False)
    
    if len(filtered_data[gne_mask & pan_asian_mask]) > 0:
        filtered_data = filtered_data[~(gne_mask & pan_asian_mask)]
    
    if len(filtered_data) == 0:
        return {
            'grp_sum': 0,
            'frequency_avg': np.nan,
            'reach1_avg': np.nan,
            'reach3_avg': np.nan
        }
    
    # Aggregate the data for this month
    available_cols = filtered_data.columns.tolist()
    tv_metric_columns = ['National GRP', 'Frequency', 'Reach 1+', 'Reach 3+']
    missing_cols = [col for col in tv_metric_columns if col not in available_cols]
    if missing_cols:
        logger.warning(f"Missing TV metric columns in month-specific function: {missing_cols}")
    
    grp_sum = filtered_data['National GRP'].dropna().sum() if 'National GRP' in available_cols else 0
    
    freq_values = filtered_data['Frequency'].dropna() if 'Frequency' in available_cols else pd.Series(dtype=float)
    frequency_avg = freq_values.mean() if len(freq_values) > 0 else np.nan
    
    reach1_values = filtered_data['Reach 1+'].dropna() if 'Reach 1+' in available_cols else pd.Series(dtype=float)
    reach1_avg = reach1_values.mean() if len(reach1_values) > 0 else np.nan
    
    reach3_values = filtered_data['Reach 3+'].dropna() if 'Reach 3+' in available_cols else pd.Series(dtype=float)
    reach3_avg = reach3_values.mean() if len(reach3_values) > 0 else np.nan
    
    return {
        'grp_sum': grp_sum,
        'frequency_avg': frequency_avg,
        'reach1_avg': reach1_avg,
        'reach3_avg': reach3_avg
    }

def load_and_prepare_data(excel_path):
    """Load raw Lumina data and prepare it in the format expected by the presentation generator."""
    logger.info(f"Loading raw Lumina data from: {excel_path}")
    try:
        # Load raw data
        raw_df = pd.read_excel(excel_path, header=0)
        logger.info(f"Loaded {len(raw_df)} rows from raw data")
        
        # Check if Month column exists, if not extract from Plan Start Date
        if 'Month' not in raw_df.columns or raw_df['Month'].isna().all():
            if '**Flight Start Date' in raw_df.columns:
                logger.info("Month column missing or empty. Extracting month from Flight Start Date...")
                # Convert to datetime and extract month name
                raw_df['**Flight Start Date'] = pd.to_datetime(raw_df['**Flight Start Date'])
                raw_df['Month'] = raw_df['**Flight Start Date'].dt.strftime('%b')  # Short month names (Jan, Feb, etc.)
                logger.info(f"Successfully extracted months from {len(raw_df[raw_df['Month'].notna()])} rows")
            else:
                logger.error("No Month column or Flight Start Date column found")
                return None
        
        # Media Type column should already exist - no need to create it
        if 'Media Type' in raw_df.columns:
            logger.info(f"Media Type column found with values: {dict(raw_df['Media Type'].value_counts())}")
        else:
            logger.error("Media Type column not found in data")
            return None
        
        # Extract country from hierarchical format
        def extract_country(country_str):
            if pd.isna(country_str):
                return None
            parts = str(country_str).split(' | ')
            return parts[-1].strip()  # Take last part
        
        # Clean brand name (remove "Haleon | " prefix)
        def clean_brand(brand_str):
            if pd.isna(brand_str):
                return ""
            brand = str(brand_str)
            if " | " in brand:
                return brand.split(" | ")[-1].strip()
            return brand.strip()
        
        # Month name mapping
        month_mapping = {
            'Jan': 'Jan', 'Feb': 'Feb', 'Mar': 'Mar', 'Apr': 'Apr', 
            'May': 'May', 'Jun': 'Jun', 'Jul': 'Jul', 'Aug': 'Aug',
            'Sep': 'Sep', 'Sept': 'Sep', 'Oct': 'Oct', 'Nov': 'Nov', 'Dec': 'Dec'
        }
        
        # Process data
        processed_data = []
        
        # Filter out GNE Television rows with "Pan Asian TV" in Flight Comments
        def should_exclude_row(row):
            geography_raw = str(row['Plan - Geography'])
            media_type = row['Media Type']
            flight_comments = str(row.get('Flight Comments', '')).strip()
            
            # Exclude GNE Television campaigns with "Pan Asian TV" in Flight Comments
            # Check for 'GNE' in the raw geography string (e.g., "Global | EMEA | MEA | GNE")
            if ('GNE' in geography_raw and 
                media_type == 'Television' and 
                'Pan Asian TV' in flight_comments):
                return True
            return False
        
        # Apply filter
        logger.info(f"Applying GNE Television filter (excluding 'Pan Asian TV' from Flight Comments)...")
        initial_count = len(raw_df)
        raw_df = raw_df[~raw_df.apply(should_exclude_row, axis=1)]
        filtered_count = len(raw_df)
        logger.info(f"Filtered out {initial_count - filtered_count} rows. Remaining: {filtered_count}")
        
        # Group by Geography + Brand + Campaign + Year + Month for initial aggregation
        group_cols = ['Plan - Geography', 'Plan - Brand', '**Campaign Name(s)', 'Plan - Year', 'Month', 'Media Type']
        
        logger.info("Aggregating data by month...")
        
        for name, group in raw_df.groupby(group_cols):
            geography_raw, brand_raw, campaign, year, month_raw, media_type = name
            
            # Clean and validate data
            country = extract_country(geography_raw)
            brand = clean_brand(brand_raw)
            
            # Map month name
            month = month_mapping.get(month_raw, month_raw)
            
            # Skip if essential data is missing
            if not country or not brand or not campaign:
                continue
            
            # Aggregate financial data (always present)
            total_cost = group['*Cost to Client'].sum()
            
            # Aggregate TV metrics (only for Television)
            grp_sum = np.nan
            freq_avg = np.nan
            reach1_avg = np.nan
            reach3_avg = np.nan
            
            if media_type == 'Television':
                # Debug: Check if TV metric columns exist
                available_cols = group.columns.tolist()
                tv_metric_columns = ['National GRP', 'Frequency', 'Reach 1+', 'Reach 3+']
                missing_cols = [col for col in tv_metric_columns if col not in available_cols]
                if missing_cols:
                    logger.warning(f"Missing TV metric columns for {country}-{brand}-{campaign}: {missing_cols}")
                    logger.debug(f"Available columns: {[col for col in available_cols if any(keyword in col.lower() for keyword in ['grp', 'reach', 'freq'])]}")
                
                # Sum GRPs
                if 'National GRP' in available_cols:
                    grp_values = group['National GRP'].dropna()
                    if len(grp_values) > 0:
                        grp_sum = grp_values.sum()
                        logger.debug(f"Found {len(grp_values)} GRP values for {country}-{brand}-{campaign}, sum: {grp_sum}")
                
                # Average Frequency and Reach
                if 'Frequency' in available_cols:
                    freq_values = group['Frequency'].dropna()
                    if len(freq_values) > 0:
                        freq_avg = freq_values.mean()
                        logger.debug(f"Found {len(freq_values)} Frequency values for {country}-{brand}-{campaign}, avg: {freq_avg}")
                
                if 'Reach 1+' in available_cols:
                    reach1_values = group['Reach 1+'].dropna()
                    if len(reach1_values) > 0:
                        reach1_avg = reach1_values.mean()
                        logger.debug(f"Found {len(reach1_values)} Reach 1+ values for {country}-{brand}-{campaign}, avg: {reach1_avg}")
                
                if 'Reach 3+' in available_cols:
                    reach3_values = group['Reach 3+'].dropna()
                    if len(reach3_values) > 0:
                        reach3_avg = reach3_values.mean()
                        logger.debug(f"Found {len(reach3_values)} Reach 3+ values for {country}-{brand}-{campaign}, avg: {reach3_avg}")
            
            # Get other fields (take first non-null value)
            campaign_type = group['**Campaign Type'].dropna().iloc[0] if not group['**Campaign Type'].dropna().empty else ""
            funnel_stage = group['**Funnel Stage'].dropna().iloc[0] if not group['**Funnel Stage'].dropna().empty else ""
            
            # Store the aggregated row data
            row_data = {
                'Country': country,
                'Brand': brand,
                'Media Type': media_type,
                'Campaign Name': campaign,
                'Campaign Type': campaign_type,
                'Funnel Stage': funnel_stage,
                'Year': year,
                'Month': month,
                'Total Cost': total_cost if pd.notna(total_cost) else 0,
                'GRP': grp_sum,
                'Frequency': freq_avg,
                'Reach 1+': reach1_avg,
                'Reach 3+': reach3_avg,
            }
            
            processed_data.append(row_data)
        
        # Convert to DataFrame
        agg_df = pd.DataFrame(processed_data)
        logger.info(f"Created {len(agg_df)} month-level aggregated rows")
        
        if len(agg_df) == 0:
            logger.error("No data found after processing")
            return None
        
        # Pivot to wide format (months as columns)
        logger.info("Pivoting to wide format...")
        
        result_rows = []
        
        # Group by everything except Month to create final campaign rows
        final_group_cols = ['Country', 'Brand', 'Media Type', 'Campaign Name', 'Campaign Type', 'Funnel Stage', 'Year']
        
        for name, group in agg_df.groupby(final_group_cols):
            country, brand, media_type, campaign, campaign_type, funnel_stage, year = name
            
            # Initialize row with campaign info
            row = {
                'Country': country,
                'Brand': brand,
                'Media Type': media_type,
                'Campaign Name': campaign,
                'Campaign Type': campaign_type,
                'Funnel Stage': funnel_stage,
                'Year': year
            }
            
            # Initialize all months to 0
            months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
            for month in months:
                row[month] = 0
            
            # Fill in actual monthly data
            total_cost = 0
            total_grp = 0
            freq_values = []
            reach1_values = []
            reach3_values = []
            
            for _, month_row in group.iterrows():
                month = month_row['Month']
                if month in months:
                    # Set monthly budget
                    row[month] = month_row['Total Cost']
                    total_cost += month_row['Total Cost']
                    
                    # Accumulate TV metrics (will be averaged later)
                    if pd.notna(month_row['GRP']):
                        total_grp += month_row['GRP']
                    
                    if pd.notna(month_row['Frequency']):
                        freq_values.append(month_row['Frequency'])
                    if pd.notna(month_row['Reach 1+']):
                        reach1_values.append(month_row['Reach 1+'])
                    if pd.notna(month_row['Reach 3+']):
                        reach3_values.append(month_row['Reach 3+'])
            
            # Set totals and campaign-level averages
            row['Total Cost'] = total_cost
            row['GRP'] = total_grp if total_grp > 0 else np.nan
            row['Frequency'] = np.mean(freq_values) if freq_values else np.nan
            row['Reach 1+'] = np.mean(reach1_values) if reach1_values else np.nan
            row['Reach 3+'] = np.mean(reach3_values) if reach3_values else np.nan
            row['Flight Comments'] = ""
            
            result_rows.append(row)
        
        df = pd.DataFrame(result_rows)
        logger.info(f"Final result: {len(df)} campaigns")
        
        # Ensure columns are in expected order
        expected_columns = [
            'Country', 'Brand', 'Media Type', 'Campaign Name', 'Campaign Type', 'Funnel Stage', 'Year',
            'Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec',
            'Total Cost', 'GRP', 'Frequency', 'Reach 1+', 'Reach 3+', 'Flight Comments'
        ]
        
        # Add missing columns with default values
        for col in expected_columns:
            if col not in df.columns:
                if col in ['GRP', 'Frequency', 'Reach 1+', 'Reach 3+']:
                    df[col] = np.nan
                elif col in ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Total Cost']:
                    df[col] = 0
                else:
                    df[col] = ""
        
        # Reorder columns
        df = df[expected_columns]
        
        # Add Mapped Media Type column for compatibility
        df['Mapped Media Type'] = df['Media Type']
        
        logger.info(f"Data loaded successfully. Shape: {df.shape}")
        logger.info(f"Media types: {dict(df['Media Type'].value_counts())}")
        logger.info(f"TV campaigns with metrics: {len(df[(df['Media Type'] == 'Television') & df['GRP'].notna()])}")
        
        return df
        
    except FileNotFoundError:
        logger.error(f"Error: Excel file not found at {excel_path}")
        return None
    except Exception as e:
        logger.error(f"Error loading data: {str(e)}")
        return None
    except Exception as e:
        logger.error(f"Error loading or preparing data from {excel_path}: {e}")
        logger.error(traceback.format_exc())
        return None

def _prepare_main_table_data_detailed(df, region, masterbrand, year=None, excel_path=None):
    """
    Prepare table data for a specific region/masterbrand combination.
    Adapted from v1.5 create_campaign_table function for v2.0 template approach.
    
    Args:
        df: DataFrame with loaded Excel data
        region: Region name to filter by
        masterbrand: Masterbrand name to filter by
        
    Returns:
        tuple: (table_data, cell_metadata) where:
            - table_data: List of lists representing table rows
            - cell_metadata: Dictionary with metadata for each cell for conditional formatting
        or (None, None) if no data found
    """
    try:
        year_text = f" - {year}" if year is not None else ""
        logger.info(f"Preparing table data for {region} - {masterbrand}{year_text}")
        
        # **ENHANCED DEBUGGING**: Log data filtering process
        logger.debug(f"Total rows in DataFrame before filtering: {len(df)}")
        logger.debug(f"Unique countries in data: {sorted(df['Country'].dropna().unique())}")
        logger.debug(f"Unique masterbrands in data: {sorted(df['Brand'].dropna().unique())}")
        if year is not None:
            logger.debug(f"Unique years in data: {sorted(df['Year'].dropna().unique())}")
        
        # Month mapping for quarterly calculations
        month_mapping = {
            'January': 1, 'Jan': 1, 'JAN': 1, '1': 1, 1: 1,
            'February': 2, 'Feb': 2, 'FEB': 2, '2': 2, 2: 2,
            'March': 3, 'Mar': 3, 'MAR': 3, '3': 3, 3: 3,
            'April': 4, 'Apr': 4, 'APR': 4, '4': 4, 4: 4,
            'May': 5, 'MAY': 5, '5': 5, 5: 5,
            'June': 6, 'Jun': 6, 'JUN': 6, '6': 6, 6: 6,
            'July': 7, 'Jul': 7, 'JUL': 7, '7': 7, 7: 7,
            'August': 8, 'Aug': 8, 'AUG': 8, '8': 8, 8: 8,
            'September': 9, 'Sep': 9, 'SEP': 9, '9': 9, 9: 9,
            'October': 10, 'Oct': 10, 'OCT': 10, '10': 10, 10: 10,
            'November': 11, 'Nov': 11, 'NOV': 11, '11': 11, 11: 11,
            'December': 12, 'Dec': 12, 'DEC': 12, '12': 12
        }
        
        # Constants for table structure
        MONTHS = ['JAN', 'FEB', 'MAR', 'Q1', 'APR', 'MAY', 'JUN', 'Q2', 'JUL', 'AUG', 'SEP', 'Q3', 'OCT', 'NOV', 'DEC', 'Q4']
        GRP_DATA_IDENTIFIER = "GRPs"  # Identifier for GRP sub-rows
        
        # Filter data for the specific region, masterbrand, and year
        # Use the processed column names from load_and_prepare_data
        filter_conditions = [
            (df['Country'].astype(str).str.strip() == region),
            (df['Brand'].astype(str).str.strip() == masterbrand)
        ]
        
        if year is not None:
            filter_conditions.append(df['Year'].astype(str).str.strip() == str(year))
        
        filtered_df = df[
            filter_conditions[0] & filter_conditions[1] & (filter_conditions[2] if len(filter_conditions) > 2 else True)
        ].copy()
        
        # **ENHANCED DEBUGGING**: Log filtering results
        logger.debug(f"Rows after filtering for {region} - {masterbrand}: {len(filtered_df)}")
        
        if filtered_df.empty:
            logger.warning(f"No data found for {region} - {masterbrand}")
            # **ENHANCED DEBUGGING**: Show available combinations for troubleshooting
            available_combinations = df.groupby(['Country', 'Brand']).size().reset_index(name='count')
            logger.debug(f"Available country/masterbrand combinations:")
            for _, combo in available_combinations.iterrows():
                logger.debug(f"  {combo['Country']} - {combo['Brand']} ({combo['count']} rows)")
            return None, None
            
        logger.info(f"Found {len(filtered_df)} rows for {region} - {masterbrand}")
        
        # Get unique campaigns and sort them
        campaign_column = 'Campaign Name'
        campaigns = sorted(filtered_df[campaign_column].dropna().unique())
        
        # **ENHANCED DEBUGGING**: Log campaign details
        logger.info(f"Found {len(campaigns)} unique campaigns for {region} - {masterbrand}:")
        for i, campaign in enumerate(campaigns):
            campaign_rows = len(filtered_df[filtered_df[campaign_column] == campaign])
            logger.info(f"  {i+1}. '{campaign}' ({campaign_rows} rows)")
        
        # **ENHANCED DEBUGGING**: Check for potential data quality issues
        if any(campaign in str(campaigns) for campaign in ['Gender', 'You Did It']):
            logger.warning(f"DATA QUALITY ALERT: Suspicious campaigns found for {region} - {masterbrand}: {campaigns}")
            logger.warning("This may indicate incorrect data filtering or source data issues")
        
        # Initialize table data with headers
        table_data = [['CAMPAIGN NAME', 'BUDGET', 'GRPs', 'REACH', 'OTS', '%', 'MEDIA TYPE'] + MONTHS]
        monthly_totals = [0.0] * len(MONTHS)
        grp_monthly_totals = [0.0] * len(MONTHS)
        
        # Initialize cell metadata for conditional formatting
        # Structure: {(row_idx, col_idx): {'has_data': bool, 'media_type': str, 'value': float}}
        cell_metadata = {}
        
        # Calculate total budget for percentage calculations
        total_budget_for_percentage = filtered_df['Total Cost'].sum()
        grand_total_grp = 0  # Sum of campaign total GRPs for footer
        
        logger.debug(f"Total budget for percentage calculation: £{total_budget_for_percentage:,.2f}")
        
        for campaign_name in campaigns:
            if pd.isna(campaign_name) or str(campaign_name).strip() == '':
                continue
            
            campaign_df = filtered_df[filtered_df[campaign_column] == campaign_name]
            logger.debug(f"Processing campaign: '{campaign_name}' ({len(campaign_df)} rows)")
            
            campaign_total_budget = campaign_df['Total Cost'].sum()
            campaign_percentage = (campaign_total_budget / total_budget_for_percentage * 100) if total_budget_for_percentage > 0 else 0
            
            logger.debug(f"  Campaign budget: £{campaign_total_budget:,.2f} ({campaign_percentage:.1f}%)")
            
            # Calculate total TV GRP for this campaign
            campaign_total_tv_grp_value = 0
            if 'GRP' in campaign_df.columns:
                campaign_tv_df = campaign_df[campaign_df['Media Type'].astype(str).str.upper().isin(['TV', 'TELEVISION'])]
                if not campaign_tv_df.empty:
                    grp_sum = campaign_tv_df['GRP'].sum()
                    if not pd.isna(grp_sum):
                        campaign_total_tv_grp_value = grp_sum
            
            logger.debug(f"  Campaign TV GRPs: {campaign_total_tv_grp_value}")
            
            # Get unique media types for this campaign, prioritizing TV first
            media_types_for_campaign = sorted(campaign_df['Media Type'].fillna('').astype(str).unique())
            media_types_for_campaign.sort(key=lambda mt: 0 if normalize_media_type(str(mt)) == 'Television' else 1)
            
            logger.debug(f"  Media types: {media_types_for_campaign}")
            
            first_media_type_for_campaign = True
            
            for media_type_original in media_types_for_campaign:
                if pd.isna(media_type_original) or str(media_type_original).strip() == '':
                    continue
                
                current_media_type_normalized = normalize_media_type(str(media_type_original))
                media_df = campaign_df[campaign_df['Media Type'] == media_type_original]
                
                monthly_budget_data_formatted = ['-'] * len(MONTHS)
                monthly_grp_data_formatted = ['-'] * len(MONTHS)
                
                is_tv_current_media = current_media_type_normalized == 'Television'
                
                # For wide-format data, monthly budgets are already in columns
                # No need to group by Month column
                
                # Populate monthly budget data with quarterly totals
                q_budget_total = 0
                current_row_idx = len(table_data)  # Track current row index for metadata
                
                for i, month_name in enumerate(MONTHS):
                    month_col_idx = i + 7  # Month columns start after first 7 columns
                    
                    if month_name.startswith('Q'):  # Q1, Q2, Q3, Q4
                        formatted_budget = format_number(q_budget_total, is_budget=True, is_monthly_column=True)
                        monthly_budget_data_formatted[i] = formatted_budget
                        monthly_totals[i] += q_budget_total
                        # Add metadata for quarter columns - FIXED: align with display logic
                        cell_metadata[(current_row_idx, month_col_idx)] = {
                            'has_data': not is_empty_formatted_value(formatted_budget),
                            'media_type': current_media_type_normalized,
                            'value': q_budget_total
                        }
                        q_budget_total = 0
                    else:
                        # For wide-format data, get budget directly from month column
                        # Convert month name to proper case (Jan, Feb, Mar, etc.)
                        month_col_name = month_name.capitalize() if len(month_name) == 3 else month_name
                        if month_col_name in media_df.columns:
                            budget_val = media_df[month_col_name].sum()
                            formatted_budget = format_number(budget_val, is_budget=True, is_monthly_column=True)
                            monthly_budget_data_formatted[i] = formatted_budget
                            monthly_totals[i] += budget_val
                            q_budget_total += budget_val
                            # Add metadata for month columns - FIXED: align with display logic
                            cell_metadata[(current_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(formatted_budget),
                                'media_type': current_media_type_normalized,
                                'value': budget_val
                            }
                        else:
                            logger.warning(f"Month column {month_name} not found in data.")
                            cell_metadata[(current_row_idx, month_col_idx)] = {
                                'has_data': False,
                                'media_type': current_media_type_normalized,
                                'value': 0
                            }
                
                # Process GRP data for TV media types
                if is_tv_current_media and 'GRP' in df.columns:
                    # For wide-format data, GRP is in a single column, not monthly
                    q_grp_total = 0
                    grp_row_idx = current_row_idx + 1  # GRP row comes after main media row
                    
                    for i, month_name in enumerate(MONTHS):
                        month_col_idx = i + 7  # Month columns start after first 7 columns
                        
                        if month_name.startswith('Q'):  # Q1, Q2, Q3, Q4
                            formatted_grp = format_number(q_grp_total, is_grp=True)
                            monthly_grp_data_formatted[i] = formatted_grp
                            grp_monthly_totals[i] += q_grp_total
                            # Add metadata for GRP quarter columns - FIXED: align with display logic
                            cell_metadata[(grp_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(formatted_grp),
                                'media_type': 'GRPs',
                                'value': q_grp_total
                            }
                            q_grp_total = 0
                        else:
                            # Use month-specific aggregation from raw data
                            metrics = get_month_specific_tv_metrics(excel_path, region, masterbrand, campaign_name, year, month_name)
                            grp_val = metrics['grp_sum']
                            
                            
                            if grp_val > 0:
                                formatted_grp = format_number(grp_val, is_grp=True)
                                monthly_grp_data_formatted[i] = formatted_grp
                                grp_monthly_totals[i] += grp_val
                                q_grp_total += grp_val
                            else:
                                formatted_grp = '-'
                                monthly_grp_data_formatted[i] = formatted_grp
                            
                            # Add metadata for GRP month columns - FIXED: align with display logic
                            cell_metadata[(grp_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(formatted_grp),
                                'media_type': 'GRPs',
                                'value': grp_val
                            }
                
                # CRITICAL FIX: Display campaign info for ALL media types, not just first
                # This resolves the issue where Digital campaigns appeared with blank names
                campaign_name_display = str(campaign_name).upper()
                campaign_budget_display = format_number(campaign_total_budget, is_budget=True)
                campaign_percentage_display = format_number(campaign_percentage, is_percentage=True)
                
                # Display total GRP for the campaign only on the primary TV row
                campaign_total_grp_display = format_number(campaign_total_tv_grp_value, is_grp=True) if first_media_type_for_campaign and is_tv_current_media else ""
                if first_media_type_for_campaign and is_tv_current_media:
                    grand_total_grp += campaign_total_tv_grp_value  # Accumulate for footer
                
                # Calculate TOTAL REACH and TOTAL FREQ for TV campaigns
                total_reach_display = ""
                total_freq_display = ""
                if first_media_type_for_campaign and is_tv_current_media:
                    # Calculate TOTAL REACH as average of all Reach@1+ values for this campaign
                    # Get all monthly Reach@1+ values for this campaign from raw data
                    all_reach1_values = []
                    total_grps_for_campaign = 0
                    
                    # Collect all monthly reach@1+ values and GRPs for this campaign
                    for month in ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']:
                        metrics = get_month_specific_tv_metrics(excel_path, region, masterbrand, campaign_name, year, month)
                        if not pd.isna(metrics['reach1_avg']) and metrics['reach1_avg'] > 0:
                            all_reach1_values.append(metrics['reach1_avg'])
                        if not pd.isna(metrics['grp_sum']) and metrics['grp_sum'] > 0:
                            total_grps_for_campaign += metrics['grp_sum']
                    
                    # Calculate TOTAL REACH as average of all monthly Reach@1+ values
                    if all_reach1_values:
                        total_reach_avg = sum(all_reach1_values) / len(all_reach1_values)
                        reach_pct = total_reach_avg * 100
                        total_reach_display = f"{reach_pct:.0f}%"
                        
                        # Calculate TOTAL FREQ as Total GRPs ÷ Total Average Reach@1+
                        if total_reach_avg > 0 and total_grps_for_campaign > 0:
                            # Fix frequency calculation: GRPs are already percentage, reach is decimal
                            # Convert reach from decimal to percentage for calculation
                            total_freq = total_grps_for_campaign / (total_reach_avg * 100)
                            # Format as whole number with thousands separator
                            total_freq_display = f"{int(round(total_freq)):,}"
                        else:
                            total_freq_display = "-"
                    else:
                        total_reach_display = "-"
                        total_freq_display = "-"
                else:
                    total_reach_display = "-"
                    total_freq_display = "-"

                # Create primary media type row
                primary_row_data = [
                    campaign_name_display,
                    campaign_budget_display,
                    campaign_total_grp_display,
                    total_reach_display,
                    total_freq_display,
                    campaign_percentage_display,
                    current_media_type_normalized
                ] + monthly_budget_data_formatted
                table_data.append(primary_row_data)
                
                # Add GRP sub-row for TV media types
                if is_tv_current_media:
                    grp_sub_row_data = [
                        "",  # Campaign Name blank
                        "",  # Budget blank
                        "",  # TV GRPs (campaign total) blank
                        "",  # TOTAL REACH blank
                        "",  # TOTAL FREQ blank
                        "",  # % blank
                        GRP_DATA_IDENTIFIER  # Media Type is GRP_DATA_IDENTIFIER
                    ] + monthly_grp_data_formatted
                    table_data.append(grp_sub_row_data)
                    
                    # Add Reach@1+ sub-row for TV campaigns
                    # Calculate month-by-month Reach 1+ values
                    monthly_reach_data = []
                    reach_row_idx = len(table_data)  # Get the row index for metadata
                    
                    # Track quarterly values for average calculation
                    q1_values, q2_values, q3_values, q4_values = [], [], [], []
                    current_quarter_values = q1_values
                    
                    for i, month_name in enumerate(MONTHS):
                        month_col_idx = i + 7  # Month columns start after first 7 columns
                        
                        if month_name.startswith('Q'):  # Q1, Q2, Q3, Q4
                            # Calculate quarterly average from accumulated values
                            if current_quarter_values:
                                quarterly_avg = sum(current_quarter_values) / len(current_quarter_values)
                                # Convert decimal average to percentage (0.70 -> 70%)
                                quarterly_pct = quarterly_avg * 100
                                reach1_formatted = f"{quarterly_pct:.0f}%" if quarterly_pct > 0 else "-"
                            else:
                                reach1_formatted = "-"
                                quarterly_avg = 0
                            
                            monthly_reach_data.append(reach1_formatted)
                            # Add cell metadata for coloring - FIXED: align with display logic
                            cell_metadata[(reach_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(reach1_formatted),
                                'media_type': 'Reach@1+',
                                'value': quarterly_avg
                            }
                            
                            # Switch to next quarter
                            if month_name == 'Q1':
                                current_quarter_values = q2_values
                            elif month_name == 'Q2':
                                current_quarter_values = q3_values
                            elif month_name == 'Q3':
                                current_quarter_values = q4_values
                        else:
                            # Use month-specific aggregation from raw data
                            metrics = get_month_specific_tv_metrics(excel_path, region, masterbrand, campaign_name, year, month_name)
                            reach1_avg = metrics['reach1_avg']
                            
                            if not pd.isna(reach1_avg) and reach1_avg > 0:
                                # Convert decimal to percentage (0.70 -> 70%)
                                reach1_pct = reach1_avg * 100
                                reach1_formatted = f"{reach1_pct:.0f}%"
                                month_reach1_avg = reach1_avg
                                current_quarter_values.append(month_reach1_avg)
                            else:
                                reach1_formatted = "-"
                                month_reach1_avg = 0
                            
                            monthly_reach_data.append(reach1_formatted)
                            # Add cell metadata for coloring - FIXED: align with display logic
                            cell_metadata[(reach_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(reach1_formatted),
                                'media_type': 'Reach@1+',
                                'value': month_reach1_avg
                            }
                    
                    reach_sub_row_data = [
                        "",  # Campaign Name blank
                        "-",  # Budget shows dash
                        "-",  # TV GRPs shows dash
                        "-",  # TOTAL REACH shows dash
                        "-",  # TOTAL FREQ shows dash
                        "-",  # % shows dash
                        "Reach@1+"  # Media Type column
                    ] + monthly_reach_data
                    table_data.append(reach_sub_row_data)
                    
                    # Add OTS@1+ sub-row for TV campaigns
                    # Calculate month-by-month OTS@1+ values (GRPs ÷ Reach@1+)
                    monthly_freq_data = []
                    freq_row_idx = len(table_data)  # Get the row index for metadata
                    
                    # Track quarterly values for average calculation
                    q1_values, q2_values, q3_values, q4_values = [], [], [], []
                    current_quarter_values = q1_values
                    
                    for i, month_name in enumerate(MONTHS):
                        month_col_idx = i + 7  # Month columns start after first 7 columns
                        
                        if month_name.startswith('Q'):  # Q1, Q2, Q3, Q4
                            # Calculate quarterly average from accumulated values
                            if current_quarter_values:
                                quarterly_avg = sum(current_quarter_values) / len(current_quarter_values)
                                freq_formatted = f"{quarterly_avg:.0f}" if quarterly_avg > 0 else "-"
                            else:
                                freq_formatted = "-"
                                quarterly_avg = 0
                            
                            monthly_freq_data.append(freq_formatted)
                            # Add cell metadata for coloring - FIXED: align with display logic
                            cell_metadata[(freq_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(freq_formatted),
                                'media_type': 'OTS@1+',
                                'value': quarterly_avg
                            }
                            
                            # Switch to next quarter
                            if month_name == 'Q1':
                                current_quarter_values = q2_values
                            elif month_name == 'Q2':
                                current_quarter_values = q3_values
                            elif month_name == 'Q3':
                                current_quarter_values = q4_values
                        else:
                            # Use month-specific aggregation from raw data
                            metrics = get_month_specific_tv_metrics(excel_path, region, masterbrand, campaign_name, year, month_name)
                            grp_sum = metrics['grp_sum']
                            reach1_avg = metrics['reach1_avg']
                            
                            # Calculate OTS@1+ = GRPs ÷ Reach@1+
                            if (not pd.isna(grp_sum) and grp_sum > 0 and 
                                not pd.isna(reach1_avg) and reach1_avg > 0):
                                # reach1_avg is in decimal form (0.68 = 68%), so divide by 100
                                ots1_value = grp_sum / (reach1_avg * 100)
                                freq_formatted = f"{ots1_value:.0f}"
                                month_freq_avg = ots1_value
                                current_quarter_values.append(month_freq_avg)
                            else:
                                freq_formatted = "-"
                                month_freq_avg = 0
                            
                            monthly_freq_data.append(freq_formatted)
                            # Add cell metadata for coloring - FIXED: align with display logic
                            cell_metadata[(freq_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(freq_formatted),
                                'media_type': 'OTS@1+',
                                'value': month_freq_avg
                            }
                    
                    freq_sub_row_data = [
                        "",  # Campaign Name blank
                        "-",  # Budget shows dash
                        "-",  # TV GRPs shows dash
                        "-",  # TOTAL REACH shows dash
                        "-",  # TOTAL FREQ shows dash
                        "-",  # % shows dash
                        "OTS@1+"  # Media Type column
                    ] + monthly_freq_data
                    table_data.append(freq_sub_row_data)
                    
                    # Add Reach@3+ sub-row for TV campaigns
                    # Calculate month-by-month Reach 3+ values
                    monthly_reach3_data = []
                    reach3_row_idx = len(table_data)  # Get the row index for metadata
                    
                    # Track quarterly values for average calculation
                    q1_values, q2_values, q3_values, q4_values = [], [], [], []
                    current_quarter_values = q1_values
                    
                    for i, month_name in enumerate(MONTHS):
                        month_col_idx = i + 7  # Month columns start after first 7 columns
                        
                        if month_name.startswith('Q'):  # Q1, Q2, Q3, Q4
                            # Calculate quarterly average from accumulated values
                            if current_quarter_values:
                                quarterly_avg = sum(current_quarter_values) / len(current_quarter_values)
                                # Convert decimal average to percentage (0.70 -> 70%)
                                quarterly_pct = quarterly_avg * 100
                                reach3_formatted = f"{quarterly_pct:.0f}%" if quarterly_pct > 0 else "-"
                            else:
                                reach3_formatted = "-"
                                quarterly_avg = 0
                            
                            monthly_reach3_data.append(reach3_formatted)
                            # Add cell metadata for coloring - FIXED: align with display logic
                            cell_metadata[(reach3_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(reach3_formatted),
                                'media_type': 'Reach@3+',
                                'value': quarterly_avg
                            }
                            
                            # Switch to next quarter
                            if month_name == 'Q1':
                                current_quarter_values = q2_values
                            elif month_name == 'Q2':
                                current_quarter_values = q3_values
                            elif month_name == 'Q3':
                                current_quarter_values = q4_values
                        else:
                            # Use month-specific aggregation from raw data
                            metrics = get_month_specific_tv_metrics(excel_path, region, masterbrand, campaign_name, year, month_name)
                            reach3_avg = metrics['reach3_avg']
                            
                            if not pd.isna(reach3_avg) and reach3_avg > 0:
                                # Convert decimal to percentage (0.70 -> 70%)
                                reach3_pct = reach3_avg * 100
                                reach3_formatted = f"{reach3_pct:.0f}%"
                                month_reach3_avg = reach3_avg
                                current_quarter_values.append(month_reach3_avg)
                            else:
                                reach3_formatted = "-"
                                month_reach3_avg = 0
                            
                            monthly_reach3_data.append(reach3_formatted)
                            # Add cell metadata for coloring - FIXED: align with display logic
                            cell_metadata[(reach3_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(reach3_formatted),
                                'media_type': 'Reach@3+',
                                'value': month_reach3_avg
                            }
                    
                    reach3_sub_row_data = [
                        "",  # Campaign Name blank
                        "-",  # Budget shows dash
                        "-",  # TV GRPs shows dash
                        "-",  # TOTAL REACH shows dash
                        "-",  # TOTAL FREQ shows dash
                        "-",  # % shows dash
                        "Reach@3+"  # Media Type column
                    ] + monthly_reach3_data
                    table_data.append(reach3_sub_row_data)
                    
                    # Add OTS@3+ sub-row for TV campaigns (OTS@3+ = GRPs ÷ Reach@3+)
                    # Calculate month-by-month OTS values
                    monthly_ots_data = []
                    ots_row_idx = len(table_data)  # Get the row index for metadata
                    
                    # Track quarterly values for average calculation
                    q1_values, q2_values, q3_values, q4_values = [], [], [], []
                    current_quarter_values = q1_values
                    
                    for i, month_name in enumerate(MONTHS):
                        month_col_idx = i + 7  # Month columns start after first 7 columns
                        
                        if month_name.startswith('Q'):  # Q1, Q2, Q3, Q4
                            # Calculate quarterly average from accumulated values
                            if current_quarter_values:
                                quarterly_avg = sum(current_quarter_values) / len(current_quarter_values)
                                # OTS is a regular number, not a percentage
                                ots_formatted = f"{quarterly_avg:.0f}" if quarterly_avg > 0 else "-"
                            else:
                                ots_formatted = "-"
                                quarterly_avg = 0
                            
                            monthly_ots_data.append(ots_formatted)
                            # Add cell metadata for coloring - FIXED: align with display logic
                            cell_metadata[(ots_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(ots_formatted),
                                'media_type': 'OTS@3+',
                                'value': quarterly_avg
                            }
                            
                            # Switch to next quarter
                            if month_name == 'Q1':
                                current_quarter_values = q2_values
                            elif month_name == 'Q2':
                                current_quarter_values = q3_values
                            elif month_name == 'Q3':
                                current_quarter_values = q4_values
                        else:
                            # Use month-specific aggregation from raw data
                            metrics = get_month_specific_tv_metrics(excel_path, region, masterbrand, campaign_name, year, month_name)
                            grp_sum = metrics['grp_sum']
                            reach3_avg = metrics['reach3_avg']
                            
                            # Calculate OTS@3+ = GRPs ÷ Reach@3+
                            if (not pd.isna(grp_sum) and grp_sum > 0 and 
                                not pd.isna(reach3_avg) and reach3_avg > 0):
                                # reach3_avg is in decimal form (0.15 = 15%), so divide by 100
                                ots_value = grp_sum / (reach3_avg * 100)
                                # OTS is a regular number, not a percentage
                                ots_formatted = f"{ots_value:.0f}"
                                month_ots_avg = ots_value
                                current_quarter_values.append(month_ots_avg)
                            else:
                                ots_formatted = "-"
                                month_ots_avg = 0
                            
                            monthly_ots_data.append(ots_formatted)
                            # Add cell metadata for coloring - FIXED: align with display logic
                            cell_metadata[(ots_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(ots_formatted),
                                'media_type': 'OTS@3+',
                                'value': month_ots_avg
                            }
                    
                    ots_sub_row_data = [
                        "",  # Campaign Name blank
                        "-",  # Budget shows dash
                        "-",  # TV GRPs shows dash
                        "-",  # TOTAL REACH shows dash
                        "-",  # TOTAL FREQ shows dash
                        "-",  # % shows dash
                        "OTS@3+"  # Media Type column
                    ] + monthly_ots_data
                    table_data.append(ots_sub_row_data)
                
                first_media_type_for_campaign = False
        
        # Add total row
        total_row_budget_values = [format_number(total, is_budget=True, is_monthly_column=True) for total in monthly_totals]
        
        total_row = [
            'TOTAL',
            format_number(total_budget_for_percentage, is_budget=True),
            format_number(grand_total_grp, is_grp=True),  # Sum of all campaign total GRPs for TV
            '-',  # TOTAL REACH blank for total row
            '-',  # TOTAL FREQ blank for total row
            format_number(100, is_percentage=True),
            ''  # Media type blank for total row
        ] + total_row_budget_values
        table_data.append(total_row)
        
        logger.info(f"Table data created successfully for {region} - {masterbrand}. Total rows: {len(table_data)}")
        return table_data, cell_metadata
        
    except Exception as e:
        logger.error(f"Error preparing table data for {region} - {masterbrand}: {str(e)}")
        logger.error(traceback.format_exc())
        return None, None

def _split_table_data_by_campaigns(table_data, cell_metadata):
    """
    Split table data into multiple tables if it exceeds MAX_ROWS_PER_SLIDE.
    Keeps complete campaigns together (including TV sub-rows).
    
    Args:
        table_data: List of lists representing table rows
        cell_metadata: Dictionary with metadata for each cell
        
    Returns:
        List of tuples: [(table_data_split, cell_metadata_split, is_continuation), ...]
    """
    if not table_data or len(table_data) <= MAX_ROWS_PER_SLIDE:
        # No need to split
        return [(table_data, cell_metadata, False)]
    
    logger.info(f"Table has {len(table_data)} rows, splitting needed (max: {MAX_ROWS_PER_SLIDE})")
    
    # Identify header and total row
    header_row = table_data[0]
    total_row = table_data[-1]
    
    # Group campaigns with their sub-rows
    campaigns = []
    current_campaign = []
    current_campaign_name = None
    
    # Track which rows are sub-rows
    sub_row_media_types = ['GRPs', 'Reach@1+', 'OTS@1+', 'Reach@3+', 'OTS@3+']
    
    for i in range(1, len(table_data) - 1):  # Skip header and total
        row = table_data[i]
        media_type = row[6] if len(row) > 6 else ""
        campaign_name = row[0] if row[0] else ""
        
        if media_type in sub_row_media_types:
            # This is a sub-row, add to current campaign
            if current_campaign:  # Make sure we have a campaign to add to
                current_campaign.append((i, row))
        else:
            # This is a main row
            if current_campaign and current_campaign_name:
                # Save previous campaign (including its main row and sub-rows)
                campaigns.append({
                    'name': current_campaign_name,
                    'rows': current_campaign,
                    'start_idx': current_campaign[0][0]
                })
            # Start new campaign with this main row
            current_campaign_name = campaign_name
            current_campaign = [(i, row)]
    
    # Don't forget the last campaign
    if current_campaign and current_campaign_name:
        campaigns.append({
            'name': current_campaign_name,
            'rows': current_campaign,
            'start_idx': current_campaign[0][0]
        })
    
    
    # Now split campaigns into slides
    splits = []
    current_split = [header_row]
    current_split_indices = [0]  # Track original indices for metadata
    current_row_count = 1  # Start with header
    split_number = 0
    
    # Track totals for subtotals
    running_total_budget = 0
    running_total_grp = 0
    subtotal_values = None  # Initialize for carried forward rows
    all_split_rows = []  # Track all rows across splits for final total
    
    for campaign in campaigns:
        campaign_row_count = len(campaign['rows'])
        
        # Check if adding this campaign would exceed the limit
        # Reserve 1 row for subtotal/total
        if current_row_count + campaign_row_count + 1 > MAX_ROWS_PER_SLIDE:
            # Complete current split with total/subtotal
            if split_number > 0 or len(splits) > 0:
                # For intermediate splits, use SUBTOTAL
                subtotal_values = _calculate_subtotal_for_split(current_split[1:])  # Exclude header
                subtotal_row = ['SUBTOTAL'] + subtotal_values
                current_split.append(subtotal_row)
                current_split_indices.append(-1)  # Special index for subtotal
            else:
                # For first split when there will be more, use TOTAL
                total_values = _calculate_subtotal_for_split(current_split[1:])  # Exclude header
                total_row_split = ['TOTAL'] + total_values
                current_split.append(total_row_split)
                current_split_indices.append(-1)  # Special index for total
            
            # Save current split
            split_metadata = _extract_metadata_for_indices(cell_metadata, current_split_indices, len(current_split))
            splits.append((current_split, split_metadata, split_number > 0))
            
            # Start new split
            split_number += 1
            current_split = [header_row]
            current_split_indices = [0]
            current_row_count = 1
            
            # Add carried forward subtotal if configured
            if SHOW_CARRIED_SUBTOTAL and splits and subtotal_values:
                # Make sure subtotal_values has the right number of columns
                if len(subtotal_values) == 22:  # Missing campaign name column
                    carried_row = ['CARRIED FORWARD'] + subtotal_values
                else:
                    logger.warning(f"Unexpected subtotal_values length: {len(subtotal_values)}")
                    carried_row = ['CARRIED FORWARD'] + subtotal_values
                current_split.append(carried_row)
                current_split_indices.append(-2)  # Special index for carried forward
                current_row_count += 1
        
        # Add campaign to current split
        for idx, row in campaign['rows']:
            current_split.append(row)
            current_split_indices.append(idx)
            all_split_rows.append(row)  # Track for final total
        current_row_count += campaign_row_count
    
    # Handle the last split
    if len(splits) > 0:
        # This is a continuation slide, calculate grand total from all rows
        grand_total_values = _calculate_subtotal_for_split(all_split_rows)
        grand_total_row = ['TOTAL'] + grand_total_values
        current_split.append(grand_total_row)
        current_split_indices.append(-1)
    else:
        # This is the only slide, add original total row
        current_split.append(total_row)
        current_split_indices.append(len(table_data) - 1)
    
    # Extract metadata for final split
    split_metadata = _extract_metadata_for_indices(cell_metadata, current_split_indices, len(current_split))
    splits.append((current_split, split_metadata, split_number > 0))
    
    logger.info(f"Table split into {len(splits)} slides")
    return splits


def _calculate_subtotal_for_split(rows):
    """Calculate subtotal values for a split section."""
    # Initialize subtotal values
    subtotal_values = []
    
    # Define sub-row media types to exclude from totals
    sub_row_media_types = ['GRPs', 'Reach@1+', 'OTS@1+', 'Reach@3+', 'OTS@3+']
    
    # Skip first column (campaign name), process budget column
    total_budget = 0
    for row in rows:
        # Skip sub-rows when calculating budget totals
        media_type = row[6] if len(row) > 6 else ""
        if media_type in sub_row_media_types:
            continue
            
        if row[1] and row[1] != '-':
            # Extract numeric value from budget string
            budget_str = row[1].replace('£', '').replace('K', '000').replace('M', '000000').replace(',', '')
            try:
                total_budget += float(budget_str)
            except:
                pass
    
    # Format budget
    if total_budget >= 1000000:
        subtotal_values.append(f"£{total_budget/1000000:.1f}M")
    else:
        subtotal_values.append(f"£{int(total_budget/1000)}K")
    
    # Handle other columns (simplified for now)
    # TV GRPs, TOTAL REACH, TOTAL FREQ, %
    subtotal_values.extend(['-', '-', '-', '-'])
    
    # Media Type column
    subtotal_values.append('')
    
    # Monthly columns - sum up values
    for col_idx in range(7, 23):  # Monthly columns
        monthly_total = 0
        for row in rows:
            # Skip sub-rows when calculating monthly totals
            media_type = row[6] if len(row) > 6 else ""
            if media_type in sub_row_media_types:
                continue
                
            if col_idx < len(row) and row[col_idx] and row[col_idx] != '-':
                try:
                    value_str = row[col_idx].replace('£', '').replace('K', '000').replace('M', '000000').replace(',', '')
                    monthly_total += float(value_str)
                except:
                    pass
        
        if monthly_total > 0:
            if monthly_total >= 1000000:
                subtotal_values.append(f"£{monthly_total/1000000:.1f}M")
            else:
                subtotal_values.append(f"£{int(monthly_total/1000)}K")
        else:
            subtotal_values.append('-')
    
    return subtotal_values


def _extract_metadata_for_indices(original_metadata, indices, new_row_count):
    """Extract metadata for specific row indices and remap to new positions."""
    new_metadata = {}
    
    for new_row_idx, original_idx in enumerate(indices):
        if original_idx < 0:
            # Special indices for subtotal/carried forward rows
            continue
            
        # Copy metadata for this row
        for (orig_row, col), meta in original_metadata.items():
            if orig_row == original_idx:
                new_metadata[(new_row_idx, col)] = meta
    
    return new_metadata


def _prepare_media_type_chart_data_detailed(df, region, masterbrand, year=None):
    """
    Prepare pie chart data for media type budget distribution.
    
    Args:
        df: DataFrame with loaded Excel data
        region: Region name to filter by
        masterbrand: Masterbrand name to filter by
        
    Returns:
        dict: Dictionary with media types as keys and budget values as values
    """
    try:
        year_text = f" - {year}" if year is not None else ""
        logger.info(f"Preparing media type chart data for {region} - {masterbrand}{year_text}")
        
        # Filter data for the specific region, masterbrand, and year
        filter_conditions = [
            (df['Country'].astype(str).str.strip() == region),
            (df['Brand'].astype(str).str.strip() == masterbrand)
        ]
        
        if year is not None:
            filter_conditions.append(df['Year'].astype(str).str.strip() == str(year))
        
        filtered_df = df[
            filter_conditions[0] & filter_conditions[1] & (filter_conditions[2] if len(filter_conditions) > 2 else True)
        ].copy()
        
        if filtered_df.empty:
            logger.warning(f"No data found for media type chart: {region} - {masterbrand}")
            return None
            
        # Group by mapped media type and sum budgets
        media_type_budgets = {}
        
        # Get media types and their corresponding budgets
        for media_type_orig in filtered_df['Media Type'].unique():
            if pd.isna(media_type_orig):
                continue
                
            media_type_normalized = normalize_media_type(str(media_type_orig))
            media_df = filtered_df[filtered_df['Media Type'] == media_type_orig]
            total_budget = media_df['Total Cost'].sum()
            
            if media_type_normalized in media_type_budgets:
                media_type_budgets[media_type_normalized] += total_budget
            else:
                media_type_budgets[media_type_normalized] = total_budget
        
        # Filter out zero budgets
        media_type_budgets = {k: v for k, v in media_type_budgets.items() if v > 0}
        
        logger.info(f"Media type chart data prepared: {media_type_budgets}")
        return media_type_budgets
        
    except Exception as e:
        logger.error(f"Error preparing media type chart data for {region} - {masterbrand}: {str(e)}")
        logger.error(traceback.format_exc())
        return None

def _prepare_funnel_chart_data_detailed(df, region, masterbrand, year=None):
    """
    Prepare pie chart data for funnel stage budget distribution.
    
    Args:
        df: DataFrame with loaded Excel data
        region: Region name to filter by
        masterbrand: Masterbrand name to filter by
        
    Returns:
        dict: Dictionary with funnel stages as keys and budget values as values
    """
    try:
        year_text = f" - {year}" if year is not None else ""
        logger.info(f"Preparing funnel chart data for {region} - {masterbrand}{year_text}")
        
        # Filter data for the specific region, masterbrand, and year
        filter_conditions = [
            (df['Country'].astype(str).str.strip() == region),
            (df['Brand'].astype(str).str.strip() == masterbrand)
        ]
        
        if year is not None:
            filter_conditions.append(df['Year'].astype(str).str.strip() == str(year))
        
        filtered_df = df[
            filter_conditions[0] & filter_conditions[1] & (filter_conditions[2] if len(filter_conditions) > 2 else True)
        ].copy()
        
        if filtered_df.empty:
            logger.warning(f"No data found for funnel chart: {region} - {masterbrand}")
            return None
        
        # Group by funnel stage and sum budgets
        funnel_budgets = {}
        
        if 'Funnel Stage' in df.columns:
            for funnel_stage in filtered_df['Funnel Stage'].unique():
                if pd.isna(funnel_stage) or str(funnel_stage).strip() == '':
                    continue
                    
                funnel_df = filtered_df[filtered_df['Funnel Stage'] == funnel_stage]
                total_budget = funnel_df['Total Cost'].sum()
                
                if total_budget > 0:
                    funnel_budgets[str(funnel_stage)] = total_budget
        
        logger.info(f"Funnel chart data prepared: {funnel_budgets}")
        return funnel_budgets
        
    except Exception as e:
        logger.error(f"Error preparing funnel chart data for {region} - {masterbrand}: {str(e)}")
        logger.error(traceback.format_exc())
        return None

def _prepare_campaign_type_chart_data(df, region, masterbrand, year=None):
    """
    Prepare pie chart data for campaign type budget distribution.
    
    Args:
        df: DataFrame with loaded Excel data
        region: Region name to filter by
        masterbrand: Masterbrand name to filter by
        
    Returns:
        dict: Dictionary with campaign types as keys and budget values as values
    """
    try:
        year_text = f" - {year}" if year is not None else ""
        logger.info(f"Preparing campaign type chart data for {region} - {masterbrand}{year_text}")
        
        # Filter data for the specific region, masterbrand, and year
        filter_conditions = [
            (df['Country'].astype(str).str.strip() == region),
            (df['Brand'].astype(str).str.strip() == masterbrand)
        ]
        
        if year is not None:
            filter_conditions.append(df['Year'].astype(str).str.strip() == str(year))
        
        filtered_df = df[
            filter_conditions[0] & filter_conditions[1] & (filter_conditions[2] if len(filter_conditions) > 2 else True)
        ].copy()
        
        if filtered_df.empty:
            logger.warning(f"No data found for campaign type chart: {region} - {masterbrand}")
            return None
        
        # Group by campaign type and sum budgets
        campaign_type_budgets = {}
        
        if 'Campaign Type' in df.columns:
            for campaign_type in filtered_df['Campaign Type'].unique():
                if pd.isna(campaign_type) or str(campaign_type).strip() == '':
                    continue
                    
                campaign_df = filtered_df[filtered_df['Campaign Type'] == campaign_type]
                total_budget = campaign_df['Total Cost'].sum()
                
                if total_budget > 0:
                    campaign_type_budgets[str(campaign_type)] = total_budget
        
        logger.info(f"Campaign type chart data prepared: {campaign_type_budgets}")
        return campaign_type_budgets
        
    except Exception as e:
        logger.error(f"Error preparing campaign type chart data for {region} - {masterbrand}: {str(e)}")
        logger.error(traceback.format_exc())
        return None

def set_title_text_detailed(title_shape, title_text, template_prs):
    logger.debug(f"Attempting to set title text '{title_text}' for shape '{title_shape.name}' using detailed method.")
    # --- Locate the reference title placeholder and its text_frame from the template --- #
    template_title_ref_shape = None
    # Priority: 1. Named 'TitlePlaceholder' on first slide master
    if template_prs.slide_masters:
        master = template_prs.slide_masters[0]
        for placeholder in master.placeholders:
            if placeholder.name == "TitlePlaceholder":
                template_title_ref_shape = placeholder
                logger.debug(f"Found reference title placeholder '{placeholder.name}' by name in slide master.")
                break
            elif placeholder.is_placeholder and placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                template_title_ref_shape = placeholder
                logger.debug(f"Found reference title placeholder type '{placeholder.placeholder_format.type}' in slide master.")
                break # Take the first one found
    
    # Priority: 2. Named 'TitlePlaceholder' on first slide layout (if not found in master)
    if not template_title_ref_shape and template_prs.slide_layouts:
        layout = template_prs.slide_layouts[0] # Assuming first layout is relevant
        for placeholder in layout.placeholders:
            if placeholder.name == "TitlePlaceholder":
                template_title_ref_shape = placeholder
                logger.debug(f"Found reference title placeholder '{placeholder.name}' by name in slide layout '{layout.name}'.")
                break
            elif placeholder.is_placeholder and placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                template_title_ref_shape = placeholder
                logger.debug(f"Found reference title placeholder type '{placeholder.placeholder_format.type}' in slide layout '{layout.name}'.")
                break

    # Priority: 3. Named 'TitlePlaceholder' or actual title placeholder on the first actual slide (if not in master/layout)
    if not template_title_ref_shape and template_prs.slides:
        actual_slide = template_prs.slides[0]
        if actual_slide.has_title and actual_slide.shapes.title:
            template_title_ref_shape = actual_slide.shapes.title
            logger.debug(f"Using actual title shape from template's first slide: '{template_title_ref_shape.name}'")
        else:
            for shape_on_slide in actual_slide.shapes:
                if shape_on_slide.name == "TitlePlaceholder":
                    template_title_ref_shape = shape_on_slide
                    logger.debug(f"Found named shape 'TitlePlaceholder' on template's first slide.")
                    break

    if not template_title_ref_shape or not template_title_ref_shape.has_text_frame:
        logger.warning("Could not find a suitable reference title placeholder with a text frame in the template. Applying basic text only.")
        title_shape.text = title_text
        return

    template_tf = template_title_ref_shape.text_frame
    logger.debug(f"Using template text_frame from shape: '{template_title_ref_shape.name}' (ID: {template_title_ref_shape.shape_id}) for styling.")

    # --- Copy Shape Fill --- 
    if template_title_ref_shape:
        source_shape_fill = template_title_ref_shape.fill
        target_shape_fill = title_shape.fill

        # Check if source_shape_fill.type is None first, as it can be for inherited fills
        if source_shape_fill.type is None or source_shape_fill.type == MSO_FILL_TYPE.BACKGROUND:
            target_shape_fill.background() # Makes it transparent / no fill
            logger.debug("Template Title - Source fill is NONE or BACKGROUND. Applied no fill to target.")
        elif source_shape_fill.type == MSO_FILL_TYPE.SOLID:
            target_shape_fill.solid()
            if source_shape_fill.fore_color.type == MSO_COLOR_TYPE.RGB:
                target_shape_fill.fore_color.rgb = source_shape_fill.fore_color.rgb
                logger.debug(f"Template Title - Source fill is SOLID RGB: {source_shape_fill.fore_color.rgb}. Applied to target.")
            elif source_shape_fill.fore_color.type == MSO_COLOR_TYPE.SCHEME:
                target_shape_fill.fore_color.theme_color = source_shape_fill.fore_color.theme_color
                if hasattr(source_shape_fill.fore_color, 'brightness') and source_shape_fill.fore_color.brightness is not None:
                    target_shape_fill.fore_color.brightness = source_shape_fill.fore_color.brightness
                if hasattr(source_shape_fill.fore_color, 'tint') and source_shape_fill.fore_color.tint is not None:
                    target_shape_fill.fore_color.tint = source_shape_fill.fore_color.tint
                if hasattr(source_shape_fill.fore_color, 'shade') and source_shape_fill.fore_color.shade is not None:
                    target_shape_fill.fore_color.shade = source_shape_fill.fore_color.shade
                logger.debug(f"Template Title - Source fill is SOLID SCHEME Color. Applied to target.")
            else:
                logger.warning(f"Template Title - Source fill is SOLID but color type {source_shape_fill.fore_color.type} not handled for fill. Fill color not fully copied.")
        # Add more fill types like GRADIENT, PICTURE if necessary later
        else:
            logger.warning(f"Template Title - Source fill type {source_shape_fill.type} not handled. Fill not copied.")
    else:
        logger.warning("Could not find reference title placeholder shape to copy fill.")
    # --- End of Shape Fill Copying ---

    # --- Apply Vertical Anchor from template_tf to title_shape.text_frame --- #
    source_va = template_tf.vertical_anchor
    logger.debug(f"Template Title - Read source vertical_anchor: {source_va} (Enum: {MSO_VERTICAL_ANCHOR(source_va) if source_va is not None else 'None'})")
    if source_va is not None:
        title_shape.text_frame.vertical_anchor = source_va
        logger.debug(f"Applied vertical_anchor '{source_va}' to title_shape.text_frame. New VA: {title_shape.text_frame.vertical_anchor}")
    else:
        logger.debug("Source vertical_anchor is None. Vertical alignment on new title will use default (likely TOP).")

    # --- Clear existing content and add new paragraph for the title text --- #
    title_shape.text_frame.clear()
    p = title_shape.text_frame.add_paragraph()
    p.text = title_text
    logger.debug(f"Set title text to: '{title_text}' in new paragraph.")

    # --- Apply Paragraph and Font Styling from template_tf to the new paragraph 'p' and its first run --- #
    if template_tf.paragraphs:
        ref_para = template_tf.paragraphs[0]  # Use first paragraph of template for style reference
        logger.debug(f"Reference paragraph from template: Level={ref_para.level}, SpaceBefore={ref_para.space_before}, SpaceAfter={ref_para.space_after}, LineSpacing={ref_para.line_spacing}")

        # 1. Apply Paragraph Alignment (Horizontal)
        if ref_para.alignment is not None:
            p.alignment = ref_para.alignment
            logger.debug(f"Template Title - Source Paragraph Alignment: {ref_para.alignment} (Enum: {PP_ALIGN(ref_para.alignment) if ref_para.alignment is not None else 'None'}). Applied to new paragraph.")
        else:
            logger.debug("Template Title - Source paragraph alignment is None. Horizontal alignment will use default.")
        
        # Apply other paragraph properties if needed (e.g., level, spacing - though these are often complex)
        # p.level = ref_para.level # Be cautious with level, can affect master styles

        # 2. Apply Font Styling from the first run of the reference paragraph
        if ref_para.runs:
            ref_run = ref_para.runs[0]  # Reference run from template
            
            if p.runs: # The new paragraph 'p' should have one run after p.text = title_text
                target_run = p.runs[0]

                font_attrs_log = []
                # Font Name
                if ref_run.font.name:
                    target_run.font.name = ref_run.font.name
                    font_attrs_log.append(f"Name='{ref_run.font.name}'")
                # Font Size
                if ref_run.font.size:
                    target_run.font.size = ref_run.font.size
                    font_attrs_log.append(f"Size={ref_run.font.size}")
                # Bold
                if ref_run.font.bold is not None:
                    target_run.font.bold = ref_run.font.bold
                    font_attrs_log.append(f"Bold={ref_run.font.bold}")
                # Italic
                if ref_run.font.italic is not None:
                    target_run.font.italic = ref_run.font.italic
                    font_attrs_log.append(f"Italic={ref_run.font.italic}")
                # Underline
                if ref_run.font.underline is not None:
                    target_run.font.underline = ref_run.font.underline
                    font_attrs_log.append(f"Underline={ref_run.font.underline}")
                # Color
                if ref_run.font.color.type:
                    if ref_run.font.color.type == MSO_COLOR_TYPE.RGB:
                        target_run.font.color.rgb = ref_run.font.color.rgb
                        font_attrs_log.append(f"ColorRGB='{ref_run.font.color.rgb}'")
                    elif ref_run.font.color.type == MSO_COLOR_TYPE.SCHEME:
                        target_run.font.color.theme_color = ref_run.font.color.theme_color
                        target_run.font.color.brightness = ref_run.font.color.brightness
                        font_attrs_log.append(f"ColorTheme='{ref_run.font.color.theme_color}', Brightness='{ref_run.font.color.brightness}'")
                    # Add other color types if necessary
                
                logger.debug(f"Template Title - Source Font Details: {', '.join(font_attrs_log)}")
                applied_font_attrs_log = [
                    f"Name='{target_run.font.name}'", f"Size={target_run.font.size}", 
                    f"Bold={target_run.font.bold}", f"Italic={target_run.font.italic}", 
                    f"Underline={target_run.font.underline}"
                ]
                if target_run.font.color.type == MSO_COLOR_TYPE.RGB:
                    applied_font_attrs_log.append(f"ColorRGB='{target_run.font.color.rgb}'")
                elif target_run.font.color.type == MSO_COLOR_TYPE.SCHEME:
                    applied_font_attrs_log.append(f"ColorTheme='{target_run.font.color.theme_color}', Brightness='{target_run.font.color.brightness}'")
                logger.debug(f"Applied Font to new run: {', '.join(applied_font_attrs_log)}")

            else:
                logger.warning("Title paragraph has no runs after setting text. Cannot apply font styling.")
        else:
            logger.debug("Template Title - Reference paragraph in template has no runs. Cannot copy font styling.")
    else:
        logger.debug("Template Title - Reference text_frame in template has no paragraphs. Cannot copy paragraph/font styling.")

    logger.debug(f"Finished setting title for shape '{title_shape.name}'. Text frame word_wrap: {title_shape.text_frame.word_wrap}, auto_size: {title_shape.text_frame.auto_size}")

def _add_and_style_table(slide, table_data, cell_metadata, template_slide=None):
    """
    Add a table to the slide and apply styling based on template or v1.5 specifications.
    
    Args:
        slide: The PowerPoint slide to add the table to
        table_data: List of lists representing table rows
        cell_metadata: Dictionary with metadata for conditional formatting
        template_slide: Optional template slide to copy table styling from
        
    Returns:
        bool: True if table was created successfully, False otherwise
    """
    try:
        if not table_data or len(table_data) < 2:
            logger.warning("No table data provided or insufficient data for table creation")
            return False
        
        rows = len(table_data)
        cols = len(table_data[0])
        
        logger.info(f"Creating table with {rows} rows and {cols} columns")
        
        # Table positioning using centralized coordinate system
        table_pos = get_element_position('main_table')
        if not table_pos:
            logger.error("Failed to get table position coordinates")
            return False
        
        # Try to drop into our pre-styled placeholder first
        table_shape = None
        
        # ENHANCED DEBUGGING: Commented out to reduce log file size
        # Uncomment for debugging placeholder issues
        """
        logger.info(f"=== PLACEHOLDER DEBUGGING for slide {slide} ===")
        logger.info(f"Total placeholders found: {len(slide.placeholders)}")
        
        for i, ph in enumerate(slide.placeholders):
            ph_name = getattr(ph, "name", "NO_NAME")
            ph_type = None
            ph_idx = None
            
            try:
                ph_type = ph.placeholder_format.type if hasattr(ph, 'placeholder_format') else "NO_FORMAT"
                ph_idx = ph.placeholder_format.idx if hasattr(ph, 'placeholder_format') else "NO_IDX"
            except:
                ph_type = "ERROR_GETTING_TYPE"
                ph_idx = "ERROR_GETTING_IDX"
            
            logger.info(f"  Placeholder {i}: Name='{ph_name}', Type={ph_type}, Idx={ph_idx}")
            
            # Check if this matches our criteria
            name_match = ph_name == TABLE_PLACEHOLDER_NAME
            type_match = ph_type == PP_PLACEHOLDER.TABLE
            logger.info(f"    Name match ('{ph_name}' == '{TABLE_PLACEHOLDER_NAME}'): {name_match}")
            logger.info(f"    Type match ({ph_type} == {PP_PLACEHOLDER.TABLE}): {type_match}")
        
        logger.info(f"=== END PLACEHOLDER DEBUGGING ===")
        """
        
        for ph in slide.placeholders:
            if (ph.placeholder_format.type == PP_PLACEHOLDER.TABLE
                and getattr(ph, "name", "") == TABLE_PLACEHOLDER_NAME):
                try:
                    table_shape = ph.insert_table(rows, cols)
                    logger.info(f"Inserted table into placeholder '{TABLE_PLACEHOLDER_NAME}'")
                except Exception as e:
                    logger.warning(f"Placeholder.insert_table() failed: {e}")
                break

        if not table_shape:
            # fallback for slides/layouts without our named placeholder
            logger.warning(f"Placeholder '{TABLE_PLACEHOLDER_NAME}' missing—using add_table() fallback")
            table_shape = slide.shapes.add_table(
                rows, cols,
                table_pos['left'], table_pos['top'],
                table_pos['width'], table_pos['height']
            )

        # rename & grab the Table object
        table_shape.name = SHAPE_NAME_TABLE
        table = table_shape.table
        
        # CRITICAL FIX: Set proper shape name to match template
        table_shape.name = SHAPE_NAME_TABLE  # "MainDataTable"
        
        # CRITICAL ROW HEIGHT FIX: PRECISION TARGET VALUES for pixel-perfect compliance
        # Based on diagnostic measurements - applying optimized heights for final compactness
        ROW_HEIGHT_HEADER = Pt(10.0)    # Target: 0.140" (reduced from 0.167") - 16% reduction
        ROW_HEIGHT_BODY = Pt(8.5)       # Target: 0.120" (reduced from 0.139") - 14% reduction  
        ROW_HEIGHT_SUBTOTAL = Pt(10.0)  # Target: 0.140" (reduced from 0.167") - 16% reduction
        
        # DIAGNOSTIC LOGGING: Report TARGET row heights for pixel-perfect compliance
        logger.info(f"=== PRECISION TARGET: Implementing Optimized Row Heights ===")
        logger.info(f"Header Row Height: {ROW_HEIGHT_HEADER} ({ROW_HEIGHT_HEADER.inches:.3f} inches) - TARGET")
        logger.info(f"Body Row Height: {ROW_HEIGHT_BODY} ({ROW_HEIGHT_BODY.inches:.3f} inches) - TARGET")
        logger.info(f"Subtotal Row Height: {ROW_HEIGHT_SUBTOTAL} ({ROW_HEIGHT_SUBTOTAL.inches:.3f} inches) - TARGET")
        
        for i, row in enumerate(table.rows):
            if i == 0:
                # Header row
                row.height = ROW_HEIGHT_HEADER
                logger.debug(f"DIAGNOSTIC: Applied header height {ROW_HEIGHT_HEADER} to row {i}")
            elif i == len(table_data) - 1:
                # Subtotal row (last row)
                row.height = ROW_HEIGHT_SUBTOTAL
                logger.debug(f"DIAGNOSTIC: Applied subtotal height {ROW_HEIGHT_SUBTOTAL} to row {i}")
            else:
                # Body rows (data rows)
                row.height = ROW_HEIGHT_BODY
                logger.debug(f"DIAGNOSTIC: Applied body height {ROW_HEIGHT_BODY} to row {i}")
            
            # CRITICAL FIX: Use EXACTLY height rule to prevent auto-expansion
            if TABLE_HEIGHT_RULE_AVAILABLE:
                try:
                    row.height_rule = WD_ROW_HEIGHT_RULE.AT_LEAST # Changed from EXACTLY
                    logger.debug(f"Applied AT_LEAST height rule to row {i}")
                except Exception as e:
                    logger.debug(f"Could not set exact height rule for row {i}: {e}")
            else:
                logger.debug(f"Height rule not available - row {i} height set but may auto-expand")
            
            # CRITICAL: UNIVERSAL CELL ALIGNMENT & MINIMAL PADDING - Apply to ALL cells
            for j, cell in enumerate(row.cells):
                # DISABLE text wrapping and auto-sizing FIRST
                cell.text_frame.word_wrap = False
                cell.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                
                # CRITICAL: UNIVERSAL VERTICAL CENTERING for ALL cells
                cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                
                # CRITICAL: UNIVERSAL HORIZONTAL CENTERING + MINIMAL PARAGRAPH SPACING for ALL paragraphs in ALL cells
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    
                    # **NEW CRITICAL STEP**: MINIMAL PARAGRAPH SPACING for true vertical compactness
                    # This controls spacing *within* the text block itself, vital for true vertical compactness
                    paragraph.line_spacing = 1.0  # Single line spacing (no extra spacing between lines)
                    paragraph.space_before = Pt(0)  # No space before the paragraph
                    paragraph.space_after = Pt(0)   # No space after the paragraph
                    
                    # **CRITICAL FIX**: FORCE Calibri font on ALL paragraphs/runs during initial setup
                    # This prevents any default font inheritance that could cause Roboto 18pt issues
                    if not paragraph.runs:
                        run = paragraph.add_run()  # Create run if none exists
                    
                    for run in paragraph.runs:
                        run.font.name = DEFAULT_FONT_NAME  # FORCE Calibri
                        # Set appropriate size based on row type
                        if i == 0:  # Header row
                            run.font.size = FONT_SIZE_HEADER   # 8pt
                        else:  # Body/subtotal rows
                            run.font.size = FONT_SIZE_BODY    # 7pt
                        run.font.color.rgb = CLR_BLACK
                
                # CRITICAL: ABSOLUTE MINIMAL cell padding (0pt = 0 EMU) for ultra-compact layout
                absolute_minimal_margin = 0  # 0pt in EMU - absolute minimum for maximum density
                try:
                    cell.margin_left = 0
                    cell.margin_right = 0
                    cell.margin_top = 0
                    cell.margin_bottom = 0
                except Exception as margin_error:
                    # Fallback to OXML method for margins
                    try:
                        cell._tc.tcPr.marL = 0  # Left margin = 0 EMU
                        cell._tc.tcPr.marR = 0  # Right margin = 0 EMU
                        cell._tc.tcPr.marT = 0  # Top margin = 0 EMU
                        cell._tc.tcPr.marB = 0  # Bottom margin = 0 EMU
                    except Exception as oxml_margin_error:
                        logger.debug(f"Could not set zero margins for cell ({i},{j}): {oxml_margin_error}")
        
        # REMOVED: Fixed table height constraint that was causing oversized rows
        # table_shape.height = Inches(2.340)   # REMOVED - was forcing fixed height
        
        # Ensure exact positioning
        table_shape.top = Inches(0.812)      # Maintain exact top position
        
        logger.info(f"Table created with individual row height constraints (no fixed total height)")
        
        # ENHANCED TABLE: Adjust column widths for new TOTAL REACH and TOTAL FREQ columns
        column_widths = [
            Inches(0.65),   # Campaign Name (reduced to allow wider key columns)
            Inches(0.50),   # Budget (increased to fit "BUDGET" on one line)
            Inches(0.35),   # TV GRPs (reduced)
            Inches(0.43),   # TOTAL REACH (increased to fit "REACH" on one line)
            Inches(0.35),   # TOTAL FREQ (reduced)
            Inches(0.40),   # % (reduced)
            Inches(0.72),   # Media Type (reduced)
        ] + [Inches(0.375)] * 16  # 16 months (JAN-DEC + Q1-Q4) - kept for single-line display
        # Total width calculation: adjusted to maintain table width
        # NEW Total: 0.65 + 0.50 + 0.35 + 0.43 + 0.35 + 0.40 + 0.72 + (16 × 0.375) = 9.40" (within bounds)

        # Apply column widths if we have enough columns
        for i, width in enumerate(column_widths[:cols]):
            table.columns[i].width = width
        
        # Populate table with data
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_value in enumerate(row_data):
                if col_idx < cols:  # Safety check
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value) if cell_value is not None else ""
                    
                    # Apply cell styling with metadata
                    _apply_table_cell_styling(cell, row_idx, col_idx, table_data, cell_metadata)
        
        # CRITICAL FIX: Apply borders to the entire table after all cells are styled
        _apply_table_borders(table)
        
        # CRITICAL FIX: Apply internal borders for first 7 columns only
        # _apply_internal_table_borders(table, rows) # Intentionally removed
        
        logger.info(f"Table created successfully with individual row height constraints")
        
        # DIAGNOSTIC LOGGING: Commented out to reduce log file size
        # Uncomment for debugging table structure issues
        """
        logger.info(f"=== DIAGNOSTIC: Final Table Structure Analysis ===")
        logger.info(f"Total table rows created: {len(table.rows)}")
        logger.info(f"Total table columns: {len(table.columns)}")
        
        # Sample a few rows to report their actual applied heights
        sample_rows = min(5, len(table.rows))  # Sample first 5 rows or all if fewer
        for i in range(sample_rows):
            row = table.rows[i]
            row_type = "HEADER" if i == 0 else ("SUBTOTAL" if i == len(table_data) - 1 else "BODY")
            height_pt = row.height.pt if hasattr(row.height, 'pt') else "Unknown"
            height_inches = row.height.inches if hasattr(row.height, 'inches') else "Unknown"
            logger.info(f"Row {i} ({row_type}): Height = {height_pt}pt ({height_inches:.3f}\")")
            
            # Report cell alignment for first few cells in this row
            for j in range(min(3, len(row.cells))):  # Sample first 3 cells
                cell = row.cells[j]
                v_anchor = cell.text_frame.vertical_anchor if hasattr(cell.text_frame, 'vertical_anchor') else "Unknown"
                h_align = cell.text_frame.paragraphs[0].alignment if cell.text_frame.paragraphs else "Unknown"
                logger.info(f"  Cell ({i},{j}): V_Anchor={v_anchor}, H_Align={h_align}")
        
        logger.info(f"=== END DIAGNOSTIC: Table Structure Analysis ===")
        """
        
        return True
        
    except Exception as e:
        logger.error(f"Error creating table: {str(e)}")
        logger.error(traceback.format_exc())
        return False

def _apply_table_cell_styling(cell, row_idx, col_idx, table_data, cell_metadata):
    """
    Apply styling to a table cell based on its position and content.
    CRITICAL FIX: Ensures ALL cells (including empty ones) get explicit Calibri font to eradicate "Roboto 18pt" defaults.
    Targets OXML structure from template for vertical centering and margins.
    
    Args:
        cell: The table cell to style
        row_idx: Row index (0-based)
        col_idx: Column index (0-based)
        table_data: Full table data for context
        cell_metadata: Dictionary with metadata for conditional formatting
    """
    # Constants for styling based on template analysis
    MARGIN_LEFT_RIGHT_PT = 3.6 # Approx 45720 EMU, for API calls
    # Top/Bottom margins will be omitted to match template behavior

    # CRITICAL FIX: Import qn and OxmlElement for direct OXML manipulation
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    try:
        # Import required styling constants (already at module level, but good for clarity if function was standalone)
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
        from pptx.util import Pt
        
        original_cell_text = str(table_data[row_idx][col_idx]) if col_idx < len(table_data[row_idx]) else ""
        logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Original text='{original_cell_text}'")

        # **CRITICAL STEP 1**: CELL-LEVEL VERTICAL ANCHOR (API Attempt)
        try:
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: API cell.vertical_anchor set to MIDDLE. Current: {cell.vertical_anchor}")
        except Exception as e_cell_anchor:
            logger.warning(f"CELL STYLING [{row_idx},{col_idx}]: Failed to set cell.vertical_anchor via API: {e_cell_anchor}. Will rely on OXML.")

        # **CRITICAL STEP 2**: CELL MARGINS (API Attempt - L/R only, T/B omitted)
        try:
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: API setting cell margins: L/R={MARGIN_LEFT_RIGHT_PT}pt, T/B=Default")
            cell.margin_left = Pt(MARGIN_LEFT_RIGHT_PT)
            cell.margin_right = Pt(MARGIN_LEFT_RIGHT_PT)
            # Explicitly DO NOT SET cell.margin_top and cell.margin_bottom via API
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Margins after API set: L={cell.margin_left.pt if cell.margin_left else 'Default'}pt, R={cell.margin_right.pt if cell.margin_right else 'Default'}pt, T(api)='Omitted', B(api)='Omitted'")
        except Exception as margin_error:
            logger.warning(f"CELL STYLING [{row_idx},{col_idx}]: Error setting L/R margins via API: {margin_error}. Will rely on OXML for L/R as well.")

        # **CRITICAL STEP 3**: TEXT FRAME BASIC SETUP
        text_frame = cell.text_frame
        text_frame.clear()  # Clear existing content AND PARAGRAPHS before adding new
        text_frame.word_wrap = False
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        
        # **NEW CRITICAL STEP 3.5**: SET ANCHOR ON BODYPR and ZERO TEXT_FRAME MARGINS
        try:
            # Explicitly set text_frame margins to 0 to influence lIns, tIns, rIns, bIns on bodyPr
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Text_frame margins set to 0. L={text_frame.margin_left}, T={text_frame.margin_top}, R={text_frame.margin_right}, B={text_frame.margin_bottom}")

            # Access bodyPr (should be created by setting text_frame margins if not present)
            # and forcefully set its anchor attribute.
            if hasattr(text_frame._txBody, 'bodyPr') and text_frame._txBody.bodyPr is not None:
                bodyPr = text_frame._txBody.bodyPr
            else:
                # If bodyPr somehow still doesn't exist, we might need to add it.
                # This is less common if margins were set, but as a fallback:
                from pptx.oxml.xmlchemy import OxmlElement
                from pptx.oxml.ns import qn
                bodyPr = OxmlElement("a:bodyPr")
                text_frame._txBody.insert_element_before(bodyPr, 'a:lstStyle', 'a:p') # Insert before list style or paragraphs
                logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: OXML bodyPr element created as it was missing.")

            bodyPr.set('anchor', 'ctr') # Anchor text body itself to center
            # EXPLICITLY ZERO OUT INSETS on bodyPr to remove any internal padding
            bodyPr.set('lIns', '0')
            bodyPr.set('tIns', '0')
            bodyPr.set('rIns', '0')
            bodyPr.set('bIns', '0')
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: OXML bodyPr attributes directly set: anchor='ctr', lIns='0', tIns='0', rIns='0', bIns='0'.")
            
            # Verify what the text_frame.vertical_anchor API reports now
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: text_frame.vertical_anchor after bodyPr.set: {text_frame.vertical_anchor}")

        except Exception as bodyPr_error:
            logger.error(f"CELL STYLING [{row_idx},{col_idx}]: Critical error setting OXML bodyPr or text_frame margins: {bodyPr_error} {traceback.format_exc()}")

        # **CRITICAL STEP 4**: OXML DIRECT SETTINGS (tcPr anchor and margins)
        try:
            tcPr = cell._tc.get_or_add_tcPr()
            
            # Ensure vertical anchor on tcPr is "ctr" by directly setting the attribute
            # tcPr.anchor = "ctr" # This was causing ValueError due to enum mismatch
            tcPr.set('anchor', 'ctr') # Direct attribute setting
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: OXML tcPr attribute 'anchor' directly set to 'ctr'.")

            # Add vert="horz" as seen in the template's OXML
            tcPr.set('vert', 'horz')
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: OXML tcPr attribute 'vert' directly set to 'horz'.")

            # Set specific L/R margins using MARGIN_EMU_LR
            tcPr.marL = MARGIN_EMU_LR
            tcPr.marR = MARGIN_EMU_LR
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: OXML tcPr.marL & tcPr.marR set to {MARGIN_EMU_LR}.")

            # Robustly remove marT and marB elements if they exist, using xpath
            marT_elements = tcPr.xpath('./a:marT') # xpath returns a list of matching elements
            for el_marT in marT_elements:
                tcPr.remove(el_marT)
                logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: OXML removed existing a:marT element.")
            
            marB_elements = tcPr.xpath('./a:marB')
            for el_marB in marB_elements:
                tcPr.remove(el_marB)
                logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: OXML removed existing a:marB element.")
            
            # Verify final OXML state for margins and anchor
            final_marL_oxml = tcPr.marL if hasattr(tcPr, 'marL') and tcPr.marL is not None else "NotSet"
            final_marR_oxml = tcPr.marR if hasattr(tcPr, 'marR') and tcPr.marR is not None else "NotSet"
            final_marT_oxml_exists = bool(tcPr.xpath('./a:marT'))
            final_marB_oxml_exists = bool(tcPr.xpath('./a:marB'))
            final_anchor_oxml = tcPr.anchor if hasattr(tcPr, 'anchor') and tcPr.anchor is not None else "NotSet"
            final_vert_oxml = tcPr.get('vert') if hasattr(tcPr, 'get') else "NotSet" # Check if .get exists
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Final OXML tcPr: anchor='{final_anchor_oxml}', vert='{final_vert_oxml}', marL='{final_marL_oxml}', marR='{final_marR_oxml}', marT_exists={final_marT_oxml_exists}, marB_exists={final_marB_oxml_exists}")

        except Exception as oxml_error:
            logger.error(f"CELL STYLING [{row_idx},{col_idx}]: Critical error setting OXML tcPr properties: {oxml_error} {traceback.format_exc()}")

        # **CRITICAL STEP 5**: PARAGRAPH PROPERTIES - ULTRA-COMPACT LINE SPACING FIX
        paragraphs_before_fix = len(text_frame.paragraphs)
        
        if text_frame.paragraphs:
            # Use the existing paragraph (left by text_frame.clear())
            p = text_frame.paragraphs[0]
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Reusing existing paragraph after clear(). Paragraphs count: {paragraphs_before_fix}")
        else:
            # Fallback: if no paragraphs exist (unexpected), add one
            p = text_frame.add_paragraph()
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: No paragraphs found after clear(), added new one. Paragraphs count: {len(text_frame.paragraphs)}")
        
        p.alignment = PP_ALIGN.CENTER
        
        # Get paragraph properties for direct OXML manipulation
        pPr = p._p.get_or_add_pPr()
        
        # CRITICAL FIX 1: Remove ALL existing spacing elements first
        for spacing_element in ['a:spcBef', 'a:spcAft', 'a:lnSpc']:
            existing = pPr.find(qn(spacing_element))
            if existing is not None:
                pPr.remove(existing)
                logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Removed existing {spacing_element}")
        
        # CRITICAL FIX 2: Remove default paragraph properties that might add space
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is not None:
            pPr.remove(defRPr)
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Removed default paragraph properties (a:defRPr)")
        
        # CRITICAL FIX 3: Set MINIMAL line spacing using EXACT point measurement
        # This is the key fix - use point-based line spacing that equals font size
        lnSpc = OxmlElement('a:lnSpc')
        spcPts_line = OxmlElement('a:spcPts')
        
        # Determine font size for this cell to set line spacing accordingly
        if row_idx == 0:  # Header row
            font_size_pt = FONT_SIZE_HEADER.pt  # 8pt
            line_spacing_pt = int(font_size_pt * 100)  # 800 (8pt in hundredths)
        else:  # Body/subtotal rows
            font_size_pt = FONT_SIZE_BODY.pt    # 7pt  
            line_spacing_pt = int(font_size_pt * 100)  # 700 (7pt in hundredths)
        
        spcPts_line.set('val', str(line_spacing_pt))  # Line height = font size
        lnSpc.append(spcPts_line)
        pPr.append(lnSpc)
        
        # CRITICAL FIX 4: Explicitly zero out before/after spacing
        spcBef = OxmlElement('a:spcBef')
        spcPts_before = OxmlElement('a:spcPts')
        spcPts_before.set('val', '0')
        spcBef.append(spcPts_before)
        pPr.append(spcBef)
        
        spcAft = OxmlElement('a:spcAft')
        spcPts_after = OxmlElement('a:spcPts')
        spcPts_after.set('val', '0')
        spcAft.append(spcPts_after)
        pPr.append(spcAft)
        
        logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: ULTRA-COMPACT spacing applied: font_size={font_size_pt}pt, line_spacing={line_spacing_pt/100}pt, spcBef=0, spcAft=0")
        
        # **CRITICAL STEP 6**: DETERMINE PROCESSED CELL TEXT & APPLY "-" RULE
        processed_cell_text = ""
        is_empty_cell = False  # Track if this is an empty cell that needs a dash
        
        if row_idx == 0:
            # HEADER ROW: Calibri 8pt Bold ALL CAPS
            if original_cell_text and str(original_cell_text).strip() and str(original_cell_text).strip() not in ["0", "0.0", "0.00", "0.000", "£0K", "0K", "-", "–", "0.0%"]:
                processed_cell_text = str(original_cell_text).upper()
            else:
                processed_cell_text = "-"  # Use dash for empty header cells
                is_empty_cell = True
            
        elif row_idx == len(table_data) - 1:
            # SUBTOTAL ROW: Calibri 7pt Bold ALL CAPS
            if original_cell_text and str(original_cell_text).strip() and str(original_cell_text).strip() not in ["0", "0.0", "0.00", "0.000", "£0K", "0K", "-", "–", "0.0%"]:
                processed_cell_text = str(original_cell_text).upper()
            else:
                processed_cell_text = "-"  # Use dash for empty subtotal cells
                is_empty_cell = True
            
        else:
            # BODY ROWS: Apply "-" rule for empty values in ALL columns
            if (not original_cell_text or 
                str(original_cell_text).strip() == "" or
                str(original_cell_text).strip() in ["0", "0.0", "0.00", "0.000", "£0K", "0K", "-", "–", "0.0%"]):
                processed_cell_text = "-"  # Use regular dash
                is_empty_cell = True
            else:
                # Process non-empty cells
                if col_idx == 0 and original_cell_text and str(original_cell_text).strip():
                    # Campaign name: FULL CAPS
                    processed_cell_text = str(original_cell_text).upper()
                elif col_idx < 7 and original_cell_text and str(original_cell_text).strip():
                    # Other first-seven columns: UPPER CASE
                    processed_cell_text = str(original_cell_text).upper()
                else:
                    # Month/quarter columns: keep as is
                    processed_cell_text = str(original_cell_text)
        
        # **CRITICAL STEP 7**: SET TEXT AND IMMEDIATELY APPLY EXPLICIT FONT FORMATTING
        p.text = processed_cell_text
        
        # **CRITICAL STEP 8**: ENSURE RUN EXISTS AND FORCE CALIBRI FONT
        # After setting text, PowerPoint may recreate runs, so we must ensure proper formatting
        if not p.runs:
            run = p.add_run()  # Create run if none exists
            run.text = processed_cell_text  # Ensure text is in the run
        else:
            run = p.runs[0]
        
        # **UNIVERSAL FONT ENFORCEMENT** - Apply to ALL cells to eradicate "Roboto 18pt"
        if row_idx == 0:
            # HEADER ROW: Calibri 8pt Bold
            run.font.name = DEFAULT_FONT_NAME  # FORCE Calibri
            run.font.size = FONT_SIZE_HEADER   # 8pt
            run.font.bold = True
            run.font.color.rgb = CLR_BLACK
            
        elif row_idx == len(table_data) - 1 or (row_idx > 0 and table_data[row_idx][0] in ['SUBTOTAL', 'CARRIED FORWARD']):
            # SUBTOTAL/TOTAL/CARRIED FORWARD ROW: Calibri 7pt Bold
            run.font.name = DEFAULT_FONT_NAME  # FORCE Calibri
            run.font.size = FONT_SIZE_BODY     # 7pt
            run.font.bold = True
            run.font.color.rgb = CLR_BLACK
            
        else:
            # BODY ROWS: Calibri 7pt, Bold for first 7 columns, Normal for months
            run.font.name = DEFAULT_FONT_NAME  # FORCE Calibri
            run.font.size = FONT_SIZE_BODY     # 7pt
            
            # Apply light grey color for empty cells, black for others
            if is_empty_cell:
                run.font.color.rgb = CLR_LIGHT_GRAY_TEXT  # Light grey for dashes
            else:
                run.font.color.rgb = CLR_BLACK
            
            if col_idx < 7:
                # First 7 columns: BOLD
                run.font.bold = True
            else:
                # Month/quarter columns: NORMAL weight
                run.font.bold = False
        
        # **CRITICAL STEP 9**: DOUBLE-CHECK AND RE-ENFORCE FONT PROPERTIES
        # This prevents any inheritance of unwanted defaults
        if run.font.name != DEFAULT_FONT_NAME:
            run.font.name = DEFAULT_FONT_NAME
            logger.debug(f"Re-enforced Calibri font for cell ({row_idx},{col_idx})")
        
        if not run.font.size:
            run.font.size = FONT_SIZE_BODY if row_idx != 0 else FONT_SIZE_HEADER
            logger.debug(f"Re-enforced font size for cell ({row_idx},{col_idx})")
        
        # **CRITICAL STEP 9.5**: RUN-LEVEL OXML FIX for BASELINE SPACING
        # This addresses potential baseline issues that could cause text to appear too low
        try:
            rPr = run._r.get_or_add_rPr()
            
            # Remove any default character properties that might affect baseline
            defRPr_run = rPr.find(qn('a:defRPr'))
            if defRPr_run is not None:
                rPr.remove(defRPr_run)
                logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Removed run-level default properties")
            
            # Ensure baseline is at normal position (remove any superscript/subscript offset)
            baseline = rPr.find(qn('a:baseline'))
            if baseline is not None:
                rPr.remove(baseline)
            
            # Add explicit baseline of 0 to ensure text sits at true baseline
            baseline = OxmlElement('a:baseline')
            baseline.set('val', '0')  # 0 = normal baseline position
            rPr.append(baseline)
            
            # Remove any character spacing adjustments that might affect positioning
            spc = rPr.find(qn('a:spc'))
            if spc is not None:
                rPr.remove(spc)
                logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Removed character spacing")
            
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Run-level baseline fix applied")
            
        except Exception as run_oxml_error:
            logger.warning(f"CELL STYLING [{row_idx},{col_idx}]: Error applying run-level OXML fixes: {run_oxml_error}")
        
        # **CRITICAL STEP 10**: APPLY CELL BACKGROUND COLORS BASED ON ROW TYPE AND POSITION
        if row_idx == 0:
            # Header row - light grey background for first 7 columns and quarterly columns, green for months
            if col_idx < 7:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_TABLE_GRAY
            elif col_idx in [10, 14, 18, 22]:  # Q1, Q2, Q3, Q4 header cells
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_TABLE_GRAY
            else:
                # Month columns - green background
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_HEADER_GREEN
        
        elif row_idx == len(table_data) - 1:
            # Total row - QA checklist subtotal grey #D9D9D9
            cell.fill.solid()
            cell.fill.fore_color.rgb = CLR_SUBTOTAL_GRAY
        
        elif table_data[row_idx][0] in ['SUBTOTAL', 'CARRIED FORWARD']:
            # Subtotal or carried forward rows - same styling as total row
            cell.fill.solid()
            cell.fill.fore_color.rgb = CLR_SUBTOTAL_GRAY
        
        else:
            # Data rows with conditional coloring
            cell_key = (row_idx, col_idx)
            
            if col_idx < 7:
                # First 7 columns: white background always
                cell.fill.background()
                
            elif col_idx in [10, 14, 18, 22]:  # Q1, Q2, Q3, Q4 columns
                # Quarter columns get same grey fill as TOTAL row
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_SUBTOTAL_GRAY
                
            else:
                # Month columns with conditional coloring
                if cell_key in cell_metadata:
                    cell_meta = cell_metadata[cell_key]
                    cell_value = cell_meta.get('value', 0)
                    has_data = cell_meta.get('has_data', False)
                    
                    try:
                        numeric_value_for_coloring = float(cell_value) if cell_value else 0 
                    except (ValueError, TypeError):
                        numeric_value_for_coloring = 0
                    
                    # FINAL FIX: Only use has_data flag since it now correctly reflects display status
                    if has_data:
                        # Apply media type color based on value
                        media_type = cell_meta['media_type']
                        
                        if media_type == "Television":
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = CLR_TELEVISION
                        elif media_type == "Digital":
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = CLR_DIGITAL
                        elif media_type == "OOH":
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = CLR_OOH
                        elif media_type == "GRPs":
                            # GRP cells always use TELEVISION color
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = CLR_TELEVISION
                        elif media_type in ["Reach@1+", "OTS@1+", "Reach@3+", "OTS@3+"]:
                            # New sub-rows use same green as GRP cells
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = CLR_TELEVISION
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = CLR_OTHER
                    else:
                        # No significant data: white background
                        cell.fill.background()
                else:
                    # No metadata: white background
                    cell.fill.background()
        
        # **FINAL ENFORCEMENT**: Re-apply ALL centering and spacing controls
        # This ensures nothing gets overridden by PowerPoint defaults
        # text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE # Old: text_frame level
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE # New: cell-level anchor
        logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Final enforcement - cell.vertical_anchor: {cell.vertical_anchor}")
        
        text_frame.word_wrap = False
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        
        # logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Final enforcement - text_frame.vertical_anchor: {text_frame.vertical_anchor}") # Old log

        # Apply final enforcement to ALL paragraphs in the cell
        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.line_spacing = 1.0
            # paragraph.space_before = Pt(0) # API call replaced by OXML
            # paragraph.space_after = Pt(0)  # API call replaced by OXML

            # OXML: Ensure zero paragraph spacing for final enforcement
            pPr_final = paragraph._p.get_or_add_pPr()
            spcBef_final = pPr_final.find(qn('a:spcBef'))
            if spcBef_final is not None:
                pPr_final.remove(spcBef_final)
            spcBef_final = OxmlElement('a:spcBef')
            spcPts_before_final = OxmlElement('a:spcPts')
            spcPts_before_final.set('val', '0')
            spcBef_final.append(spcPts_before_final)
            pPr_final.append(spcBef_final)

            spcAft_final = pPr_final.find(qn('a:spcAft'))
            if spcAft_final is not None:
                pPr_final.remove(spcAft_final)
            spcAft_final = OxmlElement('a:spcAft')
            spcPts_after_final = OxmlElement('a:spcPts')
            spcPts_after_final.set('val', '0')
            spcAft_final.append(spcPts_after_final)
            pPr_final.append(spcAft_final)
            
            logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Final enforcement - Paragraph {para_idx} properties: Alignment={paragraph.alignment}, LineSpacing={paragraph.line_spacing}. OXML spcBef/spcAft forced to 0.")
            
            # Ensure all runs in all paragraphs have proper font (final check)
            for run_idx, cell_run in enumerate(paragraph.runs):
                # Determine expected font properties based on row_idx and col_idx again
                expected_font_name = DEFAULT_FONT_NAME
                expected_bold = False
                expected_color_rgb = CLR_BLACK # Default to black

                if row_idx == 0: # Header
                    expected_font_size = FONT_SIZE_HEADER
                    expected_bold = True
                elif row_idx == len(table_data) - 1: # Subtotal
                    expected_font_size = FONT_SIZE_BODY 
                    expected_bold = True
                else: # Body
                    expected_font_size = FONT_SIZE_BODY
                    if col_idx < 5: # First 5 columns bold
                        expected_bold = True
                    
                    # Color for dashes
                    # Re-check if the processed_cell_text (which should be in the run now) is a dash
                    if cell_run.text == "-": # Check the actual run text
                         expected_color_rgb = CLR_LIGHT_GRAY_TEXT
                
                if cell_run.font.name != expected_font_name:
                    cell_run.font.name = expected_font_name
                    logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Final run {run_idx} font name re-enforced to {expected_font_name}")
                if cell_run.font.size != expected_font_size:
                    cell_run.font.size = expected_font_size
                    logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Final run {run_idx} font size re-enforced to {expected_font_size}")
                if cell_run.font.bold != expected_bold:
                    cell_run.font.bold = expected_bold
                    logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Final run {run_idx} font bold re-enforced to {expected_bold}")
                
                # Check and set color, carefully
                current_run_color_rgb = None
                if cell_run.font.color.type == MSO_COLOR_TYPE.RGB:
                    current_run_color_rgb = cell_run.font.color.rgb
                
                if current_run_color_rgb != expected_color_rgb:
                    cell_run.font.color.rgb = expected_color_rgb
                    logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Final run {run_idx} font color re-enforced to RGB({expected_color_rgb.r},{expected_color_rgb.g},{expected_color_rgb.b})")

        # One last time for the text_frame itself
        # text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE # Old: text_frame level
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE # New: cell-level anchor
        logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Absolutely final re-assertion of cell.vertical_anchor: {cell.vertical_anchor}")
        
        # **CRITICAL FINAL STEP**: Clean up any stray empty paragraphs and log final count
        paragraphs_final_count = len(text_frame.paragraphs)
        if paragraphs_final_count > 1:
            # Remove any extra empty paragraphs beyond the first one
            paragraphs_to_remove = []
            for i in range(1, len(text_frame.paragraphs)):
                para = text_frame.paragraphs[i]
                if not para.text.strip():  # Empty paragraph
                    paragraphs_to_remove.append(i)
            
            # Remove from highest index to lowest to avoid index shifting
            for para_idx in reversed(paragraphs_to_remove):
                # Note: Direct paragraph removal from text_frame.paragraphs is not supported
                # This is a safety check - the fix should prevent multiple paragraphs
                logger.warning(f"CELL STYLING [{row_idx},{col_idx}]: Found {len(paragraphs_to_remove)} extra empty paragraphs - this should not happen with the reuse fix")
        
        logger.debug(f"CELL STYLING [{row_idx},{col_idx}]: Final paragraph count: {len(text_frame.paragraphs)} (should be 1 for pixel-perfect alignment)")
        
        logger.debug(f"UNIVERSAL FORMATTING: Applied explicit Calibri font, centering, and spacing to cell ({row_idx},{col_idx}) - Text: '{processed_cell_text}'")
        
    except Exception as e:
        logger.error(f"Error styling cell ({row_idx},{col_idx}): {e}")
        logger.error(traceback.format_exc())

def _apply_table_borders(table):
    """
    Apply consistent borders to all table cells.
    Uses #BFBFBF color with 0.75pt width for all internal and external borders.
    
    Args:
        table: The PowerPoint table object to apply borders to
    """
    try:
        from pptx.oxml.ns import qn
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        
        # Border specifications from QA checklist
        border_color = CLR_TABLE_GRAY  # #BFBFBF (191, 191, 191)
        border_width = Pt(0.75)       # 0.75pt width
        
        logger.debug(f"Applying borders to table with {len(table.rows)} rows and {len(table.columns)} columns")
        
        # Apply borders to all cells
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                try:
                    # Access the cell's border properties
                    tc = cell._tc  # Access the underlying table cell element
                    tcPr = tc.get_or_add_tcPr()  # Table cell properties
                    
                    # Create border elements if they don't exist
                    borders = ['top', 'left', 'bottom', 'right']
                    
                    for border_name in borders:
                        border_element_name = f'{border_name}Border'
                        border_element = tcPr.find(qn(f'a:{border_element_name}'))
                        
                        if border_element is None:
                            border_element = tcPr.add(qn(f'a:{border_element_name}'))
                        
                        # Set border width (in EMUs - English Metric Units)
                        # 0.75pt = 0.75 * 12700 EMUs
                        border_width_emu = int(0.75 * 12700)
                        border_element.set('w', str(border_width_emu))
                        
                        # Set border color
                        solidFill = border_element.find(qn('a:solidFill'))
                        if solidFill is None:
                            solidFill = border_element.add(qn('a:solidFill'))
                        
                        srgbClr = solidFill.find(qn('a:srgbClr'))
                        if srgbClr is None:
                            srgbClr = solidFill.add(qn('a:srgbClr'))
                        
                        # Convert RGB to hex format
                        hex_color = f"{border_color.r:02X}{border_color.g:02X}{border_color.b:02X}"
                        srgbClr.set('val', hex_color)
                        
                except Exception as cell_border_error:
                    logger.debug(f"Could not apply borders to cell ({row_idx}, {col_idx}): {cell_border_error}")
                    # Fallback: try using the newer border API if available
                    try:
                        from pptx.table import _Cell
                        if hasattr(cell, 'border'):
                            # This is for newer python-pptx versions
                            cell.border.top.color.rgb = border_color
                            cell.border.left.color.rgb = border_color  
                            cell.border.bottom.color.rgb = border_color
                            cell.border.right.color.rgb = border_color
                            cell.border.top.width = border_width
                            cell.border.left.width = border_width
                            cell.border.bottom.width = border_width
                            cell.border.right.width = border_width
                    except Exception as fallback_error:
                        logger.debug(f"Fallback border method also failed for cell ({row_idx}, {col_idx}): {fallback_error}")
        
        logger.info(f"Table borders applied successfully with #BFBFBF color and 0.75pt width")
        return True
        
    except Exception as e:
        logger.warning(f"Error applying table borders: {str(e)}")
        return False

def _ensure_font_consistency(font, target_font_name, target_size, target_bold, target_color):
    """
    Ensure font consistency by applying target properties.
    
    Args:
        font: The font object to modify
        target_font_name: Target font name (e.g., 'Calibri')
        target_size: Target font size (e.g., Pt(7))
        target_bold: Target bold setting (True/False)
        target_color: Target color (RGB)
    """
    try:
        if target_font_name:
            font.name = target_font_name
        if target_size:
            font.size = target_size
        if target_bold is not None:
            font.bold = target_bold
        if target_color:
            font.color.rgb = target_color
    except Exception as e:
        logger.debug(f"Error setting font properties: {e}")

def _prepare_main_table_data(df, region, masterbrand):
    """
    Prepare table data and cell metadata for the main data table.
    
    Args:
        df: DataFrame with the Excel data
        region: Region filter
        masterbrand: Masterbrand filter
        
    Returns:
        tuple: (table_data, cell_metadata) or (None, None) if no data
    """
    try:
        # Filter data for this region/masterbrand combination
        filtered_df = df[
            (df['Region'] == region) &
            (df['Global Masterbrand'] == masterbrand)
        ].copy()
        
        if filtered_df.empty:
            logger.warning(f"No data found for {region} - {masterbrand}")
            return None, None
        
        # Prepare table structure
        table_data = []
        cell_metadata = {}
        
        # Header row
        header_row = ['Campaign Name', 'Budget', 'TV GRPs', 'TOTAL REACH', 'TOTAL FREQ', '%', 'Media Type'] + \
                    ['JAN', 'FEB', 'MAR', 'APR', 'MAY', 'JUN', 
                     'JUL', 'AUG', 'SEP', 'OCT', 'NOV', 'DEC'] + \
                    ['Q1', 'Q2', 'Q3', 'Q4']
        table_data.append(header_row)
        
        # Process each row of data
        row_idx = 1
        for _, row in filtered_df.iterrows():
            # Build data row
            campaign_name = str(row.get('Campaign Name', ''))
            budget = f"${row.get('Budget', 0):,.0f}" if pd.notna(row.get('Budget')) else "$0"
            tv_grps = str(row.get('TV GRPs', '')) if pd.notna(row.get('TV GRPs')) else ""
            percentage = f"{row.get('Percentage', 0):.1f}%" if pd.notna(row.get('Percentage')) else "0.0%"
            media_type = str(row.get('Media Type', ''))
            
            data_row = [campaign_name, budget, tv_grps, percentage, media_type]
            
            # Add monthly data
            months = ['JAN', 'FEB', 'MAR', 'APR', 'MAY', 'JUN', 
                     'JUL', 'AUG', 'SEP', 'OCT', 'NOV', 'DEC']
            
            for col_idx, month in enumerate(months, start=5):
                month_value = row.get(month, 0)
                if pd.notna(month_value) and month_value != 0:
                    formatted_value = f"${month_value:,.0f}"
                    data_row.append(formatted_value)
                    # Store metadata for conditional formatting - align with display
                    cell_metadata[(row_idx, col_idx)] = {
                        'value': month_value,
                        'media_type': media_type,
                        'has_data': not is_empty_formatted_value(formatted_value)
                    }
                else:
                    data_row.append("")
                    cell_metadata[(row_idx, col_idx)] = {
                        'value': 0,
                        'media_type': media_type,
                        'has_data': False
                    }
            
            # Add quarterly data
            quarters = ['Q1', 'Q2', 'Q3', 'Q4']
            for col_idx, quarter in enumerate(quarters, start=17):
                quarter_value = row.get(quarter, 0)
                if pd.notna(quarter_value) and quarter_value != 0:
                    data_row.append(f"${quarter_value:,.0f}")
                else:
                    data_row.append("–")
            
            table_data.append(data_row)
            row_idx += 1
        
        # Add totals row
        total_budget = filtered_df['Budget'].sum() if 'Budget' in filtered_df.columns else 0
        totals_row = [
            'TOTAL',
            f"${total_budget:,.0f}",
            '',
            '',
            ''  # Media type blank for total row
        ] + \
        [f"${month_total:,.0f}" if month_total > 0 else "" for month_total in [filtered_df[month].sum() for month in months]] + \
        [f"${quarter_total:,.0f}" if quarter_total > 0 else "–" for quarter_total in [filtered_df[quarter].sum() for quarter in quarters]]
        
        table_data.append(totals_row)
        
        logger.info(f"Prepared table data with {len(table_data)} rows for {region} - {masterbrand}")
        return table_data, cell_metadata
        
    except Exception as e:
        logger.error(f"Error preparing table data: {str(e)}")
        logger.error(traceback.format_exc())
        return None, None

def _prepare_funnel_chart_data(df, region, masterbrand):
    """
    Prepare data for the funnel chart.
    
    Args:
        df: DataFrame with the Excel data
        region: Region filter
        masterbrand: Masterbrand filter
        
    Returns:
        list: Chart data or None if no data
    """
    try:
        # Filter data for this region/masterbrand combination
        filtered_df = df[
            (df['Region'] == region) &
            (df['Global Masterbrand'] == masterbrand)
        ].copy()
        
        if filtered_df.empty:
            return None
        
        # Group by funnel stage and sum budgets
        if 'Funnel Stage' in filtered_df.columns and 'Budget' in filtered_df.columns:
            funnel_data = filtered_df.groupby('Funnel Stage')['Budget'].sum().to_dict()
            return [(stage, budget) for stage, budget in funnel_data.items() if budget > 0]
        
        return None
        
    except Exception as e:
        logger.error(f"Error preparing funnel chart data: {str(e)}")
        return None

def _prepare_media_type_chart_data(df, region, masterbrand):
    """
    Prepare data for the media type chart.
    
    Args:
        df: DataFrame with the Excel data
        region: Region filter
        masterbrand: Masterbrand filter
        
    Returns:
        list: Chart data or None if no data
    """
    try:
        # Filter data for this region/masterbrand combination
        filtered_df = df[
            (df['Region'] == region) &
            (df['Global Masterbrand'] == masterbrand)
        ].copy()
        
        if filtered_df.empty:
            return None
        
        # Group by media type and sum budgets
        if 'Media Type' in filtered_df.columns and 'Budget' in filtered_df.columns:
            media_data = filtered_df.groupby('Media Type')['Budget'].sum().to_dict()
            return [(media_type, budget) for media_type, budget in media_data.items() if budget > 0]
        
        return None
        
    except Exception as e:
        logger.error(f"Error preparing media type chart data: {str(e)}")
        return None

def _prepare_brand_chart_data(df, region, masterbrand):
    """
    Prepare data for the brand chart.
    
    Args:
        df: DataFrame with the Excel data
        region: Region filter
        masterbrand: Masterbrand filter
        
    Returns:
        list: Chart data or None if no data
    """
    try:
        # Filter data for this region/masterbrand combination
        filtered_df = df[
            (df['Region'] == region) &
            (df['Global Masterbrand'] == masterbrand)
        ].copy()
        
        if filtered_df.empty:
            return None
        
        # Group by brand and sum budgets
        if 'Brand' in filtered_df.columns and 'Budget' in filtered_df.columns:
            brand_data = filtered_df.groupby('Brand')['Budget'].sum().to_dict()
            return [(brand, budget) for brand, budget in brand_data.items() if budget > 0]
        
        return None
        
    except Exception as e:
        logger.error(f"Error preparing brand chart data: {str(e)}")
        return None

def _add_pie_chart(slide, chart_data, chart_title, position_info, chart_name=None):
    """
    Add a pie chart to the slide at the specified position.
    
    Args:
        slide: The PowerPoint slide to add the chart to
        chart_data: Dictionary with labels as keys and values as values
        chart_title: Title for the chart
        position_info: Dictionary with 'left', 'top', 'width', 'height' in inches
        chart_name: Optional name for the chart shape (for template compliance)
        
    Returns:
        bool: True if chart was created successfully, False otherwise
    """
    try:
        if not chart_data:
            logger.warning(f"No chart data provided for {chart_title}")
            return False
        
        logger.info(f"Creating pie chart: {chart_title}")
        
        # Calculate total for percentages
        total_value = sum(chart_data.values())
        if total_value <= 0:
            logger.warning(f"Invalid total value for chart {chart_title}: {total_value}")
            return False
        
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = list(chart_data.keys())
        chart_data_obj.add_series('Budget', list(chart_data.values()))
        
        # Position and size (position_info already contains Inches objects from get_element_position)
        left = position_info['left']
        top = position_info['top']
        width = position_info['width']
        height = position_info['height']
        
        # Add chart to slide
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, left, top, width, height, chart_data_obj
        )
        
        # CRITICAL FIX: Set proper shape name to match template (fixes "Chart 14" issue)
        if chart_name:
            chart_shape.name = chart_name
            logger.debug(f"Chart shape name set to: {chart_name}")
        
        chart = chart_shape.chart
        
        # Set chart title
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        
        # PIXEL-PERFECT: Style the title with Calibri 8pt non-bold (smaller than before)
        title_font = chart.chart_title.text_frame.paragraphs[0].runs[0].font
        
        # PIXEL-PERFECT: Use enhanced font consistency enforcement for smaller chart titles
        _ensure_font_consistency(
            title_font, 
            target_font_name=DEFAULT_FONT_NAME,
            target_size=FONT_SIZE_CHART_TITLE,  # 8pt for chart titles
            target_bold=False,  # NON-BOLD for pixel-perfect compliance
            target_color=CLR_BLACK
        )
        
        # FINAL QA: Disable AutoFit for chart title
        chart.chart_title.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        
        # Configure chart appearance
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        # PIXEL-PERFECT: Enhanced legend font consistency - smaller font
        try:
            legend_font = chart.legend.font
            _ensure_font_consistency(
                legend_font,
                target_font_name=DEFAULT_FONT_NAME,
                target_size=FONT_SIZE_CHART_LABELS,  # 6pt for legend
                target_bold=False,
                target_color=CLR_BLACK
            )
            logger.debug(f"Chart legend font consistency applied for '{chart_title}'")
        except Exception as legend_font_error:
            logger.debug(f"Could not apply legend font consistency: {legend_font_error}")
        
        # Apply colors based on data type
        series = chart.series[0]
        
        # Color mapping for different chart types
        color_mapping = {
            # Media type colors
            'Television': CLR_TELEVISION,
            'Digital': CLR_DIGITAL, 
            'OOH': CLR_OOH,
            'Other': CLR_OTHER,
            # Funnel stage colors (using existing palette)
            'Awareness': CLR_TELEVISION,
            'Consideration': CLR_BLACK,
            'Purchase': CLR_OOH,
            # Campaign type colors (using existing palette)
            'Always On': CLR_TELEVISION,
            'Brand': CLR_DIGITAL,
            'Product': CLR_OOH
        }
        
        # Apply colors to pie slices
        for i, (label, value) in enumerate(chart_data.items()):
            if i < len(series.points):
                point = series.points[i]
                
                # Get color for this label, default to a rotation of existing colors
                if label in color_mapping:
                    color = color_mapping[label]
                else:
                    # Default color rotation
                    colors = [CLR_TELEVISION, CLR_DIGITAL, CLR_OOH, CLR_OTHER]
                    color = colors[i % len(colors)]
                
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = color
        
        # Enable data labels with percentages
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.show_percentage = True
        data_labels.show_value = False
        
        # PIXEL-PERFECT: Set percentage format to one decimal place (e.g., "45.2%")
        try:
            data_labels.number_format = "0.0%"  # One decimal place format
            logger.debug(f"Applied one decimal place format to chart data labels for '{chart_title}'")
        except Exception as format_error:
            logger.debug(f"Could not set decimal format for data labels: {format_error}")
        
        # PRIORITY 3: Enhanced data labels font consistency
        try:
            data_labels_font = data_labels.font
            _ensure_font_consistency(
                data_labels_font,
                target_font_name=DEFAULT_FONT_NAME,
                target_size=FONT_SIZE_CHART_LABELS,  # 6pt for data labels
                target_bold=False,
                target_color=CLR_BLACK
            )
            logger.debug(f"Chart data labels font consistency applied for '{chart_title}'")
        except Exception as data_labels_font_error:
            logger.debug(f"Could not apply data labels font consistency: {data_labels_font_error}")
        
        data_labels.font.name = DEFAULT_FONT_NAME  # Force Calibri
        data_labels.font.size = FONT_SIZE_CHART_LABELS    # 6pt as per FINAL QA specification (was 7pt)
        data_labels.font.color.rgb = CLR_BLACK
        data_labels.font.bold = False
        
        logger.info(f"Chart '{chart_title}' created successfully with {len(chart_data)} data points")
        return True
        
    except Exception as e:
        logger.error(f"Error creating pie chart '{chart_title}': {str(e)}")
        logger.error(traceback.format_exc())
        return False

def _populate_slide_content(new_slide, prs, combination_row, slide_title_suffix, 
                          split_table_data, split_metadata, split_idx, df, excel_path):
    """Populate a single slide with all content (title, table, charts, comments)."""
    
    # Copy Title Placeholder from template FIRST
    template_title_shape = _get_shape_by_name(prs.slides[0], SHAPE_NAME_TITLE)
    if template_title_shape:
        title_text = f"{combination_row[0]} – {combination_row[1]} ({combination_row[2]}){slide_title_suffix}"
        copied_title_shape = _copy_text_box(template_title_shape, new_slide, new_name=SHAPE_NAME_TITLE, new_text=title_text)
        if copied_title_shape:
            # QA FIX: Ensure correct font for title
            if copied_title_shape.has_text_frame and copied_title_shape.text_frame.paragraphs:
                for paragraph in copied_title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # PRIORITY 3: Enhanced font consistency for titles
                        _ensure_font_consistency(
                            run.font,
                            target_font_name=DEFAULT_FONT_NAME,
                            target_size=FONT_SIZE_TITLE,
                            target_bold=True,
                            target_color=CLR_WHITE
                        )
            logger.info(f"Title copied and populated for slide.")
    
    # Copy Static Text Elements (Comments Title and Box)
    template_comments_title_shape = _get_shape_by_name(prs.slides[0], SHAPE_NAME_COMMENTS_TITLE)
    if template_comments_title_shape:
        copied_comments_title = _copy_text_box(template_comments_title_shape, new_slide, new_name=SHAPE_NAME_COMMENTS_TITLE, new_text="COMMENTS")
        if copied_comments_title:
            # Position comments title using centralized coordinates
            title_pos = get_element_position('comments_title')
            if title_pos:
                copied_comments_title.left = title_pos['left']
                copied_comments_title.top = title_pos['top']
                copied_comments_title.width = title_pos['width']
                copied_comments_title.height = title_pos['height']
    
    # Copy Comments Box
    template_comments_box_shape = _get_shape_by_name(prs.slides[0], SHAPE_NAME_COMMENTS_BOX)
    if template_comments_box_shape:
        copied_comments_box = _copy_text_box(template_comments_box_shape, new_slide, new_name=SHAPE_NAME_COMMENTS_BOX, new_text="")
        if copied_comments_box:
            # Position comments box using centralized coordinates
            box_pos = get_element_position('comments_box')
            if box_pos:
                copied_comments_box.left = box_pos['left']
                copied_comments_box.top = box_pos['top']
                copied_comments_box.width = box_pos['width']
                copied_comments_box.height = box_pos['height']
    
    logger.info("Skipping legend copying - legend elements are part of slide master")
    
    # Create and populate the main data table
    logger.info(f"Creating table for {combination_row[0]} - {combination_row[1]} - {combination_row[2]}{slide_title_suffix}")
    
    # Use the split table data
    table_success = _add_and_style_table(new_slide, split_table_data, split_metadata, prs.slides[0])
    if table_success:
        logger.info(f"Table created successfully for slide")
    else:
        logger.warning(f"Failed to create table for slide")
    
    # Create and add the three pie charts
    if SHOW_CHARTS_ON_SPLITS == "all" or (split_idx == 0 and SHOW_CHARTS_ON_SPLITS == "first_only"):
        logger.info(f"Creating charts for {combination_row[0]} - {combination_row[1]} - {combination_row[2]}")
        
        # Chart positioning using centralized coordinate system
        chart_positions = [
            get_element_position('chart_1'),
            get_element_position('chart_2'),
            get_element_position('chart_3')
        ]
        
        # Validate chart positions
        if not all(chart_positions):
            logger.error("Failed to get chart position coordinates")
            chart_positions = []
        
        # 1. Funnel Chart
        funnel_data = _prepare_funnel_chart_data_detailed(df, combination_row[0], combination_row[1], combination_row[2])
        if funnel_data:
            funnel_success = _add_pie_chart(new_slide, funnel_data, "BUDGET BY FUNNEL STAGE", chart_positions[0], SHAPE_NAME_FUNNEL_CHART)
            if funnel_success:
                logger.info(f"Funnel chart created successfully for slide")
        
        # 2. Media Type Chart  
        media_type_data = _prepare_media_type_chart_data_detailed(df, combination_row[0], combination_row[1], combination_row[2])
        if media_type_data:
            media_success = _add_pie_chart(new_slide, media_type_data, "BUDGET BY MEDIA TYPE", chart_positions[1], SHAPE_NAME_MEDIA_TYPE_CHART)
            if media_success:
                logger.info(f"Media type chart created successfully for slide")
        
        # 3. Campaign Type Chart
        campaign_type_data = _prepare_campaign_type_chart_data(df, combination_row[0], combination_row[1], combination_row[2])
        if campaign_type_data:
            campaign_success = _add_pie_chart(new_slide, campaign_type_data, "BUDGET BY CAMPAIGN TYPE", chart_positions[2], SHAPE_NAME_CAMPAIGN_TYPE_CHART)
            if campaign_success:
                logger.info(f"Campaign type chart created successfully for slide")
                
        logger.info(f"Chart creation completed for slide")

def create_presentation(template_path, excel_path, output_path):
    """Creates a PowerPoint presentation based on a template and Excel data."""
    logger.info(f"Starting presentation creation using template: {template_path}")
    try:
        prs = Presentation(template_path)
        if not prs.slides:
            logger.error("Template presentation has no slides. Cannot proceed.")
            return False

        logger.debug("--- Available Slide Layouts in Template ---")
        for i, layout in enumerate(prs.slide_layouts):
            layout_name = layout.name if hasattr(layout, 'name') else 'Unknown Name'
            placeholder_count = len(layout.placeholders) if hasattr(layout, 'placeholders') else 0
            logger.debug(f"  Layout Index {i}: {layout_name} (Placeholders: {placeholder_count})")
        logger.debug("--- End Available Slide Layouts ---")

        # Log template slide structure
        logger.debug("--- Template Slide (First Slide) Structure ---")
        if prs.slides:
            template_slide = prs.slides[0]
            logger.debug(f"  Template slide has {len(template_slide.shapes)} shapes:")
            for i, shape in enumerate(template_slide.shapes):
                shape_name = shape.name if hasattr(shape, 'name') else 'Unnamed'
                shape_type = shape.shape_type if hasattr(shape, 'shape_type') else 'Unknown Type'
                logger.debug(f"    Shape {i}: Name='{shape_name}', Type={shape_type}")
        else:
            logger.debug("  No slides found in template!")
        logger.debug("--- End Template Slide Structure ---")

        df = load_and_prepare_data(excel_path)
        if df is None or df.empty:
            logger.error("Failed to load or prepare data. Aborting presentation creation.")
            return False

        # The load_and_prepare_data() function returns a processed DataFrame with these column names
        country_col_name = 'Country'
        brand_col_name = 'Brand'
        year_col_name = 'Year'
        
        unique_combinations = df[[country_col_name, brand_col_name, year_col_name]].drop_duplicates().values.tolist()
        logger.info(f"Found {len(unique_combinations)} unique Country/Global Masterbrand/Year combinations.")

        # Calculate total investment for each combination
        combinations_with_investment = []
        for combination in unique_combinations:
            country, brand, year = combination
            # Filter data for this specific combination
            combination_data = df[
                (df[country_col_name] == country) & 
                (df[brand_col_name] == brand) & 
                (df[year_col_name] == year)
            ]
            # Calculate total investment (sum of Total Cost for this combination)
            total_investment = combination_data['Total Cost'].fillna(0).sum()
            
            combinations_with_investment.append((country, brand, year, total_investment))
        
        # Group by market (country) and calculate total market investment
        market_investments = {}
        for country, brand, year, investment in combinations_with_investment:
            if country not in market_investments:
                market_investments[country] = {'total': 0, 'combinations': []}
            market_investments[country]['total'] += investment
            market_investments[country]['combinations'].append((country, brand, year, investment))
        
        # Sort markets by total investment (highest first)
        sorted_markets = sorted(market_investments.items(), key=lambda x: x[1]['total'], reverse=True)
        
        # Build final sorted list: markets ordered by total, brands within market ordered by individual investment
        unique_combinations = []
        for market, data in sorted_markets:
            # Sort combinations within this market by individual investment
            market_combos = sorted(data['combinations'], key=lambda x: x[3], reverse=True)
            # Add to final list (without investment values)
            unique_combinations.extend([(c, b, y) for c, b, y, _ in market_combos])
        
        # Log the sorted order with market totals
        logger.info("Slides will be generated grouped by market, ordered by total market investment:")
        for i, (market, data) in enumerate(sorted_markets[:10]):  # Log top 10 markets
            logger.info(f"  {i+1}. {market}: £{data['total']:,.0f} total")
            for combo in sorted(data['combinations'], key=lambda x: x[3], reverse=True)[:3]:  # Show top 3 brands
                logger.info(f"      - {combo[1]} ({combo[2]}): £{combo[3]:,.0f}")
            if len(data['combinations']) > 3:
                logger.info(f"      ... and {len(data['combinations']) - 3} more brands")
        if len(sorted_markets) > 10:
            logger.info(f"  ... and {len(sorted_markets) - 10} more markets")

        if not unique_combinations:
            logger.warning("No unique Country/Global Masterbrand combinations found in the data.")
            # Decide if an empty presentation should be saved or an error returned
            # For now, let's save an empty presentation (after removing template slide)

        current_market = None
        for idx, combination_row in enumerate(unique_combinations):
            # Check if we're starting a new market
            if combination_row[0] != current_market:
                current_market = combination_row[0]
                # Fix market name display
                display_market_name = "Morocco" if current_market == "MOR" else current_market
                # Add a market delimiter slide with black background
                try:
                    # Try to use a blank slide layout 
                    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
                    delimiter_slide = prs.slides.add_slide(blank_layout)
                except:
                    delimiter_slide = prs.slides.add_slide(prs.slide_layouts[0])
                
                # Add full black background
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Create a black rectangle covering the entire slide
                from pptx.enum.shapes import MSO_SHAPE
                black_bg = delimiter_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=0,
                    top=0,
                    width=slide_width,
                    height=slide_height
                )
                
                # Make the background solid black
                black_bg.fill.solid()
                black_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure black
                black_bg.line.fill.background()  # Remove border
                
                # Create a text box in the center of the slide for the market name
                text_box = delimiter_slide.shapes.add_textbox(
                    left=Inches(1),
                    top=int(slide_height * 0.45),  # Center vertically
                    width=slide_width - Inches(2),
                    height=Inches(1.5)
                )
                
                text_frame = text_box.text_frame
                text_frame.text = str(display_market_name).upper() if display_market_name else "UNKNOWN"  # Make market name full caps
                text_frame.word_wrap = False
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                # Format the text - large white font
                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(36)  # Large font (30-40 range)
                        run.font.bold = True
                        run.font.name = DEFAULT_FONT_NAME
                        run.font.color.rgb = RGBColor(255, 255, 255)  # Pure white
                logger.info(f"Added market delimiter slide for: {display_market_name}")
            
            logger.info(f"Processing combination {idx+1}/{len(unique_combinations)}: {combination_row[0]} - {combination_row[1]} - {combination_row[2]}")
            
            # First, prepare the table data to check if splitting is needed
            table_result = _prepare_main_table_data_detailed(df, combination_row[0], combination_row[1], combination_row[2], excel_path)
            
            if table_result[0] is None:
                logger.warning(f"No table data generated for {combination_row[0]} - {combination_row[1]}")
                continue
            
            table_data, cell_metadata = table_result
            
            # Split table if needed
            table_splits = _split_table_data_by_campaigns(table_data, cell_metadata)
            
            # Create a slide for each split
            for split_idx, (split_table_data, split_metadata, is_continuation) in enumerate(table_splits):
                # Add slide number to title if there are multiple splits
                if len(table_splits) > 1:
                    slide_title_suffix = f" ({split_idx + 1} of {len(table_splits)})"
                else:
                    slide_title_suffix = ""
                
                new_slide = prs.slides.add_slide(prs.slide_layouts[0])
                # logger.debug(f"Added new slide for {combination_row[0]} - {combination_row[1]} - {combination_row[2]}{slide_title_suffix}")
                
                # Populate this slide with content immediately
                _populate_slide_content(
                    new_slide, prs, combination_row, slide_title_suffix, 
                    split_table_data, split_metadata, split_idx, df, excel_path
                )

                # Diagnostic: Log all shapes on the new slide (commented out to reduce log size)
            # This entire block is commented out to reduce log file size
            # Re-enable for debugging shape-related issues
            """
            logger.debug(f"--- Shapes on new_slide (Layout: {prs.slide_layouts[0].name}, Index: {idx+1}) ---")
            if hasattr(new_slide, 'shapes'):
                logger.debug(f"  new_slide.shapes attribute exists. Number of shapes found: {len(new_slide.shapes)}")
                if new_slide.shapes:
                    for i, shape in enumerate(new_slide.shapes):
                        shape_name = shape.name if hasattr(shape, 'name') else 'Unnamed Shape'
                        shape_type = shape.shape_type if hasattr(shape, 'shape_type') else 'Unknown Type'
                        is_placeholder = shape.is_placeholder if hasattr(shape, 'is_placeholder') else False
                        placeholder_type = None
                        placeholder_idx = None
                        if is_placeholder and hasattr(shape, 'placeholder_format'):
                            ph_format = shape.placeholder_format
                            if hasattr(ph_format, 'type'):
                                placeholder_type = ph_format.type
                            if hasattr(ph_format, 'idx'):
                                placeholder_idx = ph_format.idx
                        logger.debug(f"  Shape {i}: Name='{shape_name}', Type={shape_type}, IsPlaceholder={is_placeholder}, PlaceholderType={placeholder_type}, PlaceholderIdx={placeholder_idx}")
                else:
                    logger.debug("  new_slide.shapes collection is empty (it exists but len is 0).")
            else:
                logger.debug("  new_slide.shapes attribute is missing.")
            logger.debug(f"--- End Shapes on new_slide (Index: {idx+1}) ---")
                """

        # ORIGINAL CONTENT POPULATION CODE MOVED TO _populate_slide_content() FUNCTION
        # (This section was deleted to prevent empty slides)
        # Remove the original template slide (slide 1) which shows "Region / Global Masterbrand"
        if len(prs.slides) > 1: # Only remove if we have other slides
            logger.info(f"Removing template slide from presentation (total slides: {len(prs.slides)}).")
            try:
                # The template slide is always the first slide added (index 0)
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
                logger.info("Template slide removed successfully.")
            except Exception as e:
                logger.warning(f"Could not remove template slide: {e}")
        else:
            logger.warning(f"Not removing template slide - only {len(prs.slides)} slides found.")

        # CRITICAL FIX: Ensure output directory exists and add proper save error handling
        from pathlib import Path
        output_path_obj = Path(output_path)
        
        # Ensure output directory exists
        output_path_obj.parent.mkdir(parents=True, exist_ok=True)
        logger.info(f"Output directory ensured: {output_path_obj.parent}")
        
        # Ensure .pptx extension
        if not output_path.lower().endswith('.pptx'):
            output_path = output_path + '.pptx'
            logger.info(f"Added .pptx extension: {output_path}")

        # Save with proper error handling
        try:
            prs.save(output_path)
            logger.info(f"Presentation saved to {output_path}")
            
            # Verify file creation and get size
            file_size = os.path.getsize(output_path)
            logger.info(f"File verified: {file_size:,} bytes")
            return True
        except Exception as e:
            logger.exception(f"❌ Save failed: {e}")
            return False

        return True

    except Exception as e:
        logger.error(f"An error occurred during presentation creation: {e}")
        logger.error(traceback.format_exc())
        return False

def _apply_internal_table_borders(table, total_rows):
    """
    Apply specific internal borders for the first 7 columns only (excluding header and total rows).
    Uses direct OXML manipulation to add #BFBFBF borders with 0.75pt thickness.
    
    Args:
        table: The PowerPoint table object
        total_rows: Total number of rows in the table
    """
    try:
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        
        # Border specifications - same as external borders for consistency
        border_color_rgb = (191, 191, 191) # CLR_TABLE_GRAY is RGBColor(191, 191, 191)
        border_width_emu = int(0.75 * 12700)  # 0.75pt in EMUs
        hex_color = f"{border_color_rgb[0]:02X}{border_color_rgb[1]:02X}{border_color_rgb[2]:02X}"
        
        logger.debug(f"Applying internal borders to first 7 columns for rows 1 to {total_rows-2}")
        
        # Apply internal borders only to body rows (excluding header row 0 and total row)
        for row_idx in range(1, total_rows - 1):  # Skip header (0) and total (last) rows
            row = table.rows[row_idx]
            
            for col_idx in range(7):  # First 7 columns only
                cell = row.cells[col_idx]
                
                try:
                    # Access the cell's table cell properties
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    
                    # Add bottom border for horizontal lines (apply to all qualifying cells)
                    if row_idx < total_rows - 2:  # Don't add bottom border to row just above total
                        bottom_border = tcPr.find(qn('a:lnB'))
                        if bottom_border is None:
                            bottom_border = tcPr.add(qn('a:lnB'))
                        
                        # Set border properties
                        bottom_border.set('w', str(border_width_emu))
                        bottom_border.set('cap', 'flat')
                        bottom_border.set('cmpd', 'sng')
                        
                        # Add solid fill with color
                        solidFill = bottom_border.find(qn('a:solidFill'))
                        if solidFill is None:
                            solidFill = bottom_border.add(qn('a:solidFill'))
                        
                        srgbClr = solidFill.find(qn('a:srgbClr'))
                        if srgbClr is None:
                            srgbClr = solidFill.add(qn('a:srgbClr'))
                        srgbClr.set('val', hex_color)
                        
                        # Add preset dash (solid line)
                        prstDash = bottom_border.find(qn('a:prstDash'))
                        if prstDash is None:
                            prstDash = bottom_border.add(qn('a:prstDash'))
                        prstDash.set('val', 'solid')
                    
                    # Add right border for vertical lines (apply to first 6 columns only)
                    if col_idx < 6:  # Columns 0, 1, 2, 3, 4, 5 get right borders
                        right_border = tcPr.find(qn('a:lnR'))
                        if right_border is None:
                            right_border = tcPr.add(qn('a:lnR'))
                        
                        # Set border properties
                        right_border.set('w', str(border_width_emu))
                        right_border.set('cap', 'flat')
                        right_border.set('cmpd', 'sng')
                        
                        # Add solid fill with color
                        solidFill = right_border.find(qn('a:solidFill'))
                        if solidFill is None:
                            solidFill = right_border.add(qn('a:solidFill'))
                        
                        srgbClr = solidFill.find(qn('a:srgbClr'))
                        if srgbClr is None:
                            srgbClr = solidFill.add(qn('a:srgbClr'))
                        srgbClr.set('val', hex_color)
                        
                        # Add preset dash (solid line)
                        prstDash = right_border.find(qn('a:prstDash'))
                        if prstDash is None:
                            prstDash = right_border.add(qn('a:prstDash'))
                        prstDash.set('val', 'solid')
                    
                except Exception as cell_border_error:
                    logger.debug(f"Could not apply internal borders to cell ({row_idx}, {col_idx}): {cell_border_error}")
        
        logger.info(f"Internal table borders applied successfully for first 7 columns with #BFBFBF color and 0.75pt width")
        return True
        
    except Exception as e:
        logger.warning(f"Error applying internal table borders: {str(e)}")
        return False

def _verify_file_exists(label, path_str, extra_search=()):
    """
    Return an absolute path to an existing file.

    extra_search – optional directories (relative to project root) to probe
                   if the user gave a bare filename.
    """
    from pathlib import Path
    p = Path(path_str).expanduser()
    
    # Try the path as given first
    if p.is_file():
        return str(p.resolve())
    
    # If the file doesn't exist at the given path and we have fallback directories,
    # and if it looks like a bare filename (no directory separators), search fallbacks
    if extra_search and '/' not in str(p) and '\\' not in str(p):
        project_root = Path(__file__).parent.parent  # Go up from scripts/ to project root
        for alt in extra_search:
            candidate = project_root / alt / p.name
            if candidate.is_file():
                return str(candidate.resolve())

    # If we get here, the file wasn't found anywhere
    p_resolved = p.resolve()
    raise FileNotFoundError(f"❌ {label} file not found: '{p_resolved}'")

def _unit_test__no_orphan_self():
    import inspect, re, pathlib
    src = pathlib.Path(__file__).read_text(encoding="utf-8")
    tree = ast.parse(src, filename=str(pathlib.Path(__file__)))
    offender_lines = []
    for node in ast.walk(tree):
        # Set parent attributes for all nodes to allow easy traversal up the tree
        for child in ast.iter_child_nodes(node):
            child.parent = node
            
        if isinstance(node, ast.Attribute) and isinstance(node.value, ast.Name) and node.value.id == 'self':
            # Ascend until we hit a FunctionDef, ClassDef, or Module
            current_node = node
            legit = False
            is_in_class_scope = False # Track if we are within a ClassDef scope at all
            
            while hasattr(current_node, 'parent'): # Check if parent attribute exists
                parent = current_node.parent
                if isinstance(parent, ast.FunctionDef):
                    # Check if 'self' is the first argument of this function
                    if parent.args.args and parent.args.args[0].arg == 'self':
                        # Now, ensure this FunctionDef is itself within a ClassDef
                        # Keep ascending from the FunctionDef to find a ClassDef
                        func_parent = parent
                        while hasattr(func_parent, 'parent'):
                            func_parent = func_parent.parent
                            if isinstance(func_parent, ast.ClassDef):
                                legit = True # 'self' is in a method of a class
                                break
                            if isinstance(func_parent, ast.Module): # Reached top level without finding a class
                                break # Stop ascent once a FunctionDef is found and checked
                
                if isinstance(parent, ast.ClassDef):
                    is_in_class_scope = True # We are inside a class, but 'self' might be at class level (e.g. self.x = 10)
                    # If 'self' is used directly under ClassDef (not in a method), it's an error unless it's part of an assignment in __init__ or other method
                    # The current check correctly flags 'self.attribute' if not in a method like 'def meth(self, ...):'
                    # If we hit a ClassDef before a qualifying FunctionDef that has 'self' as first arg, it's likely a class variable access using 'self', which is wrong.
                    break # Stop ascent if ClassDef is found before a qualifying FunctionDef
                
                if isinstance(parent, ast.Module):
                    break # Reached top level
                current_node = parent
            
            if not legit:
                # Additional check: if it's in a class scope but not in a method, it's an error
                # e.g. class MyClass: self.x = 10 (invalid)
                # vs class MyClass: def __init__(self): self.x = 10 (valid)
                # The 'legit' flag already covers the valid case. If not legit, it's an offender.
                offender_lines.append(node.lineno)

    if offender_lines:
        # Remove duplicates and sort
        unique_offender_lines = sorted(list(set(offender_lines)))
        raise RuntimeError(f" orphan 'self.' detected outside a method or in an invalid class context at lines: {unique_offender_lines} – aborting build")

_unit_test__no_orphan_self()

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Generate PowerPoint presentation from Excel data using a template.')
    parser.add_argument('--excel', help='Path to the input Excel file (absolute or relative, must exist).')
    parser.add_argument('--template', help='Path to the PowerPoint template file (absolute or relative, must exist).')
    parser.add_argument('--output', help='Path for the output PowerPoint file. Defaults to Output_Presentation_v2_[timestamp].pptx')
    parser.add_argument('--log', help='Path for the log file. Defaults to logs/excel_to_ppt_v2_log_[timestamp].log')
    parser.add_argument("--list-templates", action="store_true",
        help="Show all *.pptx files in 'data/input/templates/', 'templates/', ., or .. directories and exit")

    args = parser.parse_args()

    # Handle --list-templates flag
    if args.list_templates:
        from glob import glob
        here = Path(__file__).parent.parent.parent  # project root
        TEMPLATE_FALLBACK_DIRS = ("input/templates", "data/input/templates", ".", "..")
        print("Available template files:")
        found_any = False
        for d in TEMPLATE_FALLBACK_DIRS:
            for f in glob(str(here / d / "*.pptx")):
                print(f"  {Path(f).relative_to(here)}")
                found_any = True
        if not found_any:
            print("  No .pptx files found in data/input/templates/, templates/, ., or .. directories")
        sys.exit(0)

    # Check required arguments for normal operation
    if not args.excel or not args.template:
        parser.error("--excel and --template are required when not using --list-templates")

    # Initialize logging later after we create output directory
    # logger = setup_logging()
    logger = None  # Will be set up after output directory is created

    # --- Path Validation and Logging --- 
    excel_input_path_str = args.excel # No longer .name, and no .close() needed
    
    TEMPLATE_FALLBACK_DIRS = ("input/templates", "data/input/templates", ".", "..")

    try:
        template_path = _verify_file_exists("Template", args.template, TEMPLATE_FALLBACK_DIRS)
        excel_path    = _verify_file_exists("Excel", excel_input_path_str)  # no fallbacks; XLSX is user-supplied
    except FileNotFoundError as e:
        # Log the error using the existing logger setup if available, or print to stderr
        if 'logger' in globals() and logger is not None:
            logger.error(f"File validation failed: {str(e)}")
        else:
            print(f"Pre-logging ERROR: File validation failed: {str(e)}", file=sys.stderr)
        print(f"FATAL ERROR: {str(e)}", file=sys.stderr) # Ensure message goes to console
        sys.exit(2) # Exit code 2 for "bad invocation" distinct from "script bug"

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # Create timestamped output subfolder
    output_subfolder = os.path.join("output", f"run_{timestamp}")
    os.makedirs(output_subfolder, exist_ok=True)
    
    # Configure File Handler for logging in the output subfolder
    if args.log:
        log_file_path = args.log
    else:
        log_file_path = os.path.join(output_subfolder, f"automation_log_{timestamp}.log")
    
    # Ensure logger is fully configured before first use post-validation
    # If setup_logging() was called earlier, this might reconfigure or add handlers.
    # Assuming logger is configured by setup_logging() called near the top or on first use.
    # The critical part is the print to stderr and sys.exit(2).

    # Log validated paths (logger should be configured by now if setup_logging is called early)
    # If setup_logging is called after arg parsing, these specific logs might only go to console if logger isn't ready.
    # Best practice: Initialize logger once, early. Assuming 'logger = setup_logging()' is effective globally.
    
    # Initialize logging with the log file in the output subfolder
    logger = setup_logging(log_path_base=os.path.join(output_subfolder, "automation_log"))

    logger.info(f"Script execution started with arguments: {args}") # Moved this log after path validation
    logger.info(f"Using template: {template_path}")
    logger.info(f"Using Excel   : {excel_path}")

    # File logging is already configured in setup_logging()
    logger.info(f"Logging to: {log_file_path}")

    # Set output filename in the timestamped subfolder
    if args.output:
        # If user specified output, use their filename but in our subfolder
        output_filename = os.path.join(output_subfolder, os.path.basename(args.output))
    else:
        output_filename = os.path.join(output_subfolder, f"AMP_Presentation_{timestamp}.pptx")
    
    logger.info(f"Output will be saved to: {output_filename}")
    logger.info(f"Log file will be saved to: {log_file_path}")

    success = create_presentation(
        template_path=template_path, # Use validated absolute path
        excel_path=excel_path,       # Use validated absolute path
        output_path=output_filename
    )

    if success:
        logger.info(f"Presentation generation completed successfully. Output: {output_filename}")
        # Display the location prominently for easy access
        print("\n" + "="*80)
        print("✅ PRESENTATION GENERATED SUCCESSFULLY!")
        print("="*80)
        print(f"📁 LOCATION: {output_filename}")
        print("="*80 + "\n")
        sys.exit(0)
    else:
        logger.error("Presentation generation failed.")
        sys.exit(1)
