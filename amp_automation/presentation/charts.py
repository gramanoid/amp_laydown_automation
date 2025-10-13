"""Chart rendering helpers."""

from __future__ import annotations

import logging
from dataclasses import dataclass
from typing import Mapping, Sequence

import pandas as pd
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from amp_automation.presentation.tables import ensure_font_consistency
from amp_automation.utils.media import normalize_media_type

logger = logging.getLogger("amp_automation.presentation.charts")

__all__ = [
    "ChartStyleContext",
    "prepare_funnel_chart_data",
    "prepare_media_type_chart_data",
    "prepare_campaign_type_chart_data",
    "add_pie_chart",
]


@dataclass(slots=True)
class ChartStyleContext:
    """Styling parameters applied when rendering charts."""

    font_name: str
    title_font_size: Pt
    label_font_size: Pt
    font_color: RGBColor
    legend_position: XL_LEGEND_POSITION = XL_LEGEND_POSITION.BOTTOM
    data_label_number_format: str = "0.0%"


def prepare_funnel_chart_data(
    df: pd.DataFrame,
    region: str,
    masterbrand: str,
    year: str | int | None = None,
) -> dict[str, float] | None:
    """Aggregate budgets by funnel stage for the specified filters."""

    try:
        filters = [
            df["Country"].astype(str).str.strip() == str(region).strip(),
            df["Brand"].astype(str).str.strip() == str(masterbrand).strip(),
        ]
        if year is not None:
            filters.append(df["Year"].astype(str).str.strip() == str(year).strip())

        mask = filters[0]
        for condition in filters[1:]:
            mask &= condition

        filtered = df.loc[mask].copy()
        if filtered.empty:
            logger.warning("No data found for funnel chart: %s - %s", region, masterbrand)
            return None

        budgets: dict[str, float] = {}
        if "Funnel Stage" not in filtered.columns:
            logger.warning("Missing 'Funnel Stage' column for funnel chart")
            return None

        for stage in filtered["Funnel Stage"].dropna().unique():
            stage_str = str(stage).strip()
            if not stage_str:
                continue
            total_budget = filtered.loc[filtered["Funnel Stage"] == stage, "Total Cost"].sum()
            if total_budget > 0:
                budgets[stage_str] = float(total_budget)

        return budgets or None

    except Exception as exc:  # pragma: no cover - defensive logging
        logger.error(
            "Error preparing funnel chart data for %s - %s: %s",
            region,
            masterbrand,
            exc,
        )
        return None


def prepare_media_type_chart_data(
    df: pd.DataFrame,
    region: str,
    masterbrand: str,
    year: str | int | None = None,
) -> dict[str, float] | None:
    """Aggregate budgets by normalized media type for the specified filters."""

    try:
        filters = [
            df["Country"].astype(str).str.strip() == str(region).strip(),
            df["Brand"].astype(str).str.strip() == str(masterbrand).strip(),
        ]
        if year is not None:
            filters.append(df["Year"].astype(str).str.strip() == str(year).strip())

        mask = filters[0]
        for condition in filters[1:]:
            mask &= condition

        filtered = df.loc[mask].copy()
        if filtered.empty:
            logger.warning("No data found for media type chart: %s - %s", region, masterbrand)
            return None

        if "Media Type" not in filtered.columns or "Total Cost" not in filtered.columns:
            logger.warning("Missing media type columns for chart data")
            return None

        budgets: dict[str, float] = {}
        for raw_media_type in filtered["Media Type"].dropna().unique():
            normalized = normalize_media_type(raw_media_type)
            media_total = filtered.loc[
                filtered["Media Type"] == raw_media_type, "Total Cost"
            ].sum()
            if media_total <= 0:
                continue
            budgets[normalized] = budgets.get(normalized, 0.0) + float(media_total)

        return budgets or None

    except Exception as exc:  # pragma: no cover - defensive logging
        logger.error(
            "Error preparing media type chart data for %s - %s: %s",
            region,
            masterbrand,
            exc,
        )
        return None


def prepare_campaign_type_chart_data(
    df: pd.DataFrame,
    region: str,
    masterbrand: str,
    year: str | int | None = None,
) -> dict[str, float] | None:
    """Aggregate budgets by campaign type for the specified filters."""

    try:
        filters = [
            df["Country"].astype(str).str.strip() == str(region).strip(),
            df["Brand"].astype(str).str.strip() == str(masterbrand).strip(),
        ]
        if year is not None:
            filters.append(df["Year"].astype(str).str.strip() == str(year).strip())

        mask = filters[0]
        for condition in filters[1:]:
            mask &= condition

        filtered = df.loc[mask].copy()
        if filtered.empty:
            logger.warning("No data found for campaign type chart: %s - %s", region, masterbrand)
            return None

        if "Campaign Type" not in filtered.columns or "Total Cost" not in filtered.columns:
            logger.warning("Missing campaign type columns for chart data")
            return None

        budgets: dict[str, float] = {}
        for campaign_type in filtered["Campaign Type"].dropna().unique():
            campaign_str = str(campaign_type).strip()
            if not campaign_str:
                continue
            total_budget = filtered.loc[
                filtered["Campaign Type"] == campaign_type, "Total Cost"
            ].sum()
            if total_budget > 0:
                budgets[campaign_str] = float(total_budget)

        return budgets or None

    except Exception as exc:  # pragma: no cover - defensive logging
        logger.error(
            "Error preparing campaign type chart data for %s - %s: %s",
            region,
            masterbrand,
            exc,
        )
        return None


def add_pie_chart(
    slide,
    chart_data: Mapping[str, float],
    chart_title: str,
    position: Mapping[str, object],
    style: ChartStyleContext,
    color_mapping: Mapping[str, RGBColor],
    default_colors: Sequence[RGBColor],
    *,
    chart_name: str | None = None,
) -> bool:
    """Render a pie chart with consistent styling."""

    try:
        if not chart_data:
            logger.warning("No chart data provided for %s", chart_title)
            return False

        total = float(sum(chart_data.values()))
        if total <= 0:
            logger.warning("Invalid total value for chart %s: %s", chart_title, total)
            return False

        categories = list(chart_data.keys())
        values = [float(value) for value in chart_data.values()]

        data = CategoryChartData()
        data.categories = categories
        data.add_series("Budget", values)

        try:
            graphic_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                position["left"],
                position["top"],
                position["width"],
                position["height"],
                data,
            )
        except KeyError as exc:
            logger.error("Missing chart position key %s for %s", exc, chart_title)
            return False

        if chart_name:
            graphic_frame.name = chart_name

        chart = graphic_frame.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title

        title_run = chart.chart_title.text_frame.paragraphs[0].runs[0]
        ensure_font_consistency(
            title_run.font,
            style.font_name,
            style.title_font_size,
            False,
            style.font_color,
        )
        chart.chart_title.text_frame.auto_size = MSO_AUTO_SIZE.NONE

        chart.has_legend = True
        chart.legend.position = style.legend_position
        try:
            ensure_font_consistency(
                chart.legend.font,
                style.font_name,
                style.label_font_size,
                False,
                style.font_color,
            )
        except Exception as exc:  # pragma: no cover - defensive logging
            logger.debug("Unable to style legend font for %s: %s", chart_title, exc)

        series = chart.series[0]

        for index, category in enumerate(categories):
            try:
                point = series.points[index]
            except IndexError:
                continue

            color = color_mapping.get(category)
            if color is None and default_colors:
                color = default_colors[index % len(default_colors)]
            if color is None:
                continue

            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color

        series.has_data_labels = True
        labels = series.data_labels
        labels.show_percentage = True
        labels.show_value = False

        if style.data_label_number_format:
            try:
                labels.number_format = style.data_label_number_format
            except Exception as exc:  # pragma: no cover - defensive logging
                logger.debug(
                    "Unable to apply number format '%s' for %s: %s",
                    style.data_label_number_format,
                    chart_title,
                    exc,
                )

        try:
            ensure_font_consistency(
                labels.font,
                style.font_name,
                style.label_font_size,
                False,
                style.font_color,
            )
            labels.font.bold = False
        except Exception as exc:  # pragma: no cover - defensive logging
            logger.debug("Unable to style data labels for %s: %s", chart_title, exc)

        logger.info("Chart '%s' created successfully with %s data points", chart_title, len(categories))
        return True

    except Exception as exc:  # pragma: no cover - defensive logging
        logger.error("Error creating pie chart '%s': %s", chart_title, exc)
        return False
