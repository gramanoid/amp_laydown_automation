"""Table construction helpers."""

from __future__ import annotations

import logging
import traceback
from dataclasses import dataclass, field
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


ALIGNMENT_LOOKUP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def _resolve_alignment(value) -> PP_ALIGN:
    if isinstance(value, PP_ALIGN):
        return value
    if isinstance(value, str):
        return ALIGNMENT_LOOKUP.get(value.lower(), PP_ALIGN.CENTER)
    if isinstance(value, int):
        try:
            return PP_ALIGN(value)
        except ValueError:
            return PP_ALIGN.CENTER
    return PP_ALIGN.CENTER

logger = logging.getLogger("amp_automation.presentation.tables")

__all__ = [
    "ensure_font_consistency",
    "apply_table_borders",
    "CellStyleContext",
    "TableLayout",
    "style_table_cell",
    "add_and_style_table",
]


def ensure_font_consistency(font, target_font_name: Optional[str], target_size, target_bold: Optional[bool], target_color):
    """Apply the target font properties to the given run font."""

    try:
        if target_font_name:
            font.name = target_font_name
        if target_size:
            font.size = target_size
        if target_bold is not None:
            font.bold = target_bold
        if target_color:
            font.color.rgb = target_color
    except Exception as exc:  # pragma: no cover - defensive logging
        logger.debug("Error setting font properties: %s", exc)


def apply_table_borders(table, border_color: RGBColor, border_width_pt: float = 0.75) -> bool:
    """Apply consistent borders to the provided PowerPoint table."""

    try:
        border_width = Pt(border_width_pt)
        border_width_emu = int(border_width_pt * 12700)
        hex_color = str(border_color)
        if hex_color.startswith("0x"):
            hex_color = hex_color[2:]
        hex_color = hex_color.upper().zfill(6)

        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                try:
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()

                    for edge in ("lnT", "lnL", "lnB", "lnR"):
                        border = tcPr.find(qn(f"a:{edge}"))
                        if border is None:
                            border = OxmlElement(f"a:{edge}")
                            tcPr.append(border)
                        # Clear any existing children/attributes so we start from a clean state
                        border.attrib.clear()
                        for child in list(border):
                            border.remove(child)

                        border.set("w", str(border_width_emu))

                        solid_fill = OxmlElement("a:solidFill")
                        srgb = OxmlElement("a:srgbClr")
                        srgb.set("val", hex_color)
                        solid_fill.append(srgb)
                        border.append(solid_fill)

                        dash = OxmlElement("a:prstDash")
                        dash.set("val", "solid")
                        border.append(dash)

                except Exception as cell_error:  # pragma: no cover - defensive log
                    logger.debug("Border styling failed for cell (%s,%s): %s", row_idx, col_idx, cell_error)

        logger.info("Table borders applied successfully")
        return True

    except Exception as exc:  # pragma: no cover - defensive log
        logger.warning("Error applying table borders: %s", exc)
        return False


def _scale_table_height(table, target_height_emu: int | None) -> None:
    """Scale table row heights proportionally to reach the target total height."""

    if not target_height_emu:
        return

    current_height = sum(row.height for row in table.rows)
    if not current_height:
        return

    scale = target_height_emu / current_height
    if abs(scale - 1.0) < 1e-6:
        return

    for row in table.rows:
        row.height = int(row.height * scale)


@dataclass(slots=True)
class CellStyleContext:
    margin_left_right_pt: float
    margin_emu_lr: int
    default_font_name: str
    font_size_header: Pt
    font_size_body: Pt
    font_size_body_compact: Pt
    color_black: RGBColor
    color_light_gray_text: RGBColor
    color_table_gray: RGBColor
    color_header_green: RGBColor
    color_subtotal_gray: RGBColor
    color_tv: RGBColor
    color_digital: RGBColor
    color_ooh: RGBColor
    color_other: RGBColor
    column_alignment: dict[int, object] = field(default_factory=dict)
    word_wrap_columns: set[int] = field(default_factory=set)
    shrink_to_fit_columns: set[int] = field(default_factory=set)
    uppercase_columns: set[int] = field(default_factory=set)
    dual_line_labels: dict[str, list[str]] = field(default_factory=dict)


@dataclass(slots=True)
class TableLayout:
    """Presentation geometry and rules for constructing a table."""

    placeholder_name: str
    shape_name: str
    position: dict[str, object]
    row_height_header: Pt
    row_height_body: Pt
    row_height_subtotal: Pt
    column_widths: list
    top_override: object | None = None
    height_rule_available: bool = False
    height_rule_value: object | None = None


def style_table_cell(
    cell,
    row_idx: int,
    col_idx: int,
    table_data: list[list[str]],
    cell_metadata: dict[tuple[int, int], dict[str, object]],
    context: CellStyleContext,
    logger: logging.Logger,
) -> None:
    """Apply styling to a specific table cell."""

    MARGIN_LEFT_RIGHT_PT = context.margin_left_right_pt
    MARGIN_EMU_LR = context.margin_emu_lr
    DEFAULT_FONT_NAME = context.default_font_name
    FONT_SIZE_HEADER = context.font_size_header
    FONT_SIZE_BODY = context.font_size_body
    FONT_SIZE_BODY_COMPACT = context.font_size_body_compact
    FONT_SIZE_MONTHLY_TOTAL = Pt(6.5)
    CLR_BLACK = context.color_black
    CLR_LIGHT_GRAY_TEXT = context.color_light_gray_text
    CLR_TABLE_GRAY = context.color_table_gray
    CLR_HEADER_GREEN = context.color_header_green
    CLR_SUBTOTAL_GRAY = context.color_subtotal_gray
    CLR_TELEVISION = context.color_tv
    CLR_DIGITAL = context.color_digital
    CLR_OOH = context.color_ooh
    CLR_OTHER = context.color_other
    wrap_from_config = col_idx in context.word_wrap_columns
    shrink_to_fit = col_idx in context.shrink_to_fit_columns
    use_compact_font = False
    MONTH_HEADER_COLUMNS = set(range(3, 15))
    CLR_WHITE = RGBColor(255, 255, 255)
    SUBTOTAL_LABELS = {"SUBTOTAL", "CARRIED FORWARD", "MONTHLY TOTAL (£ 000)", "GRAND TOTAL"}

    normalized_row_label = ""
    if 0 <= row_idx < len(table_data) and table_data[row_idx]:
        normalized_row_label = " ".join(str(table_data[row_idx][0]).replace("\xa0", " ").split()).upper()
    is_monthly_total_row = normalized_row_label == "MONTHLY TOTAL (£ 000)"
    is_total_row = row_idx == len(table_data) - 1 or (row_idx > 0 and normalized_row_label in SUBTOTAL_LABELS)

    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
    alignment = PP_ALIGN.CENTER

    try:
        original_cell_text = str(table_data[row_idx][col_idx]) if col_idx < len(table_data[row_idx]) else ""
        logger.debug("CELL STYLING [%s,%s]: Original text='%s'", row_idx, col_idx, original_cell_text)

        try:
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            logger.debug("CELL STYLING [%s,%s]: API cell.vertical_anchor set to MIDDLE. Current: %s", row_idx, col_idx, cell.vertical_anchor)
        except Exception as e_cell_anchor:
            logger.warning(
                "CELL STYLING [%s,%s]: Failed to set cell.vertical_anchor via API: %s. Will rely on OXML.",
                row_idx,
                col_idx,
                e_cell_anchor,
            )

        try:
            logger.debug(
                "CELL STYLING [%s,%s]: API setting cell margins: L/R=%spt, T/B=Default",
                row_idx,
                col_idx,
                MARGIN_LEFT_RIGHT_PT,
            )
            cell.margin_left = Pt(MARGIN_LEFT_RIGHT_PT)
            cell.margin_right = Pt(MARGIN_LEFT_RIGHT_PT)
            logger.debug(
                "CELL STYLING [%s,%s]: Margins after API set: L=%spt, R=%spt, T(api)='Omitted', B(api)='Omitted'",
                row_idx,
                col_idx,
                cell.margin_left.pt if cell.margin_left else "Default",
                cell.margin_right.pt if cell.margin_right else "Default",
            )
        except Exception as margin_error:
            logger.warning(
                "CELL STYLING [%s,%s]: Error setting L/R margins via API: %s. Will rely on OXML for L/R as well.",
                row_idx,
                col_idx,
                margin_error,
            )

        text_frame = cell.text_frame
        text_frame.clear()
        if use_compact_font and processed_cell_text:
            processed_cell_text = (
                processed_cell_text.replace(" ", "\u2060 ")
                .replace("-", "\u2011")
            )

        text_frame.word_wrap = wrap_from_config
        text_frame.auto_size = MSO_AUTO_SIZE.NONE

        try:
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            logger.debug(
                "CELL STYLING [%s,%s]: Text_frame margins set to 0. L=%s, T=%s, R=%s, B=%s",
                row_idx,
                col_idx,
                text_frame.margin_left,
                text_frame.margin_top,
                text_frame.margin_right,
                text_frame.margin_bottom,
            )

            if hasattr(text_frame._txBody, "bodyPr") and text_frame._txBody.bodyPr is not None:
                bodyPr = text_frame._txBody.bodyPr
            else:
                bodyPr = OxmlElement("a:bodyPr")
                text_frame._txBody.insert_element_before(bodyPr, "a:lstStyle", "a:p")
                logger.debug(
                    "CELL STYLING [%s,%s]: OXML bodyPr element created as it was missing.",
                    row_idx,
                    col_idx,
                )

            bodyPr.set("anchor", "ctr")
            bodyPr.set("lIns", "0")
            bodyPr.set("tIns", "0")
            bodyPr.set("rIns", "0")
            bodyPr.set("bIns", "0")
            logger.debug(
                "CELL STYLING [%s,%s]: OXML bodyPr attributes directly set: anchor='ctr', lIns='0', tIns='0', rIns='0', bIns='0'.",
                row_idx,
                col_idx,
            )
            logger.debug(
                "CELL STYLING [%s,%s]: text_frame.vertical_anchor after bodyPr.set: %s",
                row_idx,
                col_idx,
                text_frame.vertical_anchor,
            )

        except Exception as bodyPr_error:
            logger.error(
                "CELL STYLING [%s,%s]: Critical error setting OXML bodyPr or text_frame margins: %s %s",
                row_idx,
                col_idx,
                bodyPr_error,
                traceback.format_exc(),
            )

        try:
            tcPr = cell._tc.get_or_add_tcPr()
            tcPr.set("anchor", "ctr")
            logger.debug(
                "CELL STYLING [%s,%s]: OXML tcPr attribute 'anchor' directly set to 'ctr'.",
                row_idx,
                col_idx,
            )
            tcPr.set("vert", "horz")
            logger.debug(
                "CELL STYLING [%s,%s]: OXML tcPr attribute 'vert' directly set to 'horz'.",
                row_idx,
                col_idx,
            )
            tcPr.marL = MARGIN_EMU_LR
            tcPr.marR = MARGIN_EMU_LR
            logger.debug(
                "CELL STYLING [%s,%s]: OXML tcPr.marL & tcPr.marR set to %s.",
                row_idx,
                col_idx,
                MARGIN_EMU_LR,
            )

            marT_elements = tcPr.xpath("./a:marT")
            for el_marT in marT_elements:
                tcPr.remove(el_marT)
                logger.debug(
                    "CELL STYLING [%s,%s]: OXML removed existing a:marT element.",
                    row_idx,
                    col_idx,
                )

            marB_elements = tcPr.xpath("./a:marB")
            for el_marB in marB_elements:
                tcPr.remove(el_marB)
                logger.debug(
                    "CELL STYLING [%s,%s]: OXML removed existing a:marB element.",
                    row_idx,
                    col_idx,
                )

            final_marL_oxml = tcPr.marL if hasattr(tcPr, "marL") and tcPr.marL is not None else "NotSet"
            final_marR_oxml = tcPr.marR if hasattr(tcPr, "marR") and tcPr.marR is not None else "NotSet"
            final_marT_oxml_exists = bool(tcPr.xpath("./a:marT"))
            final_marB_oxml_exists = bool(tcPr.xpath("./a:marB"))
            final_anchor_oxml = tcPr.anchor if hasattr(tcPr, "anchor") and tcPr.anchor is not None else "NotSet"
            final_vert_oxml = tcPr.get("vert") if hasattr(tcPr, "get") else "NotSet"
            logger.debug(
                "CELL STYLING [%s,%s]: Final OXML tcPr: anchor='%s', vert='%s', marL='%s', marR='%s', marT_exists=%s, marB_exists=%s",
                row_idx,
                col_idx,
                final_anchor_oxml,
                final_vert_oxml,
                final_marL_oxml,
                final_marR_oxml,
                final_marT_oxml_exists,
                final_marB_oxml_exists,
            )

        except Exception as oxml_error:
            logger.error(
                "CELL STYLING [%s,%s]: Critical error setting OXML tcPr properties: %s %s",
                row_idx,
                col_idx,
                oxml_error,
                traceback.format_exc(),
            )

        paragraphs_before_fix = len(text_frame.paragraphs)

        if text_frame.paragraphs:
            p = text_frame.paragraphs[0]
            logger.debug(
                "CELL STYLING [%s,%s]: Reusing existing paragraph after clear(). Paragraphs count: %s",
                row_idx,
                col_idx,
                paragraphs_before_fix,
            )
        else:
            p = text_frame.add_paragraph()
            logger.debug(
                "CELL STYLING [%s,%s]: No paragraphs found after clear(), added new one. Paragraphs count: %s",
                row_idx,
                col_idx,
                len(text_frame.paragraphs),
            )

        p.alignment = alignment
        pPr = p._p.get_or_add_pPr()

        for spacing_element in ["a:spcBef", "a:spcAft", "a:lnSpc"]:
            existing = pPr.find(qn(spacing_element))
            if existing is not None:
                pPr.remove(existing)
                logger.debug(
                    "CELL STYLING [%s,%s]: Removed existing %s",
                    row_idx,
                    col_idx,
                    spacing_element,
                )

        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is not None:
            pPr.remove(defRPr)
            logger.debug(
                "CELL STYLING [%s,%s]: Removed default paragraph properties (a:defRPr)",
                row_idx,
                col_idx,
            )

        lnSpc = OxmlElement("a:lnSpc")
        spcPts_line = OxmlElement("a:spcPts")
        if row_idx == 0:
            font_size_pt = FONT_SIZE_HEADER.pt
        else:
            font_size_pt = FONT_SIZE_BODY.pt
        line_spacing_pt = int(font_size_pt * 100)
        spcPts_line.set("val", str(line_spacing_pt))
        lnSpc.append(spcPts_line)
        pPr.append(lnSpc)

        spcBef = OxmlElement("a:spcBef")
        spcPts_before = OxmlElement("a:spcPts")
        spcPts_before.set("val", "0")
        spcBef.append(spcPts_before)
        pPr.append(spcBef)

        spcAft = OxmlElement("a:spcAft")
        spcPts_after = OxmlElement("a:spcPts")
        spcPts_after.set("val", "0")
        spcAft.append(spcPts_after)
        pPr.append(spcAft)

        logger.debug(
            "CELL STYLING [%s,%s]: ULTRA-COMPACT spacing applied: font_size=%spt, line_spacing=%spt, spcBef=0, spcAft=0",
            row_idx,
            col_idx,
            font_size_pt,
            line_spacing_pt / 100,
        )

        processed_cell_text = ""
        is_empty_cell = False

        if row_idx == 0:
            if (
                original_cell_text
                and str(original_cell_text).strip()
                and str(original_cell_text).strip()
                not in ["0", "0.0", "0.00", "0.000", "£0K", "0K", "-", "–", "0.0%"]
            ):
                processed_cell_text = str(original_cell_text).upper()
            else:
                processed_cell_text = "-"
                is_empty_cell = True

        elif row_idx == len(table_data) - 1:
            if (
                original_cell_text
                and str(original_cell_text).strip()
                and str(original_cell_text).strip()
                not in ["0", "0.0", "0.00", "0.000", "£0K", "0K", "-", "–", "0.0%"]
            ):
                processed_cell_text = str(original_cell_text).upper()
            else:
                processed_cell_text = "-"
                is_empty_cell = True

        else:
            if (
                not original_cell_text
                or str(original_cell_text).strip() == ""
                or str(original_cell_text).strip() in ["0", "0.0", "0.00", "0.000", "£0K", "0K", "-", "–", "0.0%"]
            ):
                processed_cell_text = "-"
                is_empty_cell = True
            else:
                if col_idx in (0, 1) and str(original_cell_text).strip():
                    processed_cell_text = str(original_cell_text).upper()
                else:
                    processed_cell_text = str(original_cell_text)

        if processed_cell_text and processed_cell_text.strip().endswith("%"):
            pct_value = processed_cell_text.strip()[:-1]
            try:
                numeric_pct = float(pct_value)
                if abs(numeric_pct - round(numeric_pct)) < 1e-6:
                    processed_cell_text = f"{int(round(numeric_pct))}%"
            except ValueError:
                pass

        if processed_cell_text not in ("", "-", "-") and col_idx in context.uppercase_columns:
            processed_cell_text = processed_cell_text.upper()

        normalized_key = (
            processed_cell_text.replace("\r", " ").replace("\n", " ").strip().upper()
            if processed_cell_text
            else ""
        )
        if normalized_key and normalized_key in context.dual_line_labels:
            processed_cell_text = "\r".join(context.dual_line_labels[normalized_key])

        processed_cell_text = (
            processed_cell_text.replace("\r\n", "\r").replace("\n", "\r") if processed_cell_text else processed_cell_text
        )

        normalized_for_compact = (
            "".join(ch for ch in processed_cell_text if ch.isalnum()) if processed_cell_text else ""
        )
        compact_length = len(normalized_for_compact)
        compact_columns = {0, 1}
        use_compact_font = (
            row_idx > 0
            and col_idx in compact_columns
            and compact_length >= 10
        )

        if use_compact_font and processed_cell_text:
            processed_cell_text = processed_cell_text.replace(" ", "\u00A0").replace("-", "\u2011")

        text_frame.word_wrap = wrap_from_config
        text_frame.auto_size = MSO_AUTO_SIZE.NONE

        body_font_size = FONT_SIZE_BODY
        empty_font_name = "Verdana"
        empty_font_size = Pt(6)

        if not processed_cell_text or processed_cell_text == "-":
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = "-"
            run.font.name = empty_font_name
            run.font.size = empty_font_size
            run.font.bold = False
            run.font.color.rgb = CLR_LIGHT_GRAY_TEXT
            paragraph.font.name = empty_font_name
            paragraph.font.size = empty_font_size
            processed_cell_text = "-"
        else:
            if not p.runs:
                run = p.add_run()
            else:
                run = p.runs[0]
            run.text = processed_cell_text

        if row_idx == 0:
            run.font.name = DEFAULT_FONT_NAME
            run.font.size = FONT_SIZE_HEADER
            run.font.bold = True
            header_color = CLR_WHITE if col_idx in MONTH_HEADER_COLUMNS else CLR_BLACK
            run.font.color.rgb = header_color

        elif is_total_row:
            run.font.name = DEFAULT_FONT_NAME
            run.font.size = FONT_SIZE_MONTHLY_TOTAL if is_monthly_total_row else FONT_SIZE_HEADER
            run.font.bold = True
            run.font.color.rgb = CLR_BLACK

        else:
            if processed_cell_text == "-":
                run.font.name = empty_font_name
                run.font.size = empty_font_size
                run.font.bold = False
                run.font.color.rgb = CLR_LIGHT_GRAY_TEXT
            else:
                run.font.name = DEFAULT_FONT_NAME
                run.font.size = body_font_size
                run.font.color.rgb = CLR_BLACK
                run.font.bold = col_idx < 3 or col_idx in (15, 16)  # CAMPAIGN/MEDIA/METRICS + TOTAL/GRPs columns


        if processed_cell_text != "-" and run.font.name != DEFAULT_FONT_NAME:
            run.font.name = DEFAULT_FONT_NAME
            logger.debug("Re-enforced Calibri font for cell (%s,%s)", row_idx, col_idx)

        if processed_cell_text != "-" and not run.font.size:
            run.font.size = FONT_SIZE_BODY if row_idx != 0 else FONT_SIZE_HEADER
            logger.debug("Re-enforced font size for cell (%s,%s)", row_idx, col_idx)

        try:
            rPr = run._r.get_or_add_rPr()
            defRPr_run = rPr.find(qn("a:defRPr"))
            if defRPr_run is not None:
                rPr.remove(defRPr_run)
                logger.debug("CELL STYLING [%s,%s]: Removed run-level default properties", row_idx, col_idx)

            baseline = rPr.find(qn("a:baseline"))
            if baseline is not None:
                rPr.remove(baseline)

            baseline = OxmlElement("a:baseline")
            baseline.set("val", "0")
            rPr.append(baseline)

            spc = rPr.find(qn("a:spc"))
            if spc is not None:
                rPr.remove(spc)
                logger.debug("CELL STYLING [%s,%s]: Removed character spacing", row_idx, col_idx)

            logger.debug("CELL STYLING [%s,%s]: Run-level baseline fix applied", row_idx, col_idx)

        except Exception as run_oxml_error:
            logger.warning(
                "CELL STYLING [%s,%s]: Error applying run-level OXML fixes: %s",
                row_idx,
                col_idx,
                run_oxml_error,
            )

        subtotal_labels = {"SUBTOTAL", "CARRIED FORWARD", "MONTHLY TOTAL (\u00a3 000)", "GRAND TOTAL"}
        row_label = str(table_data[row_idx][0]).strip().upper() if row_idx < len(table_data) else ""

        def _apply_rgb_fill(target_cell, rgb_color):
            target_cell.fill.solid()
            target_cell.fill.fore_color.rgb = rgb_color

        def _apply_theme_fill(target_cell, theme_color, brightness=None):
            target_cell.fill.solid()
            target_cell.fill.fore_color.theme_color = theme_color
            if brightness is not None:
                target_cell.fill.fore_color.brightness = brightness
            if hasattr(target_cell.fill.fore_color, "tint"):
                target_cell.fill.fore_color.tint = None
            if hasattr(target_cell.fill.fore_color, "shade"):
                target_cell.fill.fore_color.shade = None

        def _apply_base_background(target_cell):
            _apply_theme_fill(target_cell, MSO_THEME_COLOR_INDEX.BACKGROUND_1, brightness=0)

        highlight_alias_map = {
            "GRPS": "TELEVISION",
            "REACH": "TELEVISION",
            "REACH@1+": "TELEVISION",
            "REACH1+": "TELEVISION",
            "OTS": "TELEVISION",
            "OTS@3+": "TELEVISION",
        }

        if row_idx == 0:
            header_text2_cols = {0, 1, 2, 15, 16, 17}
            if col_idx in header_text2_cols:
                _apply_theme_fill(cell, MSO_THEME_COLOR_INDEX.TEXT_2, brightness=-0.1)
            else:
                _apply_theme_fill(cell, MSO_THEME_COLOR_INDEX.TEXT_1, brightness=0.35)

        elif row_idx == len(table_data) - 1 or row_label in subtotal_labels:
            _apply_rgb_fill(cell, CLR_SUBTOTAL_GRAY)

        else:
            cell_key = (row_idx, col_idx)
            total_col_idx = len(table_data[row_idx]) - 3 if row_idx < len(table_data) and len(table_data[row_idx]) >= 3 else None

            def resolve_media_type() -> str | None:
                if cell_key in cell_metadata:
                    return cell_metadata[cell_key].get("media_type")

                # Check column 2 (METRICS) for special metric names with assigned colors
                metrics_value = table_data[row_idx][2] if row_idx < len(table_data) and len(table_data[row_idx]) > 2 else None
                if metrics_value:
                    metrics_normalized = str(metrics_value).strip().upper()
                    metric_to_media = {
                        "REACH@1+": "TELEVISION",
                        "OTS@3+": "TELEVISION",
                        "META REACH": "DIGITAL",
                        "TT REACH": "DIGITAL",
                    }
                    if metrics_normalized in metric_to_media:
                        return metric_to_media[metrics_normalized]

                # Fall back to column 1 (media type)
                media_value = table_data[row_idx][1] if row_idx < len(table_data) and len(table_data[row_idx]) > 1 else None
                media_value = str(media_value).strip() if media_value else ""
                if media_value and media_value != "-":
                    return media_value
                if row_idx > 0:
                    prev_row = table_data[row_idx - 1]
                    if len(prev_row) > 1:
                        prev_media = str(prev_row[1]).strip()
                        if prev_media and prev_media != "-":
                            return prev_media
                return None

            media_type = resolve_media_type()

            def apply_media_highlight() -> bool:
                if not media_type:
                    return False

                normalized = media_type.strip().upper()
                normalized = highlight_alias_map.get(normalized, normalized)
                if normalized in {"TELEVISION", "TV"}:
                    _apply_rgb_fill(cell, CLR_TELEVISION)
                    return True
                if normalized == "DIGITAL":
                    _apply_rgb_fill(cell, CLR_DIGITAL)
                    return True
                if normalized == "OOH":
                    _apply_rgb_fill(cell, CLR_OOH)
                    return True
                if normalized in {"RADIO", "CINEMA", "PRINT", "OTHER"}:
                    _apply_rgb_fill(cell, CLR_OTHER)
                    return True
                if normalized == "SUBTOTAL":
                    _apply_rgb_fill(cell, CLR_SUBTOTAL_GRAY)
                    return True
                return False

            base_applied = False
            if col_idx in (0, 1):
                _apply_base_background(cell)
                base_applied = True
            elif total_col_idx is not None and col_idx >= total_col_idx:
                _apply_base_background(cell)
                base_applied = True

            if col_idx == 2:
                if not apply_media_highlight() and not base_applied:
                    _apply_base_background(cell)
            elif 3 <= col_idx <= 14:
                cell_meta = cell_metadata.get(cell_key, {})
                if cell_meta.get("has_data"):
                    if not apply_media_highlight() and not base_applied:
                        _apply_base_background(cell)
                else:
                    cell.fill.background()
            elif not base_applied:
                cell.fill.background()

        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        logger.debug(
            "CELL STYLING [%s,%s]: Final enforcement - cell.vertical_anchor: %s",
            row_idx,
            col_idx,
            cell.vertical_anchor,
        )

        text_frame.word_wrap = wrap_from_config
        text_frame.auto_size = MSO_AUTO_SIZE.NONE

        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            paragraph.alignment = alignment
            paragraph.line_spacing = 1.0

            pPr_final = paragraph._p.get_or_add_pPr()
            spcBef_final = pPr_final.find(qn("a:spcBef"))
            if spcBef_final is not None:
                pPr_final.remove(spcBef_final)
            spcBef_final = OxmlElement("a:spcBef")
            spcPts_before_final = OxmlElement("a:spcPts")
            spcPts_before_final.set("val", "0")
            spcBef_final.append(spcPts_before_final)
            pPr_final.append(spcBef_final)

            spcAft_final = pPr_final.find(qn("a:spcAft"))
            if spcAft_final is not None:
                pPr_final.remove(spcAft_final)
            spcAft_final = OxmlElement("a:spcAft")
            spcPts_after_final = OxmlElement("a:spcPts")
            spcPts_after_final.set("val", "0")
            spcAft_final.append(spcPts_after_final)
            pPr_final.append(spcAft_final)

            logger.debug(
                "CELL STYLING [%s,%s]: Final enforcement - Paragraph %s properties: Alignment=%s, LineSpacing=%s. OXML spcBef/spcAft forced to 0.",
                row_idx,
                col_idx,
                para_idx,
                paragraph.alignment,
                paragraph.line_spacing,
            )

            for run_idx, cell_run in enumerate(paragraph.runs):
                expected_font_name = DEFAULT_FONT_NAME
                expected_font_size = body_font_size if use_compact_font else FONT_SIZE_BODY
                expected_bold = False
                expected_color_rgb = CLR_BLACK

                if row_idx == 0:
                    expected_font_size = FONT_SIZE_HEADER
                    expected_bold = True
                    if col_idx in MONTH_HEADER_COLUMNS:
                        expected_color_rgb = CLR_WHITE
                elif is_total_row:
                    expected_font_size = FONT_SIZE_MONTHLY_TOTAL if is_monthly_total_row else FONT_SIZE_HEADER
                    expected_bold = True
                else:
                    if col_idx < 3:
                        expected_bold = True

                if cell_run.text.strip() in ("", "-"):
                    expected_font_name = empty_font_name
                    expected_font_size = empty_font_size
                    expected_bold = False
                    expected_color_rgb = CLR_LIGHT_GRAY_TEXT

                if cell_run.font.name != expected_font_name:
                    cell_run.font.name = expected_font_name
                    logger.debug(
                        "CELL STYLING [%s,%s]: Final run %s font name re-enforced to %s",
                        row_idx,
                        col_idx,
                        run_idx,
                        expected_font_name,
                    )
                if cell_run.font.size != expected_font_size:
                    cell_run.font.size = expected_font_size
                    logger.debug(
                        "CELL STYLING [%s,%s]: Final run %s font size re-enforced to %s",
                        row_idx,
                        col_idx,
                        run_idx,
                        expected_font_size,
                    )
                if cell_run.font.bold != expected_bold:
                    cell_run.font.bold = expected_bold
                    logger.debug(
                        "CELL STYLING [%s,%s]: Final run %s font bold re-enforced to %s",
                        row_idx,
                        col_idx,
                        run_idx,
                        expected_bold,
                    )

                current_run_color_rgb = None
                if cell_run.font.color.type == MSO_COLOR_TYPE.RGB:
                    current_run_color_rgb = cell_run.font.color.rgb

                if current_run_color_rgb != expected_color_rgb:
                    cell_run.font.color.rgb = expected_color_rgb
                    logger.debug(
                        "CELL STYLING [%s,%s]: Final run %s font color re-enforced",
                        row_idx,
                        col_idx,
                        run_idx,
                    )

        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        logger.debug(
            "CELL STYLING [%s,%s]: Absolutely final re-assertion of cell.vertical_anchor: %s",
            row_idx,
            col_idx,
            cell.vertical_anchor,
        )

        paragraphs_final_count = len(text_frame.paragraphs)
        if paragraphs_final_count > 1:
            paragraphs_to_remove = []
            for i in range(1, len(text_frame.paragraphs)):
                para = text_frame.paragraphs[i]
                if not para.text.strip():
                    paragraphs_to_remove.append(i)

            for para_idx in reversed(paragraphs_to_remove):
                logger.warning(
                    "CELL STYLING [%s,%s]: Found %s extra empty paragraphs - this should not happen with the reuse fix",
                    row_idx,
                    col_idx,
                    len(paragraphs_to_remove),
                )

        logger.debug(
            "CELL STYLING [%s,%s]: Final paragraph count: %s (should be 1 for pixel-perfect alignment)",
            row_idx,
            col_idx,
            len(text_frame.paragraphs),
        )

        logger.debug(
            "UNIVERSAL FORMATTING: Applied explicit Calibri font, centering, and spacing to cell (%s,%s) - Text: '%s'",
            row_idx,
            col_idx,
            processed_cell_text,
        )

    except Exception as exc:
        logger.error("Error styling cell (%s,%s): %s", row_idx, col_idx, exc)
        logger.error(traceback.format_exc())


def add_and_style_table(
    slide: Slide,
    table_data: list[list[str]],
    cell_metadata: dict[tuple[int, int], dict[str, object]],
    layout: TableLayout,
    style_context: CellStyleContext,
    logger: logging.Logger,
) -> bool:
    """Create a table on the slide and apply styling using shared helpers."""

    if not table_data or not table_data[0]:
        logger.warning("No table data provided or insufficient data for table creation")
        return False

    rows = len(table_data)
    cols = len(table_data[0])

    if not layout.position:
        logger.error("Table layout position is not defined")
        return False

    position = layout.position
    table_shape = None

    if layout.placeholder_name:
        for placeholder in slide.placeholders:
            if (
                placeholder.placeholder_format.type == PP_PLACEHOLDER.TABLE
                and getattr(placeholder, "name", "") == layout.placeholder_name
            ):
                try:
                    table_shape = placeholder.insert_table(rows, cols)
                    logger.info("Inserted table into placeholder '%s'", layout.placeholder_name)
                except Exception as exc:
                    logger.warning("Placeholder.insert_table() failed: %s", exc)
                break

    if table_shape is None:
        logger.warning(
            "Placeholder '%s' missing—using add_table() fallback",
            layout.placeholder_name,
        )
        table_shape = slide.shapes.add_table(
            rows,
            cols,
            position["left"],
            position["top"],
            position["width"],
            position["height"],
        )

    table_shape.name = layout.shape_name
    table = table_shape.table

    if layout.top_override is not None:
        table_shape.top = layout.top_override

    logger.info("Creating table with %s rows and %s columns", rows, cols)
    logger.info("=== PRECISION TARGET: Implementing Optimized Row Heights ===")
    logger.info(
        "Header Row Height: %s (%0.3f inches) - TARGET",
        layout.row_height_header,
        layout.row_height_header.inches,
    )
    logger.info(
        "Body Row Height: %s (%0.3f inches) - TARGET",
        layout.row_height_body,
        layout.row_height_body.inches,
    )
    logger.info(
        "Subtotal Row Height: %s (%0.3f inches) - TARGET",
        layout.row_height_subtotal,
        layout.row_height_subtotal.inches,
    )

    for row_index, table_row in enumerate(table.rows):
        if row_index == 0:
            table_row.height = layout.row_height_header
        elif row_index == rows - 1:
            table_row.height = layout.row_height_subtotal
        else:
            table_row.height = layout.row_height_body

        if layout.height_rule_available and layout.height_rule_value is not None:
            try:
                table_row.height_rule = layout.height_rule_value
            except Exception as exc:
                logger.debug("Could not set height rule for row %s: %s", row_index, exc)

    target_height_emu = None
    if layout.position and layout.position.get("height") is not None:
        height_value = layout.position.get("height")
        if hasattr(height_value, "emu"):
            target_height_emu = int(height_value.emu)
        else:
            try:
                target_height_emu = int(height_value)
            except (TypeError, ValueError):
                target_height_emu = None
    if target_height_emu is None:
        target_height_emu = getattr(table_shape, "height", None)

    _scale_table_height(table, target_height_emu)

    if layout.column_widths:
        for col_index, width in enumerate(layout.column_widths[:cols]):
            try:
                table.columns[col_index].width = width
            except Exception as exc:
                logger.debug("Could not set width for column %s: %s", col_index, exc)

    for row_index, row_data in enumerate(table_data):
        for col_index, cell_value in enumerate(row_data[:cols]):
            cell = table.cell(row_index, col_index)
            cell.text = "" if cell_value is None else str(cell_value)
            style_table_cell(
                cell,
                row_index,
                col_index,
                table_data,
                cell_metadata,
                style_context,
                logger,
            )

    apply_table_borders(table, style_context.color_table_gray)
    logger.info("Table created successfully with individual row height constraints")

    empty_font_name = "Verdana"
    empty_font_size = Pt(6)
    for row in table.rows:
        for cell in row.cells:
            if not cell.text or not cell.text.strip() or cell.text.strip() == "-":
                text_frame = cell.text_frame
                text_frame.clear()
                text_frame.text = "\u200b"
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = empty_font_name
                paragraph.font.size = empty_font_size
                run = paragraph.runs[0]
                run.font.name = empty_font_name
                run.font.size = empty_font_size
                run.font.bold = False
                run.font.color.rgb = style_context.color_light_gray_text

    return True
