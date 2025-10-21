"""Helpers for cloning template shapes and tables for pixel-accurate slides."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.slide import Slide
from pptx.oxml.ns import qn
from copy import deepcopy


_REL_ATTRS = {qn("r:embed"), qn("r:link"), qn("r:id")}


def _clone_element(source_shape, target_slide: Slide):
    xml = deepcopy(source_shape.element)
    source_part = source_shape.part
    target_part = target_slide.part

    for node in xml.iter():
        for attr in list(node.attrib):
            if attr in _REL_ATTRS:
                r_id = node.attrib.get(attr)
                if not r_id:
                    continue
                rel = source_part.rels.get(r_id)
                if rel is None:
                    continue
                if rel.is_external:
                    new_rid = target_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                else:
                    new_rid = target_part.relate_to(rel.target_part, rel.reltype)
                node.set(attr, new_rid)

    target_slide.shapes._spTree.append(xml)
    return target_slide.shapes[-1]


class TemplateCloneError(RuntimeError):
    """Raised when required template shapes cannot be cloned."""


def load_template(template_path: Path) -> Presentation:
    return Presentation(template_path)


def clone_template_table(template_slide: Slide, target_slide: Slide, table_name: str) -> Slide:
    for shape in template_slide.shapes:
        if getattr(shape, "name", "") == table_name and getattr(shape, "has_table", False):
            return _clone_element(shape, target_slide)
    raise TemplateCloneError(f"Table '{table_name}' not found on template slide")


def clone_template_shape(template_slide: Slide, target_slide: Slide, shape_name: str):
    for shape in template_slide.shapes:
        if getattr(shape, "name", "") == shape_name:
            return _clone_element(shape, target_slide)
    raise TemplateCloneError(f"Shape '{shape_name}' not found on template slide")
