LEGEND_GROUPS_CONFIG: dict[str, object] = {}
SUMMARY_TILE_CONFIG: dict[str, object] = {}
FONT_FAMILY_LEGEND = "Calibri"

LEGEND_COLOR_WIDTH_IN = 24.84 / 96.0
LEGEND_COLOR_HEIGHT_IN = 13.37 / 96.0
LEGEND_TEXT_GAP_IN = 0.05
LEGEND_MIN_TEXT_WIDTH_IN = 0.25

def _apply_configured_position(shape, position_config):
    if not shape or not position_config:
        return

    try:
        if "left_inches" in position_config:
            shape.left = Inches(float(position_config["left_inches"]))
        if "top_inches" in position_config:
            shape.top = Inches(float(position_config["top_inches"]))
        if "width_inches" in position_config:
            shape.width = Inches(float(position_config["width_inches"]))
        if "height_inches" in position_config:
            shape.height = Inches(float(position_config["height_inches"]))
    except Exception as exc:
        logger.debug("Unable to apply configured position for shape '%s': %s", getattr(shape, "name", "<unnamed>"), exc)


def _prune_paragraph_runs(paragraph):
    if not paragraph.runs:
        return paragraph.add_run()
    primary_run = paragraph.runs[0]
    for extra in paragraph.runs[1:]:
        paragraph._p.remove(extra._r)
    return primary_run


def _set_shape_text(shape, template_shape, text):
    if not shape or not getattr(shape, "has_text_frame", False):
        return

    template_tf = template_shape.text_frame if template_shape and getattr(template_shape, "has_text_frame", False) else None
    text_frame = shape.text_frame
    text_frame.clear()

    if template_tf:
        text_frame.word_wrap = template_tf.word_wrap
        try:
            text_frame.auto_size = template_tf.auto_size
        except AttributeError:
            pass
        for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            try:
                setattr(text_frame, attr, getattr(template_tf, attr))
            except Exception:
                continue

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[idx] if idx < len(text_frame.paragraphs) else text_frame.add_paragraph()
        template_para = template_tf.paragraphs[idx] if template_tf and idx < len(template_tf.paragraphs) else None

        if template_para:
            paragraph.alignment = template_para.alignment
            paragraph.level = template_para.level
            paragraph.line_spacing = template_para.line_spacing
            paragraph.space_before = template_para.space_before
            paragraph.space_after = template_para.space_after
        elif paragraph.alignment is None:
            paragraph.alignment = PP_ALIGN.CENTER

        run = _prune_paragraph_runs(paragraph)
        run.text = line

        if template_para and template_para.runs:
            template_run = template_para.runs[0]
            run.font.name = template_run.font.name
            run.font.size = template_run.font.size
            run.font.bold = template_run.font.bold
            run.font.italic = template_run.font.italic
            template_color = template_run.font.color
            if template_color.type == MSO_COLOR_TYPE.RGB:
                run.font.color.rgb = template_color.rgb
            elif template_color.type == MSO_COLOR_TYPE.SCHEME:
                run.font.color.theme_color = template_color.theme_color
                if getattr(template_color, "brightness", None) is not None:
                    run.font.color.brightness = template_color.brightness
                if getattr(template_color, "tint", None) is not None:
                    run.font.color.tint = template_color.tint
                if getattr(template_color, "shade", None) is not None:
                    run.font.color.shade = template_color.shade
        else:
            run.font.name = FONT_FAMILY_LEGEND
            run.font.size = FONT_SIZE_BODY

    while len(text_frame.paragraphs) > len(lines):
        extra_paragraph = text_frame.paragraphs[len(lines)]
        text_frame._txBody.remove(extra_paragraph._p)


def _get_child_shape(parent_shape, child_name):
    if not parent_shape or not hasattr(parent_shape, "shapes"):
        return None
    for child in parent_shape.shapes:
        if getattr(child, "name", "") == child_name:
            return child
    return None


def _remove_shape_by_name(container, shape_name):
    """Remove a shape (or nested shape) with the provided name from the container."""
    if not shape_name or not container or not hasattr(container, "shapes"):
        return False
    removed = False
    for shape in list(container.shapes):
        name = getattr(shape, "name", "")
        if name == shape_name:
            element = shape._element
            parent = element.getparent()
            if parent is not None:
                parent.remove(element)
            removed = True
            continue
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            if _remove_shape_by_name(shape, shape_name):
                removed = True
    return removed


def _populate_summary_tiles(slide, template_slide, df, combination_row, excel_path):
    if not SUMMARY_TILE_CONFIG:
        return

    for section in ("quarter_budgets", "media_share", "funnel_share"):
        for tile in SUMMARY_TILE_CONFIG.get(section, {}).values():
            shape_name = tile.get("shape") if isinstance(tile, dict) else None
            if not shape_name:
                continue
            if not any(getattr(shape, "name", "") == shape_name for shape in slide.shapes):
                clone_template_shape(template_slide, slide, shape_name)

    footer_cfg = SUMMARY_TILE_CONFIG.get("footer_notes", {})
    footer_shape = footer_cfg.get("shape") if isinstance(footer_cfg, dict) else None
    if footer_shape and not any(getattr(shape, "name", "") == footer_shape for shape in slide.shapes):
        clone_template_shape(template_slide, slide, footer_shape)

    market, brand, year = combination_row
    combo_filter = (
        (df["Country"].astype(str).str.strip() == str(market).strip())
        & (df["Brand"].astype(str).str.strip() == str(brand).strip())
        & (df["Year"].astype(str).str.strip() == str(year).strip())
    )
    subset = df.loc[combo_filter].copy()
    if subset.empty:
        logger.warning("Summary tiles: no data for %s - %s (%s)", market, brand, year)
        return

    total_cost = subset["Total Cost"].sum()

    _populate_quarter_tiles(slide, template_slide, subset)
    _populate_media_share_tiles(slide, template_slide, subset, total_cost)
    _populate_funnel_share_tiles(slide, template_slide, subset, total_cost)
    _populate_footer(slide, template_slide, excel_path)


def _populate_quarter_tiles(slide, template_slide, subset):
    months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    month_values = {month: float(subset[month].sum()) for month in months}

    for quarter_key, config in SUMMARY_TILE_CONFIG.get("quarter_budgets", {}).items():
        shape_name = config.get("shape")
        if not shape_name:
            continue
        shape = next((s for s in slide.shapes if getattr(s, "name", "") == shape_name and getattr(s, "has_text_frame", False)), None)
        if not shape:
            logger.warning("Quarter tile shape '%s' missing", shape_name)
            continue
        template_shape = _get_shape_by_name(template_slide, shape_name)
        _apply_configured_position(shape, config.get("position"))

        quarter_months = {
            "q1": ("Jan", "Feb", "Mar"),
            "q2": ("Apr", "May", "Jun"),
            "q3": ("Jul", "Aug", "Sep"),
            "q4": ("Oct", "Nov", "Dec"),
        }.get(quarter_key.lower(), ())
        value = sum(month_values.get(month, 0.0) for month in quarter_months)
        formatted = _format_tile_value(config, value)
        prefix = config.get("prefix", "")

        _set_shape_text(shape, template_shape, f"{prefix}{formatted}")


def _populate_media_share_tiles(slide, template_slide, subset, total_cost):
    media_group = subset.groupby("Mapped Media Type")["Total Cost"].sum()

    for media_key, config in SUMMARY_TILE_CONFIG.get("media_share", {}).items():
        shape_name = config.get("shape")
        if not shape_name:
            continue

        shape = next((s for s in slide.shapes if getattr(s, "name", "") == shape_name and getattr(s, "has_text_frame", False)), None)
        if not shape:
            logger.warning("Media share shape '%s' missing", shape_name)
            continue
        template_shape = _get_shape_by_name(template_slide, shape_name)
        _apply_configured_position(shape, config.get("position"))

        normalized_label = media_key.capitalize()
        if media_key.lower() == "television":
            lookup_key = "Television"
        elif media_key.lower() == "digital":
            lookup_key = "Digital"
        else:
            lookup_key = "Other"

        value = float(media_group.get(lookup_key, 0.0))
        formatted = _format_percentage_tile(config, value, total_cost)
        label = config.get("label", normalized_label)

        _set_shape_text(shape, template_shape, f"{label}: {formatted}")


def _populate_funnel_share_tiles(slide, template_slide, subset, total_cost):
    funnel_group = subset.groupby("Funnel Stage")["Total Cost"].sum()

    for funnel_key, config in SUMMARY_TILE_CONFIG.get("funnel_share", {}).items():
        shape_name = config.get("shape")
        if not shape_name:
            continue
        shape = next((s for s in slide.shapes if getattr(s, "name", "") == shape_name and getattr(s, "has_text_frame", False)), None)
        if not shape:
            logger.warning("Funnel share shape '%s' missing", shape_name)
            continue
        template_shape = _get_shape_by_name(template_slide, shape_name)
        _apply_configured_position(shape, config.get("position"))

        lookup_key = {
            "awareness": "Awareness",
            "consideration": "Consideration",
            "purchase": "Purchase",
        }.get(funnel_key.lower(), funnel_key)

        value = float(funnel_group.get(lookup_key, 0.0))
        formatted = _format_percentage_tile(config, value, total_cost)
        label = config.get("label", lookup_key[:3].upper())

        _set_shape_text(shape, template_shape, f"{label}: {formatted}")


def _apply_legend_color(shape, media_key):
    if not shape or not hasattr(shape, "fill"):
        return

    color_lookup = {
        "television": CLR_TELEVISION,
        "tv": CLR_TELEVISION,
        "digital": CLR_DIGITAL,
        "ooh": CLR_OOH,
        "other": CLR_OTHER,
        "radio": CLR_OTHER,
        "cinema": CLR_OTHER,
        "print": CLR_OTHER,
    }

    color = color_lookup.get(str(media_key).strip().lower())
    if not color:
        return

    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        if hasattr(shape, "line"):
            shape.line.fill.background()
    except Exception as exc:
        logger.debug("Unable to apply legend color for '%s': %s", media_key, exc)


def _create_legend_color_shape(slide, shape_name, position):
    if not position:
        return None

    try:
        left = Inches(float(position.get("left_inches", 0.0)))
        top = Inches(float(position.get("top_inches", 0.0)))
        total_height = Inches(float(position.get("height_inches", LEGEND_COLOR_HEIGHT_IN)))
    except Exception as exc:
        logger.debug("Invalid legend position config for '%s': %s", shape_name, exc)
        return None

    color_width = Inches(LEGEND_COLOR_WIDTH_IN)
    color_height = Inches(LEGEND_COLOR_HEIGHT_IN)
    color_top = top + max((total_height - color_height) / 2, 0)

    color_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left,
        color_top,
        color_width,
        color_height,
    )
    color_shape.name = shape_name or f"LegendColor_{len(slide.shapes)}"
    color_shape.line.fill.background()
    return color_shape


def _create_legend_text_shape(slide, shape_name, position):
    if not position:
        return None

    try:
        left = Inches(float(position.get("left_inches", 0.0)))
        top = Inches(float(position.get("top_inches", 0.0)))
        total_width = Inches(float(position.get("width_inches", LEGEND_COLOR_WIDTH_IN + LEGEND_TEXT_GAP_IN + LEGEND_MIN_TEXT_WIDTH_IN)))
        total_height = Inches(float(position.get("height_inches", LEGEND_COLOR_HEIGHT_IN)))
    except Exception as exc:
        logger.debug("Invalid legend position config for '%s': %s", shape_name, exc)
        return None

    text_left = left + Inches(LEGEND_COLOR_WIDTH_IN + LEGEND_TEXT_GAP_IN)
    text_width = max(total_width - (text_left - left), Inches(LEGEND_MIN_TEXT_WIDTH_IN))

    text_shape = slide.shapes.add_textbox(text_left, top, text_width, total_height)
    text_shape.name = shape_name or f"LegendText_{len(slide.shapes)}"
    return text_shape


def _set_legend_text(shape, template_shape, text):
    if template_shape and getattr(template_shape, "has_text_frame", False):
        _set_shape_text(shape, template_shape, text)
        return

    if not shape or not getattr(shape, "has_text_frame", False):
        return

    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_FAMILY_LEGEND
    run.font.size = FONT_SIZE_LEGEND
    run.font.bold = True


def _ensure_legend_shapes(slide, template_slide):
    if not LEGEND_GROUPS_CONFIG:
        return
    for media_key, legend_cfg in LEGEND_GROUPS_CONFIG.items():
        group_name = legend_cfg.get("group_shape")
        color_shape_name = legend_cfg.get("color_shape")
        text_shape_name = legend_cfg.get("text_shape")
        position = legend_cfg.get("position", {})

        template_group_shape = _get_shape_by_name(template_slide, group_name) if (template_slide and group_name) else None
        template_color_shape = None
        template_text_shape = None

        if template_group_shape:
            if color_shape_name:
                template_color_shape = _get_child_shape(template_group_shape, color_shape_name)
            if text_shape_name:
                template_text_shape = _get_child_shape(template_group_shape, text_shape_name)
        else:
            if template_slide and color_shape_name:
                template_color_shape = _get_shape_by_name(template_slide, color_shape_name)
            if template_slide and text_shape_name:
                template_text_shape = _get_shape_by_name(template_slide, text_shape_name)

        template_has_entry = any(
            shape is not None for shape in (template_group_shape, template_color_shape, template_text_shape)
        )

        if template_slide is not None and not template_has_entry:
            for candidate in (group_name, color_shape_name, text_shape_name):
                _remove_shape_by_name(slide, candidate)
            continue

        group_shape = None
        if group_name:
            group_shape = _get_shape_by_name(slide, group_name)
            if not group_shape and template_group_shape is not None:
                try:
                    group_shape = clone_template_shape(template_slide, slide, group_name)
                except TemplateCloneError as exc:
                    logger.debug("Unable to clone legend group '%s': %s", group_name, exc)
                    group_shape = None
            if group_shape:
                _apply_configured_position(group_shape, legend_cfg.get("position"))

        if color_shape_name:
            if group_shape:
                color_shape = _get_child_shape(group_shape, color_shape_name)
            else:
                color_shape = _get_shape_by_name(slide, color_shape_name)
                if color_shape is None and template_color_shape is not None:
                    try:
                        color_shape = clone_template_shape(template_slide, slide, color_shape_name)
                    except TemplateCloneError:
                        color_shape = _create_legend_color_shape(slide, color_shape_name, position)
                if color_shape is None and template_color_shape is None:
                    color_shape = _create_legend_color_shape(slide, color_shape_name, position)
                elif color_shape and position:
                    _apply_configured_position(color_shape, position)
            if color_shape:
                _apply_legend_color(color_shape, media_key)

        if text_shape_name:
            if group_shape:
                text_shape = _get_child_shape(group_shape, text_shape_name)
            else:
                text_shape = _get_shape_by_name(slide, text_shape_name)
                if text_shape is None and template_text_shape is not None:
                    try:
                        text_shape = clone_template_shape(template_slide, slide, text_shape_name)
                    except TemplateCloneError:
                        text_shape = _create_legend_text_shape(slide, text_shape_name, position)
                if text_shape is None and template_text_shape is None:
                    text_shape = _create_legend_text_shape(slide, text_shape_name, position)
                elif text_shape and position:
                    _apply_configured_position(text_shape, position)
            if text_shape:
                if template_text_shape and getattr(template_text_shape, "has_text_frame", False):
                    template_text = template_text_shape.text_frame.text
                else:
                    template_text = str(media_key).upper()
                _set_legend_text(text_shape, template_text_shape, template_text)


def _populate_footer(slide, template_slide, excel_path):
    config = SUMMARY_TILE_CONFIG.get("footer_notes", {})
    shape_name = config.get("shape")
    if not shape_name:
        return
    shape = next((s for s in slide.shapes if getattr(s, "name", "") == shape_name and getattr(s, "has_text_frame", False)), None)
    if not shape:
        logger.warning("Footer shape '%s' missing", shape_name)
        return
    template_shape = _get_shape_by_name(template_slide, shape_name)
    _apply_configured_position(shape, config.get("position"))

    import re

    text = config.get("default_text", "")
    if config.get("append_date"):
        stamp_raw = _extract_export_date(excel_path, config.get("append_date_format", "%d_%m_%y"))
        if stamp_raw:
            normalized = re.sub(r"[^0-9]", "", stamp_raw)
            if not normalized:
                normalized = stamp_raw.replace("_", "").strip()
            if normalized:
                text = text.replace("DD_MM_YY", normalized).replace("DDMMYY", normalized)
                text = text.replace(f"{normalized}_Lumina", f"{normalized} Lumina")

    _set_shape_text(shape, template_shape, text)


def _extract_export_date(excel_path, output_format):
    from datetime import datetime
    import re

    path_str = str(excel_path)
    match = re.search(r"(\d{4})_(\d{2})_(\d{2})", path_str)
    if match:
        year, month, day = match.groups()
        try:
            dt = datetime(int(year), int(month), int(day))
            return dt.strftime(output_format)
        except ValueError:
            logger.debug("Failed to parse export date from filename '%s'", path_str)

    try:
        file_time = datetime.fromtimestamp(Path(excel_path).stat().st_mtime)
        return file_time.strftime(output_format)
    except Exception as exc:
        logger.debug("Falling back to current timestamp for export date: %s", exc)
        return datetime.now().strftime(output_format)


def _format_tile_value(config, value):
    scale = float(config.get("scale", 1.0))
    fmt = config.get("number_format", "{value}")
    try:
        return fmt.format(value=value * scale)
    except Exception as exc:
        logger.warning("Failed to format tile value %s with format %s: %s", value, fmt, exc)
        return str(value * scale)


def _format_percentage_tile(config, value, total):
    if total <= 0:
        pct = 0.0
    else:
        pct = value / total
    return _format_tile_value(dict(config, scale=config.get("scale", 100)), pct)
def _ensure_shape_on_slide(slide, template_slide, shape_name):
    shape = _get_shape_by_name(slide, shape_name)
    if shape:
        logger.debug("Shape '%s' already present on slide", shape_name)
        return shape

    template_shape = _get_shape_by_name(template_slide, shape_name)
    if not template_shape:
        logger.warning("Template missing shape '%s'", shape_name)
        return None

    if template_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        new_shape = _copy_shape(template_shape, slide, new_name=shape_name)
    elif template_shape.has_text_frame:
        new_shape = _copy_text_box(template_shape, slide, new_name=shape_name)
    else:
        logger.warning("Unable to clone shape '%s' of unsupported type %s", shape_name, template_shape.shape_type)
        new_shape = None
    if new_shape:
        logger.debug("Cloned shape '%s' onto slide", shape_name)
    return new_shape


def _compose_title_text(
    combination_row: tuple[str, str, int],
    slide_title_suffix: str,
) -> str:
    market, brand, year = combination_row
    market_display = _normalize_market_name(market)
    title_format = presentation_config.get("title", {}).get("format", "{market} - {brand}")
    return title_format.format(market=market_display, brand=brand, year=year) + slide_title_suffix


def _apply_title(slide, template_slide, combination_row, slide_title_suffix):
    shape_name = presentation_config.get("title", {}).get("shape") or SHAPE_NAME_TITLE
    title_shape = next((s for s in slide.shapes if getattr(s, "name", "") == shape_name and getattr(s, "has_text_frame", False)), None)

    title_text = _compose_title_text(combination_row, slide_title_suffix)

    if not title_shape:
        logger.warning("Title shape '%s' not available on slide", shape_name)
        return title_text

    text_frame = title_shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = title_text

    if paragraph.runs:
        run = paragraph.runs[0]
        _ensure_font_consistency(
            run.font,
            DEFAULT_FONT_NAME,
            FONT_SIZE_TITLE,
            True,
            CLR_WHITE,
        )

    return title_text


def _normalize_market_name(raw_market: str) -> str:
    if raw_market == "MOR":
        return "Morocco"
    return raw_market


def _normalize_row_label(value: object) -> str:
    if value is None:
        return ""
    return " ".join(str(value).replace("\xa0", " ").split()).upper()


def _derive_campaign_blocks_for_table(table_data: list[list[str]]) -> list[tuple[int, int]]:
    """Return [(campaign_start_row, monthly_total_row)] for each campaign block."""
    if not table_data or len(table_data) <= 1:
        return []

    blocks: list[tuple[int, int]] = []
    idx = 1  # skip header
    total_label = "MONTHLY TOTAL (£ 000)"
    skip_labels = {"", "-", total_label, "CARRIED FORWARD", "GRAND TOTAL"}

    while idx < len(table_data) - 1:
        label = _normalize_row_label(table_data[idx][0] if idx < len(table_data) and table_data[idx] else "")
        if label in skip_labels:
            idx += 1
            continue

        start_idx = idx
        search_idx = idx + 1
        monthly_idx = None

        while search_idx < len(table_data):
            search_label = _normalize_row_label(
                table_data[search_idx][0] if search_idx < len(table_data) and table_data[search_idx] else ""
            )
            if search_label == total_label:
                monthly_idx = search_idx
                break
            if search_label and search_label not in {"", "-", total_label}:
                # Reached next campaign without finding monthly total; abort this block.
                break
            search_idx += 1

        if monthly_idx is None:
            idx += 1
            continue

        blocks.append((start_idx, monthly_idx))
        idx = monthly_idx + 1

    return blocks


def _apply_campaign_cell_merges(table, table_data: list[list[str]]) -> None:
    """Merge campaign name column cells and monthly total headers to prevent row inflation."""
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.text import MSO_VERTICAL_ANCHOR

    blocks = _derive_campaign_blocks_for_table(table_data)
    if not blocks:
        return

    for start_idx, monthly_idx in blocks:
        try:
            if start_idx >= len(table.rows) or monthly_idx >= len(table.rows):
                continue

            merge_end = max(start_idx, monthly_idx - 1)
            if merge_end > start_idx:
                top_cell = table.cell(start_idx, 0)
                for merge_row in range(start_idx + 1, merge_end + 1):
                    table.cell(merge_row, 0).text = ""
                merged_cell = top_cell.merge(table.cell(merge_end, 0))
                merged_label = str(table_data[start_idx][0] or "")
                merged_cell.text = merged_label
                try:
                    merged_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                except Exception:
                    pass
                try:
                    merged_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                except Exception:
                    pass

            if monthly_idx >= len(table.rows):
                continue

            left_cell = table.cell(monthly_idx, 0)
            label_text = str(left_cell.text or "")
            for col_idx in (1, 2):
                if col_idx < len(table.columns):
                    table.cell(monthly_idx, col_idx).text = ""
            merged_total = left_cell.merge(table.cell(monthly_idx, 2))
            merged_total.text = label_text
            try:
                merged_total.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            except Exception:
                pass
            try:
                merged_total.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            except Exception:
                pass
        except Exception as exc:
            logger.debug(
                "Skipping campaign cell merge for rows %s-%s: %s",
                start_idx,
                monthly_idx,
                exc,
            )


def _clear_comments(slide):
    if not comments_config.get("enabled", False):
        return

    for shape_name in (SHAPE_NAME_COMMENTS_TITLE, SHAPE_NAME_COMMENTS_BOX):
        shape = _get_shape_by_name(slide, shape_name)
        if shape and shape.has_text_frame:
            shape.text_frame.text = ""
"""
Excel to PowerPoint Automation Script - GEOGRAPHY VERSION
This version uses the 'Plan - Geography' column instead of '**Country' column
to extract market/country information from the hierarchical geography string.

Changes from FINAL_PRODUCTION.py:
- Uses 'Plan - Geography' column (column K) instead of '**Country'
- Same extraction logic: splits by " | " and takes the last part
- Example: "Global | EMEA | MEA | Pakistan" → "Pakistan"
"""
import os
import logging
import traceback
from datetime import datetime
import pandas as pd
import numpy as np
import ast, pathlib, inspect, textwrap
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from amp_automation.presentation.template_geometry import (
    TEMPLATE_V4_COLUMN_WIDTHS_EMU,
    TEMPLATE_V4_COLUMN_WIDTHS_INCHES,
    TEMPLATE_V4_ROW_HEIGHT_BODY_EMU,
    TEMPLATE_V4_ROW_HEIGHT_BODY_INCHES,
    TEMPLATE_V4_ROW_HEIGHT_HEADER_EMU,
    TEMPLATE_V4_ROW_HEIGHT_HEADER_INCHES,
    TEMPLATE_V4_ROW_HEIGHT_TRAILER_EMU,
    TEMPLATE_V4_TABLE_BOUNDS,
    TEMPLATE_V4_TABLE_LEFT_EMU,
    TEMPLATE_V4_TABLE_TOP_EMU,
    TEMPLATE_V4_TABLE_WIDTH_EMU,
    TEMPLATE_V4_TABLE_HEIGHT_EMU,
)
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL, MSO_FILL_TYPE

from amp_automation.config import Config, load_master_config
from amp_automation.data import (
    get_month_specific_tv_metrics,
    load_and_prepare_data as modular_load_and_prepare_data,
)
from amp_automation.presentation.charts import (
    ChartStyleContext,
    add_pie_chart as presentation_add_pie_chart,
    prepare_campaign_type_chart_data,
    prepare_funnel_chart_data,
    prepare_media_type_chart_data,
)
from amp_automation.presentation.tables import (
    CellStyleContext,
    apply_table_borders,
    ensure_font_consistency as _ensure_font_consistency,
    style_table_cell,
)
from amp_automation.presentation.template_clone import TemplateCloneError, clone_template_shape, clone_template_table
from amp_automation.utils.media import normalize_media_type
from amp_automation.tooling import autopptx_adapter, aspose_converter, docstrange_validator
from amp_automation.tooling.autopptx_adapter import SlidePayload

REQUIRED_SHAPE_NAMES: set[str] = set()
CHANGES_INFO_RELTYPE = "http://schemas.microsoft.com/office/2016/11/relationships/changesInfo"

def _rgb_color(config_value, fallback):
    try:
        r, g, b = config_value
        return RGBColor(int(r), int(g), int(b))
    except (TypeError, ValueError):  # pragma: no cover - defensive fallback
        return RGBColor(*fallback)


def _coord_from_config(config_mapping, fallback):
    return {
        "left": float(config_mapping.get("left_inches", config_mapping.get("left", fallback["left"]))),
        "top": float(config_mapping.get("top_inches", config_mapping.get("top", fallback["top"]))),
        "width": float(config_mapping.get("width_inches", config_mapping.get("width", fallback["width"]))),
        "height": float(config_mapping.get("height_inches", config_mapping.get("height", fallback["height"]))),
}


def _parse_alignment_map(mapping: dict | list | None) -> dict[int, object]:
    result: dict[int, object] = {}
    if not mapping:
        return result
    for key, value in mapping.items():
        key_str = str(key)
        if "-" in key_str:
            start_str, end_str = key_str.split("-", 1)
            try:
                start = int(start_str)
                end = int(end_str)
            except ValueError:
                continue
            for idx in range(start, end + 1):
                result[idx] = value
        else:
            try:
                result[int(key_str)] = value
            except ValueError:
                continue
    return result


def _build_dual_line_map(mapping: dict | None) -> dict[str, list[str]]:
    result: dict[str, list[str]] = {}
    if not mapping:
        return result
    for key, parts in mapping.items():
        key_upper = str(key).upper()
        if isinstance(parts, (list, tuple)):
            result[key_upper] = [str(part) for part in parts]
        else:
            result[key_upper] = [str(parts)]
    return result


MARGIN_EMU_LR = 45720  # Approx Pt(3.6), from template analysis for left/right cell margins
ZERO_THRESHOLD = 0.01  # Values below this (absolute) are treated as zero for display/coloring

TABLE_MONTH_ORDER = [
    "Jan",
    "Feb",
    "Mar",
    "Apr",
    "May",
    "Jun",
    "Jul",
    "Aug",
    "Sep",
    "Oct",
    "Nov",
    "Dec",
]
TABLE_MONTH_HEADERS = [month.upper() for month in TABLE_MONTH_ORDER]
TABLE_HEADER_COLUMNS = [
    "CAMPAIGN",
    "MEDIA",
    "METRICS",
    *TABLE_MONTH_HEADERS,
    "TOTAL",
    "GRPs",
    "%",
]

MEDIA_DISPLAY_ORDER = [
    "Television",
    "Digital",
    "OOH",
    "Radio",
    "Cinema",
    "Print",
    "Other",
]

MEDIA_DISPLAY_LABELS = {
    "Television": "TELEVISION",
    "Digital": "DIGITAL",
    "OOH": "OOH",
    "Radio": "RADIO",
    "Cinema": "CINEMA",
    "Print": "PRINT",
    "Other": "OTHER",
}

_CAMPAIGN_BOUNDARIES: list[tuple[int, int]] = []
TABLE_COLUMN_ALIGNMENT_MAP: dict[int, object] = {}
TABLE_WORD_WRAP_COLUMNS: set[int] = set()
TABLE_UPPERCASE_COLUMNS: set[int] = set()
TABLE_DUAL_LINE_LABELS: dict[str, list[str]] = {}
TABLE_SHRINK_TO_FIT_COLUMNS: set[int] = set()
CLONE_PIPELINE_ENABLED: bool = True
AUTOPPTX_CONFIG: dict[str, object] = {}
ASPOSE_CONFIG: dict[str, object] = {}
DOCSTRANGE_CONFIG: dict[str, object] = {}


def _normalized_media_value(raw_media: str) -> str:
    if not raw_media:
        return "Other"
    normalized = normalize_media_type(str(raw_media))
    if normalized not in MEDIA_DISPLAY_LABELS:
        return "Other"
    return normalized


def _media_display_label(media_key: str) -> str:
    return MEDIA_DISPLAY_LABELS.get(media_key, media_key.upper())


def _set_cell_metadata(
    metadata: dict[tuple[int, int], dict[str, object]],
    row_idx: int,
    col_idx: int,
    value: float | int | None,
    media_type: str,
    has_data: bool,
) -> None:
    metadata[(row_idx, col_idx)] = {
        "has_data": has_data,
        "media_type": media_type,
        "value": value or 0,
    }


def _format_budget_cell(value: float) -> str:
    formatted = format_number(value, is_budget=True, is_monthly_column=True)
    return formatted if formatted else "-"


def _format_total_budget(value: float) -> str:
    formatted = format_number(value, is_budget=True)
    return formatted if formatted else "-"


def _format_percentage_cell(value: float) -> str:
    formatted = format_number(value, is_percentage=True)
    return formatted if formatted else "0.0%"


def _collect_monthly_values(media_df: pd.DataFrame) -> list[float]:
    values: list[float] = []
    for month in TABLE_MONTH_ORDER:
        if month in media_df.columns:
            values.append(float(media_df[month].sum()))
        else:
            values.append(0.0)
    return values


def _build_campaign_monthly_total_row(
    row_idx: int,
    month_totals: list[float],
    cell_metadata: dict[tuple[int, int], dict[str, object]],
) -> list[str]:
    row: list[str] = ["MONTHLY TOTAL (£ 000)", "", ""]
    for month_idx, value in enumerate(month_totals):
        formatted = _format_budget_cell(value)
        col_idx = 3 + month_idx
        row.append(formatted)
        _set_cell_metadata(cell_metadata, row_idx, col_idx, value, "Subtotal", not is_empty_formatted_value(formatted))

    total_value = sum(month_totals)
    total_col_idx = 3 + len(TABLE_MONTH_ORDER)
    total_formatted = _format_total_budget(total_value)
    row.append(total_formatted)
    _set_cell_metadata(
        cell_metadata,
        row_idx,
        total_col_idx,
        total_value,
        "Subtotal",
        not is_empty_formatted_value(total_formatted),
    )

    row.extend(["", ""])
    return row


def _build_digital_metric_rows(
    start_row_idx: int,
    cell_metadata: dict[tuple[int, int], dict[str, object]],
) -> list[list[str]]:
    metric_labels = ["YT Reach", "META Reach", "TT Reach"]
    rows: list[list[str]] = []

    for label in metric_labels:
        row_idx = start_row_idx + len(rows)
        row: list[str] = ["-", "-", label]
        for month_idx in range(len(TABLE_MONTH_ORDER)):
            col_idx = 3 + month_idx
            row.append("-")
            _set_cell_metadata(cell_metadata, row_idx, col_idx, 0.0, "Digital", False)

        total_col_idx = 3 + len(TABLE_MONTH_ORDER)
        row.append("-")
        _set_cell_metadata(cell_metadata, row_idx, total_col_idx, 0.0, "Digital", False)

        row.extend(["", ""])
        rows.append(row)

    return rows


def _build_tv_metric_rows(
    campaign_name: str,
    region: str,
    masterbrand: str,
    year: int | None,
    excel_path: str | Path | None,
    start_row_idx: int,
    cell_metadata: dict[tuple[int, int], dict[str, object]],
) -> tuple[list[list[str]], float]:
    if not excel_path or year is None:
        return [], 0.0

    rows: list[list[str]] = []
    grp_totals: list[float] = []
    reach1_totals: list[float] = []
    freq_totals: list[float] = []

    for month in TABLE_MONTH_ORDER:
        metrics = get_month_specific_tv_metrics(
            excel_path,
            region,
            masterbrand,
            campaign_name,
            year,
            month,
        )

        grp_totals.append(float(metrics.get("grp_sum", 0.0) or 0.0))
        reach1 = metrics.get("reach1_avg", 0.0) or 0.0
        reach1_totals.append(float(reach1) * 100.0)
        freq_totals.append(float(metrics.get("frequency_avg", 0.0) or 0.0))

    grp_row_idx = start_row_idx
    grp_row: list[str] = ["-", "-", "GRPs"]
    grp_sum = 0.0
    for month_idx, value in enumerate(grp_totals):
        formatted = format_number(value, is_grp=True)
        if not formatted:
            formatted = "-"
        grp_row.append(formatted)
        col_idx = 3 + month_idx
        _set_cell_metadata(
            cell_metadata,
            grp_row_idx,
            col_idx,
            value,
            "GRPs",
            not is_empty_formatted_value(formatted),
        )
        grp_sum += value

    total_col_idx = 3 + len(TABLE_MONTH_ORDER)
    total_formatted = format_number(grp_sum, is_grp=True)
    grp_row.append(total_formatted if total_formatted else "-")
    _set_cell_metadata(
        cell_metadata,
        grp_row_idx,
        total_col_idx,
        grp_sum,
        "GRPs",
        not is_empty_formatted_value(total_formatted),
    )
    grp_row.extend([format_number(grp_sum, is_grp=True), ""])
    rows.append(grp_row)

    reach_row_idx = start_row_idx + len(rows)
    reach_row: list[str] = ["-", "-", "Reach@1+"]
    for month_idx, value in enumerate(reach1_totals):
        formatted = _format_percentage_cell(value)
        reach_row.append(formatted if formatted else "-")
        col_idx = 3 + month_idx
        _set_cell_metadata(
            cell_metadata,
            reach_row_idx,
            col_idx,
            value,
            "Reach",
            not is_empty_formatted_value(formatted),
        )

    reach_row.append("-")  # TOTAL column
    _set_cell_metadata(cell_metadata, reach_row_idx, total_col_idx, 0.0, "Reach", False)
    reach_row.extend(["", ""])
    rows.append(reach_row)

    ots_row_idx = start_row_idx + len(rows)
    ots_row: list[str] = ["-", "-", "OTS@3+"]
    for month_idx, value in enumerate(freq_totals):
        formatted = format_number(value, is_grp=False)
        if not formatted:
            formatted = "-"
        ots_row.append(formatted)
        col_idx = 3 + month_idx
        _set_cell_metadata(
            cell_metadata,
            ots_row_idx,
            col_idx,
            value,
            "OTS",
            not is_empty_formatted_value(formatted),
        )

    ots_row.append("-")
    _set_cell_metadata(cell_metadata, ots_row_idx, total_col_idx, 0.0, "OTS", False)
    ots_row.extend(["", ""])
    rows.append(ots_row)

    return rows, grp_sum


def _build_budget_row(
    campaign_label: str,
    media_label: str,
    media_key: str,
    monthly_values: list[float],
    total_cost: float,
    share_percentage: float | None,
    row_idx: int,
    cell_metadata: dict[tuple[int, int], dict[str, object]],
) -> list[str]:
    row: list[str] = [campaign_label, media_label, "£ 000"]

    for month_idx, value in enumerate(monthly_values):
        formatted = _format_budget_cell(value)
        col_idx = 3 + month_idx
        row.append(formatted)
        _set_cell_metadata(
            cell_metadata,
            row_idx,
            col_idx,
            value,
            media_key,
            not is_empty_formatted_value(formatted),
        )

    total_col_idx = 3 + len(TABLE_MONTH_ORDER)
    total_formatted = _format_total_budget(total_cost)
    row.append(total_formatted)
    _set_cell_metadata(
        cell_metadata,
        row_idx,
        total_col_idx,
        total_cost,
        media_key,
        not is_empty_formatted_value(total_formatted),
    )

    row.append("")  # GRPs column blank for budget row
    if share_percentage is not None:
        row.append(_format_percentage_cell(share_percentage))
    else:
        row.append("")

    return row


def _build_campaign_block(
    campaign_name: str,
    campaign_df: pd.DataFrame,
    base_row_idx: int,
    total_budget_for_percentage: float,
    cell_metadata: dict[tuple[int, int], dict[str, object]],
    region: str,
    masterbrand: str,
    year: int | None,
    excel_path: str | Path | None,
) -> tuple[list[list[str]], list[float], float]:
    campaign_df = campaign_df.copy()
    campaign_df["_NormalizedMedia"] = campaign_df["Mapped Media Type"].apply(_normalized_media_value)

    block_rows: list[list[str]] = []
    block_month_totals = [0.0] * len(TABLE_MONTH_ORDER)
    block_grp_total = 0.0

    campaign_total_budget = float(campaign_df["Total Cost"].sum() or 0.0)
    share_percentage = (
        (campaign_total_budget / total_budget_for_percentage) * 100.0
        if total_budget_for_percentage > 0
        else 0.0
    )

    first_media = True

    for media_key in MEDIA_DISPLAY_ORDER:
        media_mask = (
            campaign_df["_NormalizedMedia"].astype(str).str.lower()
            == media_key.lower()
        )
        media_df = campaign_df[media_mask]
        if media_df.empty:
            continue

        monthly_values = _collect_monthly_values(media_df)
        block_month_totals = [
            existing + value for existing, value in zip(block_month_totals, monthly_values)
        ]
        total_cost = float(media_df["Total Cost"].sum() or 0.0)

        row_idx = base_row_idx + len(block_rows)
        campaign_label = str(campaign_name).upper() if first_media else "-"

        row = _build_budget_row(
            campaign_label,
            _media_display_label(media_key),
            media_key,
            monthly_values,
            total_cost,
            share_percentage if first_media else None,
            row_idx,
            cell_metadata,
        )
        block_rows.append(row)

        if media_key == "Television":
            tv_rows, tv_grp_total = _build_tv_metric_rows(
                str(campaign_name),
                region,
                masterbrand,
                year,
                excel_path,
                base_row_idx + len(block_rows),
                cell_metadata,
            )
            if tv_rows:
                block_rows.extend(tv_rows)
                block_grp_total += tv_grp_total
        elif media_key == "Digital":
            digital_rows = _build_digital_metric_rows(
                base_row_idx + len(block_rows),
                cell_metadata,
            )
            if digital_rows:
                block_rows.extend(digital_rows)

        first_media = False

    total_row = _build_campaign_monthly_total_row(
        base_row_idx + len(block_rows),
        block_month_totals,
        cell_metadata,
    )
    block_rows.append(total_row)

    return block_rows, block_month_totals, block_grp_total


def _build_grand_total_row(
    monthly_totals: list[float],
    total_budget: float,
    grand_total_grp: float,
) -> list[str]:
    row: list[str] = ["GRAND TOTAL", "", ""]
    for value in monthly_totals:
        row.append(_format_budget_cell(value))

    row.append(_format_total_budget(total_budget))
    row.append(format_number(grand_total_grp, is_grp=True))
    row.append(_format_percentage_cell(100.0))
    return row


def _build_carried_forward_row(
    month_totals: list[float],
    total_budget: float,
    grp_total: float,
) -> list[str]:
    row: list[str] = ["CARRIED FORWARD", "", ""]
    for value in month_totals:
        row.append(_format_budget_cell(value))

    row.append(_format_total_budget(total_budget))
    row.append(format_number(grp_total, is_grp=True) if grp_total else "")
    row.append("")
    return row


def _build_carried_forward_metadata_values(
    month_totals: list[float],
    total_budget: float,
    grp_total: float,
) -> dict[int, dict[str, object]]:
    metadata: dict[int, dict[str, object]] = {}
    month_start_col = 3
    for idx, value in enumerate(month_totals):
        col_idx = month_start_col + idx
        metadata[col_idx] = {
            "has_data": abs(value) > ZERO_THRESHOLD,
            "media_type": "Subtotal",
            "value": value,
        }

    total_col_idx = month_start_col + len(month_totals)
    metadata[total_col_idx] = {
        "has_data": abs(total_budget) > ZERO_THRESHOLD,
        "media_type": "Subtotal",
        "value": total_budget,
    }

    grp_col_idx = total_col_idx + 1
    if abs(grp_total) > ZERO_THRESHOLD:
        metadata[grp_col_idx] = {
            "has_data": True,
            "media_type": "GRPs",
            "value": grp_total,
        }

    return metadata


def _build_grand_total_metadata_values(
    month_totals: list[float],
    total_budget: float,
    grand_total_grp: float,
) -> dict[int, dict[str, object]]:
    metadata: dict[int, dict[str, object]] = {}
    month_start_col = 3
    for idx, value in enumerate(month_totals):
        col_idx = month_start_col + idx
        metadata[col_idx] = {
            "has_data": abs(value) > ZERO_THRESHOLD,
            "media_type": "Subtotal",
            "value": value,
        }

    total_col_idx = month_start_col + len(month_totals)
    metadata[total_col_idx] = {
        "has_data": abs(total_budget) > ZERO_THRESHOLD,
        "media_type": "Subtotal",
        "value": total_budget,
    }

    grp_col_idx = total_col_idx + 1
    metadata[grp_col_idx] = {
        "has_data": abs(grand_total_grp) > ZERO_THRESHOLD,
        "media_type": "GRPs",
        "value": grand_total_grp,
    }

    pct_col_idx = grp_col_idx + 1
    metadata[pct_col_idx] = {
        "has_data": True,
        "media_type": "Subtotal",
        "value": 100.0,
    }

    return metadata


def _coerce_year(value: object) -> int | None:
    if value is None:
        return None
    try:
        return int(str(value).strip())
    except (ValueError, TypeError):
        return None


def _initialize_from_config(config: Config) -> None:
    global MASTER_CONFIG
    global presentation_config, fonts_config, font_sizes
    global colors_config, media_colors_config, ui_colors_config
    global table_config, row_heights_config, table_position_config
    global comments_config, comments_title_pos_config, comments_box_pos_config
    global title_position_config, charts_config, chart_positions_config
    global summary_tiles_config
    global TABLE_PLACEHOLDER_NAME
    global CLR_BLACK, CLR_WHITE, CLR_LIGHT_GRAY_TEXT, CLR_TABLE_GRAY, CLR_HEADER_GREEN
    global CLR_COMMENTS_GRAY, CLR_SUBTOTAL_GRAY
    global CLR_TELEVISION, CLR_DIGITAL, CLR_OOH, CLR_OTHER
    global DEFAULT_FONT_NAME, FONT_SIZE_HEADER, FONT_SIZE_BODY, FONT_SIZE_BODY_COMPACT
    global FONT_SIZE_CHART_TITLE, FONT_SIZE_CHART_LABELS
    global FONT_SIZE_TITLE, FONT_SIZE_LEGEND, FONT_SIZE_COMMENTS
    global TABLE_ROW_HEIGHT_HEADER, TABLE_ROW_HEIGHT_BODY, TABLE_ROW_HEIGHT_SUBTOTAL, TABLE_ROW_HEIGHT_TRAILER
    global TABLE_COLUMN_WIDTHS, TABLE_TOP_OVERRIDE
    global TABLE_CELL_STYLE_CONTEXT
    global CHART_STYLE_CONTEXT, CHART_COLOR_MAPPING, CHART_COLOR_CYCLE
    global MAX_ROWS_PER_SLIDE, SPLIT_STRATEGY, SHOW_CHARTS_ON_SPLITS, SHOW_CARRIED_SUBTOTAL, CONTINUATION_INDICATOR
    global ELEMENT_COORDINATES
    global REQUIRED_SHAPE_NAMES
    global LEGEND_GROUPS_CONFIG
    global SUMMARY_TILE_CONFIG
    global TABLE_COLUMN_ALIGNMENT_MAP, TABLE_WORD_WRAP_COLUMNS, TABLE_UPPERCASE_COLUMNS, TABLE_DUAL_LINE_LABELS, TABLE_SHRINK_TO_FIT_COLUMNS
    global CLONE_PIPELINE_ENABLED
    global AUTOPPTX_CONFIG, ASPOSE_CONFIG, DOCSTRANGE_CONFIG

    MASTER_CONFIG = config

    features_section = MASTER_CONFIG.get("features", {}) if isinstance(MASTER_CONFIG.get("features", {}), dict) else {}
    CLONE_PIPELINE_ENABLED = bool(features_section.get("clone_pipeline_enabled", True))

    presentation_config = MASTER_CONFIG.section("presentation")
    fonts_config = presentation_config.get("fonts", {})
    font_sizes = fonts_config.get("sizes_pt", {})
    colors_config = presentation_config.get("colors", {})
    media_colors_config = colors_config.get("media_types", {})
    ui_colors_config = colors_config.get("ui_elements", {})
    table_config = presentation_config.get("table", {})
    styling_config = table_config.get("styling", {})
    TABLE_COLUMN_ALIGNMENT_MAP = _parse_alignment_map(styling_config.get("column_alignment", {}))
    TABLE_WORD_WRAP_COLUMNS = {int(idx) for idx in styling_config.get("word_wrap_columns", [])}
    TABLE_SHRINK_TO_FIT_COLUMNS = {int(idx) for idx in styling_config.get("shrink_to_fit_columns", [])}
    TABLE_UPPERCASE_COLUMNS = {int(idx) for idx in styling_config.get("uppercase_columns", [])}
    TABLE_DUAL_LINE_LABELS = _build_dual_line_map(styling_config.get("dual_line_labels", {}))
    row_heights_config = table_config.get("row_heights", {})
    table_position_config = table_config.get("positioning", {})
    comments_config = presentation_config.get("comments", {})
    comments_title_pos_config = comments_config.get("title_positioning", {})
    comments_box_pos_config = comments_config.get("box_positioning", {})
    title_position_config = presentation_config.get("title", {}).get("positioning", {})
    layout_config = presentation_config.get("layout", {})
    LEGEND_GROUPS_CONFIG = layout_config.get("legend_groups", {})
    charts_config = presentation_config.get("charts", {})
    chart_positions_config = charts_config.get("positioning", {})
    summary_tiles_config = presentation_config.get("summary_tiles", {})
    SUMMARY_TILE_CONFIG = summary_tiles_config

    tooling_section = MASTER_CONFIG.get("tooling", {})
    tooling_config = tooling_section if isinstance(tooling_section, dict) else {}
    AUTOPPTX_CONFIG = dict(tooling_config.get("autopptx", {}))
    ASPOSE_CONFIG = dict(tooling_config.get("aspose", {}))
    DOCSTRANGE_CONFIG = dict(tooling_config.get("docstrange", {}))

    TABLE_PLACEHOLDER_NAME = table_config.get("placeholder_name", "Table Placeholder 1")

    CLR_BLACK = RGBColor(0, 0, 0)
    CLR_WHITE = RGBColor(255, 255, 255)
    CLR_LIGHT_GRAY_TEXT = _rgb_color(ui_colors_config.get("light_gray_text", {}).get("rgb"), (191, 191, 191))
    CLR_TABLE_GRAY = _rgb_color(ui_colors_config.get("table_gray", {}).get("rgb"), (191, 191, 191))
    CLR_HEADER_GREEN = _rgb_color(ui_colors_config.get("header_green", {}).get("rgb"), (56, 236, 4))
    CLR_COMMENTS_GRAY = _rgb_color(ui_colors_config.get("comments_gray", {}).get("rgb"), (242, 242, 242))
    CLR_SUBTOTAL_GRAY = _rgb_color(ui_colors_config.get("subtotal_gray", {}).get("rgb"), (217, 217, 217))

    CLR_TELEVISION = _rgb_color(media_colors_config.get("television", {}).get("rgb"), (113, 212, 141))
    CLR_DIGITAL = _rgb_color(media_colors_config.get("digital", {}).get("rgb"), (253, 242, 183))
    CLR_OOH = _rgb_color(media_colors_config.get("ooh", {}).get("rgb"), (255, 191, 0))
    CLR_OTHER = _rgb_color(media_colors_config.get("other", {}).get("rgb"), (176, 211, 255))

    DEFAULT_FONT_NAME = fonts_config.get("default_family", "Calibri")
    TABLE_FONT_NAME = fonts_config.get("table_family", DEFAULT_FONT_NAME)
    FONT_SIZE_HEADER = Pt(float(font_sizes.get("header", 7.5)))
    FONT_SIZE_BODY = Pt(float(font_sizes.get("body", 7.0)))
    FONT_SIZE_BODY_COMPACT = Pt(float(font_sizes.get("body_compact", font_sizes.get("body", 7.0))))
    FONT_SIZE_CHART_TITLE = Pt(float(font_sizes.get("chart_title", 8.0)))
    FONT_SIZE_CHART_LABELS = Pt(float(font_sizes.get("chart_labels", 6.0)))
    FONT_SIZE_TITLE = Pt(float(font_sizes.get("title", 11.0)))
    FONT_SIZE_LEGEND = Pt(float(font_sizes.get("legend", 6.0)))
    FONT_SIZE_COMMENTS = Pt(float(font_sizes.get("comments", 9.0)))
    FONT_FAMILY_LEGEND = fonts_config.get("legend_family", DEFAULT_FONT_NAME)

    header_height_emu = row_heights_config.get("header_emu")
    body_height_emu = row_heights_config.get("body_emu")
    subtotal_height_emu = row_heights_config.get("subtotal_emu")
    trailer_height_emu = row_heights_config.get("trailer_emu")

    if header_height_emu is not None:
        TABLE_ROW_HEIGHT_HEADER = Emu(int(header_height_emu))
    else:
        header_inches = float(
            row_heights_config.get("header_inches", TEMPLATE_V4_ROW_HEIGHT_HEADER_INCHES)
        )
        TABLE_ROW_HEIGHT_HEADER = Inches(header_inches)

    if body_height_emu is not None:
        TABLE_ROW_HEIGHT_BODY = Emu(int(body_height_emu))
    else:
        body_inches = float(
            row_heights_config.get("body_inches", TEMPLATE_V4_ROW_HEIGHT_BODY_INCHES)
        )
        TABLE_ROW_HEIGHT_BODY = Inches(body_inches)

    if subtotal_height_emu is not None:
        TABLE_ROW_HEIGHT_SUBTOTAL = Emu(int(subtotal_height_emu))
    else:
        subtotal_inches = float(
            row_heights_config.get("subtotal_inches", TEMPLATE_V4_ROW_HEIGHT_BODY_INCHES)
        )
        TABLE_ROW_HEIGHT_SUBTOTAL = Inches(subtotal_inches)

    if trailer_height_emu is not None:
        TABLE_ROW_HEIGHT_TRAILER = Emu(int(trailer_height_emu))
    else:
        TABLE_ROW_HEIGHT_TRAILER = Emu(TEMPLATE_V4_ROW_HEIGHT_TRAILER_EMU)

    column_widths_emu_config = table_config.get("column_widths_emu") or []
    if column_widths_emu_config:
        TABLE_COLUMN_WIDTHS = [Emu(int(width)) for width in column_widths_emu_config]
    else:
        column_widths_config = table_config.get("column_widths_inches")
        column_widths_source = (
            column_widths_config if column_widths_config else TEMPLATE_V4_COLUMN_WIDTHS_INCHES
        )
        TABLE_COLUMN_WIDTHS = [Inches(float(width)) for width in column_widths_source]

    table_top_inches = float(
        table_position_config.get("top_inches", TEMPLATE_V4_TABLE_BOUNDS.top)
    )
    TABLE_TOP_OVERRIDE = Inches(table_top_inches)

    TABLE_CELL_STYLE_CONTEXT = CellStyleContext(
        margin_left_right_pt=3.6,
        margin_emu_lr=MARGIN_EMU_LR,
        default_font_name=TABLE_FONT_NAME,
        font_size_header=FONT_SIZE_HEADER,
        font_size_body=FONT_SIZE_BODY,
        font_size_body_compact=FONT_SIZE_BODY_COMPACT,
        color_black=CLR_BLACK,
        color_light_gray_text=CLR_LIGHT_GRAY_TEXT,
        color_table_gray=CLR_TABLE_GRAY,
        color_header_green=CLR_HEADER_GREEN,
        color_subtotal_gray=CLR_SUBTOTAL_GRAY,
        color_tv=CLR_TELEVISION,
        color_digital=CLR_DIGITAL,
        color_ooh=CLR_OOH,
        color_other=CLR_OTHER,
        column_alignment=TABLE_COLUMN_ALIGNMENT_MAP,
        word_wrap_columns=TABLE_WORD_WRAP_COLUMNS,
        shrink_to_fit_columns=TABLE_SHRINK_TO_FIT_COLUMNS,
        uppercase_columns=TABLE_UPPERCASE_COLUMNS,
        dual_line_labels=TABLE_DUAL_LINE_LABELS,
    )

    CHART_STYLE_CONTEXT = ChartStyleContext(
        font_name=DEFAULT_FONT_NAME,
        title_font_size=FONT_SIZE_CHART_TITLE,
        label_font_size=FONT_SIZE_CHART_LABELS,
        font_color=CLR_BLACK,
    )

    CHART_COLOR_MAPPING = {
        "Television": CLR_TELEVISION,
        "Digital": CLR_DIGITAL,
        "OOH": CLR_OOH,
        "Other": CLR_OTHER,
        "Awareness": CLR_TELEVISION,
        "Consideration": CLR_BLACK,
        "Purchase": CLR_OOH,
        "Always On": CLR_TELEVISION,
        "Brand": CLR_DIGITAL,
        "Product": CLR_OOH,
    }

    CHART_COLOR_CYCLE = [CLR_TELEVISION, CLR_DIGITAL, CLR_OOH, CLR_OTHER]

    MAX_ROWS_PER_SLIDE = int(table_config.get("max_rows_per_slide", 32))
    SPLIT_STRATEGY = table_config.get("split_strategy", "by_campaign")
    SHOW_CHARTS_ON_SPLITS = table_config.get("show_charts_on_splits", "all")
    SHOW_CARRIED_SUBTOTAL = bool(table_config.get("show_carried_subtotal", True))
    CONTINUATION_INDICATOR = table_config.get("continuation_indicator", " (Continued)")

    ELEMENT_COORDINATES = {
        "title": _coord_from_config(
            title_position_config,
            {"left": 0.184, "top": 0.308, "width": 2.952, "height": 0.370},
        ),
        "main_table": _coord_from_config(
            table_position_config,
            TEMPLATE_V4_TABLE_BOUNDS.as_dict(),
        ),
        "comments_title": _coord_from_config(
            comments_title_pos_config,
            {"left": 1.097, "top": 3.697, "width": 0.640, "height": 0.151},
        ),
        "comments_box": _coord_from_config(
            comments_box_pos_config,
            {"left": 0.184, "top": 3.886, "width": 2.466, "height": 1.489},
        ),
        "chart_1": _coord_from_config(
            chart_positions_config.get("funnel_chart", {}),
            {"left": 2.650, "top": 3.300, "width": 2.466, "height": 2.000},
        ),
        "chart_2": _coord_from_config(
            chart_positions_config.get("media_type_chart", {}),
            {"left": 4.725, "top": 3.300, "width": 2.647, "height": 2.000},
        ),
        "chart_3": _coord_from_config(
            chart_positions_config.get("campaign_type_chart", {}),
            {"left": 6.985, "top": 3.300, "width": 2.647, "height": 2.000},
        ),
        "tv_legend_color": {"left": 6.645, "top": 0.438, "width": 0.259, "height": 0.139},
        "tv_legend_text": {"left": 6.841, "top": 0.416, "width": 0.612, "height": 0.219},
        "digital_legend_color": {"left": 7.463, "top": 0.449, "width": 0.259, "height": 0.139},
        "digital_legend_text": {"left": 7.658, "top": 0.416, "width": 0.467, "height": 0.219},
        "ooh_legend_color": {"left": 8.196, "top": 0.449, "width": 0.259, "height": 0.139},
        "ooh_legend_text": {"left": 8.392, "top": 0.416, "width": 0.393, "height": 0.219},
        "other_legend_color": {"left": 8.866, "top": 0.449, "width": 0.259, "height": 0.139},
        "other_legend_text": {"left": 9.061, "top": 0.416, "width": 0.439, "height": 0.219},
    }


_initialize_from_config(load_master_config())


def configure(config: Config) -> None:
    global REQUIRED_SHAPE_NAMES
    _initialize_from_config(config)
    REQUIRED_SHAPE_NAMES = _collect_required_shape_names()



def _collect_required_shape_names() -> set[str]:
    """Build the set of shape names the template must expose."""

    required: set[str] = {SHAPE_NAME_TABLE}

    tiles = SUMMARY_TILE_CONFIG if SUMMARY_TILE_CONFIG else {}
    for section in ("quarter_budgets", "media_share", "funnel_share"):
        for tile in tiles.get(section, {}).values():
            shape_name = tile.get("shape") if isinstance(tile, dict) else None
            if shape_name:
                required.add(str(shape_name))

    footer_cfg = tiles.get("footer_notes", {})
    footer_shape = footer_cfg.get("shape") if isinstance(footer_cfg, dict) else None
    if footer_shape:
        required.add(str(footer_shape))

    # Only require comments placeholder if comments enabled in config
    if comments_config.get("enabled", False):
        required.update({SHAPE_NAME_COMMENTS_TITLE, SHAPE_NAME_COMMENTS_BOX})

    # Charts are optional and controlled via config flag
    if charts_config.get("enabled", False):
        required.update({SHAPE_NAME_FUNNEL_CHART, SHAPE_NAME_MEDIA_TYPE_CHART, SHAPE_NAME_CAMPAIGN_TYPE_CHART})

    return required


def _validate_template_shapes(template_slide) -> bool:
    """Verify the template slide exposes all required shapes."""

    available = {getattr(shape, "name", "") for shape in template_slide.shapes}
    missing = sorted(name for name in REQUIRED_SHAPE_NAMES if name and name not in available)

    if missing:
        logger.error("Template missing required shapes: %s", ", ".join(missing))
        return False

    logger.debug("Template shape validation passed (%s shapes present).", len(REQUIRED_SHAPE_NAMES))
    return True

# Define the constant we need (EXACTLY = 1 is the standard value)
class WD_ROW_HEIGHT_RULE:
    AT_LEAST = 0  # Add this for more flexible row height
    EXACTLY = 1

TABLE_HEIGHT_RULE_AVAILABLE = True

# --- Constants for Named Shapes (from Template_Refactoring_Guide.md) ---
SHAPE_NAME_TITLE = "TitlePlaceholder"
SHAPE_NAME_TABLE = "MainDataTable"
SHAPE_NAME_COMMENTS_TITLE = "CommentsTitle"
SHAPE_NAME_COMMENTS_BOX = "CommentsBox"
SHAPE_NAME_FUNNEL_CHART = "FunnelChart"
SHAPE_NAME_MEDIA_TYPE_CHART = "MediaTypeChart"
SHAPE_NAME_CAMPAIGN_TYPE_CHART = "CampaignTypeChart"
SHAPE_NAME_TV_LEGEND_COLOR = "TelevisionLegendColor"
SHAPE_NAME_TV_LEGEND_TEXT = "TelevisionLegendText"
SHAPE_NAME_DIGITAL_LEGEND_COLOR = "DigitalLegendColor"
SHAPE_NAME_DIGITAL_LEGEND_TEXT = "DigitalLegendText"
SHAPE_NAME_OOH_LEGEND_COLOR = "OOHLegendColor"
SHAPE_NAME_OOH_LEGEND_TEXT = "OOHLegendText"
SHAPE_NAME_OTHER_LEGEND_COLOR = "OtherLegendColor"
SHAPE_NAME_OTHER_LEGEND_TEXT = "OtherLegendText"

REQUIRED_SHAPE_NAMES = _collect_required_shape_names()

# Must exactly match the placeholder name in Slide Master
TABLE_PLACEHOLDER_NAME = "Table Placeholder 1"

# --- Color Constants (QA checklist verified RGB values) ---
CLR_BLACK = RGBColor(0, 0, 0)
CLR_WHITE = RGBColor(255, 255, 255)
CLR_LIGHT_GRAY_TEXT = _rgb_color(ui_colors_config.get("light_gray_text", {}).get("rgb"), (191, 191, 191))
CLR_TABLE_GRAY = _rgb_color(ui_colors_config.get("table_gray", {}).get("rgb"), (191, 191, 191))
CLR_HEADER_GREEN = _rgb_color(ui_colors_config.get("header_green", {}).get("rgb"), (56, 236, 4))
CLR_COMMENTS_GRAY = _rgb_color(ui_colors_config.get("comments_gray", {}).get("rgb"), (242, 242, 242))
CLR_SUBTOTAL_GRAY = _rgb_color(ui_colors_config.get("subtotal_gray", {}).get("rgb"), (217, 217, 217))

CLR_TELEVISION = _rgb_color(media_colors_config.get("television", {}).get("rgb"), (113, 212, 141))
CLR_DIGITAL = _rgb_color(media_colors_config.get("digital", {}).get("rgb"), (253, 242, 183))
CLR_OOH = _rgb_color(media_colors_config.get("ooh", {}).get("rgb"), (255, 191, 0))
CLR_OTHER = _rgb_color(media_colors_config.get("other", {}).get("rgb"), (176, 211, 255))

# --- PIXEL-PERFECT FONT CONSTANTS ---
DEFAULT_FONT_NAME = fonts_config.get("default_family", "Calibri")
FONT_SIZE_HEADER = Pt(float(font_sizes.get("header", 7.5)))
FONT_SIZE_BODY = Pt(float(font_sizes.get("body", 7.0)))
FONT_SIZE_BODY_COMPACT = Pt(
    float(font_sizes.get("body_compact", font_sizes.get("body", 7.0))))
FONT_SIZE_CHART_TITLE = Pt(float(font_sizes.get("chart_title", 8.0)))
FONT_SIZE_CHART_LABELS = Pt(float(font_sizes.get("chart_labels", 6.0)))

TABLE_ROW_HEIGHT_HEADER = Pt(float(row_heights_config.get("header_inches", 0.139)) * 72)
TABLE_ROW_HEIGHT_BODY = Pt(float(row_heights_config.get("body_inches", 0.118)) * 72)
TABLE_ROW_HEIGHT_SUBTOTAL = Pt(float(row_heights_config.get("subtotal_inches", 0.139)) * 72)
_table_column_config = table_config.get("column_widths_inches")
_table_column_source = (
    _table_column_config if _table_column_config else TEMPLATE_V4_COLUMN_WIDTHS_INCHES
)
TABLE_COLUMN_WIDTHS = [Inches(float(width)) for width in _table_column_source]

TABLE_TOP_OVERRIDE = Inches(
    float(table_position_config.get("top_inches", TEMPLATE_V4_TABLE_BOUNDS.top))
)

TABLE_CELL_STYLE_CONTEXT = CellStyleContext(
    margin_left_right_pt=3.6,
    margin_emu_lr=MARGIN_EMU_LR,
    default_font_name=DEFAULT_FONT_NAME,
    font_size_header=FONT_SIZE_HEADER,
    font_size_body=FONT_SIZE_BODY,
    font_size_body_compact=FONT_SIZE_BODY_COMPACT,
    color_black=CLR_BLACK,
    color_light_gray_text=CLR_LIGHT_GRAY_TEXT,
    color_table_gray=CLR_TABLE_GRAY,
    color_header_green=CLR_HEADER_GREEN,
    color_subtotal_gray=CLR_SUBTOTAL_GRAY,
    color_tv=CLR_TELEVISION,
    color_digital=CLR_DIGITAL,
    color_ooh=CLR_OOH,
    color_other=CLR_OTHER,
    column_alignment=TABLE_COLUMN_ALIGNMENT_MAP,
    word_wrap_columns=TABLE_WORD_WRAP_COLUMNS,
    shrink_to_fit_columns=TABLE_SHRINK_TO_FIT_COLUMNS,
    uppercase_columns=TABLE_UPPERCASE_COLUMNS,
    dual_line_labels=TABLE_DUAL_LINE_LABELS,
)

CHART_STYLE_CONTEXT = ChartStyleContext(
    font_name=DEFAULT_FONT_NAME,
    title_font_size=FONT_SIZE_CHART_TITLE,
    label_font_size=FONT_SIZE_CHART_LABELS,
    font_color=CLR_BLACK,
)

CHART_COLOR_MAPPING = {
    "Television": CLR_TELEVISION,
    "Digital": CLR_DIGITAL,
    "OOH": CLR_OOH,
    "Other": CLR_OTHER,
    "Awareness": CLR_TELEVISION,
    "Consideration": CLR_BLACK,
    "Purchase": CLR_OOH,
    "Always On": CLR_TELEVISION,
    "Brand": CLR_DIGITAL,
    "Product": CLR_OOH,
}

CHART_COLOR_CYCLE = [CLR_TELEVISION, CLR_DIGITAL, CLR_OOH, CLR_OTHER]

FONT_SIZE_TITLE = Pt(float(font_sizes.get("title", 11.0)))
FONT_SIZE_LEGEND = Pt(float(font_sizes.get("legend", 6.0)))
FONT_SIZE_COMMENTS = Pt(float(font_sizes.get("comments", 9.0)))

# --- TABLE SPLITTING CONSTANTS ---
# Template main table exposes 35 physical rows (header + 32 body slots + carried-forward + slide grand total).
# Keep the default body-row limit at 32 so cloned slides stay pixel-perfect while still reserving space
# for the carried-forward subtotal and slide-level GRAND TOTAL rows we append during splitting.
MAX_ROWS_PER_SLIDE = int(table_config.get("max_rows_per_slide", 32))
SPLIT_STRATEGY = table_config.get("split_strategy", "by_campaign")
SMART_PAGINATION_ENABLED = bool(table_config.get("smart_pagination_enabled", False))
SHOW_CHARTS_ON_SPLITS = table_config.get("show_charts_on_splits", "all")
SHOW_CARRIED_SUBTOTAL = bool(table_config.get("show_carried_subtotal", True))
CONTINUATION_INDICATOR = table_config.get("continuation_indicator", " (Continued)")

# --- GEOSPATIAL 2D COORDINATE SYSTEM ---
# Precise positioning for 10" × 5.625" PowerPoint slide (FINAL QA VERIFIED)
# Canvas: 16:9 aspect ratio with origin (0,0) at top-left corner
# All coordinates verified against "Egypt - Centrum" reference slide + Final QA checklist
ELEMENT_COORDINATES = {
    "title": _coord_from_config(
        title_position_config,
        {"left": 0.184, "top": 0.308, "width": 2.952, "height": 0.370},
    ),
        "main_table": _coord_from_config(
            table_position_config,
            TEMPLATE_V4_TABLE_BOUNDS.as_dict(),
        ),
    "comments_title": _coord_from_config(
        comments_title_pos_config,
        {"left": 1.097, "top": 3.697, "width": 0.640, "height": 0.151},
    ),
    "comments_box": _coord_from_config(
        comments_box_pos_config,
        {"left": 0.184, "top": 3.886, "width": 2.466, "height": 1.489},
    ),
    "chart_1": _coord_from_config(
        chart_positions_config.get("funnel_chart", {}),
        {"left": 2.650, "top": 3.300, "width": 2.466, "height": 2.000},
    ),
    "chart_2": _coord_from_config(
        chart_positions_config.get("media_type_chart", {}),
        {"left": 4.725, "top": 3.300, "width": 2.647, "height": 2.000},
    ),
    "chart_3": _coord_from_config(
        chart_positions_config.get("campaign_type_chart", {}),
        {"left": 6.985, "top": 3.300, "width": 2.647, "height": 2.000},
    ),
    "tv_legend_color": {"left": 6.645, "top": 0.438, "width": 0.259, "height": 0.139},
    "tv_legend_text": {"left": 6.841, "top": 0.416, "width": 0.612, "height": 0.219},
    "digital_legend_color": {"left": 7.463, "top": 0.449, "width": 0.259, "height": 0.139},
    "digital_legend_text": {"left": 7.658, "top": 0.416, "width": 0.467, "height": 0.219},
    "ooh_legend_color": {"left": 8.196, "top": 0.449, "width": 0.259, "height": 0.139},
    "ooh_legend_text": {"left": 8.392, "top": 0.416, "width": 0.393, "height": 0.219},
    "other_legend_color": {"left": 8.866, "top": 0.449, "width": 0.259, "height": 0.139},
    "other_legend_text": {"left": 9.061, "top": 0.416, "width": 0.439, "height": 0.219},
}

def get_element_position(element_name):
    """
    Get positioning coordinates for a named element in inches.
    
    Args:
        element_name: Key from ELEMENT_COORDINATES dictionary
        
    Returns:
        dict: Dictionary with left, top, width, height in Inches objects
    """
    if element_name not in ELEMENT_COORDINATES:
        logger.error(f"Element '{element_name}' not found in coordinate system")
        return None
    
    coords = ELEMENT_COORDINATES[element_name]
    return {
        'left': Inches(coords['left']),
        'top': Inches(coords['top']), 
        'width': Inches(coords['width']),
        'height': Inches(coords['height'])
    }

logger = logging.getLogger("amp_automation.legacy")

# --- CORRECTED Excel Column Indices for YOUR ACTUAL FILE ---
# Updated to match YOUR BULK_PLAN_EXPORT_2025_08_25.xlsx structure
COLUMN_GEOGRAPHY = 10              # Plan - Geography (extract country from hierarchy) - NO CHANGE
COLUMN_GLOBAL_MASTERBRAND = 17     # Plan - Brand - NO CHANGE
COLUMN_MEDIA_TYPE = 20             # Media Type - NO CHANGE
COLUMN_CAMPAIGN_NAMES = 83         # **Campaign Name(s) - WAS 85
COLUMN_CAMPAIGN_TYPE = 84          # **Campaign Type - WAS 86
COLUMN_FUNNEL_STAGE = 95           # **Funnel Stage - WAS 97
COLUMN_YEAR = 15                   # Plan - Year - NO CHANGE
COLUMN_FLIGHT_START_DATE = 110     # **Flight Start Date (for month extraction) - NO CHANGE
COLUMN_NET_COST = 71               # *Net Cost (budget data) - WAS 73
COLUMN_NATIONAL_GRP = 55           # GRP - WAS 56

# Additional Lumina columns for enhanced TV data (CORRECTED)
COLUMN_REACH_1PLUS = 104           # Reach 1+ - WAS 106
COLUMN_REACH_3PLUS = 64            # Reach 3+ - NO CHANGE
COLUMN_FREQUENCY = 105             # Frequency - WAS 107

# Legacy column references (kept for backwards compatibility)
COLUMN_COUNTRY = 10                # Now uses COLUMN_GEOGRAPHY
COLUMN_MONTH_START_DATE = 110      # Now uses COLUMN_FLIGHT_START_DATE

# Actual column name for GRP data, used if reading by name
COLUMN_NATIONAL_GRP_NAME = 'National GRP [Current]'

# --- Data Processing Constants (from v1.5) ---
MEDIA_TYPE_MAPPING = {
    'Television': 'TV',
    'TV': 'TV',
    'Digital': 'Digital',
    'OOH': 'OOH',
    'Other': 'Other',
    'Print': 'Other', # Group Print under Other
    'Radio': 'Other', # Group Radio under Other
    'Cinema': 'Other' # Group Cinema under Other
}

MONTH_MAP = {
    1: 'Jan', 2: 'Feb', 3: 'Mar', 4: 'Apr', 5: 'May', 6: 'Jun',
    7: 'Jul', 8: 'Aug', 9: 'Sep', 10: 'Oct', 11: 'Nov', 12: 'Dec'
}

EXPECTED_COLUMNS = [
    COLUMN_GLOBAL_MASTERBRAND,
    COLUMN_MEDIA_TYPE,
    COLUMN_CAMPAIGN_NAMES,
    COLUMN_COUNTRY,
    COLUMN_NET_COST,
    COLUMN_MONTH_START_DATE # Used to derive month
]

# --- Helper Functions ---
def _get_shape_by_name(slide, name):
    """Find a shape on a slide by its name."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    logger.debug("Shape with name '%s' not found on slide.", name)
    return None

def _copy_text_box(source_shape, target_slide, new_name=None, new_text=None):
    """Copies a source text box shape to a target slide, with optional new name and text."""
    if not source_shape.has_text_frame:
        logger.warning(f"Source shape '{source_shape.name}' does not have a text frame. Cannot copy as text box.")
        return None

    left, top, width, height = source_shape.left, source_shape.top, source_shape.width, source_shape.height
    new_shape = target_slide.shapes.add_textbox(left, top, width, height)
    new_shape.name = new_name if new_name else source_shape.name

    # Copy text frame properties
    source_tf = source_shape.text_frame
    new_tf = new_shape.text_frame

    new_tf.word_wrap = source_tf.word_wrap
    # MSO_AUTO_SIZE enum: NONE = 0, SHAPE_TO_FIT_TEXT = 1, TEXT_TO_FIT_SHAPE = 2
    # Ensure auto_size is handled correctly; direct copy might be problematic if enum values differ or not available.
    # For simplicity, let's try direct copy if it's simple. Otherwise, might need specific handling.
    try:
        new_tf.auto_size = source_tf.auto_size
    except AttributeError:
        logger.debug(f"Could not directly copy auto_size for shape '{new_shape.name}'. Defaulting.")
        # Default or no action if direct copy fails

    if new_text is not None:
        new_tf.text = new_text
    else:
        new_tf.text = source_tf.text

    # Copy paragraph formatting for the first paragraph (if text exists)
    if source_tf.paragraphs and new_tf.paragraphs:
        source_para = source_tf.paragraphs[0]
        new_para = new_tf.paragraphs[0]
        new_para.alignment = source_para.alignment
        
        if source_para.runs and new_para.runs:
            source_run = source_para.runs[0]
            new_run = new_para.runs[0]
            new_run.font.name = source_run.font.name
            new_run.font.size = source_run.font.size
            if source_run.font.color.type == MSO_COLOR_TYPE.RGB:
                 if hasattr(source_run.font.color, 'rgb'):
                    new_run.font.color.rgb = source_run.font.color.rgb
            new_run.font.bold = source_run.font.bold
            new_run.font.italic = source_run.font.italic
        elif not new_para.runs and new_text == "": # Handle empty text box case for font styling
            # If new_text is empty, add a run to apply font style from source
            if source_para.runs:
                source_run = source_para.runs[0]
                new_run = new_para.add_run()
                new_run.font.name = source_run.font.name
                new_run.font.size = source_run.font.size
                if source_run.font.color.type == MSO_COLOR_TYPE.RGB:
                     if hasattr(source_run.font.color, 'rgb'):
                        new_run.font.color.rgb = source_run.font.color.rgb
                new_run.font.bold = source_run.font.bold
                new_run.font.italic = source_run.font.italic
                new_para.text = "" # Clear the run's default text if any

    logger.debug(f"Copied text box '{source_shape.name}' to '{new_shape.name}' on target slide.")
    return new_shape

def _copy_shape(source_shape, target_slide, new_name=None):
    """Copies a basic source auto shape to a target slide, with optional new name."""
    if source_shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        logger.warning(f"Source shape '{source_shape.name}' is not an AUTO_SHAPE (type: {source_shape.shape_type}). Skipping copy.")
        return None

    left, top, width, height = source_shape.left, source_shape.top, source_shape.width, source_shape.height
    # Add the shape with the same auto_shape_type
    new_shape = target_slide.shapes.add_shape(
        source_shape.auto_shape_type, left, top, width, height
    )
    new_shape.name = new_name if new_name else source_shape.name

    # Copy fill properties
    source_fill = source_shape.fill
    new_fill = new_shape.fill
    # Assuming solid fill for legend color boxes
    new_fill.solid()
    if hasattr(source_fill.fore_color, 'rgb'):
        new_fill.fore_color.rgb = source_fill.fore_color.rgb
    else:
        logger.debug(f"Source shape '{source_shape.name}' fill fore_color has no rgb. Defaulting fill.")

    # Copy line properties
    source_line = source_shape.line
    new_line = new_shape.line
    if hasattr(source_line.color, 'rgb') and source_line.color.type != 0: # Check if color is set and not NO_FILL implicitly by type
        new_line.color.rgb = source_line.color.rgb
    elif source_line.fill.type == 0: # NO_FILL (0)
        new_line.fill.background() # Or new_line.fill.solid() with transparency, or just leave as default if that's no line
        logger.debug("Source shape '{source_shape.name}' line has NO_FILL. Applying background fill to line of new shape.")
    else:
        logger.debug(f"Source shape '{source_shape.name}' line color has no rgb or complex fill. Defaulting line color.")
    
    new_line.width = source_line.width # Width is in EMUs, direct copy is fine

    if hasattr(source_shape, "text") and hasattr(new_shape, "text"):
        new_shape.text = source_shape.text

    logger.debug(f"Copied shape '{source_shape.name}' to '{new_shape.name}' on target slide.")
    return new_shape

def is_empty_formatted_value(formatted_value):
    """Check if a formatted value represents empty/zero data that should not be displayed"""
    if not formatted_value:
        return True
    
    # Standard empty indicators
    if formatted_value in ["-", "", "0.0%"]:
        return True
    
    # Budget-specific empty indicators
    if formatted_value in ["£0K", "£0", "£0.00K"]:
        return True
        
    return False

def format_number(value, is_budget=False, is_percentage=False, is_grp=False, is_monthly_column=False):
    """Formats numbers for display in tables, with improved accuracy for budgets and enhanced zero handling."""
    if pd.isna(value):
        if is_percentage:
            return "0.0%"
        return ""  # Return empty string for NaN values to trigger "-" display
    
    numeric_value = 0
    try:
        numeric_value = float(value)
    except ValueError:
        if is_percentage:
            return "0.0%"
        return ""  # Return empty string for invalid values to trigger "-" display

    if numeric_value == 0 or abs(numeric_value) < ZERO_THRESHOLD:  # Enhanced zero detection including near-zero values
        if is_percentage:
            return "0.0%"
        # For budget and GRP values, return empty string to trigger "-" display in cell styling
        if is_budget or is_grp:
            return ""  # This will trigger "-" display in _apply_table_cell_styling
        return ""

    # Handle percentages first
    if is_percentage:
        formatted_pct = f"{numeric_value:.1f}%"
        if formatted_pct.endswith(".0%"):
            rounded = int(round(numeric_value))
            return f"{rounded}%"
        return formatted_pct

    # Handle GRPs with K suffix for thousands
    if is_grp:
        abs_value = abs(numeric_value)
        if abs_value >= 1_000:
            formatted_val = numeric_value / 1_000.0
            # Show one decimal place for accuracy (e.g., 23.5K)
            return f"{formatted_val:.1f}K"
        else:
            # For values less than 1000, show as integer
            return f"{int(numeric_value)}"

    # IMPROVED BUDGET FORMATTING: More accurate representation
    if is_budget:
        abs_value = abs(numeric_value)

        # Values >= 1M expressed in millions
        if abs_value >= 1_000_000:
            formatted_val = numeric_value / 1_000_000.0
            # Keep a single decimal only when materially helpful; otherwise use whole numbers
            rounded = round(formatted_val, 1 if abs_value >= 5_000_000 else 0)
            if rounded == int(rounded):
                return f"£{int(rounded)}M"
            return f"£{rounded:.1f}M"

        # Values >= 1K expressed in thousands (always whole numbers to minimize width)
        if abs_value >= 1_000:
            formatted_val = numeric_value / 1_000.0
            rounded = round(formatted_val)
            return f"£{int(rounded)}K"

        # Values < 1K show actual pounds when positive for monthly columns
        if is_monthly_column and abs_value > 0:
            return f"£{int(round(numeric_value))}"

        formatted_val = numeric_value / 1_000.0
        if is_monthly_column:
            return f"£{int(round(formatted_val))}K"
        return f"£{formatted_val:.2f}K"
    
    # Non-budget formatting (fallback)
    abs_value = abs(numeric_value)
    if abs_value >= 1_000_000:
        suffix = 'M'
        divisor = 1_000_000.0
    elif abs_value >= 1_000:
        suffix = 'K'
        divisor = 1_000.0
    else:
        suffix = 'K'
        divisor = 1_000.0

    formatted_val = numeric_value / divisor
    formatted_str = f"{formatted_val:.1f}{suffix}"
    return formatted_str

def load_and_prepare_data(excel_path):  # backward compatibility proxy
    active_logger = logger if 'logger' in globals() and logger is not None else logging.getLogger("amp_automation.data")
    data_set = modular_load_and_prepare_data(excel_path, MASTER_CONFIG, active_logger)
    return data_set.frame

def _prepare_main_table_data_detailed(df, region, masterbrand, year=None, excel_path=None):
    """Prepare detailed table data for a region/masterbrand/year combination."""

    global _CAMPAIGN_BOUNDARIES

    try:
        year_text = f" - {year}" if year is not None else ""
        logger.info("Preparing table data for %s - %s%s", region, masterbrand, year_text)

        filter_mask = (
            (df["Country"].astype(str).str.strip() == str(region).strip())
            & (df["Brand"].astype(str).str.strip() == str(masterbrand).strip())
        )

        if year is not None:
            filter_mask &= df["Year"].astype(str).str.strip() == str(year).strip()

        subset = df.loc[filter_mask].copy()
        logger.debug("Rows after filtering: %s", len(subset))

        if subset.empty:
            logger.warning("No data found for %s - %s%s", region, masterbrand, year_text)
            _CAMPAIGN_BOUNDARIES = []
            return None, None

        table_rows: list[list[str]] = [TABLE_HEADER_COLUMNS.copy()]
        cell_metadata: dict[tuple[int, int], dict[str, object]] = {}
        monthly_totals = [0.0] * len(TABLE_MONTH_ORDER)
        total_budget = float(subset["Total Cost"].sum() or 0.0)
        grand_total_grp = 0.0
        campaign_boundaries: list[tuple[int, int]] = []

        campaign_sort_info: list[tuple[str, int, str]] = []
        for campaign_name, campaign_group in subset.groupby("Campaign Name"):
            if pd.isna(campaign_name):
                continue
            normalized_media = {
                _normalized_media_value(media)
                for media in campaign_group["Mapped Media Type"].dropna()
            }
            try:
                first_media_index = next(
                    idx
                    for idx, media in enumerate(MEDIA_DISPLAY_ORDER)
                    if media in normalized_media
                )
            except StopIteration:
                first_media_index = len(MEDIA_DISPLAY_ORDER)
            campaign_sort_info.append(
                (
                    campaign_name,
                    first_media_index,
                    str(campaign_name).upper(),
                )
            )

        campaign_names = [
            name
            for name, _, _ in sorted(
                campaign_sort_info,
                key=lambda item: (item[1], item[2]),
            )
        ]

        coerced_year = _coerce_year(year)

        for campaign_name in campaign_names:
            campaign_df = subset[subset["Campaign Name"] == campaign_name]
            if campaign_df.empty:
                continue

            base_row_idx = len(table_rows)
            block_rows, block_month_totals, block_grp_total = _build_campaign_block(
                campaign_name,
                campaign_df,
                base_row_idx,
                total_budget,
                cell_metadata,
                region,
                masterbrand,
                coerced_year,
                excel_path,
            )

            if not block_rows:
                continue

            table_rows.extend(block_rows)
            monthly_totals = [
                total + addition
                for total, addition in zip(monthly_totals, block_month_totals)
            ]
            grand_total_grp += block_grp_total
            campaign_boundaries.append((base_row_idx, len(table_rows) - 1))

        grand_total_row = _build_grand_total_row(
            monthly_totals,
            total_budget,
            grand_total_grp,
        )
        table_rows.append(grand_total_row)

        _CAMPAIGN_BOUNDARIES = campaign_boundaries

        logger.info(
            "Table data created for %s - %s%s with %s rows",
            region,
            masterbrand,
            year_text,
            len(table_rows),
        )
        return table_rows, cell_metadata

    except Exception as exc:
        logger.error(
            "Error preparing table data for %s - %s%s: %s",
            region,
            masterbrand,
            year_text,
            exc,
        )
        logger.error(traceback.format_exc())
        _CAMPAIGN_BOUNDARIES = []
        return None, None

        """
                
                # Process GRP data for TV media types
                if is_tv_current_media and 'GRP' in df.columns:
                    # For wide-format data, GRP is in a single column, not monthly
                    q_grp_total = 0
                    grp_row_idx = current_row_idx + 1  # GRP row comes after main media row
                    
                    for i, month_name in enumerate(MONTHS):
                        month_col_idx = i + 7  # Month columns start after first 7 columns
                        
                        if month_name.startswith('Q'):  # Q1, Q2, Q3, Q4
                            formatted_grp = format_number(q_grp_total, is_grp=True)
                            monthly_grp_data_formatted[i] = formatted_grp
                            grp_monthly_totals[i] += q_grp_total
                            # Add metadata for GRP quarter columns - FIXED: align with display logic
                            cell_metadata[(grp_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(formatted_grp),
                                'media_type': 'GRPs',
                                'value': q_grp_total
                            }
                            q_grp_total = 0
                        else:
                            # Use month-specific aggregation from raw data
                            metrics = get_month_specific_tv_metrics(excel_path, region, masterbrand, campaign_name, year, month_name)
                            grp_val = metrics['grp_sum']
                            
                            
                            if grp_val > 0:
                                formatted_grp = format_number(grp_val, is_grp=True)
                                monthly_grp_data_formatted[i] = formatted_grp
                                grp_monthly_totals[i] += grp_val
                                q_grp_total += grp_val
                            else:
                                formatted_grp = '-'
                                monthly_grp_data_formatted[i] = formatted_grp
                            
                            # Add metadata for GRP month columns - FIXED: align with display logic
                            cell_metadata[(grp_row_idx, month_col_idx)] = {
                                'has_data': not is_empty_formatted_value(formatted_grp),
                                'media_type': 'GRPs',
                                'value': grp_val
                            }
                
                # CRITICAL FIX: Display campaign info for ALL media types, not just first
                # This resolves the issue where Digital campaigns appeared with blank names
                campaign_name_display = str(campaign_name).upper()
                campaign_budget_display = format_number(campaign_total_budget, is_budget=True)
                campaign_percentage_display = format_number(campaign_percentage, is_percentage=True)
                
                # Display total GRP for the campaign only on the primary TV row
                campaign_total_grp_display = format_number(campaign_total_tv_grp_value, is_grp=True) if first_media_type_for_campaign and is_tv_current_media else ""
                if first_media_type_for_campaign and is_tv_current_media:
                    grand_total_grp += campaign_total_tv_grp_value  # Accumulate for footer
                
                # Calculate TOTAL REACH and TOTAL FREQ for TV campaigns
                total_reach_display = ""
                total_freq_display = ""
                if first_media_type_for_campaign and is_tv_current_media:

        """

def _split_table_data_by_campaigns(table_data, cell_metadata):
    """Split the main table into continuation-friendly chunks respecting row limits."""

    if not table_data:
        return []

    header_idx = 0
    grand_total_idx = len(table_data) - 1

    if grand_total_idx <= 0:
        return [(table_data, cell_metadata, False)]

    body_row_count = grand_total_idx - 1
    if body_row_count <= MAX_ROWS_PER_SLIDE:
        return [(table_data, cell_metadata, False)]

    campaign_boundaries = _CAMPAIGN_BOUNDARIES or [(1, grand_total_idx - 1)]
    splits: list[tuple[list[list[str]], dict[tuple[int, int], dict[str, object]], bool]] = []

    current_indices: list[int] = [header_idx]
    current_body_count = 0
    slide_contains_continuation = False

    cumulative_months = [0.0] * len(TABLE_MONTH_ORDER)
    cumulative_total = 0.0
    cumulative_grp = 0.0

    month_start_col = 3
    total_col_idx = month_start_col + len(TABLE_MONTH_ORDER)
    grp_col_idx = total_col_idx + 1

    def _row_label(idx: int) -> str:
        row = table_data[idx]
        return str(row[0]).strip().upper() if row else ""

    def _media_label(idx: int) -> str:
        row = table_data[idx]
        if len(row) < 2:
            return ""
        return str(row[1]).strip().upper()

    def _is_media_header(idx: int) -> bool:
        label = _row_label(idx)
        media = _media_label(idx)
        if label in {"", "-"}:
            return bool(media and media != "-")
        if label in {"MONTHLY TOTAL (� 000)", "GRAND TOTAL", "CARRIED FORWARD"}:
            return False
        return True

    def _determine_block_length(start_idx: int, campaign_end_idx: int) -> int:
        """Return the number of contiguous rows that represent a media block (incl. metrics)."""

        if start_idx > campaign_end_idx:
            return 0

        label = _row_label(start_idx)
        if label == "MONTHLY TOTAL (� 000)":
            return 1

        length = 1
        idx = start_idx + 1
        while idx <= campaign_end_idx:
            next_label = _row_label(idx)
            if next_label == "MONTHLY TOTAL (� 000)":
                length += 1
                idx += 1
                break
            if _is_media_header(idx):
                break
            if next_label == "GRAND TOTAL":
                break
            length += 1
            idx += 1
        return length

    def accumulate_for_rows(row_indices: list[int]) -> tuple[list[float], float, float]:
        month_totals = [0.0] * len(TABLE_MONTH_ORDER)
        total_budget = 0.0
        grp_total = 0.0

        for row_idx in row_indices:
            if row_idx >= len(table_data) or row_idx <= 0 or row_idx == grand_total_idx:
                continue

            for month_idx in range(len(TABLE_MONTH_ORDER)):
                col_idx = month_start_col + month_idx
                meta = cell_metadata.get((row_idx, col_idx))
                if not meta:
                    continue
                media_type = meta.get("media_type")
                if media_type in {"Subtotal", "GrandTotal"}:
                    continue
                if meta.get("has_data"):
                    value = float(meta.get("value", 0.0) or 0.0)
                    month_totals[month_idx] += value

            total_meta = cell_metadata.get((row_idx, total_col_idx))
            if total_meta:
                media_type = total_meta.get("media_type")
                if media_type not in {"Subtotal", "GrandTotal"} and total_meta.get("has_data"):
                    total_budget += float(total_meta.get("value", 0.0) or 0.0)

            grp_meta = cell_metadata.get((row_idx, grp_col_idx))
            if grp_meta and grp_meta.get("has_data"):
                media_type = grp_meta.get("media_type")
                if media_type in {"Television", "GRPs"}:
                    grp_total += float(grp_meta.get("value", 0.0) or 0.0)

        return month_totals, total_budget, grp_total

    def finalize_split(has_remaining: bool) -> None:
        nonlocal current_indices, current_body_count, cumulative_months, cumulative_total, cumulative_grp, slide_contains_continuation

        if len(current_indices) <= 1:
            return

        body_indices = [
            idx for idx in current_indices if idx > 0 and idx != grand_total_idx
        ]
        slide_months, slide_total, slide_grp = accumulate_for_rows(body_indices)
        cumulative_months = [prev + inc for prev, inc in zip(cumulative_months, slide_months)]
        cumulative_total += slide_total
        cumulative_grp += slide_grp

        output_indices = current_indices.copy()
        inserted_meta_rows: dict[int, dict[int, dict[str, object]]] = {}

        if has_remaining and SHOW_CARRIED_SUBTOTAL:
            output_indices.append(-1)

        output_rows: list[list[str]] = []
        for new_row_idx, original_idx in enumerate(output_indices):
            if original_idx == -1:
                carried_months = list(cumulative_months)
                carried_row = _build_carried_forward_row(
                    carried_months,
                    cumulative_total,
                    cumulative_grp,
                )
                output_rows.append(carried_row)
                inserted_meta_rows[new_row_idx] = _build_carried_forward_metadata_values(
                    carried_months,
                    cumulative_total,
                    cumulative_grp,
                )
            else:
                output_rows.append(table_data[original_idx])

        output_metadata = _extract_metadata_for_indices(
            cell_metadata,
            output_indices,
            len(output_rows),
            inserted_meta_rows,
        )

        per_slide_row = _build_grand_total_row(slide_months, slide_total, slide_grp)
        output_rows.append(per_slide_row)
        per_slide_row_idx = len(output_rows) - 1
        per_slide_metadata = _build_grand_total_metadata_values(
            slide_months,
            slide_total,
            slide_grp,
        )
        for col_idx, meta in per_slide_metadata.items():
            output_metadata[(per_slide_row_idx, col_idx)] = meta

        splits.append((output_rows, output_metadata, len(splits) > 0))
        slide_contains_continuation = False
        current_indices = [header_idx]
        current_body_count = 0

    for boundary_idx, (start, end) in enumerate(campaign_boundaries):
        # Smart pagination: Calculate campaign row count and check if it fits
        campaign_row_count = end - start + 1

        if SMART_PAGINATION_ENABLED and campaign_row_count <= MAX_ROWS_PER_SLIDE:
            # Small campaign - check if it fits on current slide
            if current_body_count > 0 and current_body_count + campaign_row_count > MAX_ROWS_PER_SLIDE:
                # Doesn't fit - finalize current slide and start fresh
                logger.debug(
                    f"Smart pagination: Campaign with {campaign_row_count} rows doesn't fit on current slide "
                    f"({current_body_count} rows used). Starting fresh slide."
                )
                finalize_split(has_remaining=True)
                slide_contains_continuation = False

            # Add entire campaign to current slide
            logger.debug(
                f"Smart pagination: Adding campaign with {campaign_row_count} rows to slide "
                f"(current: {current_body_count}, new total: {current_body_count + campaign_row_count})"
            )
            current_indices.extend(range(start, end + 1))
            current_body_count += campaign_row_count
        else:
            # Large campaign (>MAX_ROWS_PER_SLIDE) or smart pagination disabled - use block-by-block splitting
            if SMART_PAGINATION_ENABLED and campaign_row_count > MAX_ROWS_PER_SLIDE:
                # Finalize current slide before starting large campaign
                if current_body_count > 0:
                    logger.debug(
                        f"Smart pagination: Large campaign with {campaign_row_count} rows detected. "
                        f"Finalizing current slide first."
                    )
                    finalize_split(has_remaining=True)
                    slide_contains_continuation = False

            idx = start
            while idx <= end:
                block_length = _determine_block_length(idx, end)
                if block_length <= 0:
                    idx += 1
                    continue

                if (
                    current_body_count > 0
                    and current_body_count + block_length > MAX_ROWS_PER_SLIDE
                ):
                    finalize_split(has_remaining=True)
                    slide_contains_continuation = False
                    continue

                if block_length > MAX_ROWS_PER_SLIDE:
                    logger.warning(
                        "Media block starting at row %s spans %s rows exceeding MAX_ROWS_PER_SLIDE=%s; forcing slide overflow.",
                        idx,
                        block_length,
                        MAX_ROWS_PER_SLIDE,
                    )

                if current_body_count == 0 and idx > start:
                    slide_contains_continuation = True

                current_indices.extend(range(idx, idx + block_length))
                current_body_count += block_length
                idx += block_length

        if slide_contains_continuation and current_body_count > 0:
            finalize_split(has_remaining=boundary_idx < len(campaign_boundaries) - 1)
            slide_contains_continuation = False
            continue

    finalize_split(has_remaining=False)

    return splits

def _calculate_subtotal_for_split(rows):
    """Calculate subtotal values for a split section."""
    # Initialize subtotal values
    subtotal_values = []
    
    # Define sub-row media types to exclude from totals
    sub_row_media_types = ['GRPs', 'Reach@1+', 'OTS@1+', 'Reach@3+', 'OTS@3+']
    
    # Skip first column (campaign name), process budget column
    total_budget = 0
    for row in rows:
        # Skip sub-rows when calculating budget totals
        media_type = row[6] if len(row) > 6 else ""
        if media_type in sub_row_media_types:
            continue
            
        if row[1] and row[1] != '-':
            # Extract numeric value from budget string
            budget_str = row[1].replace('£', '').replace('K', '000').replace('M', '000000').replace(',', '')
            try:
                total_budget += float(budget_str)
            except:
                pass
    
    # Format budget
    if total_budget >= 1000000:
        subtotal_values.append(f"£{total_budget/1000000:.1f}M")
    else:
        subtotal_values.append(f"£{int(total_budget/1000)}K")
    
    # Handle other columns (simplified for now)
    # TV GRPs, TOTAL REACH, TOTAL FREQ, %
    subtotal_values.extend(['-', '-', '-', '-'])
    
    # Media Type column
    subtotal_values.append('')
    
    # Monthly columns - sum up values
    for col_idx in range(7, 23):  # Monthly columns
        monthly_total = 0
        for row in rows:
            # Skip sub-rows when calculating monthly totals
            media_type = row[6] if len(row) > 6 else ""
            if media_type in sub_row_media_types:
                continue
                
            if col_idx < len(row) and row[col_idx] and row[col_idx] != '-':
                try:
                    value_str = row[col_idx].replace('£', '').replace('K', '000').replace('M', '000000').replace(',', '')
                    monthly_total += float(value_str)
                except:
                    pass
        
        if monthly_total > 0:
            if monthly_total >= 1000000:
                subtotal_values.append(f"£{monthly_total/1000000:.1f}M")
            else:
                subtotal_values.append(f"£{int(monthly_total/1000)}K")
        else:
            subtotal_values.append('-')
    
    return subtotal_values


def _extract_metadata_for_indices(original_metadata, indices, new_row_count, inserted_metadata: dict[int, dict[int, dict[str, object]]] | None = None):
    """Extract metadata for specific row indices and remap to new positions."""
    new_metadata = {}
    inserted_metadata = inserted_metadata or {}
    
    for new_row_idx, original_idx in enumerate(indices):
        if original_idx < 0:
            # Special indices for inserted rows (e.g., carried forward)
            if new_row_idx in inserted_metadata:
                for col_idx, meta in inserted_metadata[new_row_idx].items():
                    new_metadata[(new_row_idx, col_idx)] = meta
            continue
            
        # Copy metadata for this row
        for (orig_row, col), meta in original_metadata.items():
            if orig_row == original_idx:
                new_metadata[(new_row_idx, col)] = meta
    
    return new_metadata


def _prepare_media_type_chart_data_detailed(df, region, masterbrand, year=None):
    """Compatibility wrapper for modular media type chart data preparation."""

    return prepare_media_type_chart_data(df, region, masterbrand, year)

def _prepare_funnel_chart_data_detailed(df, region, masterbrand, year=None):
    """Compatibility wrapper for modular funnel chart data preparation."""

    return prepare_funnel_chart_data(df, region, masterbrand, year)

def _prepare_campaign_type_chart_data(df, region, masterbrand, year=None):
    """Compatibility wrapper for modular campaign type chart data preparation."""

    return prepare_campaign_type_chart_data(df, region, masterbrand, year)

def set_title_text_detailed(title_shape, title_text, template_prs):
    logger.debug(f"Attempting to set title text '{title_text}' for shape '{title_shape.name}' using detailed method.")
    # --- Locate the reference title placeholder and its text_frame from the template --- #
    template_title_ref_shape = None
    # Priority: 1. Named 'TitlePlaceholder' on first slide master
    if template_prs.slide_masters:
        master = template_prs.slide_masters[0]
        for placeholder in master.placeholders:
            if placeholder.name == "TitlePlaceholder":
                template_title_ref_shape = placeholder
                logger.debug(f"Found reference title placeholder '{placeholder.name}' by name in slide master.")
                break
            elif placeholder.is_placeholder and placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                template_title_ref_shape = placeholder
                logger.debug(f"Found reference title placeholder type '{placeholder.placeholder_format.type}' in slide master.")
                break # Take the first one found
    
    # Priority: 2. Named 'TitlePlaceholder' on first slide layout (if not found in master)
    if not template_title_ref_shape and template_prs.slide_layouts:
        layout = template_prs.slide_layouts[0] # Assuming first layout is relevant
        for placeholder in layout.placeholders:
            if placeholder.name == "TitlePlaceholder":
                template_title_ref_shape = placeholder
                logger.debug(f"Found reference title placeholder '{placeholder.name}' by name in slide layout '{layout.name}'.")
                break
            elif placeholder.is_placeholder and placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                template_title_ref_shape = placeholder
                logger.debug(f"Found reference title placeholder type '{placeholder.placeholder_format.type}' in slide layout '{layout.name}'.")
                break

    # Priority: 3. Named 'TitlePlaceholder' or actual title placeholder on the first actual slide (if not in master/layout)
    if not template_title_ref_shape and template_prs.slides:
        actual_slide = template_prs.slides[0]
        if actual_slide.has_title and actual_slide.shapes.title:
            template_title_ref_shape = actual_slide.shapes.title
            logger.debug(f"Using actual title shape from template's first slide: '{template_title_ref_shape.name}'")
        else:
            for shape_on_slide in actual_slide.shapes:
                if shape_on_slide.name == "TitlePlaceholder":
                    template_title_ref_shape = shape_on_slide
                    logger.debug(f"Found named shape 'TitlePlaceholder' on template's first slide.")
                    break

    if not template_title_ref_shape or not template_title_ref_shape.has_text_frame:
        logger.warning("Could not find a suitable reference title placeholder with a text frame in the template. Applying basic text only.")
        title_shape.text = title_text
        return

    template_tf = template_title_ref_shape.text_frame
    logger.debug(f"Using template text_frame from shape: '{template_title_ref_shape.name}' (ID: {template_title_ref_shape.shape_id}) for styling.")

    # --- Copy Shape Fill --- 
    if template_title_ref_shape:
        source_shape_fill = template_title_ref_shape.fill
        target_shape_fill = title_shape.fill

        # Check if source_shape_fill.type is None first, as it can be for inherited fills
        if source_shape_fill.type is None or source_shape_fill.type == MSO_FILL_TYPE.BACKGROUND:
            target_shape_fill.background() # Makes it transparent / no fill
            logger.debug("Template Title - Source fill is NONE or BACKGROUND. Applied no fill to target.")
        elif source_shape_fill.type == MSO_FILL_TYPE.SOLID:
            target_shape_fill.solid()
            if source_shape_fill.fore_color.type == MSO_COLOR_TYPE.RGB:
                target_shape_fill.fore_color.rgb = source_shape_fill.fore_color.rgb
                logger.debug(f"Template Title - Source fill is SOLID RGB: {source_shape_fill.fore_color.rgb}. Applied to target.")
            elif source_shape_fill.fore_color.type == MSO_COLOR_TYPE.SCHEME:
                target_shape_fill.fore_color.theme_color = source_shape_fill.fore_color.theme_color
                if hasattr(source_shape_fill.fore_color, 'brightness') and source_shape_fill.fore_color.brightness is not None:
                    target_shape_fill.fore_color.brightness = source_shape_fill.fore_color.brightness
                if hasattr(source_shape_fill.fore_color, 'tint') and source_shape_fill.fore_color.tint is not None:
                    target_shape_fill.fore_color.tint = source_shape_fill.fore_color.tint
                if hasattr(source_shape_fill.fore_color, 'shade') and source_shape_fill.fore_color.shade is not None:
                    target_shape_fill.fore_color.shade = source_shape_fill.fore_color.shade
                logger.debug(f"Template Title - Source fill is SOLID SCHEME Color. Applied to target.")
            else:
                logger.warning(f"Template Title - Source fill is SOLID but color type {source_shape_fill.fore_color.type} not handled for fill. Fill color not fully copied.")
        # Add more fill types like GRADIENT, PICTURE if necessary later
        else:
            logger.warning(f"Template Title - Source fill type {source_shape_fill.type} not handled. Fill not copied.")
    else:
        logger.warning("Could not find reference title placeholder shape to copy fill.")
    # --- End of Shape Fill Copying ---

    # --- Apply Vertical Anchor from template_tf to title_shape.text_frame --- #
    source_va = template_tf.vertical_anchor
    logger.debug(f"Template Title - Read source vertical_anchor: {source_va} (Enum: {MSO_VERTICAL_ANCHOR(source_va) if source_va is not None else 'None'})")
    if source_va is not None:
        title_shape.text_frame.vertical_anchor = source_va
        logger.debug(f"Applied vertical_anchor '{source_va}' to title_shape.text_frame. New VA: {title_shape.text_frame.vertical_anchor}")
    else:
        logger.debug("Source vertical_anchor is None. Vertical alignment on new title will use default (likely TOP).")

    # --- Clear existing content and add new paragraph for the title text --- #
    title_shape.text_frame.clear()
    p = title_shape.text_frame.add_paragraph()
    p.text = title_text
    logger.debug(f"Set title text to: '{title_text}' in new paragraph.")

    # --- Apply Paragraph and Font Styling from template_tf to the new paragraph 'p' and its first run --- #
    if template_tf.paragraphs:
        ref_para = template_tf.paragraphs[0]  # Use first paragraph of template for style reference
        logger.debug(f"Reference paragraph from template: Level={ref_para.level}, SpaceBefore={ref_para.space_before}, SpaceAfter={ref_para.space_after}, LineSpacing={ref_para.line_spacing}")

        # 1. Apply Paragraph Alignment (Horizontal)
        if ref_para.alignment is not None:
            p.alignment = ref_para.alignment
            logger.debug(f"Template Title - Source Paragraph Alignment: {ref_para.alignment} (Enum: {PP_ALIGN(ref_para.alignment) if ref_para.alignment is not None else 'None'}). Applied to new paragraph.")
        else:
            logger.debug("Template Title - Source paragraph alignment is None. Horizontal alignment will use default.")
        
        # Apply other paragraph properties if needed (e.g., level, spacing - though these are often complex)
        # p.level = ref_para.level # Be cautious with level, can affect master styles

        # 2. Apply Font Styling from the first run of the reference paragraph
        if ref_para.runs:
            ref_run = ref_para.runs[0]  # Reference run from template
            
            if p.runs: # The new paragraph 'p' should have one run after p.text = title_text
                target_run = p.runs[0]

                font_attrs_log = []
                # Font Name
                if ref_run.font.name:
                    target_run.font.name = ref_run.font.name
                    font_attrs_log.append(f"Name='{ref_run.font.name}'")
                # Font Size
                if ref_run.font.size:
                    target_run.font.size = ref_run.font.size
                    font_attrs_log.append(f"Size={ref_run.font.size}")
                # Bold
                if ref_run.font.bold is not None:
                    target_run.font.bold = ref_run.font.bold
                    font_attrs_log.append(f"Bold={ref_run.font.bold}")
                # Italic
                if ref_run.font.italic is not None:
                    target_run.font.italic = ref_run.font.italic
                    font_attrs_log.append(f"Italic={ref_run.font.italic}")
                # Underline
                if ref_run.font.underline is not None:
                    target_run.font.underline = ref_run.font.underline
                    font_attrs_log.append(f"Underline={ref_run.font.underline}")
                # Color
                if ref_run.font.color.type:
                    if ref_run.font.color.type == MSO_COLOR_TYPE.RGB:
                        target_run.font.color.rgb = ref_run.font.color.rgb
                        font_attrs_log.append(f"ColorRGB='{ref_run.font.color.rgb}'")
                    elif ref_run.font.color.type == MSO_COLOR_TYPE.SCHEME:
                        target_run.font.color.theme_color = ref_run.font.color.theme_color
                        target_run.font.color.brightness = ref_run.font.color.brightness
                        font_attrs_log.append(f"ColorTheme='{ref_run.font.color.theme_color}', Brightness='{ref_run.font.color.brightness}'")
                    # Add other color types if necessary
                
                logger.debug(f"Template Title - Source Font Details: {', '.join(font_attrs_log)}")
                applied_font_attrs_log = [
                    f"Name='{target_run.font.name}'", f"Size={target_run.font.size}", 
                    f"Bold={target_run.font.bold}", f"Italic={target_run.font.italic}", 
                    f"Underline={target_run.font.underline}"
                ]
                if target_run.font.color.type == MSO_COLOR_TYPE.RGB:
                    applied_font_attrs_log.append(f"ColorRGB='{target_run.font.color.rgb}'")
                elif target_run.font.color.type == MSO_COLOR_TYPE.SCHEME:
                    applied_font_attrs_log.append(f"ColorTheme='{target_run.font.color.theme_color}', Brightness='{target_run.font.color.brightness}'")
                logger.debug(f"Applied Font to new run: {', '.join(applied_font_attrs_log)}")

            else:
                logger.warning("Title paragraph has no runs after setting text. Cannot apply font styling.")
        else:
            logger.debug("Template Title - Reference paragraph in template has no runs. Cannot copy font styling.")
    else:
        logger.debug("Template Title - Reference text_frame in template has no paragraphs. Cannot copy paragraph/font styling.")

    logger.debug(f"Finished setting title for shape '{title_shape.name}'. Text frame word_wrap: {title_shape.text_frame.word_wrap}, auto_size: {title_shape.text_frame.auto_size}")

def _add_table_row(table):
    """
    Add a row to an existing table by cloning the last row's XML structure.

    Python-pptx doesn't provide a public API for adding rows to existing tables.
    This function uses XML manipulation to clone the last row and append it.

    Args:
        table: A python-pptx Table object

    Returns:
        The newly added row object

    Raises:
        ValueError: If the table has no rows to clone from
    """
    from copy import deepcopy

    num_rows = len(table.rows)
    if num_rows == 0:
        raise ValueError("Cannot add row to table with no existing rows")

    # Get the last row's XML element directly from the table rows
    # Use positive indexing since _RowCollection may not support negative indexing
    last_row = table.rows[num_rows - 1]
    last_tr = last_row._tr
    new_tr = deepcopy(last_tr)

    # Clear text content from all cells in the new row
    from pptx.oxml.ns import qn
    for tc in new_tr.findall(qn('a:tc')):
        # Find text elements and clear them
        text_frame = tc.find(qn('a:txBody'))
        if text_frame is not None:
            for paragraph in text_frame.findall(qn('a:p')):
                for run in paragraph.findall(qn('a:r')):
                    for text_elem in run.findall(qn('a:t')):
                        text_elem.text = ''

    # Append the new row to the table
    tbl = table._tbl
    tbl.append(new_tr)

    # Return the new row (access via positive indexing)
    return table.rows[len(table.rows) - 1]

def _populate_cloned_table(table_shape, table_data, cell_metadata):
    table = table_shape.table
    rows_needed = len(table_data)
    if rows_needed == 0:
        logger.warning("No table data supplied for cloned table population")
        return False

    cols_needed = len(table_data[0])
    if len(table.columns) < cols_needed:
        logger.warning(
            "Cloned table has fewer columns (%s) than required (%s)",
            len(table.columns),
            cols_needed,
        )
        return False

    # Debug: Log current table state
    logger.debug(f"Cloned table has {len(table.rows)} rows, need {rows_needed} rows")

    # Ensure sufficient row capacity by appending rows as needed.
    while len(table.rows) < rows_needed:
        logger.debug(f"Adding row {len(table.rows) + 1}/{rows_needed}")
        _add_table_row(table)

    for col_idx, width in enumerate(TABLE_COLUMN_WIDTHS[: len(table.columns)]):
        try:
            table.columns[col_idx].width = width
        except Exception as exc:
            logger.debug("Unable to enforce column width for column %s: %s", col_idx, exc)

    header_height = TABLE_ROW_HEIGHT_HEADER
    body_height = TABLE_ROW_HEIGHT_BODY
    subtotal_height = TABLE_ROW_HEIGHT_SUBTOTAL
    trailer_height = TABLE_ROW_HEIGHT_TRAILER

    def _apply_row_height(target_row, target_height, row_index, *, lock_exact=False):
        """Assign height and optionally pin it via hRule='exact'."""
        target_row.height = target_height
        if lock_exact:
            try:
                target_row._tr.set("hRule", "exact")
            except Exception as exc:  # pragma: no cover - defensive logging
                logger.debug("Unable to lock row %s height rule: %s", row_index, exc)
        else:
            if target_row._tr.get("hRule") is not None:
                try:
                    del target_row._tr.attrib["hRule"]
                except Exception as exc:  # pragma: no cover - defensive logging
                    logger.debug("Unable to clear row %s height rule: %s", row_index, exc)

    subtotal_labels = {"SUBTOTAL", "CARRIED FORWARD", "MONTHLY TOTAL (£ 000)", "GRAND TOTAL"}

    for row_idx in range(rows_needed):
        row = table.rows[row_idx]
        row_data = table_data[row_idx]
        raw_label = str(row_data[0]) if row_data else ""
        label = " ".join(raw_label.split()).upper()
        is_blank_row = all(
            (value is None or str(value).strip() in ("", "-"))
            for value in row_data
        )

        if row_idx == 0:
            _apply_row_height(row, header_height, row_idx, lock_exact=True)
        elif row_idx == rows_needed - 1 and is_blank_row:
            _apply_row_height(row, trailer_height, row_idx)
        elif label in subtotal_labels:
            _apply_row_height(row, subtotal_height, row_idx, lock_exact=True)
        else:
            _apply_row_height(row, body_height, row_idx, lock_exact=True)

        for col_idx in range(cols_needed):
            value = row_data[col_idx] if col_idx < len(row_data) else ""
            cell = table.cell(row_idx, col_idx)
            cell.text = "" if value is None else str(value)
            style_table_cell(
                cell,
                row_idx,
                col_idx,
                table_data,
                cell_metadata,
                TABLE_CELL_STYLE_CONTEXT,
                logger,
            )

    _apply_campaign_cell_merges(table, table_data)

    # Clear any unused rows beyond data to avoid stray template content.
    total_rows = len(table.rows)
    for row_idx in range(rows_needed, total_rows):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
        try:
            _apply_row_height(table.rows[row_idx], trailer_height, row_idx)
        except Exception as exc:
            logger.debug("Unable to collapse unused row %s height: %s", row_idx, exc)

    if total_rows > rows_needed:
        logger.debug("Removing %s trailing empty rows from cloned table", total_rows - rows_needed)
        for row_idx in range(total_rows - 1, rows_needed - 1, -1):
            tr = table.rows[row_idx]._tr
            table._tbl.remove(tr)

    apply_table_borders(table, TABLE_CELL_STYLE_CONTEXT.color_table_gray)
    try:
        table_shape.left = Inches(TEMPLATE_V4_TABLE_BOUNDS.left)
        table_shape.top = Inches(TEMPLATE_V4_TABLE_BOUNDS.top)
        table_shape.width = Inches(TEMPLATE_V4_TABLE_BOUNDS.width)
        table_shape.height = Inches(TEMPLATE_V4_TABLE_BOUNDS.height)
    except Exception as exc:
        logger.debug("Unable to snap table geometry to template bounds: %s", exc)
    return True


def _add_and_style_table(slide, table_data, cell_metadata, template_slide=None):
    table_pos = get_element_position('main_table')
    if not table_pos:
        logger.error("Failed to get table position coordinates")
        return False

    if template_slide is not None:
        try:
            cloned_table_shape = clone_template_table(template_slide, slide, SHAPE_NAME_TABLE)
            if cloned_table_shape:
                success = _populate_cloned_table(cloned_table_shape, table_data, cell_metadata)
                if success:
                    return True
                logger.error("Populating cloned table failed; aborting table population")
        except Exception as exc:  # pragma: no cover - cloning fallback
            logger.error("Template cloning failed: %s", exc)

    logger.error("Template cloning not available; unable to populate MainDataTable")
    return False

def _prepare_main_table_data(
    df: pd.DataFrame,
    region: str,
    masterbrand: str,
    year: int | None = None,
    excel_path: str | Path | None = None,
):
    """
    Compatibility wrapper returning the detailed table structure used by both pipelines.

    The legacy AutoPPTX workflow expects the same table payload as the clone pipeline so
    we delegate to :func:`_prepare_main_table_data_detailed` while keeping the public
    helper signature alive for existing call sites.
    """

    return _prepare_main_table_data_detailed(
        df,
        region,
        masterbrand,
        year,
        excel_path,
    )

def _prepare_funnel_chart_data(df, region, masterbrand):
    """
    Prepare data for the funnel chart.
    
    Args:
        df: DataFrame with the Excel data
        region: Region filter
        masterbrand: Masterbrand filter
        
    Returns:
        list: Chart data or None if no data
    """
    try:
        # Filter data for this region/masterbrand combination
        filtered_df = df[
            (df['Region'] == region) &
            (df['Global Masterbrand'] == masterbrand)
        ].copy()
        
        if filtered_df.empty:
            return None
        
        # Group by funnel stage and sum budgets
        if 'Funnel Stage' in filtered_df.columns and 'Budget' in filtered_df.columns:
            funnel_data = filtered_df.groupby('Funnel Stage')['Budget'].sum().to_dict()
            return [(stage, budget) for stage, budget in funnel_data.items() if budget > 0]
        
        return None
        
    except Exception as e:
        logger.error(f"Error preparing funnel chart data: {str(e)}")
        return None

def _prepare_media_type_chart_data(df, region, masterbrand):
    """
    Prepare data for the media type chart.
    
    Args:
        df: DataFrame with the Excel data
        region: Region filter
        masterbrand: Masterbrand filter
        
    Returns:
        list: Chart data or None if no data
    """
    try:
        # Filter data for this region/masterbrand combination
        filtered_df = df[
            (df['Region'] == region) &
            (df['Global Masterbrand'] == masterbrand)
        ].copy()
        
        if filtered_df.empty:
            return None
        
        # Group by media type and sum budgets
        if 'Media Type' in filtered_df.columns and 'Budget' in filtered_df.columns:
            media_data = filtered_df.groupby('Media Type')['Budget'].sum().to_dict()
            return [(media_type, budget) for media_type, budget in media_data.items() if budget > 0]
        
        return None
        
    except Exception as e:
        logger.error(f"Error preparing media type chart data: {str(e)}")
        return None

def _prepare_brand_chart_data(df, region, masterbrand):
    """
    Prepare data for the brand chart.
    
    Args:
        df: DataFrame with the Excel data
        region: Region filter
        masterbrand: Masterbrand filter
        
    Returns:
        list: Chart data or None if no data
    """
    try:
        # Filter data for this region/masterbrand combination
        filtered_df = df[
            (df['Region'] == region) &
            (df['Global Masterbrand'] == masterbrand)
        ].copy()
        
        if filtered_df.empty:
            return None
        
        # Group by brand and sum budgets
        if 'Brand' in filtered_df.columns and 'Budget' in filtered_df.columns:
            brand_data = filtered_df.groupby('Brand')['Budget'].sum().to_dict()
            return [(brand, budget) for brand, budget in brand_data.items() if budget > 0]
        
        return None
        
    except Exception as e:
        logger.error(f"Error preparing brand chart data: {str(e)}")
        return None

def _add_pie_chart(slide, chart_data, chart_title, position_info, chart_name=None):
    """Compatibility wrapper for modular pie chart generation."""

    if not position_info:
        logger.error("Failed to get chart position coordinates")
        return False

    return presentation_add_pie_chart(
        slide,
        chart_data,
        chart_title,
        position_info,
        CHART_STYLE_CONTEXT,
        CHART_COLOR_MAPPING,
        CHART_COLOR_CYCLE,
        chart_name=chart_name,
    )

def _populate_slide_content(new_slide, prs, combination_row, slide_title_suffix, 
                          split_table_data, split_metadata, split_idx, df, excel_path):
    """Populate a single slide with all content (title, table, charts, comments)."""

    template_slide = prs.slides[0]

    title_text = _apply_title(new_slide, template_slide, combination_row, slide_title_suffix)
    _clear_comments(new_slide)
    
    # Create and populate the main data table
    logger.info(f"Creating table for {combination_row[0]} - {combination_row[1]} - {combination_row[2]}{slide_title_suffix}")
    
    table_success = _add_and_style_table(new_slide, split_table_data, split_metadata, prs.slides[0])
    if table_success:
        logger.info(f"Table created successfully for slide")
    else:
        logger.warning(f"Failed to create table for slide")
    
    _populate_summary_tiles(new_slide, template_slide, df, combination_row, excel_path)
    _ensure_legend_shapes(new_slide, template_slide)

    return {
        "title": title_text,
        "subtitle": slide_title_suffix.strip() if slide_title_suffix else None,
        "table": split_table_data,
        "notes": None,
    }

def create_presentation(template_path, excel_path, output_path):
    """Creates a PowerPoint presentation based on a template and Excel data."""
    logger.info(f"Starting presentation creation using template: {template_path}")
    try:
        prs = Presentation(template_path)
        if not prs.slides:
            logger.error("Template presentation has no slides. Cannot proceed.")
            return False

        logger.debug("--- Available Slide Layouts in Template ---")
        for i, layout in enumerate(prs.slide_layouts):
            layout_name = layout.name if hasattr(layout, 'name') else 'Unknown Name'
            placeholder_count = len(layout.placeholders) if hasattr(layout, 'placeholders') else 0
            logger.debug(f"  Layout Index {i}: {layout_name} (Placeholders: {placeholder_count})")
        logger.debug("--- End Available Slide Layouts ---")

        # Log template slide structure
        logger.debug("--- Template Slide (First Slide) Structure ---")
        if prs.slides:
            template_slide = prs.slides[0]
            logger.debug(f"  Template slide has {len(template_slide.shapes)} shapes:")
            for i, shape in enumerate(template_slide.shapes):
                shape_name = shape.name if hasattr(shape, 'name') else 'Unnamed'
                shape_type = shape.shape_type if hasattr(shape, 'shape_type') else 'Unknown Type'
                logger.debug(f"    Shape {i}: Name='{shape_name}', Type={shape_type}")
        else:
            logger.debug("  No slides found in template!")
        logger.debug("--- End Template Slide Structure ---")

        if CLONE_PIPELINE_ENABLED and not _validate_template_shapes(prs.slides[0]):
            logger.error("Template validation failed; aborting presentation generation.")
            return False

        df = load_and_prepare_data(excel_path)
        if df is None or df.empty:
            logger.error("Failed to load or prepare data. Aborting presentation creation.")
            return False

        # The load_and_prepare_data() function returns a processed DataFrame with these column names
        country_col_name = 'Country'
        brand_col_name = 'Brand'
        year_col_name = 'Year'
        
        unique_combinations_raw = df[[country_col_name, brand_col_name, year_col_name]].drop_duplicates().values.tolist()
        logger.info(f"Found {len(unique_combinations_raw)} unique Country/Global Masterbrand/Year combinations.")

        # Calculate total investment for each combination
        combinations_with_investment = []
        for combination in unique_combinations_raw:
            country, brand, year = combination
            # Filter data for this specific combination
            combination_data = df[
                (df[country_col_name] == country) & 
                (df[brand_col_name] == brand) & 
                (df[year_col_name] == year)
            ]
            # Calculate total investment (sum of Total Cost for this combination)
            total_investment = combination_data['Total Cost'].fillna(0).sum()
            
            combinations_with_investment.append((country, brand, year, total_investment))
        
        # Group by market (country) and calculate total market investment
        market_investments = {}
        for country, brand, year, investment in combinations_with_investment:
            if country not in market_investments:
                market_investments[country] = {'total': 0, 'combinations': []}
            market_investments[country]['total'] += investment
            market_investments[country]['combinations'].append((country, brand, year, investment))
        
        # Sort markets by total investment (highest first)
        sorted_markets = sorted(market_investments.items(), key=lambda x: x[1]['total'], reverse=True)
        
        # Build final sorted list: markets ordered by total, brands within market ordered by individual investment
        ordered_combinations: list[tuple[str, str, int]] = []
        for market, data in sorted_markets:
            # Sort combinations within this market by individual investment
            market_combos = sorted(data['combinations'], key=lambda x: x[3], reverse=True)
            # Add to final list (without investment values)
            ordered_combinations.extend([(c, b, y) for c, b, y, _ in market_combos])
        
        # Log the sorted order with market totals
        logger.info("Slides will be generated grouped by market, ordered by total market investment:")
        for i, (market, data) in enumerate(sorted_markets[:10]):  # Log top 10 markets
            logger.info(f"  {i+1}. {market}: £{data['total']:,.0f} total")
            for combo in sorted(data['combinations'], key=lambda x: x[3], reverse=True)[:3]:  # Show top 3 brands
                logger.info(f"      - {combo[1]} ({combo[2]}): £{combo[3]:,.0f}")
            if len(data['combinations']) > 3:
                logger.info(f"      ... and {len(data['combinations']) - 3} more brands")
        if len(sorted_markets) > 10:
            logger.info(f"  ... and {len(sorted_markets) - 10} more markets")

        if not ordered_combinations:
            logger.warning("No unique Country/Global Masterbrand combinations found in the data.")
            # Decide if an empty presentation should be saved or an error returned
            # For now, let's save an empty presentation (after removing template slide)

        if not CLONE_PIPELINE_ENABLED:
            logger.info("Clone pipeline disabled via configuration; falling back to legacy AutoPPTX workflow.")
            return _generate_autopptx_only(
                template_path,
                output_path,
                df,
                ordered_combinations,
                excel_path,
            )

        autopptx_payloads: list[dict[str, object]] = []

        current_market = None
        for idx, combination_row in enumerate(ordered_combinations):
            # Check if we're starting a new market
            if combination_row[0] != current_market:
                current_market = combination_row[0]
                # Fix market name display
                display_market_name = "Morocco" if current_market == "MOR" else current_market
                # Add a market delimiter slide with black background
                try:
                    # Try to use a blank slide layout 
                    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
                    delimiter_slide = prs.slides.add_slide(blank_layout)
                except:
                    delimiter_slide = prs.slides.add_slide(prs.slide_layouts[0])
                
                # Add full black background
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Create a black rectangle covering the entire slide
                from pptx.enum.shapes import MSO_SHAPE
                black_bg = delimiter_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=0,
                    top=0,
                    width=slide_width,
                    height=slide_height
                )
                
                # Make the background solid black
                black_bg.fill.solid()
                black_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure black
                black_bg.line.fill.background()  # Remove border
                
                # Create a text box in the center of the slide for the market name
                text_box = delimiter_slide.shapes.add_textbox(
                    left=Inches(1),
                    top=int(slide_height * 0.45),  # Center vertically
                    width=slide_width - Inches(2),
                    height=Inches(1.5)
                )
                
                text_frame = text_box.text_frame
                text_frame.text = str(display_market_name).upper() if display_market_name else "UNKNOWN"  # Make market name full caps
                text_frame.word_wrap = False
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                # Format the text - large white font
                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(36)  # Large font (30-40 range)
                        run.font.bold = True
                        run.font.name = FONT_FAMILY_LEGEND
                        run.font.color.rgb = RGBColor(255, 255, 255)  # Pure white
                logger.info(f"Added market delimiter slide for: {display_market_name}")
            
            logger.info(f"Processing combination {idx+1}/{len(ordered_combinations)}: {combination_row[0]} - {combination_row[1]} - {combination_row[2]}")
            
            # First, prepare the table data to check if splitting is needed
            table_result = _prepare_main_table_data_detailed(df, combination_row[0], combination_row[1], combination_row[2], excel_path)
            
            if table_result[0] is None:
                logger.warning(f"No table data generated for {combination_row[0]} - {combination_row[1]}")
                continue
            
            table_data, cell_metadata = table_result
            
            # Split table if needed
            table_splits = _split_table_data_by_campaigns(table_data, cell_metadata)
            
            # Create a slide for each split
            for split_idx, (split_table_data, split_metadata, is_continuation) in enumerate(table_splits):
                # Add slide number to title if there are multiple splits
                if len(table_splits) > 1:
                    slide_title_suffix = f" ({split_idx + 1} of {len(table_splits)})"
                else:
                    slide_title_suffix = ""
                
                new_slide = prs.slides.add_slide(prs.slide_layouts[0])
                # logger.debug(f"Added new slide for {combination_row[0]} - {combination_row[1]} - {combination_row[2]}{slide_title_suffix}")
                
                # Populate this slide with content immediately
                payload = _populate_slide_content(
                    new_slide, prs, combination_row, slide_title_suffix, 
                    split_table_data, split_metadata, split_idx, df, excel_path
                )

                if payload:
                    autopptx_payloads.append(payload)

                # Diagnostic: Log all shapes on the new slide (commented out to reduce log size)
            # This entire block is commented out to reduce log file size
            # Re-enable for debugging shape-related issues
            """
            logger.debug(f"--- Shapes on new_slide (Layout: {prs.slide_layouts[0].name}, Index: {idx+1}) ---")
            if hasattr(new_slide, 'shapes'):
                logger.debug(f"  new_slide.shapes attribute exists. Number of shapes found: {len(new_slide.shapes)}")
                if new_slide.shapes:
                    for i, shape in enumerate(new_slide.shapes):
                        shape_name = shape.name if hasattr(shape, 'name') else 'Unnamed Shape'
                        shape_type = shape.shape_type if hasattr(shape, 'shape_type') else 'Unknown Type'
                        is_placeholder = shape.is_placeholder if hasattr(shape, 'is_placeholder') else False
                        placeholder_type = None
                        placeholder_idx = None
                        if is_placeholder and hasattr(shape, 'placeholder_format'):
                            ph_format = shape.placeholder_format
                            if hasattr(ph_format, 'type'):
                                placeholder_type = ph_format.type
                            if hasattr(ph_format, 'idx'):
                                placeholder_idx = ph_format.idx
                        logger.debug(f"  Shape {i}: Name='{shape_name}', Type={shape_type}, IsPlaceholder={is_placeholder}, PlaceholderType={placeholder_type}, PlaceholderIdx={placeholder_idx}")
                else:
                    logger.debug("  new_slide.shapes collection is empty (it exists but len is 0).")
            else:
                logger.debug("  new_slide.shapes attribute is missing.")
            logger.debug(f"--- End Shapes on new_slide (Index: {idx+1}) ---")
                """

        # ORIGINAL CONTENT POPULATION CODE MOVED TO _populate_slide_content() FUNCTION
        # (This section was deleted to prevent empty slides)
        # Remove the original template slide (slide 1) which shows "Region / Global Masterbrand"
        if len(prs.slides) > 1: # Only remove if we have other slides
            logger.info(f"Removing template slide from presentation (total slides: {len(prs.slides)}).")
            try:
                # The template slide is always the first slide added (index 0)
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
                logger.info("Template slide removed successfully.")
            except Exception as e:
                logger.warning(f"Could not remove template slide: {e}")
        else:
            logger.warning(f"Not removing template slide - only {len(prs.slides)} slides found.")

        # Change-tracking metadata references the removed template slide and breaks COM automation
        for rel_id, rel in list(prs.part.rels.items()):
            if rel.reltype == CHANGES_INFO_RELTYPE:
                try:
                    prs.part.drop_rel(rel_id)
                    logger.info("Removed stale changesInfo relationship %s", rel_id)
                except Exception as exc:
                    logger.warning(f"Failed to drop changesInfo relationship {rel_id}: {exc}")

        # CRITICAL FIX: Ensure output directory exists and add proper save error handling
        from pathlib import Path
        output_path_obj = Path(output_path)
        
        # Ensure output directory exists
        output_path_obj.parent.mkdir(parents=True, exist_ok=True)
        logger.info(f"Output directory ensured: {output_path_obj.parent}")
        
        # Ensure .pptx extension
        if not output_path.lower().endswith('.pptx'):
            output_path = output_path + '.pptx'
            logger.info(f"Added .pptx extension: {output_path}")

        # Save with proper error handling
        try:
            prs.save(output_path)
            logger.info(f"Presentation saved to {output_path}")
            
            # Verify file creation and get size
            file_size = os.path.getsize(output_path)
            logger.info(f"File verified: {file_size:,} bytes")

            _run_autopptx_pipeline(
                template_path,
                output_path,
                autopptx_payloads,
            )

            _run_aspose_exports(output_path)
            _run_docstrange_comparison(output_path, template_path)
            return True
        except Exception as e:
            logger.exception(f"❌ Save failed: {e}")
            return False

    except Exception as e:
        logger.error(f"An error occurred during presentation creation: {e}")
        logger.error(traceback.format_exc())
        return False


def _generate_autopptx_only(
    template_path: str | Path,
    output_path: str | Path,
    df: pd.DataFrame,
    ordered_combinations: list[tuple[str, str, int]],
    excel_path: str | Path | None,
) -> bool:
    if not AUTOPPTX_CONFIG.get("enabled", True):
        logger.warning("tooling.autopptx.enabled is false; proceeding because AutoPPTX was invoked explicitly.")

    if not autopptx_adapter.autopptx_available():
        logger.error("AutoPPTX dependency is unavailable; install 'autopptx' to enable the legacy workflow.")
        return False

    slide_payloads: list[SlidePayload] = []
    for combination_row in ordered_combinations:
        table_result = _prepare_main_table_data(
            df,
            combination_row[0],
            combination_row[1],
            combination_row[2],
            excel_path,
        )
        if not table_result:
            logger.warning(
                "Skipping combination %s - unable to prepare table data for legacy AutoPPTX pipeline",
                combination_row,
            )
            continue

        table_data, cell_metadata = table_result
        table_splits = _split_table_data_by_campaigns(table_data, cell_metadata)
        for split_idx, (split_table_data, _, _is_continuation) in enumerate(table_splits):
            suffix = f" ({split_idx + 1} of {len(table_splits)})" if len(table_splits) > 1 else ""
            title_text = _compose_title_text(combination_row, suffix)
            subtitle = suffix.strip() if suffix else None

            normalized_table = [
                ["" if cell is None else str(cell) for cell in row]
                for row in split_table_data
            ]

            slide_payloads.append(
                SlidePayload(
                    title=title_text,
                    subtitle=subtitle,
                    tables=[normalized_table],
                    notes=None,
                )
            )

    if not slide_payloads:
        logger.warning("No slide payloads generated for AutoPPTX workflow; presentation will not be created.")
        return False

    output_path_obj = Path(output_path)
    output_path_obj.parent.mkdir(parents=True, exist_ok=True)

    base_slide_index = int(AUTOPPTX_CONFIG.get("base_slide_index", 0))
    table_font_name = AUTOPPTX_CONFIG.get("table_font_name")
    table_font_size = AUTOPPTX_CONFIG.get("table_font_size")

    autopptx_adapter.generate_presentation(
        template_path,
        slide_payloads,
        output_path_obj,
        base_slide_index=base_slide_index,
        table_font_name=table_font_name,
        table_font_size=table_font_size,
    )
    logger.info("AutoPPTX presentation saved to %s", output_path_obj)
    return True


def _run_autopptx_pipeline(template_path: str | Path, main_output: str | Path, payloads: list[dict[str, object]]) -> None:
    if not AUTOPPTX_CONFIG.get("enabled"):
        return

    if not payloads:
        logger.warning("AutoPPTX requested but no slide payloads were captured; skipping generation")
        return

    if not autopptx_adapter.autopptx_available():  # pragma: no cover - optional dependency
        logger.warning("AutoPPTX package not available. Install 'autopptx' to enable placeholder rendering")
        return

    try:
        output_suffix = AUTOPPTX_CONFIG.get("output_suffix", "_autopptx")
        base_slide_index = int(AUTOPPTX_CONFIG.get("base_slide_index", 0))
        table_font_name = AUTOPPTX_CONFIG.get("table_font_name")
        table_font_size = AUTOPPTX_CONFIG.get("table_font_size")

        main_output_path = Path(main_output)
        output_path = main_output_path.with_stem(main_output_path.stem + output_suffix)

        slide_payloads: list[SlidePayload] = []
        for payload in payloads:
            table_rows = payload.get("table") or []
            tables = None
            if table_rows:
                normalized = [
                    ["" if cell is None else str(cell) for cell in row]
                    for row in table_rows
                ]
                tables = [normalized]

            slide_payloads.append(
                SlidePayload(
                    title=str(payload.get("title") or ""),
                    subtitle=payload.get("subtitle") or None,
                    tables=tables,
                    notes=payload.get("notes") or None,
                )
            )

        autopptx_adapter.generate_presentation(
            template_path,
            slide_payloads,
            output_path,
            base_slide_index=base_slide_index,
            table_font_name=table_font_name,
            table_font_size=int(table_font_size) if table_font_size else None,
        )
    except Exception as exc:  # pragma: no cover - optional integration
        logger.error("AutoPPTX generation failed: %s", exc)


def _run_aspose_exports(main_output: str | Path) -> None:
    if not ASPOSE_CONFIG.get("enabled"):
        return

    export_formats = ASPOSE_CONFIG.get("export_formats") or []
    if not export_formats:
        logger.warning("Aspose export enabled but no export_formats configured; skipping")
        return

    main_output_path = Path(main_output)

    output_target = ASPOSE_CONFIG.get("output_dir") or ASPOSE_CONFIG.get("output_subdir")
    if output_target:
        output_dir = Path(output_target)
        if not output_dir.is_absolute():
            output_dir = main_output_path.parent / output_dir
    else:
        output_dir = main_output_path.parent / "aspose_exports"

    try:
        aspose_converter.export_with_aspose(
            main_output_path,
            export_formats,
            output_dir,
            client_id=ASPOSE_CONFIG.get("client_id"),
            client_secret=ASPOSE_CONFIG.get("client_secret"),
            base_url=ASPOSE_CONFIG.get(
                "base_url",
                aspose_converter.ASPOSE_DEFAULT_BASE_URL,
            ),
            timeout=int(ASPOSE_CONFIG.get("timeout_seconds", 120)),
        )
    except aspose_converter.AsposeConfigurationError as exc:  # pragma: no cover - optional integration
        logger.warning("Aspose export skipped: %s", exc)
    except Exception as exc:  # pragma: no cover - optional integration
        logger.error("Aspose export failed: %s", exc)


def _run_docstrange_comparison(main_output: str | Path, template_path: str | Path) -> None:
    if not DOCSTRANGE_CONFIG.get("enabled"):
        return

    command = DOCSTRANGE_CONFIG.get("command")
    output_dir_cfg = (
        DOCSTRANGE_CONFIG.get("output_dir")
        or DOCSTRANGE_CONFIG.get("output_subdir")
        or "docstrange"
    )

    main_output_path = Path(main_output)
    output_dir = Path(output_dir_cfg)
    if not output_dir.is_absolute():
        output_dir = main_output_path.parent / output_dir

    reference = DOCSTRANGE_CONFIG.get("reference_template") or template_path
    output_format = DOCSTRANGE_CONFIG.get("format", "markdown")

    try:
        docstrange_validator.compare_presentations(
            main_output_path,
            reference,
            output_dir,
            command=command,
            output_format=output_format,
        )
    except docstrange_validator.DocStrangeNotAvailable as exc:  # pragma: no cover - optional integration
        logger.warning("DocStrange comparison skipped: %s", exc)
    except Exception as exc:  # pragma: no cover - optional integration
        logger.error("DocStrange comparison failed: %s", exc)

def _apply_internal_table_borders(table, total_rows):
    """
    Apply specific internal borders for the first 7 columns only (excluding header and total rows).
    Uses direct OXML manipulation to add #BFBFBF borders with 0.75pt thickness.
    
    Args:
        table: The PowerPoint table object
        total_rows: Total number of rows in the table
    """
    try:
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        
        # Border specifications - same as external borders for consistency
        border_color_rgb = (191, 191, 191) # CLR_TABLE_GRAY is RGBColor(191, 191, 191)
        border_width_emu = int(0.75 * 12700)  # 0.75pt in EMUs
        hex_color = f"{border_color_rgb[0]:02X}{border_color_rgb[1]:02X}{border_color_rgb[2]:02X}"
        
        logger.debug(f"Applying internal borders to first 7 columns for rows 1 to {total_rows-2}")
        
        # Apply internal borders only to body rows (excluding header row 0 and total row)
        for row_idx in range(1, total_rows - 1):  # Skip header (0) and total (last) rows
            row = table.rows[row_idx]
            
            for col_idx in range(7):  # First 7 columns only
                cell = row.cells[col_idx]
                
                try:
                    # Access the cell's table cell properties
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    
                    # Add bottom border for horizontal lines (apply to all qualifying cells)
                    if row_idx < total_rows - 2:  # Don't add bottom border to row just above total
                        bottom_border = tcPr.find(qn('a:lnB'))
                        if bottom_border is None:
                            bottom_border = tcPr.add(qn('a:lnB'))
                        
                        # Set border properties
                        bottom_border.set('w', str(border_width_emu))
                        bottom_border.set('cap', 'flat')
                        bottom_border.set('cmpd', 'sng')
                        
                        # Add solid fill with color
                        solidFill = bottom_border.find(qn('a:solidFill'))
                        if solidFill is None:
                            solidFill = bottom_border.add(qn('a:solidFill'))
                        
                        srgbClr = solidFill.find(qn('a:srgbClr'))
                        if srgbClr is None:
                            srgbClr = solidFill.add(qn('a:srgbClr'))
                        srgbClr.set('val', hex_color)
                        
                        # Add preset dash (solid line)
                        prstDash = bottom_border.find(qn('a:prstDash'))
                        if prstDash is None:
                            prstDash = bottom_border.add(qn('a:prstDash'))
                        prstDash.set('val', 'solid')
                    
                    # Add right border for vertical lines (apply to first 6 columns only)
                    if col_idx < 6:  # Columns 0, 1, 2, 3, 4, 5 get right borders
                        right_border = tcPr.find(qn('a:lnR'))
                        if right_border is None:
                            right_border = tcPr.add(qn('a:lnR'))
                        
                        # Set border properties
                        right_border.set('w', str(border_width_emu))
                        right_border.set('cap', 'flat')
                        right_border.set('cmpd', 'sng')
                        
                        # Add solid fill with color
                        solidFill = right_border.find(qn('a:solidFill'))
                        if solidFill is None:
                            solidFill = right_border.add(qn('a:solidFill'))
                        
                        srgbClr = solidFill.find(qn('a:srgbClr'))
                        if srgbClr is None:
                            srgbClr = solidFill.add(qn('a:srgbClr'))
                        srgbClr.set('val', hex_color)
                        
                        # Add preset dash (solid line)
                        prstDash = right_border.find(qn('a:prstDash'))
                        if prstDash is None:
                            prstDash = right_border.add(qn('a:prstDash'))
                        prstDash.set('val', 'solid')
                    
                except Exception as cell_border_error:
                    logger.debug(f"Could not apply internal borders to cell ({row_idx}, {col_idx}): {cell_border_error}")
        
        logger.info(f"Internal table borders applied successfully for first 7 columns with #BFBFBF color and 0.75pt width")
        return True
        
    except Exception as e:
        logger.warning(f"Error applying internal table borders: {str(e)}")
        return False

def _verify_file_exists(label, path_str, extra_search=()):
    """
    Return an absolute path to an existing file.

    extra_search – optional directories (relative to project root) to probe
                   if the user gave a bare filename.
    """
    from pathlib import Path
    p = Path(path_str).expanduser()
    
    # Try the path as given first
    if p.is_file():
        return str(p.resolve())
    
    # If the file doesn't exist at the given path and we have fallback directories,
    # and if it looks like a bare filename (no directory separators), search fallbacks
    if extra_search and '/' not in str(p) and '\\' not in str(p):
        project_root = Path(__file__).parent.parent  # Go up from scripts/ to project root
        for alt in extra_search:
            candidate = project_root / alt / p.name
            if candidate.is_file():
                return str(candidate.resolve())

    # If we get here, the file wasn't found anywhere
    p_resolved = p.resolve()
    raise FileNotFoundError(f"❌ {label} file not found: '{p_resolved}'")
def _unit_test__no_orphan_self():
    import inspect, re, pathlib
    src = pathlib.Path(__file__).read_text(encoding="utf-8")
    tree = ast.parse(src, filename=str(pathlib.Path(__file__)))
    offender_lines = []
    for node in ast.walk(tree):
        # Set parent attributes for all nodes to allow easy traversal up the tree
        for child in ast.iter_child_nodes(node):
            child.parent = node
            
        if isinstance(node, ast.Attribute) and isinstance(node.value, ast.Name) and node.value.id == 'self':
            # Ascend until we hit a FunctionDef, ClassDef, or Module
            current_node = node
            legit = False
            is_in_class_scope = False # Track if we are within a ClassDef scope at all
            
            while hasattr(current_node, 'parent'): # Check if parent attribute exists
                parent = current_node.parent
                if isinstance(parent, ast.FunctionDef):
                    # Check if 'self' is the first argument of this function
                    if parent.args.args and parent.args.args[0].arg == 'self':
                        # Now, ensure this FunctionDef is itself within a ClassDef
                        # Keep ascending from the FunctionDef to find a ClassDef
                        func_parent = parent
                        while hasattr(func_parent, 'parent'):
                            func_parent = func_parent.parent
                            if isinstance(func_parent, ast.ClassDef):
                                legit = True # 'self' is in a method of a class
                                break
                            if isinstance(func_parent, ast.Module): # Reached top level without finding a class
                                break # Stop ascent once a FunctionDef is found and checked
                
                if isinstance(parent, ast.ClassDef):
                    is_in_class_scope = True # We are inside a class, but 'self' might be at class level (e.g. self.x = 10)
                    # If 'self' is used directly under ClassDef (not in a method), it's an error unless it's part of an assignment in __init__ or other method
                    # The current check correctly flags 'self.attribute' if not in a method like 'def meth(self, ...):'
                    # If we hit a ClassDef before a qualifying FunctionDef that has 'self' as first arg, it's likely a class variable access using 'self', which is wrong.
                    break # Stop ascent if ClassDef is found before a qualifying FunctionDef
                
                if isinstance(parent, ast.Module):
                    break # Reached top level
                current_node = parent
            
            if not legit:
                # Additional check: if it's in a class scope but not in a method, it's an error
                # e.g. class MyClass: self.x = 10 (invalid)
                # vs class MyClass: def __init__(self): self.x = 10 (valid)
                # The 'legit' flag already covers the valid case. If not legit, it's an offender.
                offender_lines.append(node.lineno)

    if offender_lines:
        # Remove duplicates and sort
        unique_offender_lines = sorted(list(set(offender_lines)))
        raise RuntimeError(f" orphan 'self.' detected outside a method or in an invalid class context at lines: {unique_offender_lines} – aborting build")

_unit_test__no_orphan_self()


def build_presentation(template_path, excel_path, output_path):
    """Backward-compatible wrapper around ``create_presentation``."""

    return create_presentation(template_path, excel_path, output_path)
