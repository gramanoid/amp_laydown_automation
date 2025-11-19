"""
CLI entry point for post-processing operations.

This module provides a command-line interface for Python-based post-processing,
replacing slow COM automation with fast python-pptx operations.

**IMPORTANT**: Cell merges are created during deck generation (assembly.py),
NOT in post-processing. This CLI is primarily for normalization and edge case repairs.

Typical usage:
    python -m amp_automation.presentation.postprocess.cli \\
        --presentation-path "path/to/deck.pptx" \\
        --operations normalize

Edge case repairs (rarely needed):
    python -m amp_automation.presentation.postprocess.cli \\
        --presentation-path "path/to/deck.pptx" \\
        --operations normalize,merge-campaign,merge-monthly,merge-summary \\
        --slide-filter 2,3,4 \\
        --verbose

Performance: ~30 seconds for 88-slide deck (vs 10+ hours with COM)
See: docs/ARCHITECTURE_DECISION_COM_PROHIBITION.md
"""

import argparse
import logging
import sys
from pathlib import Path
from typing import List, Optional

from pptx import Presentation

from . import (
    normalize_table_layout,
    apply_blank_cell_formatting,
    normalize_table_fonts,
    delete_carried_forward_rows,
    fix_grand_total_wrapping,
    remove_pound_signs_from_totals,
    reset_primary_column_spans,
    reset_column_group,
    merge_campaign_cells,
    merge_media_cells,
    merge_percentage_cells,
    merge_monthly_total_cells,
    merge_summary_cells,
    unmerge_all_cells,
    unmerge_primary_columns,
)

logger = logging.getLogger(__name__)


class PostProcessorCLI:
    """CLI handler for presentation post-processing operations."""

    OPERATIONS = {
        "postprocess-all": "RECOMMENDED: Complete post-processing workflow (unmerge -> clean -> merge -> format)",
        "normalize": "Normalize table layout and cell formatting",
        "normalize-fonts": "Enforce Verdana fonts: 6pt body, 7pt header/bottom",
        "delete-carried-forward": "Delete CARRIED FORWARD rows from tables",
        "fix-grand-total-wrap": "Fix GRAND TOTAL row wrapping to single line",
        "remove-pound-totals": "Remove £ signs from GRAND TOTAL and MONTHLY TOTAL rows",
        "unmerge-all": "Unmerge ALL cells - clean slate before selective merging",
        "unmerge-primary": "Unmerge only primary columns (1-3) - less aggressive than unmerge-all",
        "reset-spans": "Reset column spans in primary columns (edge case repair)",
        "merge-campaign": "Merge campaign cells vertically in column 1 (apply after unmerge)",
        "merge-media": "Merge media channel cells vertically in column 2 (apply after unmerge)",
        "merge-percentage": "Merge percentage cells vertically in column 17 (apply after merge-monthly)",
        "merge-monthly": "Merge monthly total cells horizontally cols 1-3 (apply after unmerge)",
        "merge-summary": "Merge summary cells GRAND TOTAL cols 1-3 (apply after unmerge)",
    }

    # Definitive post-processing workflow (validated on 88-slide deck, 100% success rate)
    POSTPROCESS_ALL_WORKFLOW = [
        "unmerge-all",
        "delete-carried-forward",
        "merge-campaign",
        "merge-media",
        "merge-monthly",
        "merge-percentage",  # Must run AFTER merge-monthly to find gray MONTHLY TOTAL rows
        "merge-summary",
        "fix-grand-total-wrap",
        "remove-pound-totals",
        "normalize-fonts",
    ]

    def __init__(self, presentation_path: Path, slide_filter: Optional[List[int]] = None):
        self.presentation_path = presentation_path
        self.slide_filter = slide_filter
        self.prs = None

    def load_presentation(self):
        """Load the presentation file."""
        if not self.presentation_path.exists():
            raise FileNotFoundError(f"Presentation not found: {self.presentation_path}")

        logger.info(f"Loading presentation: {self.presentation_path}")
        self.prs = Presentation(str(self.presentation_path))
        logger.info(f"Loaded {len(self.prs.slides)} slides")

    def save_presentation(self):
        """Save the presentation file."""
        logger.info(f"Saving presentation: {self.presentation_path}")
        self.prs.save(str(self.presentation_path))
        logger.info("Presentation saved successfully")

    def run_operation(self, operation: str, slide_idx: int, table) -> bool:
        """
        Run a single operation on a table.

        Args:
            operation: Operation name (e.g., "normalize", "merge-campaign", "postprocess-all")
            slide_idx: Slide index (1-based)
            table: python-pptx table object

        Returns:
            True if operation succeeded, False otherwise
        """
        try:
            if operation == "postprocess-all":
                # Run complete workflow in sequence
                for sub_op in self.POSTPROCESS_ALL_WORKFLOW:
                    if not self.run_operation(sub_op, slide_idx, table):
                        logger.error(f"Slide {slide_idx} - Workflow failed at operation: {sub_op}")
                        return False
                return True
            elif operation == "normalize":
                normalize_table_layout(table)
                apply_blank_cell_formatting(table)
            elif operation == "normalize-fonts":
                normalize_table_fonts(table)
            elif operation == "delete-carried-forward":
                delete_carried_forward_rows(table)
            elif operation == "fix-grand-total-wrap":
                fix_grand_total_wrapping(table)
            elif operation == "remove-pound-totals":
                remove_pound_signs_from_totals(table)
            elif operation == "unmerge-all":
                unmerge_all_cells(table)
            elif operation == "unmerge-primary":
                unmerge_primary_columns(table, max_cols=3)
            elif operation == "reset-spans":
                reset_primary_column_spans(table, max_cols=3)
                reset_column_group(table, max_cols=3)
            elif operation == "merge-campaign":
                merge_campaign_cells(table)
            elif operation == "merge-media":
                merge_media_cells(table)
            elif operation == "merge-monthly":
                merge_monthly_total_cells(table)
            elif operation == "merge-percentage":
                merge_percentage_cells(table)
            elif operation == "merge-summary":
                merge_summary_cells(table)
            else:
                logger.warning(f"Unknown operation: {operation}")
                return False

            return True
        except Exception as e:
            logger.error(f"Slide {slide_idx} - Operation '{operation}' failed: {e}")
            return False

    def process(self, operations: List[str]) -> int:
        """
        Process all slides with the specified operations.

        Args:
            operations: List of operation names to run

        Returns:
            Exit code (0 = success, 1 = error)
        """
        self.load_presentation()

        total_operations = 0
        failed_operations = 0

        for slide_idx, slide in enumerate(self.prs.slides, start=1):
            # Apply slide filter if specified
            if self.slide_filter and slide_idx not in self.slide_filter:
                continue

            logger.info(f"Processing slide {slide_idx}")

            # Find the main table (largest table by cell count)
            tables = [shape for shape in slide.shapes if shape.has_table]
            if not tables:
                logger.debug(f"Slide {slide_idx} - No tables found")
                continue

            # Select largest table
            main_table = max(tables, key=lambda s: s.table.rows.__len__() * s.table.columns.__len__())
            table = main_table.table

            row_count = len(table.rows)
            col_count = len(table.columns)
            logger.debug(f"Slide {slide_idx} - Processing table: {row_count} rows × {col_count} columns")

            # Run each operation
            for operation in operations:
                total_operations += 1
                logger.debug(f"Slide {slide_idx} - Running: {operation}")

                if not self.run_operation(operation, slide_idx, table):
                    failed_operations += 1

        self.save_presentation()

        logger.info(f"Completed: {total_operations} operations, {failed_operations} failed")

        return 1 if failed_operations > 0 else 0


def main():
    """Main CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Post-process PowerPoint presentations with Python (replaces slow COM automation)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=f"""
Available operations:
{chr(10).join(f"  {op:15} - {desc}" for op, desc in PostProcessorCLI.OPERATIONS.items())}

IMPORTANT: Cell merges are created during generation, NOT post-processing.
           Only use merge operations to repair decks from failed generation.

Recommended usage (normalization only):
  python -m amp_automation.presentation.postprocess.cli \\
      --presentation-path deck.pptx --operations normalize

Edge case repairs (broken decks from failed generation):
  python -m amp_automation.presentation.postprocess.cli \\
      --presentation-path deck.pptx \\
      --operations normalize,merge-campaign,merge-monthly,merge-summary \\
      --slide-filter 2,3,4 --verbose

Performance: ~30 seconds for 88-slide deck (vs 10+ hours COM automation)
""",
    )

    parser.add_argument(
        "--presentation-path",
        type=Path,
        required=True,
        help="Path to the PowerPoint presentation (.pptx)",
    )

    parser.add_argument(
        "--operations",
        type=str,
        required=True,
        help=f"Comma-separated list of operations: {', '.join(PostProcessorCLI.OPERATIONS.keys())}",
    )

    parser.add_argument(
        "--slide-filter",
        type=str,
        help="Comma-separated list of slide numbers to process (1-based). If omitted, processes all slides.",
    )

    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="Enable verbose logging",
    )

    args = parser.parse_args()

    # Configure logging
    log_level = logging.DEBUG if args.verbose else logging.INFO
    logging.basicConfig(
        level=log_level,
        format="%(asctime)s - %(levelname)s - %(message)s",
        datefmt="%H:%M:%S",
    )

    # Parse operations
    operations = [op.strip() for op in args.operations.split(",")]
    invalid_ops = [op for op in operations if op not in PostProcessorCLI.OPERATIONS]
    if invalid_ops:
        logger.error(f"Invalid operations: {', '.join(invalid_ops)}")
        logger.error(f"Valid operations: {', '.join(PostProcessorCLI.OPERATIONS.keys())}")
        return 1

    # Parse slide filter
    slide_filter = None
    if args.slide_filter:
        try:
            slide_filter = [int(s.strip()) for s in args.slide_filter.split(",")]
        except ValueError:
            logger.error(f"Invalid slide filter: {args.slide_filter}")
            return 1

    # Run processor
    try:
        processor = PostProcessorCLI(args.presentation_path, slide_filter)
        return processor.process(operations)
    except Exception as e:
        logger.exception(f"Fatal error: {e}")
        return 1


if __name__ == "__main__":
    sys.exit(main())
