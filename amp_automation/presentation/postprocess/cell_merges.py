"""
Cell merge operations for campaign, monthly, and summary rows.

This module provides Python-based implementations of cell merge operations,
replacing slow COM-based PowerShell operations.
"""

import logging
import re
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# Media type colors (matching legend colors for visual consistency)
CLR_TELEVISION = RGBColor(211, 254, 201)  # Light green
CLR_DIGITAL = RGBColor(253, 242, 183)     # Light yellow
CLR_OOH = RGBColor(255, 217, 97)          # Orange/gold
CLR_OTHER = RGBColor(176, 211, 255)       # Light blue

# Font configuration (matching template requirements)
FONT_NAME = "Verdana"
CAMPAIGN_FONT_SIZE = 6  # Campaign text
MEDIA_FONT_SIZE = 6  # Media column text
MONTHLY_TOTAL_FONT_SIZE = 6  # Body text
BRAND_TOTAL_FONT_SIZE = 7  # BRAND TOTAL row (larger to stand out)
SUMMARY_FONT_SIZE = 6  # Bottom row (regular rows)


def _smart_line_break(text: str) -> str:
    """
    Intelligently break campaign names into multiple lines to prevent mid-word breaks.

    Rules:
    - Replace dashes with spaces (e.g., "FACES-CONDITION" becomes "FACES CONDITION")
    - 1 word: keep as-is
    - 2 words: split into 2 lines (1 word per line)
    - 3 words: split into 2 lines (2 words on first line, 1 word on second line)
    - 4+ words: split into 2 lines (half words on each line, rounded up on first line)

    Edge cases:
    - "FACES-CONDITION" → "FACES\nCONDITION" (explicit split at hyphen)

    Args:
        text: Campaign name text

    Returns:
        Text with newline characters inserted for optimal wrapping
    """
    if not text or text.strip() == "":
        return text

    # EDGE CASE: Explicit handling for hyphenated compound words
    # Force split at hyphen for cases like "FACES-CONDITION"
    if "-" in text:
        parts = text.split("-")
        # If we have exactly 2 parts separated by hyphen, split them
        if len(parts) == 2:
            result = f"{parts[0].strip()}\n{parts[1].strip()}"
            logger.debug(f"Smart line break (hyphen split): '{text}' -> '{result}'")
            return result

    # Replace dashes with spaces, then split by whitespace
    normalized_text = text.replace("-", " ")
    words = [word for word in normalized_text.split() if word]

    logger.debug(f"Smart line break: '{text}' -> normalized: '{normalized_text}' -> words: {words}")

    if len(words) <= 1:
        # Single word or empty - no breaking needed, but use normalized text
        result = normalized_text.strip()
    elif len(words) == 2:
        # 2 words: put each on separate line
        result = "\n".join(words)
    elif len(words) == 3:
        # 3 words: 2 on first line, 1 on second line
        result = f"{words[0]} {words[1]}\n{words[2]}"
    else:
        # 4+ words: split roughly in half
        mid_point = (len(words) + 1) // 2  # Round up
        first_line = " ".join(words[:mid_point])
        second_line = " ".join(words[mid_point:])
        result = f"{first_line}\n{second_line}"

    logger.debug(f"Smart line break result: '{result}'")
    return result


def merge_campaign_cells(table):
    """
    Merge campaign cells vertically in column 1.

    Equivalent to PowerShell function: Campaign merge operations

    This function identifies campaign rows (rows between MONTHLY TOTAL rows)
    and merges the cells in column 1 vertically to create a single cell
    spanning multiple rows for each campaign.

    Args:
        table: python-pptx table object

    Returns:
        int: Number of campaign merges performed
    """
    logger.debug("Merging campaign cells")

    row_count = len(table.rows)
    campaign_start = None
    campaign_name = None  # Store campaign name when we find it
    merges_performed = 0

    # Iterate through rows (skip header row index 0)
    for row_idx in range(1, row_count):
        cell = table.cell(row_idx, 0)  # Column 0 (CAMPAIGN column)
        cell_text = _get_cell_text(cell)
        is_gray = _has_gray_background(cell)

        # CRITICAL FIX: Some WHITE cells have "MONTHLY TOTAL\nCAMPAIGN_NAME" format
        # Extract campaign name from white cells only (preserve gray MONTHLY TOTAL cells)
        if not is_gray:
            actual_campaign_name = _extract_campaign_name(cell_text)
        else:
            # Gray cell - don't extract, use first line for detection
            actual_campaign_name = normalize_label(cell_text)

        # Check if this is a special row using the EXTRACTED campaign name (not raw text)
        # This prevents white cells with "MONTHLY TOTAL\nCAMPAIGN" from being misidentified
        is_monthly = is_monthly_total(actual_campaign_name) if actual_campaign_name else False
        is_grand = is_grand_total(actual_campaign_name) if actual_campaign_name else False

        # Track campaign start (non-empty, non-special rows, non-gray)
        # Campaign cells are WHITE cells with campaign names
        if actual_campaign_name and not is_monthly and not is_grand and not is_gray:
            if campaign_start is None:
                campaign_start = row_idx
                campaign_name = actual_campaign_name  # Save the campaign name immediately
                logger.debug(f"Found campaign start at row {row_idx}: {campaign_name}")

        # Perform merge when we hit a GRAY MONTHLY TOTAL row
        # (White cells with "MONTHLY TOTAL\nCAMPAIGN" are campaign cells, not triggers)
        if is_monthly and is_gray and campaign_start is not None:
            campaign_end = row_idx - 1

            if campaign_end > campaign_start:
                try:
                    # Merge cells vertically in column 0 (CAMPAIGN column)
                    top_cell = table.cell(campaign_start, 0)
                    bottom_cell = table.cell(campaign_end, 0)

                    # Check if already merged
                    if not _cells_are_same(top_cell, bottom_cell):
                        top_cell.merge(bottom_cell)
                        merges_performed += 1
                        logger.debug(f"Merged campaign rows {campaign_start}-{campaign_end}: {campaign_name}")

                    # Apply styling to merged cell and set cleaned campaign name with smart line breaks
                    merged_cell = table.cell(campaign_start, 0)
                    # Apply smart line breaking to prevent mid-word breaks
                    formatted_campaign_name = _smart_line_break(campaign_name)
                    _apply_cell_styling(
                        merged_cell,
                        text=formatted_campaign_name,  # Set the campaign name with smart line breaks
                        font_size=CAMPAIGN_FONT_SIZE,
                        bold=True,
                        center_align=True,
                        vertical_center=True
                    )

                except Exception as e:
                    logger.error(f"Failed to merge campaign rows {campaign_start}-{campaign_end}: {e}")

            campaign_start = None
            campaign_name = None  # Reset campaign name too

    logger.info(f"Campaign merges completed: {merges_performed} merge(s)")
    return merges_performed


def merge_media_cells(table):
    """
    Merge media channel cells vertically in column 1 (MEDIA column).

    This function identifies media channels (TELEVISION, DIGITAL, OOH, OTHER)
    and merges the cells vertically to span all metric rows for that media channel.

    Args:
        table: python-pptx table object

    Returns:
        int: Number of media channel merges performed
    """
    logger.debug("Merging media channel cells")

    row_count = len(table.rows)
    merges_performed = 0

    # Known media channels to look for
    MEDIA_CHANNELS = ["TELEVISION", "DIGITAL", "OOH", "OTHER", "RADIO", "PRINT", "CINEMA"]

    media_start = None
    media_name = None

    # Iterate through rows (skip header row index 0)
    for row_idx in range(1, row_count):
        cell = table.cell(row_idx, 1)  # Column 1 (0-indexed) is MEDIA column
        cell_text = _get_cell_text(cell).strip().upper()

        # Check if this is a media channel name
        is_media_channel = any(channel in cell_text for channel in MEDIA_CHANNELS)

        # Check if this row is empty or has dash (metric row)
        is_empty_or_dash = cell_text in ["", "-"]

        if is_media_channel:
            # If we were tracking a previous media channel, merge it now
            if media_start is not None and media_start < row_idx - 1:
                media_end = row_idx - 1
                try:
                    top_cell = table.cell(media_start, 1)
                    bottom_cell = table.cell(media_end, 1)

                    # Check if already merged
                    if not _cells_are_same(top_cell, bottom_cell):
                        top_cell.merge(bottom_cell)
                        merges_performed += 1
                        logger.debug(f"Merged {media_name} rows {media_start}-{media_end}")

                    # Apply styling to merged cell
                    merged_cell = table.cell(media_start, 1)
                    _apply_cell_styling(
                        merged_cell,
                        text=media_name,
                        font_size=MEDIA_FONT_SIZE,
                        bold=True,
                        center_align=True,
                        vertical_center=True
                    )

                except Exception as e:
                    logger.error(f"Failed to merge {media_name} rows {media_start}-{media_end}: {e}")

            # Start tracking this new media channel
            media_start = row_idx
            media_name = cell_text
            logger.debug(f"Found media channel at row {row_idx}: {media_name}")

        # Check if this row ends the current media channel (MONTHLY TOTAL or campaign name)
        campaign_cell = table.cell(row_idx, 0)  # Column 0 (CAMPAIGN)
        campaign_text = _get_cell_text(campaign_cell).strip().upper()
        # Normalize to handle non-breaking characters from styling
        normalized_campaign = normalize_label(campaign_text)
        is_special_row = (
            is_monthly_total(normalized_campaign) or
            is_grand_total(normalized_campaign)
        )

        # If we hit a special row and were tracking a media channel, merge it
        if is_special_row and media_start is not None:
            media_end = row_idx - 1

            if media_end >= media_start:
                try:
                    top_cell = table.cell(media_start, 1)
                    bottom_cell = table.cell(media_end, 1)

                    # Check if already merged
                    if not _cells_are_same(top_cell, bottom_cell):
                        top_cell.merge(bottom_cell)
                        merges_performed += 1
                        logger.debug(f"Merged {media_name} rows {media_start}-{media_end}")

                    # Apply styling to merged cell
                    merged_cell = table.cell(media_start, 1)
                    _apply_cell_styling(
                        merged_cell,
                        text=media_name,
                        font_size=MEDIA_FONT_SIZE,
                        bold=True,
                        center_align=True,
                        vertical_center=True
                    )

                except Exception as e:
                    logger.error(f"Failed to merge {media_name} rows {media_start}-{media_end}: {e}")

            media_start = None
            media_name = None

    # Handle case where table ends while still tracking a media channel
    if media_start is not None and media_start < row_count - 1:
        media_end = row_count - 1
        try:
            top_cell = table.cell(media_start, 1)
            bottom_cell = table.cell(media_end, 1)

            if not _cells_are_same(top_cell, bottom_cell):
                top_cell.merge(bottom_cell)
                merges_performed += 1
                logger.debug(f"Merged {media_name} rows {media_start}-{media_end}")

            merged_cell = table.cell(media_start, 1)
            _apply_cell_styling(
                merged_cell,
                text=media_name,
                font_size=MEDIA_FONT_SIZE,
                bold=True,
                center_align=True,
                vertical_center=True
            )

        except Exception as e:
            logger.error(f"Failed to merge {media_name} rows {media_start}-{media_end}: {e}")

    logger.info(f"Media channel merges completed: {merges_performed} merge(s)")
    return merges_performed


def _apply_colored_media_splits(cell, text: str):
    """
    Apply colored formatting to media split text in MONTHLY TOTAL rows.

    Colors media abbreviations to match their legend colors:
    - TV → Green
    - DIG → Yellow
    - OOH → Orange
    - OTH → Blue

    Example: "TOTAL - TV 38% • DIG 34% • OOH 25% • OTH 2%"

    Args:
        cell: python-pptx cell object
        text: Text containing media splits (e.g., "TOTAL - TV 38% • DIG 34%...")
    """
    try:
        text_frame = cell.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.margin_left = Pt(1)
        text_frame.margin_right = Pt(1)
        text_frame.margin_top = Pt(1)
        text_frame.margin_bottom = Pt(1)

        # Create single paragraph
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        # Define media colors
        media_colors = {
            "TV": CLR_TELEVISION,
            "DIG": CLR_DIGITAL,
            "OOH": CLR_OOH,
            "OTH": CLR_OTHER,
        }

        # Parse text and create colored runs
        # Pattern: "TOTAL - TV 38% • DIG 34% • OOH 25% • OTH 2%"
        parts = re.split(r'(TV|DIG|OOH|OTH)', text)

        for part in parts:
            if not part:
                continue

            run = p.add_run()
            run.text = part
            run.font.name = FONT_NAME
            run.font.size = Pt(MONTHLY_TOTAL_FONT_SIZE)
            run.font.bold = True

            # Apply color if this is a media abbreviation
            if part in media_colors:
                run.font.color.rgb = media_colors[part]

        # Apply vertical centering
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    except Exception as e:
        logger.error(f"Failed to apply colored media splits: {e}")
        # Fallback to regular styling
        cell.text_frame.text = text


def merge_monthly_total_cells(table):
    """
    Merge monthly total cells horizontally (columns 1-3).

    Equivalent to PowerShell function: Monthly total merge operations

    This function identifies MONTHLY TOTAL rows and merges cells horizontally
    across columns 1-3 to create a single cell for the "MONTHLY TOTAL" label.

    IMPORTANT: Only merges cells with gray background. White background cells
    with "MONTHLY TOTAL" text are campaign names and should NOT be merged.

    Args:
        table: python-pptx table object

    Returns:
        int: Number of monthly total merges performed
    """
    logger.debug("Merging monthly total cells")

    row_count = len(table.rows)
    col_count = len(table.columns)
    merges_performed = 0

    # Iterate through rows (skip header row index 0)
    for row_idx in range(1, row_count):
        cell = table.cell(row_idx, 0)  # Column 1 (0-indexed)
        cell_text = _get_cell_text(cell)

        # Check both text content AND background color
        if is_monthly_total(cell_text) and _has_gray_background(cell):
            normalized = normalize_label(cell_text)

            try:
                # Merge horizontally across columns 1-3 (0-2 in 0-indexed)
                # Merge from first column to last column in ONE operation
                left_cell = table.cell(row_idx, 0)
                right_cell = table.cell(row_idx, 2)  # Column 3 (0-indexed as 2)

                # Check if already merged
                if not _cells_are_same(left_cell, right_cell):
                    left_cell.merge(right_cell)
                    merges_performed += 1
                    logger.debug(f"Merged MONTHLY TOTAL row {row_idx} across columns 1-3")

                # Apply styling to merged cell
                merged_cell = table.cell(row_idx, 0)

                # Check if text contains media splits (new format)
                if "TV" in normalized or "DIG" in normalized or "OOH" in normalized:
                    # Apply colored media split styling
                    _apply_colored_media_splits(merged_cell, normalized)
                else:
                    # Apply regular styling (old format)
                    _apply_cell_styling(
                        merged_cell,
                        text=normalized,
                        font_size=MONTHLY_TOTAL_FONT_SIZE,
                        bold=True,
                        center_align=True,
                        vertical_center=True
                    )

            except Exception as e:
                logger.error(f"Failed to merge MONTHLY TOTAL row {row_idx}: {e}")

    logger.info(f"Monthly total merges completed: {merges_performed} merge(s)")
    return merges_performed


def merge_summary_cells(table):
    """
    Merge summary cells horizontally (columns 1-3).

    Equivalent to PowerShell function: Summary merge operations

    This function identifies GRAND TOTAL and CARRIED FORWARD rows and merges
    cells horizontally across columns 1-3 to create a single cell for the label.

    Args:
        table: python-pptx table object

    Returns:
        int: Number of summary merges performed
    """
    logger.debug("Merging summary cells (GRAND TOTAL, CARRIED FORWARD)")

    row_count = len(table.rows)
    col_count = len(table.columns)
    merges_performed = 0

    # Iterate through rows (skip header row index 0)
    for row_idx in range(1, row_count):
        cell = table.cell(row_idx, 0)  # Column 1 (0-indexed)
        cell_text = _get_cell_text(cell)

        # Check both text content AND background color
        if is_grand_total(cell_text) and _has_gray_background(cell):
            normalized = normalize_label(cell_text)

            try:
                # Merge horizontally across columns 1-3 (0-2 in 0-indexed)
                # Merge from first column to last column in ONE operation
                left_cell = table.cell(row_idx, 0)
                right_cell = table.cell(row_idx, 2)  # Column 3 (0-indexed as 2)

                # Check if already merged
                if not _cells_are_same(left_cell, right_cell):
                    left_cell.merge(right_cell)
                    merges_performed += 1
                    logger.debug(f"Merged summary row {row_idx} ({normalized}) across columns 1-3")

                # Apply styling to merged cell
                # Use 7pt for BRAND TOTAL, 6pt for other summary rows
                merged_cell = table.cell(row_idx, 0)
                is_brand_total = "BRAND" in normalized
                font_size_to_use = BRAND_TOTAL_FONT_SIZE if is_brand_total else SUMMARY_FONT_SIZE

                _apply_cell_styling(
                    merged_cell,
                    text=normalized,
                    font_size=font_size_to_use,
                    bold=True,
                    center_align=True,
                    vertical_center=True
                )

                # Apply green background if this is BRAND TOTAL
                if is_brand_total:
                    try:
                        fill = merged_cell.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(0x30, 0xea, 0x03)  # Green #30ea03
                        logger.debug(f"Applied green background to BRAND TOTAL row {row_idx}")
                    except Exception as bg_error:
                        logger.warning(f"Failed to apply green background to BRAND TOTAL: {bg_error}")

            except Exception as e:
                logger.error(f"Failed to merge summary row {row_idx} ({normalized}): {e}")

    logger.info(f"Summary merges completed: {merges_performed} merge(s)")
    return merges_performed


def merge_percentage_cells(table):
    """
    Merge percentage cells vertically in column 17 (% column).

    Similar to campaign merging, this function identifies campaign rows and merges
    the percentage cells vertically from the start of a campaign until the cell
    before the gray MONTHLY TOTAL row.

    Args:
        table: python-pptx table object

    Returns:
        int: Number of percentage merges performed
    """
    logger.debug("Merging percentage cells")

    row_count = len(table.rows)
    merge_start = None
    merges_performed = 0

    # Iterate through rows (skip header row index 0)
    for row_idx in range(1, row_count):
        cell = table.cell(row_idx, 0)  # Check CAMPAIGN column to detect row type
        cell_text = _get_cell_text(cell)
        is_gray = _has_gray_background(cell)

        # Extract campaign name from cell (handling "MONTHLY TOTAL\nCAMPAIGN" format)
        if not is_gray:
            actual_campaign_name = _extract_campaign_name(cell_text)
        else:
            actual_campaign_name = normalize_label(cell_text)

        is_monthly = is_monthly_total(actual_campaign_name) if actual_campaign_name else False
        is_grand = is_grand_total(actual_campaign_name) if actual_campaign_name else False

        # Track merge start (non-empty, non-special rows, non-gray)
        if actual_campaign_name and not is_monthly and not is_grand and not is_gray:
            if merge_start is None:
                merge_start = row_idx
                logger.debug(f"Found percentage merge start at row {row_idx}")

        # Perform merge when we hit a GRAY MONTHLY TOTAL row
        if is_monthly and is_gray and merge_start is not None:
            merge_end = row_idx - 1  # Merge until row before MONTHLY TOTAL

            if merge_end > merge_start:
                try:
                    # Clear intermediate cell contents before merging
                    for clear_row in range(merge_start + 1, merge_end + 1):
                        cell_to_clear = table.cell(clear_row, 17)
                        cell_to_clear.text = ""

                    # Merge cells in column 17 (% column)
                    top_cell = table.cell(merge_start, 17)
                    bottom_cell = table.cell(merge_end, 17)

                    # Check if already merged
                    if not _cells_are_same(top_cell, bottom_cell):
                        top_cell.merge(bottom_cell)
                        merges_performed += 1
                        logger.debug(f"Merged percentage cells rows {merge_start}-{merge_end}")

                    # Apply styling to merged cell (bold, centered, vertically centered)
                    merged_cell = table.cell(merge_start, 17)
                    _apply_cell_styling(
                        merged_cell,
                        font_size=6,
                        bold=True,
                        center_align=True,
                        vertical_center=True
                    )

                except Exception as e:
                    logger.error(f"Failed to merge percentage cells rows {merge_start}-{merge_end}: {e}")

            merge_start = None

    logger.info(f"Percentage merges completed: {merges_performed} merge(s)")
    return merges_performed


# Helper function for future implementation
def normalize_label(text: str) -> str:
    """
    Normalize label text for comparison.

    IMPORTANT: Cells may contain multiple lines (e.g., "MONTHLY TOTAL\nTELEVISION").
    We only want the first line for labels.

    Also normalizes non-breaking characters that may have been added during styling:
    - \u2011 (non-breaking hyphen) → - (regular hyphen)
    - \u00A0 (non-breaking space) → regular space

    Args:
        text: Cell text content

    Returns:
        Normalized text (uppercase, stripped whitespace, first line only, regular characters)
    """
    if not text:
        return ""
    # Take only the first line (before any newline character)
    first_line = text.split('\n')[0].split('\r')[0]
    # Normalize non-breaking characters back to regular characters
    normalized = first_line.strip().replace('\u2011', '-').replace('\u00A0', ' ')
    return normalized.upper()


def is_monthly_total(cell_text: str) -> bool:
    """
    Check if cell text represents a MONTHLY TOTAL row.

    Args:
        cell_text: Text content from cell in column 1

    Returns:
        True if this is a MONTHLY TOTAL row
    """
    normalized = normalize_label(cell_text)
    # Match both old format "MONTHLY TOTAL (£ 000)" and new format "TOTAL - TV 38%..."
    return (
        ("MONTHLY" in normalized and "TOTAL" in normalized) or
        normalized.startswith("TOTAL -") or
        normalized.startswith("TOTAL-")
    )


def is_grand_total(cell_text: str) -> bool:
    """
    Check if cell text represents a GRAND TOTAL or BRAND TOTAL row.

    Args:
        cell_text: Text content from cell in column 1

    Returns:
        True if this is a GRAND TOTAL or BRAND TOTAL row
    """
    normalized = normalize_label(cell_text)
    return ("GRAND" in normalized or "BRAND" in normalized) and "TOTAL" in normalized


# Private helper functions
def _extract_campaign_name(cell_text: str) -> str:
    """
    Extract campaign name from cell text.

    Some cells have format: "MONTHLY TOTAL ( 000)\nCAMPAIGN_NAME"
    We need to extract just the campaign name (the line after MONTHLY TOTAL).

    Args:
        cell_text: Raw cell text content

    Returns:
        Campaign name (normalized, uppercase), or empty string
    """
    if not cell_text:
        return ""

    # Split into lines
    lines = [line.strip() for line in cell_text.split('\n') if line.strip()]

    if not lines:
        return ""

    # If first line is "MONTHLY TOTAL" or "TOTAL - ...", use second line as campaign name
    first_line_normalized = lines[0].upper()
    is_monthly_total_line = (
        ("MONTHLY" in first_line_normalized and "TOTAL" in first_line_normalized) or
        first_line_normalized.startswith("TOTAL -") or
        first_line_normalized.startswith("TOTAL-")
    )

    if is_monthly_total_line:
        if len(lines) > 1:
            # Return second line (the actual campaign name)
            return lines[1].upper()
        else:
            # Just "MONTHLY TOTAL" or "TOTAL - ..." with no campaign name - not a campaign row
            return ""

    # Otherwise, first line is the campaign name
    return lines[0].upper()


def _has_gray_background(cell) -> bool:
    """
    Check if a cell has a gray background fill.

    Args:
        cell: python-pptx cell object

    Returns:
        True if cell has gray background
    """
    try:
        fill = cell.fill

        # Check if solid fill
        if fill.type == 1:  # MSO_FILL_TYPE.SOLID
            # Get RGB color
            color = fill.fore_color
            if color.type == 1:  # MSO_COLOR_TYPE.RGB
                rgb = color.rgb
                # Gray colors have R=G=B and are not pure white (255,255,255)
                # Typical gray values in templates: around (191,191,191) to (217,217,217)
                r, g, b = rgb[0], rgb[1], rgb[2]

                # Check if it's a gray color (R≈G≈B) and not white
                is_gray = (abs(r - g) <= 10 and abs(g - b) <= 10 and abs(r - b) <= 10)
                is_not_white = r < 250 and g < 250 and b < 250

                return is_gray and is_not_white

        return False
    except Exception as e:
        logger.debug(f"Error checking cell background: {e}")
        return False


def _get_cell_text(cell) -> str:
    """
    Extract text content from a table cell.

    Args:
        cell: python-pptx cell object

    Returns:
        Text content from the cell, or empty string if no text
    """
    try:
        if cell.text_frame and cell.text_frame.text:
            return cell.text_frame.text
        return ""
    except Exception:
        return ""


def _cells_are_same(cell1, cell2) -> bool:
    """
    Check if two cells are already merged (reference the same underlying cell).

    Args:
        cell1: First python-pptx cell object
        cell2: Second python-pptx cell object

    Returns:
        True if cells are already merged
    """
    try:
        # In python-pptx, merged cells share the same text_frame object
        return cell1.text_frame is cell2.text_frame
    except Exception:
        return False


def _apply_cell_styling(cell, text: str = None, font_size: int = None,
                       bold: bool = False, center_align: bool = False,
                       vertical_center: bool = False):
    """
    Apply formatting to a table cell.

    Args:
        cell: python-pptx cell object
        text: Text to set in cell (optional)
        font_size: Font size in points (optional)
        bold: Whether to make text bold
        center_align: Whether to center-align text horizontally
        vertical_center: Whether to center-align text vertically
    """
    try:
        text_frame = cell.text_frame

        # Enable word wrap
        text_frame.word_wrap = True

        # Enable text auto-fit: shrink text if needed to prevent overflow
        # This allows PowerPoint to automatically reduce font size to fit text properly
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Set margins to allow text to wrap properly within cell bounds
        # Minimal margins to maximize available space for text wrapping
        text_frame.margin_left = Pt(1)   # Minimal left margin
        text_frame.margin_right = Pt(1)  # Minimal right margin
        text_frame.margin_top = Pt(1)    # Minimal top margin
        text_frame.margin_bottom = Pt(1) # Minimal bottom margin

        # Set text if provided - split by newlines and create separate paragraphs
        if text is not None:
            # Split by newlines and filter out empty lines
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            text_frame.clear()  # Clear existing content (leaves one empty paragraph)

            for idx, line in enumerate(lines):
                # Use the existing first paragraph after clear(), don't create a new one
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = line

                # Apply paragraph formatting
                if center_align:
                    p.alignment = PP_ALIGN.CENTER

                # Apply run formatting
                if p.runs:
                    for run in p.runs:
                        run.font.name = FONT_NAME
                        if font_size is not None:
                            run.font.size = Pt(font_size)
                        if bold:
                            run.font.bold = True
        else:
            # If no text provided, apply formatting to existing text in cell
            # This is used for merged cells that already have content
            for paragraph in text_frame.paragraphs:
                if center_align:
                    paragraph.alignment = PP_ALIGN.CENTER

                for run in paragraph.runs:
                    run.font.name = FONT_NAME
                    if font_size is not None:
                        run.font.size = Pt(font_size)
                    if bold:
                        run.font.bold = True

        # Apply vertical alignment
        if vertical_center:
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    except Exception as e:
        logger.error(f"Failed to apply cell styling: {e}")
