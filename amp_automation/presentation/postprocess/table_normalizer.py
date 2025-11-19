"""
Table normalization and cell formatting operations.

This module provides Python-based implementations of table layout normalization
and blank cell formatting, replacing slow COM-based PowerShell operations.
"""

import logging
from typing import Optional

from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt

logger = logging.getLogger(__name__)

# Constants from PowerShell script
ZERO_WIDTH_SPACE = "\u200B"
BLANK_FONT_NAME = "Calibri"
BLANK_FONT_SIZE = Pt(8)


def set_cell_fixed_layout(cell):
    """
    Set fixed layout properties for a table cell.

    Equivalent to PowerShell function: Set-CellFixedLayout

    Sets:
    - AutoSize = none (no auto-sizing)
    - WordWrap = enabled
    - All margins = 0
    - VerticalAnchor = middle

    Args:
        cell: python-pptx table cell object
    """
    try:
        text_frame = cell.text_frame

        # Set auto-size and word wrap
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.word_wrap = True

        # Set margins to 0
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        # Set vertical anchor to middle
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    except Exception as e:
        logger.debug(f"Error setting fixed layout for cell: {e}")


def normalize_cell_content(text: str) -> str:
    """
    Normalize cell content by stripping whitespace.

    Args:
        text: Cell text content

    Returns:
        Normalized text with trailing whitespace removed
    """
    if not text:
        return ""
    return text.rstrip()


def ensure_blank_cell_formatting(cell):
    """
    Apply formatting to blank cells and cells containing only "-".

    Equivalent to PowerShell function: Ensure-BlankCellFormatting

    Args:
        cell: python-pptx table cell object
    """
    try:
        # First apply fixed layout
        set_cell_fixed_layout(cell)

        text_frame = cell.text_frame

        # Get normalized content
        content = normalize_cell_content(text_frame.text)
        is_dash = content == "-"
        is_blank = len(content) == 0

        # Only process blank cells or cells with "-"
        if not is_dash and not is_blank:
            return

        # Replace blank cells with zero-width space
        if is_blank:
            text_frame.text = ZERO_WIDTH_SPACE

        # Apply formatting
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # Format the paragraph
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

            # Format runs (text segments)
            for run in paragraph.runs:
                run.font.name = BLANK_FONT_NAME
                run.font.size = BLANK_FONT_SIZE
                run.font.bold = False

    except Exception as e:
        logger.debug(f"Error formatting blank cell: {e}")


def normalize_table_layout(table):
    """
    Normalize layout for all cells in a table.

    Equivalent to PowerShell function: Normalize-TableLayout

    This function iterates through all cells and applies fixed layout
    properties to ensure consistent formatting across the table.

    Args:
        table: python-pptx table object
    """
    logger.debug(f"Normalizing table layout: {len(table.rows)} rows × {len(table.columns)} columns")

    cell_count = 0
    error_count = 0

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            try:
                set_cell_fixed_layout(cell)
                cell_count += 1
            except Exception as e:
                error_count += 1
                logger.debug(f"Error normalizing cell ({row_idx},{col_idx}): {e}")

    logger.debug(f"Normalized {cell_count} cells ({error_count} errors)")


def apply_blank_cell_formatting(table):
    """
    Apply blank cell formatting to all cells in a table.

    Equivalent to PowerShell function: Apply-BlankCellFormatting

    This function iterates through all cells and applies special formatting
    to blank cells and cells containing only "-".

    Args:
        table: python-pptx table object
    """
    logger.debug(f"Applying blank cell formatting: {len(table.rows)} rows × {len(table.columns)} columns")

    cell_count = 0
    blank_count = 0
    error_count = 0

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            try:
                content = normalize_cell_content(cell.text_frame.text)
                is_blank = len(content) == 0 or content == "-"

                if is_blank:
                    blank_count += 1

                ensure_blank_cell_formatting(cell)
                cell_count += 1

            except Exception as e:
                error_count += 1
                logger.debug(f"Error formatting cell ({row_idx},{col_idx}): {e}")

    logger.debug(f"Formatted {cell_count} cells ({blank_count} blank, {error_count} errors)")


def normalize_table_fonts(table):
    """
    Normalize fonts for all cells in a table.

    Font rules:
    - Header row (row 0): Verdana, 7pt
    - BRAND TOTAL row (bottom row containing "BRAND" and "TOTAL"): Verdana, 7pt
    - Campaign column (column 0, non-header): Verdana, 6pt
    - Body rows (all others): Verdana, 6pt
    - Bottom row (other bottom rows): Verdana, 6pt

    Args:
        table: python-pptx table object

    Returns:
        dict: Statistics about font normalization
    """
    logger.debug(f"Normalizing table fonts: {len(table.rows)} rows × {len(table.columns)} columns")

    row_count = len(table.rows)
    header_cells = 0
    body_cells = 0
    bottom_cells = 0
    error_count = 0

    for row_idx in range(row_count):
        for col_idx in range(len(table.columns)):
            try:
                cell = table.cell(row_idx, col_idx)
                text_frame = cell.text_frame

                # Check if this is BRAND TOTAL row (needs 7pt)
                is_brand_total = False
                if row_idx == row_count - 1:
                    first_cell = table.cell(row_idx, 0)
                    first_cell_text = first_cell.text_frame.text.strip().upper() if first_cell.text_frame else ""
                    is_brand_total = "BRAND" in first_cell_text and "TOTAL" in first_cell_text

                # Determine font size based on row position and column
                if row_idx == 0:
                    # Header row: Verdana 7pt
                    font_size = Pt(7)
                    header_cells += 1
                elif row_idx == row_count - 1 and is_brand_total:
                    # BRAND TOTAL row: Verdana 7pt
                    font_size = Pt(7)
                    bottom_cells += 1
                else:
                    # All other rows (campaign column, body, other bottom rows): Verdana 6pt
                    font_size = Pt(6)
                    body_cells += 1

                # Check if cell is empty or whitespace only
                cell_text = text_frame.text.strip() if text_frame.text else ""

                if not cell_text or cell_text == ZERO_WIDTH_SPACE:
                    # Empty cell - set to dash with proper font
                    text_frame.clear()
                    paragraph = text_frame.paragraphs[0]
                    run = paragraph.add_run()
                    run.text = "-"
                    run.font.name = "Verdana"
                    run.font.size = font_size
                else:
                    # Cell has content - format existing runs
                    for paragraph in text_frame.paragraphs:
                        if paragraph.runs:
                            for run in paragraph.runs:
                                run.font.name = "Verdana"
                                run.font.size = font_size
                        else:
                            # Has text but no runs - create one
                            run = paragraph.add_run()
                            run.font.name = "Verdana"
                            run.font.size = font_size

            except Exception as e:
                error_count += 1
                logger.debug(f"Error normalizing font for cell ({row_idx},{col_idx}): {e}")

    total_cells = header_cells + body_cells + bottom_cells
    logger.info(f"Normalized fonts: {total_cells} cells (header: {header_cells}, body: {body_cells}, bottom: {bottom_cells}, errors: {error_count})")

    return {
        "total": total_cells,
        "header": header_cells,
        "body": body_cells,
        "bottom": bottom_cells,
        "errors": error_count
    }


def delete_carried_forward_rows(table):
    """
    Delete all CARRIED FORWARD rows from a table.

    Removes rows where column 1 contains "CARRIED FORWARD" text.

    Args:
        table: python-pptx table object

    Returns:
        int: Number of rows deleted
    """
    logger.debug(f"Deleting CARRIED FORWARD rows from table: {len(table.rows)} rows")

    rows_to_delete = []

    # First pass: identify rows to delete
    for row_idx in range(len(table.rows)):
        try:
            cell = table.cell(row_idx, 0)
            cell_text = cell.text_frame.text.strip().upper() if cell.text_frame else ""

            if "CARRIED" in cell_text and "FORWARD" in cell_text:
                rows_to_delete.append(row_idx)
                logger.debug(f"Marked row {row_idx} for deletion: CARRIED FORWARD")

        except Exception as e:
            logger.debug(f"Error checking row {row_idx}: {e}")

    # Second pass: delete rows in reverse order to preserve indices
    deleted_count = 0
    for row_idx in reversed(rows_to_delete):
        try:
            # Access the table's underlying XML to delete the row
            tbl = table._tbl
            tr = table.rows[row_idx]._tr
            tbl.remove(tr)
            deleted_count += 1
            logger.debug(f"Deleted row {row_idx}")
        except Exception as e:
            logger.error(f"Failed to delete row {row_idx}: {e}")

    logger.info(f"Deleted {deleted_count} CARRIED FORWARD row(s)")
    return deleted_count


def fix_grand_total_wrapping(table):
    """
    Fix word wrapping in GRAND TOTAL rows to prevent multi-line values.

    Uses combination approach:
    - Font: Verdana 6pt (same as body)
    - Margins: 0 (maximize horizontal space)
    - Word wrap: Disabled
    - Auto-size: Shrink to fit if needed

    Args:
        table: python-pptx table object

    Returns:
        int: Number of rows fixed
    """
    logger.debug(f"Fixing GRAND TOTAL wrapping in table: {len(table.rows)} rows")

    fixed_count = 0

    # Find GRAND TOTAL row
    for row_idx in range(len(table.rows)):
        try:
            cell = table.cell(row_idx, 0)
            cell_text = cell.text_frame.text.strip().upper() if cell.text_frame else ""

            if "GRAND" in cell_text and "TOTAL" in cell_text:
                # Fix wrapping for all cells in this row
                for col_idx in range(len(table.columns)):
                    cell = table.cell(row_idx, col_idx)
                    text_frame = cell.text_frame

                    # Disable word wrap to prevent multi-line values
                    text_frame.word_wrap = False

                    # Set margins to 0 for maximum horizontal space
                    text_frame.margin_left = 0
                    text_frame.margin_right = 0
                    text_frame.margin_top = 0
                    text_frame.margin_bottom = 0

                    # Set auto-size to shrink text if needed
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                    # Set font to Verdana 6pt (same as body)
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = "Verdana"
                            run.font.size = Pt(6)

                fixed_count += 1
                logger.debug(f"Fixed wrapping for GRAND TOTAL row {row_idx}")

        except Exception as e:
            logger.debug(f"Error fixing row {row_idx}: {e}")

    logger.info(f"Fixed wrapping for {fixed_count} GRAND TOTAL row(s)")
    return fixed_count


def remove_pound_signs_from_totals(table):
    """
    Remove pound (£) signs from GRAND TOTAL and MONTHLY TOTAL rows.
    Also applies bold and center alignment to these rows.

    Args:
        table: python-pptx table object

    Returns:
        int: Number of cells cleaned
    """
    logger.debug(f"Removing pound signs from total rows: {len(table.rows)} rows")

    cells_cleaned = 0

    for row_idx in range(len(table.rows)):
        try:
            cell = table.cell(row_idx, 0)
            cell_text = cell.text_frame.text.strip().upper() if cell.text_frame else ""

            # Check if this is GRAND TOTAL or MONTHLY TOTAL row
            is_grand_total = "GRAND" in cell_text and "TOTAL" in cell_text
            is_monthly_total = (
                ("MONTHLY" in cell_text and "TOTAL" in cell_text) or
                cell_text.startswith("TOTAL -") or
                cell_text.startswith("TOTAL-")
            )

            if is_grand_total or is_monthly_total:
                # IMPORTANT FIX (Point 6): Keep pound symbol in label cell (col 0), only remove from numeric cells
                for col_idx in range(len(table.columns)):
                    cell = table.cell(row_idx, col_idx)
                    text_frame = cell.text_frame

                    # Remove pound sign ONLY from numeric cells (col >= 3), NOT from label (col 0)
                    # This preserves "MONTHLY TOTAL (£ 000)" while cleaning numeric cells
                    if col_idx > 0 and text_frame.text and "£" in text_frame.text:
                        text_frame.text = text_frame.text.replace("£", "")
                        cells_cleaned += 1

                    # Apply vertical center alignment
                    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

                    # Apply horizontal center and bold to all paragraphs
                    for paragraph in text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER

                        for run in paragraph.runs:
                            run.font.bold = True

        except Exception as e:
            logger.debug(f"Error removing pound signs from row {row_idx}: {e}")

    logger.info(f"Removed pound signs from {cells_cleaned} cells")
    return cells_cleaned
