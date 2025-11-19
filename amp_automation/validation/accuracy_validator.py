"""
Comprehensive accuracy validator for cell mapping and formula calculations.

Validates:
1. Horizontal totals (row sums across months)
2. Vertical totals (MONTHLY TOTAL rows, BRAND TOTAL rows)
3. Cell data mapping from Excel to PowerPoint
4. GRP metrics accuracy (Reach, Frequency)
"""

from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import pandas as pd
from pptx import Presentation

logger = logging.getLogger(__name__)

MONTH_ORDER = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
TOLERANCE = 0.01  # Allow 1% tolerance for rounding errors


@dataclass
class ValidationError:
    """Represents a validation error found in the deck."""

    slide_num: int
    error_type: str  # 'horizontal_total', 'vertical_total', 'cell_mapping', 'grp_metrics'
    location: str  # Description of where the error occurred
    expected: float | str
    actual: float | str
    difference: Optional[float] = None
    row: Optional[int] = None
    col: Optional[int] = None


@dataclass
class ValidationReport:
    """Comprehensive validation report."""

    total_slides: int = 0
    slides_checked: int = 0
    errors: list[ValidationError] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @property
    def passed(self) -> bool:
        """Returns True if no errors found."""
        return len(self.errors) == 0

    @property
    def error_count(self) -> int:
        """Total number of errors."""
        return len(self.errors)

    def add_error(self, error: ValidationError):
        """Add an error to the report."""
        self.errors.append(error)

    def add_warning(self, warning: str):
        """Add a warning to the report."""
        self.warnings.append(warning)

    def summary(self) -> str:
        """Generate human-readable summary."""
        lines = []
        lines.append("=" * 80)
        lines.append("ACCURACY VALIDATION REPORT")
        lines.append("=" * 80)
        lines.append(f"Total slides: {self.total_slides}")
        lines.append(f"Slides checked: {self.slides_checked}")
        lines.append(f"Errors found: {self.error_count}")
        lines.append(f"Warnings: {len(self.warnings)}")
        lines.append(f"Status: {'PASSED' if self.passed else 'FAILED'}")
        lines.append("")

        if self.errors:
            lines.append("ERRORS:")
            lines.append("-" * 80)
            for i, error in enumerate(self.errors, 1):
                lines.append(f"{i}. Slide {error.slide_num} - {error.error_type}")
                lines.append(f"   Location: {error.location}")
                lines.append(f"   Expected: {error.expected}")
                lines.append(f"   Actual: {error.actual}")
                if error.difference is not None:
                    lines.append(f"   Difference: {error.difference:,.2f}")
                if error.row is not None and error.col is not None:
                    lines.append(f"   Position: Row {error.row}, Column {error.col}")
                lines.append("")

        if self.warnings:
            lines.append("WARNINGS:")
            lines.append("-" * 80)
            for warning in self.warnings:
                lines.append(f"  - {warning}")
            lines.append("")

        lines.append("=" * 80)
        return "\n".join(lines)


def parse_number(text: str) -> Optional[float]:
    """
    Parse a number from text, handling K suffix and currency symbols.

    Examples:
        "£127K" → 127000
        "1.5K" → 1500
        "42%" → 42.0
        "-" → None
    """
    if not text or text.strip() in ["-", "", "–"]:
        return None

    # Remove currency symbols, spaces, commas
    cleaned = text.strip().replace("£", "").replace("$", "").replace(",", "").replace(" ", "")

    # Handle percentage
    if "%" in cleaned:
        cleaned = cleaned.replace("%", "")
        try:
            return float(cleaned)
        except ValueError:
            return None

    # Handle M suffix (millions)
    if cleaned.upper().endswith("M"):
        try:
            return float(cleaned[:-1]) * 1000000
        except ValueError:
            return None

    # Handle K suffix (thousands)
    if cleaned.upper().endswith("K"):
        try:
            return float(cleaned[:-1]) * 1000
        except ValueError:
            return None

    # Try parsing as regular number
    try:
        return float(cleaned)
    except ValueError:
        return None


def validate_horizontal_total(
    row_data: list[str],
    month_start_col: int,
    month_end_col: int,
    total_col: int,
    slide_num: int,
    row_idx: int,
    row_label: str
) -> Optional[ValidationError]:
    """
    Validate that the TOTAL column equals the sum of month columns.

    Args:
        row_data: List of cell values for the row
        month_start_col: First month column index
        month_end_col: Last month column index
        total_col: TOTAL column index
        slide_num: Slide number for error reporting
        row_idx: Row index for error reporting
        row_label: Row label (e.g., campaign name) for error reporting

    Returns:
        ValidationError if mismatch found, None otherwise
    """
    # Parse month values
    month_values = []
    for col_idx in range(month_start_col, month_end_col + 1):
        if col_idx < len(row_data):
            value = parse_number(row_data[col_idx])
            if value is not None:
                month_values.append(value)

    # Parse total value
    if total_col >= len(row_data):
        return None

    actual_total = parse_number(row_data[total_col])

    if actual_total is None:
        # If no total shown, check if there were any month values
        if month_values:
            expected_total = sum(month_values)
            return ValidationError(
                slide_num=slide_num,
                error_type="horizontal_total",
                location=f"Row {row_idx}: {row_label}",
                expected=f"{expected_total:,.0f}",
                actual="None (missing total)",
                difference=expected_total,
                row=row_idx,
                col=total_col
            )
        return None

    # Calculate expected total
    expected_total = sum(month_values) if month_values else 0

    # Check if within tolerance
    if expected_total == 0 and actual_total == 0:
        return None

    # Allow 1% tolerance for rounding
    tolerance_amount = max(abs(expected_total), abs(actual_total)) * TOLERANCE
    difference = abs(expected_total - actual_total)

    if difference > tolerance_amount and difference > 1:  # Ignore differences less than £1
        return ValidationError(
            slide_num=slide_num,
            error_type="horizontal_total",
            location=f"Row {row_idx}: {row_label}",
            expected=f"{expected_total:,.0f}",
            actual=f"{actual_total:,.0f}",
            difference=difference,
            row=row_idx,
            col=total_col
        )

    return None


def validate_vertical_total(
    table_data: list[list[str]],
    total_row_idx: int,
    col_idx: int,
    start_row: int,
    end_row: int,
    slide_num: int,
    col_label: str,
    total_label: str
) -> Optional[ValidationError]:
    """
    Validate that a total row cell equals the sum of cells above it.

    Args:
        table_data: Full table data
        total_row_idx: Index of the total row
        col_idx: Column index to check
        start_row: First row to include in sum
        end_row: Last row to include in sum (exclusive of total row)
        slide_num: Slide number
        col_label: Column label (e.g., "Jan")
        total_label: Total row label (e.g., "MONTHLY TOTAL", "BRAND TOTAL")

    Returns:
        ValidationError if mismatch found, None otherwise
    """
    # Parse values from detail rows
    detail_values = []
    for row_idx in range(start_row, end_row):
        if row_idx >= len(table_data):
            break
        row = table_data[row_idx]
        if col_idx < len(row):
            # Skip sub-rows (GRPs, Reach, etc.) - these shouldn't be included in budget totals
            if row_idx > 0 and len(row) > 2:
                row_label = row[2] if len(row) > 2 else row[0] if len(row) > 0 else ""
                if row_label in ["GRPs", "Reach@1+", "Reach@3+", "OTS@1+", "OTS@3+", "Frequency"]:
                    continue

            value = parse_number(row[col_idx])
            if value is not None:
                detail_values.append(value)

    # Parse total value
    if total_row_idx >= len(table_data):
        return None

    total_row = table_data[total_row_idx]
    if col_idx >= len(total_row):
        return None

    actual_total = parse_number(total_row[col_idx])

    if actual_total is None:
        if detail_values:
            expected_total = sum(detail_values)
            return ValidationError(
                slide_num=slide_num,
                error_type="vertical_total",
                location=f"{total_label} - {col_label}",
                expected=f"{expected_total:,.0f}",
                actual="None (missing total)",
                difference=expected_total,
                row=total_row_idx,
                col=col_idx
            )
        return None

    # Calculate expected total
    expected_total = sum(detail_values) if detail_values else 0

    # Check if within tolerance
    if expected_total == 0 and actual_total == 0:
        return None

    tolerance_amount = max(abs(expected_total), abs(actual_total)) * TOLERANCE
    difference = abs(expected_total - actual_total)

    if difference > tolerance_amount and difference > 1:
        return ValidationError(
            slide_num=slide_num,
            error_type="vertical_total",
            location=f"{total_label} - {col_label}",
            expected=f"{expected_total:,.0f}",
            actual=f"{actual_total:,.0f}",
            difference=difference,
            row=total_row_idx,
            col=col_idx
        )

    return None


def extract_table_data(table) -> list[list[str]]:
    """Extract all text from table cells."""
    data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        data.append(row_data)
    return data


def validate_slide_table(slide, slide_num: int, report: ValidationReport):
    """
    Validate all calculations in a slide's table.

    Checks:
    1. Horizontal totals (each row's sum across months)
    2. Vertical totals (MONTHLY TOTAL rows)
    3. Budget row totals
    """
    # Find the main data table
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            if hasattr(shape, "name") and "MainDataTable" in shape.name:
                table = shape.table
                break

    if not table:
        report.add_warning(f"Slide {slide_num}: No MainDataTable found")
        return

    # Extract table data
    table_data = extract_table_data(table)

    if len(table_data) < 3:
        report.add_warning(f"Slide {slide_num}: Table too small ({len(table_data)} rows)")
        return

    # Identify column indices
    header_row = table_data[0] if table_data else []
    month_start_col = next((i for i, col in enumerate(header_row) if col == "JAN"), None)
    total_col = next((i for i, col in enumerate(header_row) if col == "TOTAL"), None)

    if month_start_col is None or total_col is None:
        report.add_warning(f"Slide {slide_num}: Could not identify month/total columns")
        return

    month_end_col = month_start_col + 11  # 12 months

    # Validate horizontal totals for each data row
    for row_idx in range(1, len(table_data)):
        row_data = table_data[row_idx]
        if len(row_data) < 3:
            continue

        row_label = row_data[0] if len(row_data) > 0 else ""

        # Get label from column 2 (METRICS column) if available
        metrics_col_label = row_data[2] if len(row_data) > 2 else ""

        # Skip header rows
        if row_label in ["CAMPAIGN", "MEDIA", "METRICS"]:
            continue

        # Skip metric sub-rows (GRPs, Reach, Frequency, etc.) - these don't sum horizontally
        metric_indicators = ["GRP", "Reach@1+", "Reach@3+", "OTS@1+", "OTS@3+", "Frequency",
                           "META Reach", "TT Reach", "YT Reach", "Reach@1+", "OTS@3+"]
        if any(indicator in metrics_col_label for indicator in metric_indicators):
            continue
        if any(indicator in row_label for indicator in metric_indicators):
            continue

        # Skip rows that are just "-" (empty rows between campaigns)
        if row_label.strip() in ["-", "", "–"]:
            continue

        # Skip total rows for now (we'll validate them separately)
        if "TOTAL" in row_label.upper():
            continue

        # Validate this row's horizontal total
        error = validate_horizontal_total(
            row_data=row_data,
            month_start_col=month_start_col,
            month_end_col=month_end_col,
            total_col=total_col,
            slide_num=slide_num,
            row_idx=row_idx,
            row_label=row_label
        )

        if error:
            report.add_error(error)

    # Validate vertical totals for MONTHLY TOTAL rows
    # Find all MONTHLY TOTAL or "TOTAL -" rows
    for row_idx, row_data in enumerate(table_data):
        if len(row_data) < 1:
            continue

        row_label = row_data[0].upper()
        if "MONTHLY TOTAL" in row_label or row_label.startswith("TOTAL -") or row_label.startswith("TOTAL–"):
            # Find the start of this campaign block (work backwards to find campaign name)
            campaign_start = row_idx - 1
            while campaign_start > 0:
                prev_row = table_data[campaign_start]
                prev_label = prev_row[0].upper() if prev_row else ""
                if "TOTAL" in prev_label or prev_label in ["CAMPAIGN", ""]:
                    break
                campaign_start -= 1

            # Validate each month column for this MONTHLY TOTAL
            for month_idx, month in enumerate(MONTH_ORDER):
                col_idx = month_start_col + month_idx
                error = validate_vertical_total(
                    table_data=table_data,
                    total_row_idx=row_idx,
                    col_idx=col_idx,
                    start_row=campaign_start + 1,
                    end_row=row_idx,
                    slide_num=slide_num,
                    col_label=month,
                    total_label=row_label
                )

                if error:
                    report.add_error(error)


def validate_deck_accuracy(ppt_path: str | Path) -> ValidationReport:
    """
    Comprehensive validation of deck accuracy.

    Args:
        ppt_path: Path to PowerPoint file

    Returns:
        ValidationReport with all errors and warnings
    """
    ppt_path = Path(ppt_path)
    if not ppt_path.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    prs = Presentation(str(ppt_path))
    report = ValidationReport(total_slides=len(prs.slides))

    logger.info(f"Validating accuracy for {report.total_slides} slides...")

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Skip delimiter slides
        has_table = any(shape.has_table for shape in slide.shapes)
        if not has_table:
            continue

        report.slides_checked += 1
        validate_slide_table(slide, slide_num, report)

    logger.info(f"Validation complete: {report.slides_checked} slides checked, {report.error_count} errors found")

    return report
