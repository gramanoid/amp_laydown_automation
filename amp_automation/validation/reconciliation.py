"""Automated reconciliation between generated PPT content and Lumina Excel data."""

from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence

import pandas as pd
from pptx import Presentation

from amp_automation.config.loader import Config
from amp_automation.data import load_and_prepare_data

LOGGER = logging.getLogger("amp_automation.validation.reconciliation")

MONTH_ORDER: Sequence[str] = (
    "Jan",
    "Feb",
    "Mar",
    "Apr",
    "May",
    "Jun",
    "Jul",
    "Aug",
    "Sep",
    "Oct",
    "Nov",
    "Dec",
)

QUARTER_MONTHS: Dict[str, Sequence[str]] = {
    "q1": ("Jan", "Feb", "Mar"),
    "q2": ("Apr", "May", "Jun"),
    "q3": ("Jul", "Aug", "Sep"),
    "q4": ("Oct", "Nov", "Dec"),
}

PERCENT_TOLERANCE = 0.005  # ±0.5% tolerance for share metrics
NUMBER_PATTERN = re.compile(r"-?\d+(?:[.,]\d+)?")

# Market code to display name mapping
MARKET_CODE_MAP = {
    "MOR": "MOROCCO",
    "SOUTH AFRICA": "SOUTH AFRICA",  # No mapping needed
    "KSA": "KSA",
    "GINE": "GINE",
    "EGYPT": "EGYPT",
    "TURKEY": "TURKEY",
    "PAKISTAN": "PAKISTAN",
    "KENYA": "KENYA",
    "UGANDA": "UGANDA",
    "NIGERIA": "NIGERIA",
    "MAURITIUS": "MAURITIUS",
    "FWA": "FWA",
}


@dataclass(slots=True)
class MetricComparison:
    """Comparison outcome for a single summary metric."""

    category: str
    label: str
    expected_display: str
    actual_display: str
    expected_value: Optional[float]
    actual_value: Optional[float]
    tolerance: Optional[float]
    difference: Optional[float]
    passed: bool
    notes: str = ""


@dataclass(slots=True)
class SlideReconciliation:
    """Aggregated reconciliation results for a single slide."""

    slide_index: int
    market: str
    brand: str
    year: Optional[int]
    comparisons: List[MetricComparison] = field(default_factory=list)

    @property
    def passed(self) -> bool:
        return all(comparison.passed for comparison in self.comparisons)


def generate_reconciliation_report(
    ppt_path: str | Path,
    excel_path: str | Path,
    config: Config,
    *,
    logger: Optional[logging.Logger] = None,
    data_frame: Optional[pd.DataFrame] = None,
) -> List[SlideReconciliation]:
    """Compare summary tiles in the PPT against Excel-derived expectations."""

    logger = logger or LOGGER
    ppt_path = Path(ppt_path)
    excel_path = Path(excel_path)

    if not ppt_path.is_file():
        raise FileNotFoundError(f"Presentation not found: {ppt_path}")
    if not excel_path.is_file() and data_frame is None:
        raise FileNotFoundError(f"Excel workbook not found: {excel_path}")

    presentation_cfg = config.section("presentation")
    summary_cfg = presentation_cfg.get("summary_tiles", {})
    if not summary_cfg:
        logger.warning("No summary tile configuration present; skipping reconciliation")
        return []

    if data_frame is not None:
        df = data_frame.copy()
    else:
        dataset = load_and_prepare_data(excel_path, config, logger)
        df = dataset.frame.copy()

    if df.empty:
        logger.warning("Prepared dataset is empty; no reconciliation performed")
        return []

    prs = Presentation(ppt_path)
    results: List[SlideReconciliation] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        title_text = _extract_shape_text(slide, presentation_cfg.get("title", {}).get("shape", "TitlePlaceholder"))
        if not title_text or " - " not in title_text:
            continue  # Likely a delimiter or non-data slide

        market, brand = _parse_title_tokens(title_text)
        if market is None or brand is None:
            logger.debug("Unable to parse market/brand from title '%s'", title_text)
            continue

        actual_summary = _extract_slide_summary(slide, summary_cfg)
        if not _has_summary_data(actual_summary):
            continue

        candidate_years = _candidate_years(df, market, brand)
        if not candidate_years:
            logger.warning("No dataset rows found for %s - %s", market, brand)
            results.append(
                SlideReconciliation(
                    slide_index=slide_idx,
                    market=market,
                    brand=brand,
                    year=None,
                    comparisons=_build_missing_comparisons(actual_summary),
                )
            )
            continue

        best = _select_best_year(candidate_years, df, market, brand, summary_cfg, actual_summary, logger)
        if best is None:
            logger.warning("Unable to reconcile slide %s for %s - %s", slide_idx, market, brand)
            results.append(
                SlideReconciliation(
                    slide_index=slide_idx,
                    market=market,
                    brand=brand,
                    year=None,
                    comparisons=_build_missing_comparisons(actual_summary),
                )
            )
            continue

        results.append(
            SlideReconciliation(
                slide_index=slide_idx,
                market=market,
                brand=brand,
                year=best["year"],
                comparisons=best["comparisons"],
            )
        )

    return results


def reconciliations_to_dataframe(results: Iterable[SlideReconciliation]) -> pd.DataFrame:
    """Flatten reconciliation results into a DataFrame for reporting."""

    rows: List[dict] = []
    for result in results:
        for comparison in result.comparisons:
            rows.append(
                {
                    "slide_index": result.slide_index,
                    "market": result.market,
                    "brand": result.brand,
                    "year": result.year,
                    "category": comparison.category,
                    "label": comparison.label,
                    "expected_display": comparison.expected_display,
                    "actual_display": comparison.actual_display,
                    "expected_value": comparison.expected_value,
                    "actual_value": comparison.actual_value,
                    "difference": comparison.difference,
                    "tolerance": comparison.tolerance,
                    "passed": comparison.passed,
                    "notes": comparison.notes,
                }
            )

    return pd.DataFrame(rows)


def write_reconciliation_report(results: Iterable[SlideReconciliation], output_path: str | Path) -> Path:
    """Persist reconciliation results to CSV (or empty file if no rows)."""

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    frame = reconciliations_to_dataframe(results)
    if frame.empty:
        output_path.write_text("slide_index,market,brand,year,category,label,expected_display,actual_display,expected_value,actual_value,difference,tolerance,passed,notes\n", encoding="utf-8")
    else:
        frame.to_csv(output_path, index=False)

    return output_path


# --------------------------------------------------------------------------------------
# Internal helpers
# --------------------------------------------------------------------------------------


def _normalize_market_name(df: pd.DataFrame, market: str) -> str:
    """Find the exact market name from DataFrame by case-insensitive matching and code mapping."""
    market_str = str(market).strip()
    market_lower = market_str.lower()

    # First try the market code map (e.g., "MOROCCO" -> "MOR")
    for code, display_name in MARKET_CODE_MAP.items():
        if display_name.lower() == market_lower:
            market_str = code
            break

    # Then do case-insensitive matching against the DataFrame
    market_lower = market_str.lower()
    for country in df["Country"].unique():
        if str(country).lower().strip() == market_lower:
            return str(country)
    return market_str  # Return original if no match found


def _normalize_brand_name(df: pd.DataFrame, market: str, brand: str) -> str:
    """Find the exact brand name from DataFrame by case-insensitive matching for a given market."""
    brand_lower = str(brand).lower().strip()
    market_norm = _normalize_market_name(df, market)

    market_rows = df[df["Country"].astype(str).str.strip() == str(market_norm).strip()]
    for brand_val in market_rows["Brand"].unique():
        if str(brand_val).lower().strip() == brand_lower:
            return str(brand_val)
    return brand  # Return original if no match found


def _candidate_years(df: pd.DataFrame, market: str, brand: str) -> List[int]:
    # Normalize market and brand for case-insensitive matching against DataFrame
    market_norm = _normalize_market_name(df, market)
    brand_norm = _normalize_brand_name(df, market_norm, brand)

    years = (
        df[
            (df["Country"].astype(str).str.strip() == str(market_norm).strip())
            & (df["Brand"].astype(str).str.strip() == str(brand_norm).strip())
        ]["Year"]
        .dropna()
        .unique()
        .tolist()
    )
    return sorted(int(year) for year in years)


def _select_best_year(candidate_years, df, market, brand, summary_cfg, actual_summary, logger):
    best_record = None
    best_score = (-1, float("inf"))  # (passes, total_abs_diff)

    for year in candidate_years:
        expected = _compute_expected_summary(df, market, brand, year, summary_cfg)
        if expected is None:
            logger.debug("No expected summary for %s - %s (%s)", market, brand, year)
            continue

        comparisons = _compare_summary(actual_summary, expected, summary_cfg)
        passes = sum(1 for item in comparisons if item.passed)
        diff_sum = sum(abs(item.difference) for item in comparisons if item.difference is not None and not item.passed)
        score = (passes, diff_sum)

        if score > best_score:
            best_score = score
            best_record = {"year": year, "comparisons": comparisons}

    return best_record


def _extract_slide_summary(slide, summary_cfg: dict) -> dict:
    return {
        "quarter_budgets": {
            key: _extract_shape_text(slide, config.get("shape"))
            for key, config in (summary_cfg.get("quarter_budgets", {}) or {}).items()
            if not key.startswith("_") and isinstance(config, dict)
        },
        "media_share": {
            key: _extract_shape_text(slide, config.get("shape"))
            for key, config in (summary_cfg.get("media_share", {}) or {}).items()
            if not key.startswith("_") and isinstance(config, dict)
        },
        "funnel_share": {
            key: _extract_shape_text(slide, config.get("shape"))
            for key, config in (summary_cfg.get("funnel_share", {}) or {}).items()
            if not key.startswith("_") and isinstance(config, dict)
        },
    }


def _has_summary_data(summary: dict) -> bool:
    for category in summary.values():
        for value in category.values():
            if value:
                return True
    return False


def _compute_expected_summary(df: pd.DataFrame, market: str, brand: str, year: int, summary_cfg: dict) -> Optional[dict]:
    # Normalize market and brand for case-insensitive matching against DataFrame
    market_norm = _normalize_market_name(df, market)
    brand_norm = _normalize_brand_name(df, market_norm, brand)

    subset = df[
        (df["Country"].astype(str).str.strip() == str(market_norm).strip())
        & (df["Brand"].astype(str).str.strip() == str(brand_norm).strip())
        & (df["Year"].astype(str) == str(year))
    ]

    if subset.empty:
        return None

    total_cost = float(subset["Total Cost"].sum())

    quarter_expectations = {}
    for key, config in (summary_cfg.get("quarter_budgets", {}) or {}).items():
        if key.startswith("_") or not isinstance(config, dict):
            continue
        months = QUARTER_MONTHS.get(key.lower())
        if not months:
            continue
        value = float(subset[list(months)].sum().sum())
        quarter_expectations[key] = {
            "value": value,
            "display": _format_tile_value(config, value, prefix=config.get("prefix", "")),
            "config": config,
        }

    media_expectations = {}
    media_group = subset.groupby("Mapped Media Type")["Total Cost"].sum()
    for key, config in (summary_cfg.get("media_share", {}) or {}).items():
        if key.startswith("_") or not isinstance(config, dict):
            continue
        lookup = _media_lookup_key(key)
        if lookup == "Other":
            # Aggregate OOH + Other (OOH maps to "OOH", not "Other" in config)
            value = float(media_group.get("Other", 0.0)) + float(media_group.get("OOH", 0.0))
        else:
            value = float(media_group.get(lookup, 0.0))
        proportion = 0.0 if total_cost <= 0 else value / total_cost
        media_expectations[key] = {
            "value": proportion,
            "display": _format_percentage_tile(config, value, total_cost, label=config.get("label", key.capitalize())),
            "config": config,
        }

    funnel_expectations = {}
    funnel_group = subset.groupby("Funnel Stage")["Total Cost"].sum()
    for key, config in (summary_cfg.get("funnel_share", {}) or {}).items():
        if key.startswith("_") or not isinstance(config, dict):
            continue
        lookup = _funnel_lookup_key(key)
        value = float(funnel_group.get(lookup, 0.0))
        proportion = 0.0 if total_cost <= 0 else value / total_cost
        funnel_expectations[key] = {
            "value": proportion,
            "display": _format_percentage_tile(config, value, total_cost, label=config.get("label", lookup[:3].upper())),
            "config": config,
        }

    return {
        "market": market,
        "brand": brand,
        "year": year,
        "total_cost": total_cost,
        "quarter_budgets": quarter_expectations,
        "media_share": media_expectations,
        "funnel_share": funnel_expectations,
    }


def _compare_summary(actual_summary: dict, expected_summary: dict, summary_cfg: dict) -> List[MetricComparison]:
    comparisons: List[MetricComparison] = []

    # Quarter budgets (currency)
    for key, expected in expected_summary["quarter_budgets"].items():
        config = expected["config"]
        label = config.get("label", key.upper())
        expected_value = expected["value"]
        expected_display = expected["display"]
        actual_display = actual_summary["quarter_budgets"].get(key, "")
        actual_value = _parse_display_value(actual_display, config.get("scale", 1.0), is_percentage=False)
        tolerance = _budget_tolerance(expected_value)
        difference = _compute_difference(actual_value, expected_value)
        passed, notes = _evaluate_match(actual_display, expected_display, difference, tolerance)

        comparisons.append(
            MetricComparison(
                category="quarter_budgets",
                label=label,
                expected_display=expected_display,
                actual_display=actual_display,
                expected_value=expected_value,
                actual_value=actual_value,
                tolerance=tolerance,
                difference=difference,
                passed=passed,
                notes=notes,
            )
        )

    # Media share (percentage)
    for key, expected in expected_summary["media_share"].items():
        config = expected["config"]
        label = config.get("label", key.capitalize())
        expected_value = expected["value"]
        expected_display = expected["display"]
        actual_display = actual_summary["media_share"].get(key, "")
        actual_value = _parse_display_value(actual_display, config.get("scale", 100.0), is_percentage=True)
        tolerance = PERCENT_TOLERANCE
        difference = _compute_difference(actual_value, expected_value)
        passed, notes = _evaluate_match(actual_display, expected_display, difference, tolerance)

        comparisons.append(
            MetricComparison(
                category="media_share",
                label=label,
                expected_display=expected_display,
                actual_display=actual_display,
                expected_value=expected_value,
                actual_value=actual_value,
                tolerance=tolerance,
                difference=difference,
                passed=passed,
                notes=notes,
            )
        )

    # Funnel share (percentage)
    for key, expected in expected_summary["funnel_share"].items():
        config = expected["config"]
        label = config.get("label", key.capitalize())
        expected_value = expected["value"]
        expected_display = expected["display"]
        actual_display = actual_summary["funnel_share"].get(key, "")
        actual_value = _parse_display_value(actual_display, config.get("scale", 100.0), is_percentage=True)
        tolerance = PERCENT_TOLERANCE
        difference = _compute_difference(actual_value, expected_value)
        passed, notes = _evaluate_match(actual_display, expected_display, difference, tolerance)

        comparisons.append(
            MetricComparison(
                category="funnel_share",
                label=label,
                expected_display=expected_display,
                actual_display=actual_display,
                expected_value=expected_value,
                actual_value=actual_value,
                tolerance=tolerance,
                difference=difference,
                passed=passed,
                notes=notes,
            )
        )

    return comparisons


def _extract_shape_text(slide, shape_name: Optional[str]) -> str:
    if not shape_name:
        return ""
    for shape in slide.shapes:
        if getattr(shape, "name", "") == shape_name:
            if hasattr(shape, "text"):
                return shape.text.strip()
            if getattr(shape, "has_text_frame", False):
                return shape.text_frame.text.strip()
            break
    return ""


def _parse_title_tokens(title: str) -> tuple[Optional[str], Optional[str]]:
    clean = title.strip()
    # Remove pagination markers: both "(\d+ of \d+)" and "(\d+/\d+)" formats
    clean = re.sub(r"\s*\((?:\d+\s+of\s+\d+|\d+/\d+)\)$", "", clean)
    parts = clean.split(" - ", 1)
    if len(parts) != 2:
        return None, None
    return parts[0].strip(), parts[1].strip()


def _format_tile_value(config: dict, value: float, *, prefix: str = "") -> str:
    scale = float(config.get("scale", 1.0))
    fmt = config.get("number_format", "{value}")
    try:
        rendered = fmt.format(value=value * scale)
    except Exception:
        rendered = str(value * scale)
    return f"{prefix}{rendered}" if prefix else rendered


def _format_percentage_tile(config: dict, value: float, total: float, *, label: str = "") -> str:
    if total <= 0:
        proportion = 0.0
    else:
        proportion = value / total
    display = _format_tile_value(dict(config, scale=config.get("scale", 100.0)), proportion)
    return f"{label}: {display}" if label and not display.startswith(label) else display


def _parse_display_value(display: str, scale: float, *, is_percentage: bool) -> Optional[float]:
    if not display:
        return None
    match = NUMBER_PATTERN.search(display.replace(",", ""))
    if not match:
        return None
    try:
        numeric = float(match.group(0).replace(",", ""))
    except ValueError:
        return None

    if is_percentage:
        return numeric / 100.0
    if scale == 0:
        return numeric
    return numeric / float(scale)


def _budget_tolerance(expected_value: Optional[float]) -> Optional[float]:
    if expected_value is None:
        return None
    baseline = max(abs(expected_value) * 0.005, 100.0)
    return baseline


def _compute_difference(actual_value: Optional[float], expected_value: Optional[float]) -> Optional[float]:
    if actual_value is None or expected_value is None:
        return None
    return actual_value - expected_value


def _evaluate_match(actual_display: str, expected_display: str, difference: Optional[float], tolerance: Optional[float]) -> tuple[bool, str]:
    if actual_display == expected_display:
        return True, "exact match"
    if difference is None or tolerance is None:
        return False, "value unavailable"
    if abs(difference) <= tolerance:
        return True, "within tolerance"
    return False, "difference exceeds tolerance"


def _media_lookup_key(key: str) -> str:
    lowered = key.lower()
    if lowered in {"tv", "television"}:
        return "TV"  # Mapped value from config: "Television" -> "TV"
    if lowered in {"digital", "dig."}:
        return "Digital"
    return "Other"


def _funnel_lookup_key(key: str) -> str:
    mapping = {
        "awareness": "Awareness",
        "consideration": "Consideration",
        "purchase": "Purchase",
    }
    return mapping.get(key.lower(), key.title())


def _build_missing_comparisons(actual_summary: dict) -> List[MetricComparison]:
    comparisons: List[MetricComparison] = []
    for category, metrics in actual_summary.items():
        for key, actual in metrics.items():
            comparisons.append(
                MetricComparison(
                    category=category,
                    label=key,
                    expected_display="",
                    actual_display=actual,
                    expected_value=None,
                    actual_value=None,
                    tolerance=None,
                    difference=None,
                    passed=False,
                    notes="expected data missing",
                )
            )
    return comparisons
