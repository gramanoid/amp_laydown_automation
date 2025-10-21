"""Utility helpers for generating slides via AutoPPTX."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence
import logging

from pptx import Presentation

try:  # pragma: no cover - optional dependency
    from autopptx.Text.texts import replace_bodytexts, replace_subtitle, replace_title
    from autopptx.Image.images import replace_images
    from autopptx.Table.tables import replace_tables
    from autopptx.Table.style import set_table_style
    from autopptx.Type.find import find_placeholders

    _AUTOPPTX_AVAILABLE = True
except Exception:  # pragma: no cover - optional dependency
    _AUTOPPTX_AVAILABLE = False

logger = logging.getLogger("amp_automation.autopptx")


@dataclass(slots=True)
class SlidePayload:
    """Data required to populate a single slide using AutoPPTX."""

    title: str | None = None
    subtitle: str | None = None
    bodytext: Sequence[str] | None = None
    tables: Sequence[Sequence[Sequence[str]]] | None = None
    images: Sequence[str] | None = None
    notes: str | None = None


def autopptx_available() -> bool:
    """Return ``True`` when the AutoPPTX dependency can be imported."""

    return _AUTOPPTX_AVAILABLE


def generate_presentation(
    template_path: str | Path,
    slide_payloads: Iterable[SlidePayload],
    output_path: str | Path,
    *,
    base_slide_index: int = 0,
    table_font_name: str | None = None,
    table_font_size: int | None = None,
) -> Path:
    """Render a presentation using AutoPPTX placeholder replacement.

    Parameters
    ----------
    template_path:
        Path to the template containing the slide blueprint.
    slide_payloads:
        Ordered iterable of :class:`SlidePayload` items to render.
    output_path:
        Destination path for the generated presentation.
    base_slide_index:
        Template slide index to duplicate for each payload.
    table_font_name / table_font_size:
        Optional overrides applied after table replacement via ``set_table_style``.

    Returns
    -------
    Path
        The saved presentation path for convenience.
    """

    if not _AUTOPPTX_AVAILABLE:  # pragma: no cover - import guarded
        raise RuntimeError(
            "AutoPPTX is not installed. Install 'autopptx' to enable placeholder generation."
        )

    payload_list = list(slide_payloads)
    if not payload_list:
        raise ValueError("slide_payloads must contain at least one entry")

    output_path = Path(output_path)
    template_path = Path(template_path)

    prs = Presentation(str(template_path))
    if not prs.slides:
        raise ValueError(f"Template '{template_path}' does not contain any slides")

    if base_slide_index < 0 or base_slide_index >= len(prs.slides):
        raise IndexError(f"base_slide_index {base_slide_index} out of range for template")

    base_slide = prs.slides[base_slide_index]

    # Duplicate the base slide until we have enough instances
    generated_slides: list = []
    for idx, payload in enumerate(payload_list):
        if idx == 0:
            slide = base_slide
        else:
            slide = _clone_slide(prs, base_slide)
        generated_slides.append(slide)

    # Remove surplus template slides that are not part of the generated set
    _trim_surplus_slides(prs, keep_slides=generated_slides)

    for slide, payload in zip(generated_slides, payload_list):
        _populate_slide(slide, payload, table_font_name, table_font_size)

    prs.save(str(output_path))
    logger.info("AutoPPTX generated presentation saved to %s", output_path)
    return output_path


def _clone_slide(presentation: Presentation, source_slide) -> object:
    """Duplicate *source_slide* within *presentation* using its layout."""

    slide_layout = source_slide.slide_layout
    return presentation.slides.add_slide(slide_layout)


def _trim_surplus_slides(presentation: Presentation, keep_slides: Sequence[object]) -> None:
    """Remove slides that are not present in *keep_slides*."""

    keep_identity = {id(slide) for slide in keep_slides}
    slide_ids = presentation.slides._sldIdLst

    for index in range(len(slide_ids) - 1, -1, -1):
        slide = presentation.slides[index]
        if id(slide) in keep_identity:
            continue

        r_id = slide_ids[index].rId
        presentation.part.drop_rel(r_id)
        del slide_ids[index]


def _populate_slide(
    slide,
    payload: SlidePayload,
    table_font_name: str | None,
    table_font_size: int | None,
) -> None:
    """Apply the payload to a slide using AutoPPTX helper functions."""

    if payload.title:
        replace_title(slide, payload.title)
    if payload.subtitle:
        replace_subtitle(slide, payload.subtitle)

    if payload.bodytext:
        body_seq = list(payload.bodytext)
        replace_bodytexts(slide, body_seq, distribute_to_multiple_boxes=False)

    if payload.images:
        replace_images(slide, list(payload.images))

    if payload.tables:
        table_list = [list(map(list, table)) for table in payload.tables]
        replace_tables(slide, table_list)

        if table_font_name or table_font_size:
            _apply_table_style(slide, table_font_name, table_font_size)

    if payload.notes:
        slide.notes_slide.notes_text_frame.text = payload.notes


def _apply_table_style(slide, font_name: str | None, font_size: int | None) -> None:
    """Apply uniform table styling overrides for all table placeholders."""

    table_shapes = find_placeholders(slide, "table")
    for graphic_frame in table_shapes:
        kwargs: dict[str, object] = {}
        if font_name:
            kwargs["font_name"] = font_name
        if font_size:
            kwargs["font_size"] = font_size
        if kwargs:
            set_table_style(graphic_frame, **kwargs)
