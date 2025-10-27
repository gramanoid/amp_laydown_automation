from pathlib import Path
from pptx import Presentation

src = Path(r"D:\Drive\projects\work\AMP Laydowns Automation\output\presentations\run_20251017_162901\output\GeneratedDeck_Task11.pptx")
dst = src.with_name("GeneratedDeck_Task11_resaved.pptx")

print("exists", src.exists(), "size", src.stat().st_size)

try:
    prs = Presentation(str(src))
    print("slides", len(prs.slides))
    for idx, slide in enumerate(prs.slides, 1):
        print("slide", idx, "shape_count", len(slide.shapes))
    prs.save(str(dst))
    print("resaved to", dst)
except Exception as exc:
    print("error", type(exc).__name__, exc)
