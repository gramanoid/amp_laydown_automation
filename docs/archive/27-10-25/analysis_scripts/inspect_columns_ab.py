"""Inspect columns A and B to understand structure."""

from pptx import Presentation

# Check original deck structure
prs = Presentation(r'output\presentations\run_20251024_172953\deck.pptx')
slide = prs.slides[1]  # Slide 2

print('Slide title/name:')
if slide.shapes.title:
    print(f'  {slide.shapes.title.text}')
else:
    # Try to find title in shapes
    for shape in slide.shapes:
        if shape.has_text_frame and 'RSA' in shape.text_frame.text:
            print(f'  {shape.text_frame.text}')
            break

tables = [s for s in slide.shapes if s.has_table]
if tables:
    table = max(tables, key=lambda s: s.table.rows.__len__()).table
    print(f'\nOriginal deck - First 15 rows, columns A (0) and B (1):')
    print('=' * 100)
    print(f'Row  Col A (CAMPAIGN)                              Col B (MEDIA)')
    print('=' * 100)
    for i in range(min(15, len(table.rows))):
        cell_a = table.cell(i, 0)
        cell_b = table.cell(i, 1)
        text_a = (cell_a.text_frame.text if cell_a.text_frame else '').replace('\n', '\\n')[:40]
        text_b = (cell_b.text_frame.text if cell_b.text_frame else '').replace('\n', '\\n')[:30]
        print(f'{i:2d}   {text_a:40} {text_b:30}')
