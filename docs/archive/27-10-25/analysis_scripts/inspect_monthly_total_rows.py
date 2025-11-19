#!/usr/bin/env python
"""Inspect MONTHLY TOTAL row contents."""

import sys
from pptx import Presentation

def inspect_monthly_total_rows(pptx_path: str, slide_idx: int):
    """Show content of MONTHLY TOTAL rows across all columns."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]

    # Find table
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if not table:
        print(f"No table found on slide {slide_idx + 1}")
        return

    print(f"\n{'='*100}")
    print(f"MONTHLY TOTAL Row Inspection - Slide {slide_idx + 1}")
    print(f"{'='*100}\n")

    # Check all rows
    for row_idx in range(len(table.rows)):
        cell = table.cell(row_idx, 0)
        text = cell.text_frame.text if cell.text_frame else ""

        if "MONTHLY" in text.upper() and "TOTAL" in text.upper():
            print(f"\nRow {row_idx}: MONTHLY TOTAL detected")
            print("-" * 100)

            # Show content of first 5 columns
            for col_idx in range(min(5, len(table.columns))):
                cell = table.cell(row_idx, col_idx)
                cell_text = cell.text_frame.text if cell.text_frame else ""
                # Escape newlines for display
                display_text = cell_text.replace('\n', '\\n').replace('\r', '\\r')
                print(f"  Column {col_idx + 1}: '{display_text[:80]}'")

if __name__ == "__main__":
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else "output/presentations/run_20251024_172953/deck_test_fonts.pptx"
    slide_num = int(sys.argv[2]) if len(sys.argv) > 2 else 2

    inspect_monthly_total_rows(pptx_path, slide_num - 1)
