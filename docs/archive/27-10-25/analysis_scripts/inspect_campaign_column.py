"""Diagnostic tool to inspect column 0 (CAMPAIGN column) contents."""

from pptx import Presentation

prs_path = r"output\presentations\run_20251024_172953\deck_final.pptx"
prs = Presentation(prs_path)

# Look at slide 2 (index 1)
slide = prs.slides[1]

# Find the main table
tables = [shape for shape in slide.shapes if shape.has_table]
if tables:
    main_table = max(tables, key=lambda s: s.table.rows.__len__() * s.table.columns.__len__())
    table = main_table.table

    print(f"Table size: {len(table.rows)} rows x {len(table.columns)} columns\n")
    print("Column 0 (CAMPAIGN column) contents:")
    print("=" * 80)

    for row_idx in range(len(table.rows)):
        cell = table.cell(row_idx, 0)
        text = cell.text_frame.text if cell.text_frame else ""

        # Check for merge attributes
        tc = cell._tc
        rowSpan = tc.get('rowSpan', '')
        gridSpan = tc.get('gridSpan', '')
        vMerge = tc.get('vMerge', '')
        hMerge = tc.get('hMerge', '')

        merge_info = []
        if rowSpan: merge_info.append(f"rowSpan={rowSpan}")
        if gridSpan: merge_info.append(f"gridSpan={gridSpan}")
        if vMerge: merge_info.append(f"vMerge={vMerge}")
        if hMerge: merge_info.append(f"hMerge={hMerge}")

        merge_str = " [" + ", ".join(merge_info) + "]" if merge_info else ""

        # Truncate long text
        text_display = text.replace('\n', '\\n')[:60]

        print(f"Row {row_idx:2d}: '{text_display}'{merge_str}")
