#!/usr/bin/env python
"""Diagnose merge conflicts by inspecting cell states."""

import sys
from pptx import Presentation

def diagnose_slide(pptx_path: str, slide_idx: int):
    """Diagnose merge state and row labels."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]

    # Find table
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if not table:
        print(f"No table found on slide {slide_idx + 1}")
        return

    print(f"\n{'='*80}")
    print(f"Slide {slide_idx + 1} Merge Diagnostics")
    print(f"{'='*80}")
    print(f"Table: {len(table.rows)} rows × {len(table.columns)} cols\n")

    # Check each row
    for row_idx in range(min(35, len(table.rows))):  # First 35 rows
        # Get cell text from column 1
        cell_0 = table.cell(row_idx, 0)
        text = cell_0.text_frame.text.strip() if cell_0.text_frame else ""

        # Check merge attributes in columns 1-3
        merge_info = []
        for col_idx in range(min(3, len(table.columns))):
            cell = table.cell(row_idx, col_idx)
            tc = cell._tc

            attrs = []
            if tc.get('rowSpan'):
                attrs.append(f"rSpan={tc.get('rowSpan')}")
            if tc.get('gridSpan'):
                attrs.append(f"gSpan={tc.get('gridSpan')}")
            if tc.get('vMerge'):
                attrs.append(f"vM={tc.get('vMerge')}")
            if tc.get('hMerge'):
                attrs.append(f"hM={tc.get('hMerge')}")

            if attrs:
                merge_info.append(f"C{col_idx+1}[{','.join(attrs)}]")

        # Format row info
        label = text[:30] if text else "(empty)"
        merges = " ".join(merge_info) if merge_info else "no merges"

        # Highlight special rows
        marker = ""
        if "MONTHLY" in text.upper() and "TOTAL" in text.upper():
            marker = " <-- MONTHLY TOTAL"
        elif "GRAND" in text.upper() and "TOTAL" in text.upper():
            marker = " <-- GRAND TOTAL"
        elif "CARRIED" in text.upper() and "FORWARD" in text.upper():
            marker = " <-- CARRIED FORWARD"

        print(f"Row {row_idx:2d}: {label:30s} | {merges}{marker}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python diagnose_merge_conflict.py <pptx_path> <slide_number>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_num = int(sys.argv[2])
    slide_idx = slide_num - 1

    diagnose_slide(pptx_path, slide_idx)
