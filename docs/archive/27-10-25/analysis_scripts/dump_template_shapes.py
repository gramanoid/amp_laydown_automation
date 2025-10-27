from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def dump_shapes(template_path: Path) -> None:
    prs = Presentation(template_path)
    slide = prs.slides[0]
    print(f"Slide layout: {slide.slide_layout.name}")
    for idx, shape in enumerate(slide.shapes):
        shape_type = getattr(shape, "shape_type", None)
        name = getattr(shape, "name", "")
        try:
            placeholder = shape.placeholder_format.type  # type: ignore[union-attr]
        except ValueError:
            placeholder = None
        has_table = bool(getattr(shape, "has_table", False))
        has_text = bool(getattr(shape, "has_text_frame", False))
        print(
            f"#{idx}: name={name!r}, type={shape_type}, "
            f"placeholder={placeholder}, table={has_table}, text={has_text}"
        )
        if has_table:
            table = shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            print(f"    table size: {rows} x {cols}")


if __name__ == "__main__":
    template = Path("template/Template_V4_FINAL_071025.pptx").resolve()
    dump_shapes(template)
