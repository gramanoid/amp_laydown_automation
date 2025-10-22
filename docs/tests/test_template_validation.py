from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from amp_automation.config.loader import Config
from amp_automation.validation import reconciliation as reconciliation_module


def _build_test_config(tmp_path: Path) -> Config:
    summary_tiles = {
        "quarter_budgets": {
            "q1": {"shape": "QuarterBudgetQ1", "prefix": "Q1: ", "number_format": "£{value:,.0f}K", "scale": 0.001, "label": "Q1"},
            "q2": {"shape": "QuarterBudgetQ2", "prefix": "Q2: ", "number_format": "£{value:,.0f}K", "scale": 0.001, "label": "Q2"},
            "q3": {"shape": "QuarterBudgetQ3", "prefix": "Q3: ", "number_format": "£{value:,.0f}K", "scale": 0.001, "label": "Q3"},
            "q4": {"shape": "QuarterBudgetQ4", "prefix": "Q4: ", "number_format": "£{value:,.0f}K", "scale": 0.001, "label": "Q4"},
        },
        "media_share": {
            "television": {"shape": "MediaShareTelevision", "label": "TV", "number_format": "{value:.0f}%", "scale": 100},
            "digital": {"shape": "MediaShareDigital", "label": "DIG.", "number_format": "{value:.0f}%", "scale": 100},
            "other": {"shape": "MediaShareOther", "label": "OTHER", "number_format": "{value:.0f}%", "scale": 100},
        },
        "funnel_share": {
            "awareness": {"shape": "FunnelShareAwareness", "label": "AWA", "number_format": "{value:.0f}%", "scale": 100},
            "consideration": {"shape": "FunnelShareConsideration", "label": "CON", "number_format": "{value:.0f}%", "scale": 100},
            "purchase": {"shape": "FunnelSharePurchase", "label": "PUR", "number_format": "{value:.0f}%", "scale": 100},
        },
        "footer_notes": {"shape": "FooterNotes", "default_text": ""},
    }

    payload = {
        "presentation": {
            "title": {"shape": "TitlePlaceholder", "format": "{market} - {brand}"},
            "summary_tiles": summary_tiles,
        }
    }

    return Config(data=payload, path=tmp_path / "test_config.json")


def _build_test_dataframe() -> pd.DataFrame:
    month_values = {
        "Jan": 120_000.0,
        "Feb": 80_000.0,
        "Mar": 60_000.0,
        "Apr": 90_000.0,
        "May": 40_000.0,
        "Jun": 30_000.0,
        "Jul": 25_000.0,
        "Aug": 35_000.0,
        "Sep": 20_000.0,
        "Oct": 50_000.0,
        "Nov": 45_000.0,
        "Dec": 55_000.0,
    }

    def make_row(media_type: str, funnel_stage: str, scale: float) -> dict[str, float]:
        row = {
            "Country": "KSA",
            "Brand": "Sensodyne",
            "Year": 2025,
            "Mapped Media Type": media_type,
            "Funnel Stage": funnel_stage,
        }
        for month, value in month_values.items():
            row[month] = value * scale
        row["Total Cost"] = sum(value * scale for value in month_values.values())
        return row

    rows = [
        make_row("Television", "Awareness", 1.0),
        make_row("Digital", "Consideration", 0.6),
        make_row("Other", "Purchase", 0.2),
    ]

    return pd.DataFrame(rows)


def _populate_slide_with_summary(slide, summary: dict, *, market: str, brand: str) -> None:
    textbox_specs = {
        "TitlePlaceholder": f"{market} - {brand}",
        "QuarterBudgetQ1": summary["quarter_budgets"]["q1"]["display"],
        "QuarterBudgetQ2": summary["quarter_budgets"]["q2"]["display"],
        "QuarterBudgetQ3": summary["quarter_budgets"]["q3"]["display"],
        "QuarterBudgetQ4": summary["quarter_budgets"]["q4"]["display"],
        "MediaShareTelevision": summary["media_share"]["television"]["display"],
        "MediaShareDigital": summary["media_share"]["digital"]["display"],
        "MediaShareOther": summary["media_share"]["other"]["display"],
        "FunnelShareAwareness": summary["funnel_share"]["awareness"]["display"],
        "FunnelShareConsideration": summary["funnel_share"]["consideration"]["display"],
        "FunnelSharePurchase": summary["funnel_share"]["purchase"]["display"],
        "FooterNotes": "",
    }

    for idx, (shape_name, text) in enumerate(textbox_specs.items()):
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5 + 0.3 * idx), Inches(4.0), Inches(0.25))
        textbox.name = shape_name
        textbox.text_frame.text = text


def test_reconciliation_passes_for_matching_slide(tmp_path):
    config = _build_test_config(tmp_path)
    df = _build_test_dataframe()

    summary_cfg = config.section("presentation")["summary_tiles"]
    expected_summary = reconciliation_module._compute_expected_summary(df, "KSA", "Sensodyne", 2025, summary_cfg)  # type: ignore[attr-defined]
    prs = Presentation()
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    _populate_slide_with_summary(slide, expected_summary, market="KSA", brand="Sensodyne")
    ppt_path = tmp_path / "matching.pptx"
    prs.save(ppt_path)

    results = reconciliation_module.generate_reconciliation_report(
        ppt_path,
        tmp_path / "dummy.xlsx",
        config,
        data_frame=df,
    )

    assert len(results) == 1
    assert results[0].passed
    failing = [comp for comp in results[0].comparisons if not comp.passed]
    assert not failing


def test_reconciliation_flags_mismatched_tile(tmp_path):
    config = _build_test_config(tmp_path)
    df = _build_test_dataframe()

    summary_cfg = config.section("presentation")["summary_tiles"]
    expected_summary = reconciliation_module._compute_expected_summary(df, "KSA", "Sensodyne", 2025, summary_cfg)  # type: ignore[attr-defined]
    prs = Presentation()
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    _populate_slide_with_summary(slide, expected_summary, market="KSA", brand="Sensodyne")

    # Introduce a deliberate mismatch on Q1 budget
    for shape in slide.shapes:
        if shape.name == "QuarterBudgetQ1" and shape.has_text_frame:
            shape.text_frame.text = "Q1: £0K"

    ppt_path = tmp_path / "mismatch.pptx"
    prs.save(ppt_path)

    results = reconciliation_module.generate_reconciliation_report(
        ppt_path,
        tmp_path / "dummy.xlsx",
        config,
        data_frame=df,
    )

    assert len(results) == 1
    assert not results[0].passed
    failing = [comp for comp in results[0].comparisons if not comp.passed]
    assert any(comp.category == "quarter_budgets" and comp.label == "Q1" for comp in failing)
