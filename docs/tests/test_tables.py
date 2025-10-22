"""Tests for table assembly helpers."""

from __future__ import annotations

import logging

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Emu

from amp_automation.presentation.tables import CellStyleContext
from amp_automation.presentation.template_clone import clone_template_table
from amp_automation.presentation.assembly import _populate_cloned_table


def build_context() -> CellStyleContext:
    return CellStyleContext(
        margin_left_right_pt=3.6,
        margin_emu_lr=45720,
        default_font_name="Calibri",
        font_size_header=Pt(10),
        font_size_body=Pt(9),
        color_black=RGBColor(0, 0, 0),
        color_light_gray_text=RGBColor(200, 200, 200),
        color_table_gray=RGBColor(191, 191, 191),
        color_header_green=RGBColor(0, 255, 0),
        color_subtotal_gray=RGBColor(217, 217, 217),
        color_tv=RGBColor(113, 212, 141),
        color_digital=RGBColor(253, 242, 183),
        color_ooh=RGBColor(255, 191, 0),
        color_other=RGBColor(176, 211, 255),
    )


def test_add_and_style_table_populates_cells() -> None:
    # Use a real template cloning cycle
    template_prs = Presentation("template/Template_V4_FINAL_071025.pptx")
    template_slide = template_prs.slides[0]

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    table_shape = clone_template_table(template_slide, slide, "MainDataTable")
    table_data = [
        ["Header A", "Header B", "Header C"],
        ["ROW 1", "100", "Other"],
        ["GRAND TOTAL", "100", ""],
    ]
    cell_metadata: dict[tuple[int, int], dict[str, object]] = {}

    result = _populate_cloned_table(table_shape, table_data, cell_metadata)

    assert result is True
    created_table = table_shape.table
    assert created_table.cell(0, 0).text.upper() == "HEADER A"
    assert "100" in created_table.cell(1, 1).text


def test_add_and_style_table_respects_alignment_and_dual_line() -> None:
    template_prs = Presentation("template/Template_V4_FINAL_071025.pptx")
    template_slide = template_prs.slides[0]

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    table_shape = clone_template_table(template_slide, slide, "MainDataTable")

    table_data = [
        ["Header A", "Header B", "Header C"],
        ["Row 1", "100", "Other"],
        ["GRAND TOTAL", "100", ""],
    ]
    cell_metadata: dict[tuple[int, int], dict[str, object]] = {}

    result = _populate_cloned_table(table_shape, table_data, cell_metadata)

    assert result is True
    table = table_shape.table
    cell_left = table.cell(1, 0)
    assert cell_left.text_frame.word_wrap is False
    paragraphs = cell_left.text_frame.paragraphs
    assert paragraphs[0].alignment == PP_ALIGN.CENTER
    assert cell_left.text_frame.text == "ROW 1"

    cell_right = table.cell(1, 1)
    assert cell_right.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER

    assert table.rows[0].height == Emu(161729)
    assert table.rows[1].height == Emu(99205)
    assert table.rows[2].height == Emu(99205)
