#!/usr/bin/env python3
"""
Script to unmerge horizontally merged cells in PowerPoint presentation
Requires: python-pptx library (pip install python-pptx)
"""

from pptx import Presentation

# List of merged cells to fix: (slide_index_0based, row_index_0based, content)
MERGED_CELLS_TO_FIX = [
    (1, 11, "MONTHLY TOTAL (£ 000)"),
    (1, 22, ""),
    (2, 11, "ALWAYS ON DCOMM"),
    (2, 22, ""),
    (3, 11, "GRAND TOTAL"),
    (4, 11, ""),
    (4, 22, ""),
    (6, 11, "MONTHLY TOTAL (£ 000)"),
    (6, 22, "RAMADAN ACTIVATION"),
    (10, 11, "GRAND TOTAL"),
    (12, 11, "MONTHLY TOTAL (£ 000)"),
    (12, 22, "MONTHLY TOTAL (£ 000)"),
    (12, 33, "CARRIED FORWARD"),
    (12, 34, "GRAND TOTAL"),
    (13, 11, "GRAND TOTAL"),
    (14, 11, "MONTHLY TOTAL (£ 000)"),
    (15, 11, ""),
    (15, 22, ""),
    (16, 11, "MONTHLY TOTAL (£ 000)"),
    (16, 22, "GRAND TOTAL"),
    (17, 11, "MONTHLY TOTAL (£ 000)"),
    (17, 22, "GRAND TOTAL"),
    (19, 11, "GRAND TOTAL"),
    (21, 11, "MONTHLY TOTAL (£ 000)"),
    (21, 22, ""),
    (22, 11, ""),
    (22, 22, ""),
    (23, 11, "NEXTNOVA CLINICAL REPAIR"),
    (24, 11, "GENDER"),
    (24, 22, "ALWAYS ON DCOMM"),
    (24, 33, "GRAND TOTAL"),
    (25, 11, "MONTHLY TOTAL (£ 000)"),
    (25, 22, "MASTERBRAND"),
    (25, 23, "MONTHLY TOTAL (£ 000)"),
    (26, 11, "ROCKET"),
    (26, 22, "ALWAYS ON DCOMM"),
    (26, 33, "GRAND TOTAL"),
    (27, 11, "MONTHLY TOTAL (£ 000)"),
    (28, 11, "ALWAYS ON DCOMM"),
    (28, 22, ""),
    (29, 11, ""),
    (29, 22, ""),
    (31, 11, "TIME 2.0"),
    (32, 11, "MONTHLY TOTAL (£ 000)"),
    (33, 11, "ALWAYS ON DCOMM"),
    (34, 11, "HERITAGE"),
    (35, 11, "EX REAL WORLD EVIDENCE"),
    (35, 22, "OTRIVIN PLUS"),
    (35, 23, "MONTHLY TOTAL (£ 000)"),
    (38, 11, "CLINICAL WHITE"),
    (38, 22, "ALWAYS ON DCOMM"),
    (38, 33, "GRAND TOTAL"),
    (40, 11, ""),
    (40, 22, ""),
    (41, 11, "PRONATURALS"),
    (42, 11, ""),
    (42, 22, ""),
    (43, 11, "IMMUNO VITA C MAX + ELDERBERRY"),
    (44, 11, "GRAND TOTAL"),
    (46, 11, "EX CONDITION"),
    (46, 22, ""),
    (48, 11, "GRAND TOTAL"),
    (50, 11, ""),
    (50, 22, ""),
    (50, 33, "CARRIED FORWARD"),
    (50, 34, "GRAND TOTAL"),
    (53, 11, ""),
    (53, 22, ""),
    (55, 11, "MONTHLY TOTAL (£ 000)"),
    (57, 11, "EX ‑ MINERS ARMS"),
    (60, 11, ""),
    (60, 22, ""),
    (61, 11, "LETS TREAT IT RIGHT"),
    (62, 11, "YOU DID IT"),
    (62, 22, ""),
    (63, 11, ""),
    (73, 11, "FACES‑CONDITION"),
    (73, 22, "FEEL FAMILIAR"),
    (74, 11, "MONTHLY TOTAL (£ 000)"),
    (74, 22, "ALWAYS ON DCOMM"),
    (75, 11, "MONTHLY TOTAL (£ 000)"),
    (75, 22, "GRAND TOTAL"),
    (77, 11, "EX WORLD ORAL HEALTH DAY"),
    (78, 11, "GRAND TOTAL"),
    (80, 11, "MONTHLY TOTAL (£ 000)"),
    (81, 11, ""),
    (83, 11, "HHP GLOBAL PAIN AWARENESS WEEK"),
]


def unmerge_cell_columns(input_file, output_file):
    """
    Unmerge cells that span the first 3 columns
    
    Args:
        input_file: Path to input PowerPoint file
        output_file: Path to save fixed PowerPoint file
    """
    print(f"Loading presentation: {input_file}")
    prs = Presentation(input_file)
    
    fixes_applied = 0
    
    for slide_idx, row_idx, content in MERGED_CELLS_TO_FIX:
        try:
            slide = prs.slides[slide_idx]
            
            # Find table on slide
            table = None
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    break
            
            if table is None:
                print(f"  WARNING: No table found on slide {slide_idx + 1}")
                continue
            
            # Get the cell that needs to be unmerged
            cell = table.cell(row_idx, 0)
            tc = cell._tc
            
            # Check if cell has gridSpan attribute (is merged)
            grid_span = tc.get('gridSpan')
            
            if grid_span is not None and int(grid_span) >= 3:
                # Remove the gridSpan attribute to unmerge
                tc.attrib.pop('gridSpan', None)
                
                # Ensure the content stays in column 0
                if content:
                    cell.text = content
                
                fixes_applied += 1
                print(f"  Fixed: Slide {slide_idx + 1}, Row {row_idx + 1}")
            else:
                print(f"  INFO: Slide {slide_idx + 1}, Row {row_idx + 1} - Cell not merged or span < 3")
                
        except Exception as e:
            print(f"  ERROR: Slide {slide_idx + 1}, Row {row_idx + 1} - {str(e)}")
    
    print(f"\nSaving fixed presentation to: {output_file}")
    prs.save(output_file)
    print(f"Complete! Applied {fixes_applied} fixes out of {len(MERGED_CELLS_TO_FIX)} cells")


if __name__ == "__main__":
    INPUT_FILE = "GeneratedDeck_20251022_092659.pptx"
    OUTPUT_FILE = "GeneratedDeck_20251022_092659_FIXED.pptx"
    
    unmerge_cell_columns(INPUT_FILE, OUTPUT_FILE)
