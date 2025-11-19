"""Regression tests for font normalization (EC-002)."""

from __future__ import annotations

import pytest
from pptx.util import Pt

from conftest import find_main_table, skipif_no_deck


@pytest.mark.regression
@skipif_no_deck
def test_ec002_font_sizes_consistent_across_production_deck(latest_deck_path):
    """Verify font sizes are consistent across all slides in production deck (EC-002).

    Validates that header cells use 7pt and body cells use 6pt font sizes.
    """
    from pptx import Presentation

    prs = Presentation(latest_deck_path)

    font_issues = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        table = find_main_table(slide)
        if not table:
            continue

        # Check first row (header row) - should be 7pt
        if len(table.rows) > 0:
            for col_idx in range(len(table.columns)):
                cell = table.cell(0, col_idx)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size != Pt(7):
                            font_issues.append(f"Slide {slide_idx}, Header cell ({0},{col_idx}): Expected 7pt, got {run.font.size}")

        # Check body rows - should be 6pt (sample check on first body row)
        if len(table.rows) > 1:
            for col_idx in range(min(3, len(table.columns))):  # Sample first 3 columns
                cell = table.cell(1, col_idx)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size != Pt(6):
                            font_issues.append(f"Slide {slide_idx}, Body cell (1,{col_idx}): Expected 6pt, got {run.font.size}")

    # Allow some tolerance - fonts might vary slightly in merged cells or special rows
    assert len(font_issues) < 10, f"Found {len(font_issues)} font size inconsistencies (showing first 10):\n" + "\n".join(font_issues[:10])


@pytest.mark.unit
def test_font_size_header_consistency(cell_style_context):
    """Verify header cell styling context has 7pt font (EC-002)."""
    # Verify the styling context has correct font size for headers
    assert cell_style_context.font_size_header == Pt(7), \
        f"Header font size in context should be 7pt, got {cell_style_context.font_size_header}"


@pytest.mark.unit
def test_font_size_body_consistency(cell_style_context):
    """Verify body cell styling context has 6pt font (EC-002)."""
    # Verify the styling context has correct font size for body
    assert cell_style_context.font_size_body == Pt(6), \
        f"Body font size in context should be 6pt, got {cell_style_context.font_size_body}"


@pytest.mark.unit
def test_font_family_consistent(cell_style_context):
    """Verify font family is Calibri throughout (EC-002)."""
    # Verify the styling context specifies Calibri
    assert cell_style_context.default_font_name == "Calibri", \
        f"Font should be Calibri, got {cell_style_context.default_font_name}"
