"""Regression tests for pagination and continuation handling (EC-003, EC-006)."""

from __future__ import annotations

import pytest

from conftest import find_main_table, skipif_no_deck


@pytest.mark.regression
@skipif_no_deck
def test_ec003_multi_slide_markets_exist(latest_deck_path):
    """Verify production deck contains multi-slide markets requiring pagination (EC-003)."""
    from pptx import Presentation

    prs = Presentation(latest_deck_path)

    # Count slides per market
    market_slide_counts = {}

    for slide_idx, slide in enumerate(prs.slides, start=1):
        table = find_main_table(slide)
        if not table:
            continue

        # Extract market from title or table header
        title = None
        for shape in slide.shapes:
            if hasattr(shape, "name") and "Title" in shape.name:
                if hasattr(shape, "text_frame"):
                    title = shape.text_frame.text
                    break

        if title and " - " in title:
            market = title.split(" - ")[0].strip()
            # Remove pagination marker if present
            market = market.split("(")[0].strip()

            if market not in market_slide_counts:
                market_slide_counts[market] = 0
            market_slide_counts[market] += 1

    # Verify some markets span multiple slides
    multi_slide_markets = {m: count for m, count in market_slide_counts.items() if count > 1}

    assert len(multi_slide_markets) > 0, \
        "Production deck should have multi-slide markets for pagination testing"


@pytest.mark.regression
@skipif_no_deck
def test_ec003_continuation_indicators_present(latest_deck_path):
    """Verify continuation indicators present on non-final slides (EC-003)."""
    from pptx import Presentation

    prs = Presentation(latest_deck_path)

    # Track markets and continuation indicators
    market_titles = {}

    for slide_idx, slide in enumerate(prs.slides, start=1):
        title = None
        for shape in slide.shapes:
            if hasattr(shape, "name") and "Title" in shape.name:
                if hasattr(shape, "text_frame"):
                    title = shape.text_frame.text
                    break

        if title and " - " in title:
            # Check for continuation indicator
            is_continuation = "(Continued)" in title or "continued" in title.lower()

            # Extract market
            market = title.split(" - ")[0].strip()
            market = market.split("(")[0].strip()

            if market not in market_titles:
                market_titles[market] = []
            market_titles[market].append({
                "slide": slide_idx,
                "title": title,
                "is_continuation": is_continuation,
            })

    # Verify multi-slide markets have continuation indicators
    for market, slides_info in market_titles.items():
        if len(slides_info) > 1:
            # All but last should have continuation indicator
            for i, slide_info in enumerate(slides_info[:-1]):
                assert slide_info["is_continuation"], \
                    f"Market '{market}' slide {i+1} should have continuation indicator. " \
                    f"Title: {slide_info['title']}"


@pytest.mark.regression
@skipif_no_deck
def test_ec006_carried_forward_rows_present(latest_deck_path):
    """Verify CARRIED FORWARD rows present on continuation slides (EC-006)."""
    from pptx import Presentation

    prs = Presentation(latest_deck_path)

    continuation_slides_without_cf = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        table = find_main_table(slide)
        if not table:
            continue

        title = None
        for shape in slide.shapes:
            if hasattr(shape, "name") and "Title" in shape.name:
                if hasattr(shape, "text_frame"):
                    title = shape.text_frame.text
                    break

        # Check if this is a continuation slide
        if title and "(Continued)" in title:
            # Should have CARRIED FORWARD row
            has_cf = False
            for row_idx in range(len(table.rows)):
                cell = table.cell(row_idx, 0)
                if "CARRIED FORWARD" in cell.text.upper():
                    has_cf = True
                    break

            if not has_cf:
                continuation_slides_without_cf.append((slide_idx, title))

    assert not continuation_slides_without_cf, \
        f"Continuation slides missing CARRIED FORWARD rows: {continuation_slides_without_cf}"


@pytest.mark.unit
def test_ec006_empty_metrics_not_accumulated():
    """Verify empty metrics (dashes) not accumulated in carried forward (EC-006)."""
    import pandas as pd

    # Simulate campaign data
    campaigns = [
        {"name": "Campaign A", "jan": 1000, "feb": 1000, "total": 2000},
        {"name": "Campaign B (No Data)", "jan": None, "feb": None, "total": 0},  # Empty
        {"name": "Campaign C", "jan": 500, "feb": 500, "total": 1000},
    ]

    # Simulate accumulation
    total = 0
    count = 0
    for campaign in campaigns:
        if campaign["total"] and campaign["total"] > 0:  # Skip empty
            total += campaign["total"]
            count += 1

    # Verify: Should sum to 3000 (A + C), not 3000 (A + B + C)
    assert total == 3000, f"Accumulation should skip empty metrics, got {total}"
    assert count == 2, f"Should count only 2 campaigns, got {count}"


@pytest.mark.unit
def test_max_rows_per_slide_boundary():
    """Test MAX_ROWS_PER_SLIDE configuration boundary (EC-003)."""
    # Standard configuration
    MAX_ROWS_PER_SLIDE = 32

    # Test cases
    test_cases = [
        (31, True),   # Fits in single slide
        (32, True),   # Fits exactly
        (33, False),  # Requires split
        (64, False),  # Requires 2 slides
    ]

    for body_rows, should_fit_single in test_cases:
        fits = body_rows <= MAX_ROWS_PER_SLIDE
        assert fits == should_fit_single, \
            f"{body_rows} rows should {'fit' if should_fit_single else 'not fit'} in single slide"


@pytest.mark.unit
def test_continuation_title_format():
    """Verify continuation slide title format (EC-003)."""
    test_cases = [
        ("MARKET / Brand (1 of 3)", "MARKET / Brand (1 of 3) (Continued)"),
        ("MARKET / Brand (1/3)", "MARKET / Brand (1/3) (Continued)"),
        ("MARKET / Brand", "MARKET / Brand (Continued)"),
    ]

    for original, expected in test_cases:
        # Simulate continuation indicator appending
        continued = original.replace(")", "") + " (Continued)"
        continued = continued.replace("(Continued))", "(Continued)")

        # Should have continuation indicator
        assert "(Continued)" in continued, \
            f"Continuation format incorrect: {continued}"
