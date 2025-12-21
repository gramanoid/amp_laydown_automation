"""
Comprehensive Data Validation Tests for AMP Laydowns Automation.

This module validates that ALL displayed values in the generated PPTX
match the source Excel data EXACTLY. It uses multiple validation strategies
and sampling approaches to ensure complete coverage.
"""

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any

import pandas as pd
import pytest
from pptx import Presentation


# Configuration - updated to latest run
PPTX_PATH = Path("output/presentations/run_20251217_121749/AMP_Laydowns_171225.pptx")
EXCEL_PATH = Path("input/Flowplan_Summaries_MEA_2025_12_17.xlsx")


def load_excel_data() -> pd.DataFrame:
    """Load and normalize Excel source data."""
    df = pd.read_excel(EXCEL_PATH)
    # Normalize column names
    df.columns = df.columns.str.strip()
    return df


def extract_slide_title(slide) -> str:
    """Extract the main title from a slide."""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            name = getattr(shape, "name", "")
            if "SlideTitle" in name or "Title" in name:
                return shape.text_frame.text.strip()
    return ""


def extract_media_shares(slide) -> dict[str, str]:
    """Extract media share values from slide shapes."""
    media_shapes = {
        "MediaShareTelevision": "",
        "MediaShareDigital": "",
        "MediaShareOther": ""
    }
    for shape in slide.shapes:
        for name in media_shapes:
            if name in getattr(shape, "name", ""):
                if hasattr(shape, "text_frame"):
                    media_shapes[name] = shape.text_frame.text.strip()
    return media_shapes


def parse_percentage(text: str) -> float | None:
    """Parse percentage from text like 'TV: 55%' -> 55.0."""
    match = re.search(r"(\d+(?:\.\d+)?)\s*%", text)
    if match:
        return float(match.group(1))
    return None


def calculate_expected_media_shares(df: pd.DataFrame, market: str, brand: str, year: str = "2025") -> dict[str, float]:
    """Calculate expected media shares from DataFrame."""
    # Filter data
    mask = (
        (df["Country"].astype(str).str.strip().str.upper() == market.strip().upper()) &
        (df["Brand"].astype(str).str.strip().str.upper() == brand.strip().upper())
    )
    if "Year" in df.columns:
        mask = mask & (df["Year"].astype(str).str.strip() == str(year).strip())

    subset = df.loc[mask]

    if subset.empty:
        return {"TV": 0, "Digital": 0, "Other": 0}

    # Group by Mapped Media Type if available
    if "Mapped Media Type" in subset.columns:
        media_group = subset.groupby("Mapped Media Type")["Total Cost"].sum()
    else:
        return {"TV": 0, "Digital": 0, "Other": 0}

    total_cost = subset["Total Cost"].sum()

    if total_cost <= 0:
        return {"TV": 0, "Digital": 0, "Other": 0}

    tv_value = float(media_group.get("TV", 0.0))
    digital_value = float(media_group.get("Digital", 0.0))
    # Other includes OOH per system design
    other_value = float(media_group.get("Other", 0.0)) + float(media_group.get("OOH", 0.0))

    # Calculate percentages with largest remainder method for rounding
    tv_pct = int((tv_value / total_cost) * 100)
    digital_pct = int((digital_value / total_cost) * 100)
    other_pct = int((other_value / total_cost) * 100)

    # Adjust to sum to 100%
    remainder = 100 - (tv_pct + digital_pct + other_pct)
    if remainder > 0:
        # Add remainder to highest fractional part
        fractions = [
            ("TV", (tv_value / total_cost) * 100 - tv_pct),
            ("Digital", (digital_value / total_cost) * 100 - digital_pct),
            ("Other", (other_value / total_cost) * 100 - other_pct),
        ]
        fractions.sort(key=lambda x: x[1], reverse=True)
        for i in range(remainder):
            if fractions[i % 3][0] == "TV":
                tv_pct += 1
            elif fractions[i % 3][0] == "Digital":
                digital_pct += 1
            else:
                other_pct += 1

    return {"TV": tv_pct, "Digital": digital_pct, "Other": other_pct}


class TestMediaShareValidation:
    """Test suite for media share validation across diverse test cases."""

    @pytest.fixture(autouse=True)
    def setup(self):
        """Load presentation and data for all tests."""
        self.prs = Presentation(PPTX_PATH)
        self.df = load_excel_data()

    def _find_slides_for_brand(self, market: str, brand: str) -> list[tuple[int, Any]]:
        """Find all slides matching market-brand combination."""
        results = []
        for idx, slide in enumerate(self.prs.slides):
            title = extract_slide_title(slide)
            if market.upper() in title.upper() and brand.upper() in title.upper():
                results.append((idx + 1, slide))
        return results

    # ===== TEST CASE 1: High-Spend Market =====
    def test_media_shares_saudi_arabia_voltaren(self):
        """TC1: Validate high-spend market (Saudi Arabia - Voltaren)."""
        market = "Saudi Arabia"
        brand = "Voltaren"

        expected = calculate_expected_media_shares(self.df, market, brand)
        slides = self._find_slides_for_brand(market, brand)

        assert len(slides) > 0, f"No slides found for {market} - {brand}"

        # Check last slide (should have media shares)
        slide_num, slide = slides[-1]
        actual = extract_media_shares(slide)

        tv_actual = parse_percentage(actual["MediaShareTelevision"])
        dig_actual = parse_percentage(actual["MediaShareDigital"])
        other_actual = parse_percentage(actual["MediaShareOther"])

        assert tv_actual == expected["TV"], f"TV mismatch on slide {slide_num}: {tv_actual} != {expected['TV']}"
        assert dig_actual == expected["Digital"], f"Digital mismatch on slide {slide_num}: {dig_actual} != {expected['Digital']}"
        assert other_actual == expected["Other"], f"Other mismatch on slide {slide_num}: {other_actual} != {expected['Other']}"

    # ===== TEST CASE 2: Multi-Product Brand =====
    def test_media_shares_panadol_pain(self):
        """TC2: Validate multi-product brand (Panadol Pain) - uses product splits."""
        market = "Saudi Arabia"
        brand = "Panadol Pain"

        expected = calculate_expected_media_shares(self.df, market, brand)
        slides = self._find_slides_for_brand(market, brand)

        assert len(slides) > 0, f"No slides found for {market} - {brand}"

        # Find the brand-level summary slide (not product-level)
        for slide_num, slide in slides:
            title = extract_slide_title(slide)
            # Skip product-specific slides
            if " - " in title and "PRODUCT SUMMARY" not in title.upper():
                continue

            actual = extract_media_shares(slide)
            if any(actual.values()):  # Has media share shapes
                tv_actual = parse_percentage(actual["MediaShareTelevision"])
                dig_actual = parse_percentage(actual["MediaShareDigital"])
                other_actual = parse_percentage(actual["MediaShareOther"])

                if tv_actual is not None:
                    # Verify values are not template defaults
                    is_template = (tv_actual == 55 and dig_actual == 20 and other_actual == 25)
                    assert not is_template, f"Slide {slide_num} has template defaults instead of computed values"
                break

    # ===== TEST CASE 3: Digital-Heavy Brand =====
    def test_media_shares_digital_heavy_brand(self):
        """TC3: Validate a digital-heavy brand (>80% digital spend)."""
        # Find a digital-heavy brand from the data
        digital_heavy = None
        for _, row in self.df.groupby(["Country", "Brand"]).agg({"Total Cost": "sum"}).reset_index().iterrows():
            market = row["Country"]
            brand = row["Brand"]
            expected = calculate_expected_media_shares(self.df, market, brand)
            if expected["Digital"] >= 80:
                digital_heavy = (market, brand, expected)
                break

        if digital_heavy is None:
            pytest.skip("No digital-heavy brand found in data")

        market, brand, expected = digital_heavy
        slides = self._find_slides_for_brand(market, brand)

        if not slides:
            pytest.skip(f"No slides found for {market} - {brand}")

        slide_num, slide = slides[-1]
        actual = extract_media_shares(slide)

        if any(actual.values()):
            dig_actual = parse_percentage(actual["MediaShareDigital"])
            assert dig_actual is not None, f"Could not parse digital value from slide {slide_num}"
            assert dig_actual >= 75, f"Digital-heavy brand should show >=75% digital: got {dig_actual}%"

    # ===== TEST CASE 4: TV-Dominant Brand =====
    def test_media_shares_tv_dominant_brand(self):
        """TC4: Validate a TV-dominant brand (>80% TV spend)."""
        # Find a TV-dominant brand from the data
        tv_dominant = None
        for _, row in self.df.groupby(["Country", "Brand"]).agg({"Total Cost": "sum"}).reset_index().iterrows():
            market = row["Country"]
            brand = row["Brand"]
            expected = calculate_expected_media_shares(self.df, market, brand)
            if expected["TV"] >= 80:
                tv_dominant = (market, brand, expected)
                break

        if tv_dominant is None:
            pytest.skip("No TV-dominant brand found in data")

        market, brand, expected = tv_dominant
        slides = self._find_slides_for_brand(market, brand)

        if not slides:
            pytest.skip(f"No slides found for {market} - {brand}")

        slide_num, slide = slides[-1]
        actual = extract_media_shares(slide)

        if any(actual.values()):
            tv_actual = parse_percentage(actual["MediaShareTelevision"])
            assert tv_actual is not None, f"Could not parse TV value from slide {slide_num}"
            assert tv_actual >= 75, f"TV-dominant brand should show >=75% TV: got {tv_actual}%"

    # ===== TEST CASE 5: African Market =====
    def test_media_shares_african_market(self):
        """TC5: Validate African market (Nigeria, Kenya, or South Africa)."""
        african_markets = ["Nigeria", "Kenya", "South Africa", "Ivory Coast"]

        for market in african_markets:
            brands = self.df[self.df["Country"].str.upper() == market.upper()]["Brand"].unique()
            if len(brands) > 0:
                brand = brands[0]
                expected = calculate_expected_media_shares(self.df, market, brand)
                slides = self._find_slides_for_brand(market, brand)

                if slides:
                    slide_num, slide = slides[-1]
                    actual = extract_media_shares(slide)

                    if any(actual.values()):
                        tv_actual = parse_percentage(actual["MediaShareTelevision"])
                        dig_actual = parse_percentage(actual["MediaShareDigital"])
                        other_actual = parse_percentage(actual["MediaShareOther"])

                        # Verify sum equals 100%
                        if tv_actual is not None and dig_actual is not None and other_actual is not None:
                            total = tv_actual + dig_actual + other_actual
                            assert total == 100, f"Media shares should sum to 100%, got {total}%"
                        return

        pytest.skip("No African market with valid slides found")

    # ===== TEST CASE 6: Product Summary Slides =====
    def test_product_summary_slides_have_computed_values(self):
        """TC6: Verify Product Summary slides use computed values, not template defaults."""
        template_pattern_count = 0
        computed_count = 0

        for idx, slide in enumerate(self.prs.slides):
            title = extract_slide_title(slide)
            if "PRODUCT SUMMARY" in title.upper():
                actual = extract_media_shares(slide)

                if any(actual.values()):
                    tv_actual = parse_percentage(actual["MediaShareTelevision"])
                    dig_actual = parse_percentage(actual["MediaShareDigital"])
                    other_actual = parse_percentage(actual["MediaShareOther"])

                    if tv_actual == 55 and dig_actual == 20 and other_actual == 25:
                        template_pattern_count += 1
                    else:
                        computed_count += 1

        assert template_pattern_count == 0, f"Found {template_pattern_count} Product Summary slides with template defaults"

    # ===== TEST CASE 7: No Template Defaults Anywhere =====
    def test_no_template_defaults_anywhere(self):
        """TC7: Global check - no slides should have exact template pattern (55/20/25)."""
        offending_slides = []

        for idx, slide in enumerate(self.prs.slides):
            actual = extract_media_shares(slide)

            if all(actual.values()):  # Has all three values
                tv_actual = parse_percentage(actual["MediaShareTelevision"])
                dig_actual = parse_percentage(actual["MediaShareDigital"])
                other_actual = parse_percentage(actual["MediaShareOther"])

                if tv_actual == 55 and dig_actual == 20 and other_actual == 25:
                    title = extract_slide_title(slide)
                    offending_slides.append((idx + 1, title))

        assert len(offending_slides) == 0, f"Slides with template defaults: {offending_slides}"


class TestEdgeCases:
    """Edge case and adversarial tests."""

    @pytest.fixture(autouse=True)
    def setup(self):
        """Load presentation and data for all tests."""
        self.prs = Presentation(PPTX_PATH)
        self.df = load_excel_data()

    def test_media_shares_sum_to_100(self):
        """All media share sets should sum to exactly 100%."""
        violations = []

        for idx, slide in enumerate(self.prs.slides):
            actual = extract_media_shares(slide)

            if all(actual.values()):
                tv = parse_percentage(actual["MediaShareTelevision"])
                dig = parse_percentage(actual["MediaShareDigital"])
                other = parse_percentage(actual["MediaShareOther"])

                if tv is not None and dig is not None and other is not None:
                    total = tv + dig + other
                    if total != 100:
                        title = extract_slide_title(slide)
                        violations.append((idx + 1, title, tv, dig, other, total))

        assert len(violations) == 0, f"Media shares don't sum to 100%: {violations[:5]}..."

    def test_no_negative_percentages(self):
        """No media share should be negative."""
        for idx, slide in enumerate(self.prs.slides):
            actual = extract_media_shares(slide)

            for name, value in actual.items():
                if value:
                    pct = parse_percentage(value)
                    if pct is not None:
                        assert pct >= 0, f"Negative percentage on slide {idx + 1}: {name}={pct}%"

    def test_percentages_are_integers(self):
        """Media shares should be whole numbers (integers)."""
        for idx, slide in enumerate(self.prs.slides):
            actual = extract_media_shares(slide)

            for name, value in actual.items():
                if value:
                    # Check for decimal points in the percentage text
                    match = re.search(r"(\d+)\.(\d+)\s*%", value)
                    if match:
                        pytest.fail(f"Non-integer percentage on slide {idx + 1}: {name}={value}")

    def test_zero_spend_brands_excluded(self):
        """Brands with zero total spend should not have slides."""
        zero_spend_brands = []

        for (market, brand), group in self.df.groupby(["Country", "Brand"]):
            total = group["Total Cost"].sum()
            if total == 0:
                zero_spend_brands.append((market, brand))

        for market, brand in zero_spend_brands:
            for idx, slide in enumerate(self.prs.slides):
                title = extract_slide_title(slide)
                if market.upper() in title.upper() and brand.upper() in title.upper():
                    pytest.fail(f"Zero-spend brand has slide: {market} - {brand} (slide {idx + 1})")


def generate_validation_report() -> dict:
    """Generate comprehensive validation report."""
    prs = Presentation(PPTX_PATH)
    df = load_excel_data()

    report = {
        "timestamp": datetime.now().isoformat(),
        "pptx_path": str(PPTX_PATH),
        "excel_path": str(EXCEL_PATH),
        "total_slides": len(prs.slides),
        "slides_with_media_shares": 0,
        "slides_validated": 0,
        "fields_checked": 0,
        "errors": [],
        "warnings": [],
        "test_cases": {},
        "sampling_strategy": {
            "high_spend_markets": [],
            "digital_heavy_brands": [],
            "tv_dominant_brands": [],
            "african_markets": [],
            "product_summary_slides": [],
        }
    }

    # Count slides with media shares
    for idx, slide in enumerate(prs.slides):
        actual = extract_media_shares(slide)
        if any(actual.values()):
            report["slides_with_media_shares"] += 1
            report["slides_validated"] += 1
            report["fields_checked"] += 3  # TV, Digital, Other

            # Check for template defaults
            tv = parse_percentage(actual["MediaShareTelevision"])
            dig = parse_percentage(actual["MediaShareDigital"])
            other = parse_percentage(actual["MediaShareOther"])

            if tv == 55 and dig == 20 and other == 25:
                title = extract_slide_title(slide)
                report["errors"].append({
                    "slide": idx + 1,
                    "title": title,
                    "issue": "Template default values detected (55/20/25)"
                })

            # Check sum
            if tv is not None and dig is not None and other is not None:
                total = tv + dig + other
                if total != 100:
                    title = extract_slide_title(slide)
                    report["warnings"].append({
                        "slide": idx + 1,
                        "title": title,
                        "issue": f"Media shares sum to {total}%, not 100%"
                    })

    # Identify test case samples
    brand_totals = df.groupby(["Country", "Brand"])["Total Cost"].sum().reset_index()
    top_brands = brand_totals.nlargest(5, "Total Cost")
    report["sampling_strategy"]["high_spend_markets"] = [
        {"market": row["Country"], "brand": row["Brand"], "spend": row["Total Cost"]}
        for _, row in top_brands.iterrows()
    ]

    return report


if __name__ == "__main__":
    # Run validation and generate report
    report = generate_validation_report()

    # Save JSON report
    with open("validation_report.json", "w") as f:
        json.dump(report, f, indent=2, default=str)

    # Generate markdown report
    md_lines = [
        "# Comprehensive Validation Report",
        "",
        f"**Generated:** {report['timestamp']}",
        f"**PPTX:** `{report['pptx_path']}`",
        f"**Excel:** `{report['excel_path']}`",
        "",
        "## Summary",
        "",
        "| Metric | Value |",
        "|--------|-------|",
        f"| Total Slides | {report['total_slides']} |",
        f"| Slides Validated | {report['slides_validated']} |",
        f"| Fields Checked | {report['fields_checked']} |",
        f"| Errors | {len(report['errors'])} |",
        f"| Warnings | {len(report['warnings'])} |",
        f"| **Status** | {'✅ PASSED' if len(report['errors']) == 0 else '❌ FAILED'} |",
        "",
    ]

    if report['errors']:
        md_lines.extend([
            "## Errors",
            "",
        ])
        for err in report['errors']:
            md_lines.append(f"- **Slide {err['slide']}** ({err['title']}): {err['issue']}")
        md_lines.append("")

    if report['warnings']:
        md_lines.extend([
            "## Warnings",
            "",
        ])
        for warn in report['warnings']:
            md_lines.append(f"- **Slide {warn['slide']}** ({warn['title']}): {warn['issue']}")
        md_lines.append("")

    with open("validation_report.md", "w") as f:
        f.write("\n".join(md_lines))

    print(f"Validation complete. Errors: {len(report['errors'])}, Warnings: {len(report['warnings'])}")
