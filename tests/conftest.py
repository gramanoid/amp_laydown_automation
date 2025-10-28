"""Shared pytest fixtures and utilities for test suite."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Generator

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from amp_automation.presentation.tables import CellStyleContext, TableLayout


# ============================================================================
# PATH FIXTURES
# ============================================================================


@pytest.fixture(scope="session")
def project_root() -> Path:
    """Path to project root."""
    return Path(__file__).parent.parent


@pytest.fixture(scope="session")
def template_path(project_root: Path) -> Path:
    """Path to master template."""
    return project_root / "template" / "Template_V4_FINAL_071025.pptx"


@pytest.fixture(scope="session")
def excel_path(project_root: Path) -> Path:
    """Path to test Excel data (Lumina export)."""
    return project_root / "template" / "BulkPlanData_2025_10_14.xlsx"


@pytest.fixture(scope="session")
def latest_deck_path(project_root: Path) -> Path:
    """Path to latest production deck (28-10-25)."""
    return project_root / "output" / "presentations" / "run_20251028_143931" / "AMP_Presentation_20251028_143931.pptx"


@pytest.fixture(scope="session")
def contract_path(project_root: Path) -> Path:
    """Path to structural validation contract."""
    return project_root / "config" / "presentation_contract.json"


# ============================================================================
# PRESENTATION FIXTURES
# ============================================================================


@pytest.fixture
def blank_presentation() -> Presentation:
    """Create blank presentation for testing."""
    return Presentation()


@pytest.fixture
def blank_slide(blank_presentation: Presentation) -> Presentation:
    """Create blank slide in new presentation."""
    slide_layout = blank_presentation.slide_layouts[6]  # Blank layout
    return blank_presentation.slides.add_slide(slide_layout)


# ============================================================================
# TABLE STYLING FIXTURES
# ============================================================================


@pytest.fixture
def cell_style_context() -> CellStyleContext:
    """Standard cell styling context for tests."""
    return CellStyleContext(
        margin_left_right_pt=3.6,
        margin_emu_lr=45720,
        default_font_name="Calibri",
        font_size_header=Pt(7),       # Updated for 27-10-25: 7pt headers
        font_size_body=Pt(6),         # Updated for 27-10-25: 6pt body
        font_size_body_compact=Pt(6), # Compact rows use same as body (6pt)
        color_black=RGBColor(0, 0, 0),
        color_light_gray_text=RGBColor(200, 200, 200),
        color_table_gray=RGBColor(191, 191, 191),
        color_header_green=RGBColor(0, 255, 0),
        color_subtotal_gray=RGBColor(217, 217, 217),
        color_tv=RGBColor(211, 254, 201),
        color_digital=RGBColor(253, 242, 183),
        color_ooh=RGBColor(255, 217, 97),
        color_other=RGBColor(176, 211, 255),
    )


@pytest.fixture
def table_layout() -> TableLayout:
    """Standard table layout for tests."""
    return TableLayout(
        placeholder_name="",
        shape_name="TestTable",
        position={
            "left": Inches(0.184),
            "top": Inches(1.4),
            "width": Inches(9.1),
            "height": Inches(3.3),
        },
        row_height_header=Pt(12),
        row_height_body=Pt(10.5),  # Standard from template
        row_height_subtotal=Pt(10.5),
        column_widths=[Inches(1.0), Inches(0.8), Inches(0.6), Inches(0.6), Inches(0.6)],
        top_override=None,
        height_rule_available=False,
    )


# ============================================================================
# TEST DATA FIXTURES
# ============================================================================


@pytest.fixture
def campaign_names_with_hyphens() -> list[str]:
    """Campaign names that require smart line breaking (EC-001 tests)."""
    return [
        "FACES-CONDITION",
        "LONG-CAMPAIGN-NAME",
        "SHORT",
        "MULTI-WORD-CAMPAIGN-NAME-WITH-HYPHENS",
    ]


@pytest.fixture
def market_case_variations() -> dict[str, list[str]]:
    """Market names with case variations for reconciliation tests (EC-004, EC-005)."""
    return {
        "south africa": ["south africa", "SOUTH AFRICA", "South Africa"],
        "egypt": ["egypt", "EGYPT", "Egypt"],
        "morocco": ["MOR", "MOROCCO", "Morocco"],
        "ksa": ["ksa", "KSA", "Ksa"],
    }


@pytest.fixture
def brand_case_variations() -> dict[str, list[str]]:
    """Brand names with case variations for reconciliation tests."""
    return {
        "fanta": ["fanta", "FANTA", "Fanta"],
        "sprite": ["sprite", "SPRITE", "Sprite"],
        "coca-cola": ["coca-cola", "COCA-COLA", "Coca-Cola"],
    }


@pytest.fixture
def pagination_formats() -> list[str]:
    """Pagination format variations for title parsing tests (EC-010)."""
    return [
        "SOUTH AFRICA / Fanta (1 of 3)",  # Standard format
        "SOUTH AFRICA / Fanta (1/3)",     # Alternative format
        "SOUTH AFRICA / Fanta",           # No pagination
    ]


# ============================================================================
# LOGGING FIXTURE
# ============================================================================


@pytest.fixture
def test_logger() -> logging.Logger:
    """Logger for test execution."""
    logger = logging.getLogger("test")
    logger.setLevel(logging.DEBUG)
    return logger


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================


def find_main_table(slide: Presentation) -> object | None:
    """Find main table shape in slide (first table found)."""
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def is_media_header(label: str) -> bool:
    """Check if row label is a media channel header."""
    label_upper = str(label).strip().upper()
    return label_upper in {"TELEVISION", "DIGITAL", "OOH", "OTHER"}


def is_merged_cell(table, col_idx: int, row_idx: int) -> bool:
    """Check if cell at (row_idx, col_idx) is part of a merged range."""
    try:
        cell = table.cell(row_idx, col_idx)
        # In python-pptx, merged cells have tc element with mcs (merge cell start) or vMerge
        tc = cell._tc
        grid_span = tc.get("gridSpan")
        v_merge = tc.find(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}vMerge")
        return bool(grid_span) or bool(v_merge)
    except Exception:
        return False


def extract_font_size(run) -> Pt | None:
    """Extract font size from run, handling None case."""
    if hasattr(run, 'font') and hasattr(run.font, 'size'):
        return run.font.size
    return None


# ============================================================================
# TEST SKIP MARKERS
# ============================================================================


skipif_no_deck = pytest.mark.skipif(
    not Path("output/presentations/run_20251028_143931/AMP_Presentation_20251028_143931.pptx").exists(),
    reason="Latest production deck not found"
)

skipif_no_template = pytest.mark.skipif(
    not Path("template/Template_V4_FINAL_071025.pptx").exists(),
    reason="Template not found"
)

skipif_no_excel = pytest.mark.skipif(
    not Path("template/BulkPlanData_2025_10_14.xlsx").exists(),
    reason="Excel data file not found"
)
